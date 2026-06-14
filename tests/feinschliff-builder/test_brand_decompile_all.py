"""Regression: brand_decompile_all must never clobber a hand-authored tokens.json.

Original bug: the tokens.json write ran unconditionally — `--dry-run` only
guarded the layout derivation, so a dry run still rewrote tokens.json. Worse,
an unparseable tokens.json was silently reset to `{}` and overwritten with a
minimal stub, with no backup (layouts get `layouts.bak/`; tokens got nothing).
"""
from __future__ import annotations

import importlib.util
import json
import sys
from pathlib import Path

import pytest
from pptx import Presentation


REPO_ROOT = Path(__file__).resolve().parent.parent.parent / "feinschliff-builder"
SCRIPT = REPO_ROOT / "scripts" / "brand_decompile_all.py"


def _load_script_module():
    spec = importlib.util.spec_from_file_location("brand_decompile_all", SCRIPT)
    module = importlib.util.module_from_spec(spec)
    sys.modules["brand_decompile_all"] = module
    spec.loader.exec_module(module)
    return module


def _make_source_pptx(path: Path) -> Path:
    pres = Presentation()
    pres.slides.add_slide(pres.slide_layouts[6])
    pres.save(str(path))
    return path


def _make_brand_pack(root: Path, *, tokens: str, layouts_yaml: str) -> Path:
    brand = root / "acme"
    brand.mkdir(parents=True)
    (brand / "tokens.json").write_text(tokens, encoding="utf-8")
    (brand / "verify-map.yaml").write_text(layouts_yaml, encoding="utf-8")
    return brand


def _run_main(monkeypatch, brand: Path, pptx: Path, *extra: str) -> int:
    mod = _load_script_module()
    monkeypatch.setattr(sys, "argv", [
        "brand_decompile_all.py",
        "--brand-pack", str(brand),
        "--source-pptx", str(pptx),
        *extra,
    ])
    return mod.main()


def test_dry_run_leaves_tokens_json_untouched(tmp_path, monkeypatch):
    tokens = json.dumps({"color": {"accent": {"$value": "#ff0000"}}}, indent=2) + "\n"
    brand = _make_brand_pack(tmp_path, tokens=tokens, layouts_yaml="layouts:\n  cover: 1\n")
    pptx = _make_source_pptx(tmp_path / "source.pptx")
    assert _run_main(monkeypatch, brand, pptx, "--dry-run") == 0
    assert (brand / "tokens.json").read_bytes() == tokens.encode("utf-8")
    assert not (brand / "tokens.json.bak").exists()
    assert not (brand / "layouts").exists()


def test_unparseable_tokens_json_aborts_without_modifying(tmp_path, monkeypatch):
    garbage = '{"slide": '
    brand = _make_brand_pack(tmp_path, tokens=garbage, layouts_yaml="layouts: {}\n")
    pptx = _make_source_pptx(tmp_path / "source.pptx")
    with pytest.raises(SystemExit) as excinfo:
        _run_main(monkeypatch, brand, pptx)
    assert "unparseable tokens.json" in str(excinfo.value)
    assert (brand / "tokens.json").read_text(encoding="utf-8") == garbage
    assert not (brand / "tokens.json.bak").exists()


def test_normal_run_backs_up_tokens_json_and_merges(tmp_path, monkeypatch):
    tokens = json.dumps({"color": {"accent": {"$value": "#ff0000"}}}, indent=2) + "\n"
    brand = _make_brand_pack(tmp_path, tokens=tokens, layouts_yaml="layouts: {}\n")
    pptx = _make_source_pptx(tmp_path / "source.pptx")
    assert _run_main(monkeypatch, brand, pptx) == 0
    assert (brand / "tokens.json.bak").read_bytes() == tokens.encode("utf-8")
    merged = json.loads((brand / "tokens.json").read_text(encoding="utf-8"))
    assert merged["color"]["accent"]["$value"] == "#ff0000"
    assert "width_emu" in merged["slide"]
    assert "height_emu" in merged["slide"]
