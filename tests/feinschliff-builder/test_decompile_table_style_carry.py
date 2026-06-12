"""A native-carried table must bring its referenced `<a:tblStyle>` along.

The spliced `<a:tbl>` keeps its `<a:tableStyleId>{GUID}</a:tableStyleId>`,
but the GUID points into the SOURCE deck's tableStyles.xml. Without carrying
the style, the output deck's renderer falls back to a default table style —
wrong header fill, band colours, and borders (the Shapes 'metrics' slide
rendered a maroon header instead of the source's blue).
"""
from __future__ import annotations

import base64

from lxml import etree
from pptx import Presentation

from feinschliff_builder.decompile.pptx_svg_decompile import (
    CanvasMap,
    _emit_graphic_frame,
    _source_table_style,
)

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_GUID = "{72833802-FEF1-4C79-8D5D-14CF1EAF98D9}"

_TBL_STYLE = (
    f'<a:tblStyle xmlns:a="{_A}" styleId="{_GUID}" styleName="Themed">'
    f'<a:firstRow><a:tcStyle><a:fill><a:solidFill><a:schemeClr val="accent2"/>'
    f"</a:solidFill></a:fill></a:tcStyle></a:firstRow></a:tblStyle>"
)

_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

_FRAME = f"""
  <p:graphicFrame xmlns:p="{_P}" xmlns:a="{_A}">
    <p:nvGraphicFramePr><p:cNvPr id="5" name="Table 1"/><p:cNvGraphicFramePr/><p:nvPr/></p:nvGraphicFramePr>
    <p:xfrm><a:off x="914400" y="914400"/><a:ext cx="6096000" cy="1828800"/></p:xfrm>
    <a:graphic><a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/table">
      <a:tbl>
        <a:tblPr firstRow="1"><a:tableStyleId>{_GUID}</a:tableStyleId></a:tblPr>
        <a:tblGrid><a:gridCol w="3048000"/></a:tblGrid>
        <a:tr h="370840"><a:tc><a:txBody><a:bodyPr/><a:p><a:r><a:t>Head</a:t></a:r></a:p></a:txBody><a:tcPr/></a:tc></a:tr>
      </a:tbl>
    </a:graphicData></a:graphic>
  </p:graphicFrame>
"""


def _prs_with_styled_source():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for part in prs.part.package.iter_parts():
        if str(part.partname) == "/ppt/tableStyles.xml":
            root = etree.fromstring(part.blob)
            root.append(etree.fromstring(_TBL_STYLE.encode()))
            part._blob = etree.tostring(root, xml_declaration=True,
                                        encoding="UTF-8", standalone=True)
    return prs, slide


def test_source_table_style_resolves_and_bakes():
    prs, slide = _prs_with_styled_source()
    el = _source_table_style(slide, _GUID, {"accent2": "#4D90EF"})
    assert el is not None
    srgb = el.find(f".//{{{_A}}}srgbClr")
    assert srgb is not None and srgb.get("val") == "4D90EF", (
        "schemeClr must be baked against the source theme"
    )


def test_table_carry_attaches_style_part():
    prs, slide = _prs_with_styled_source()
    frame = etree.fromstring(_FRAME.encode())
    shapes: list = []
    cmap = CanvasMap(12192000, 6858000, 1920, 1080)
    _emit_graphic_frame(frame, (0, 0), shapes, slide, cmap, {}, {})
    [shape] = shapes
    assert shape.native_xml is not None
    assert shape.native_parts, "carried table must attach its tblStyle"
    [entry] = shape.native_parts
    carried = etree.fromstring(base64.b64decode(entry["table_style"]))
    assert carried.get("styleId") == _GUID


def test_missing_style_id_carries_without_parts():
    prs, slide = _prs_with_styled_source()
    frame = etree.fromstring(_FRAME.replace(_GUID, "{0000-NOT-THERE}").encode())
    shapes: list = []
    cmap = CanvasMap(12192000, 6858000, 1920, 1080)
    _emit_graphic_frame(frame, (0, 0), shapes, slide, cmap, {}, {})
    [shape] = shapes
    assert shape.native_xml is not None
    assert not shape.native_parts
