from __future__ import annotations

import subprocess
from pathlib import Path

import pytest

pptx = pytest.importorskip("pptx")
from pptx import Presentation
from pptx.util import Inches


REPO = Path(__file__).resolve().parents[2] / "feinschliff-builder"


def _build_rough_deck(tmp_path: Path) -> Path:
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    for i in range(3):
        shape = slide.shapes.add_shape(
            1, Inches(1 + i * 1.5), Inches(2), Inches(1.2), Inches(0.8),
        )
        shape.text_frame.text = f"Box {i+1}"
    out = tmp_path / "rough.pptx"
    pres.save(str(out))
    return out


def _require_rendering_backend():
    try:
        import cairosvg  # noqa: F401
    except (ImportError, OSError) as exc:
        try:
            from playwright.sync_api import sync_playwright  # noqa: F401
        except ImportError:
            pytest.skip(f"Neither cairosvg/cairo nor playwright is available: {exc}")


def test_polish_refurbish_all_runs(tmp_path):
    _require_rendering_backend()
    src = _build_rough_deck(tmp_path)
    out = tmp_path / "polished.pptx"
    result = subprocess.run(
        [
            "uv", "run", "feinschliff", "deck", "polish",
            str(src), "-o", str(out),
            "--refurbish-all", "--brand", "feinschliff",
        ],
        capture_output=True, encoding="utf-8", cwd=REPO,
        env={
            **__import__("os").environ,
            "DYLD_FALLBACK_LIBRARY_PATH": "/opt/homebrew/opt/cairo/lib",
            "PYTHONIOENCODING": "utf-8",
            "PYTHONUTF8": "1",
        },
    )
    # The test only requires no error AND that refurbish artifacts were emitted.
    assert result.returncode == 0, result.stderr
    refurbish_dir = out.parent / "refurbished"
    assert refurbish_dir.exists()
    report = out.parent / "refurbish_report.md"
    assert report.exists()
    # At least one refurbished DSL should exist
    artifacts = list(refurbish_dir.glob("*.dsl")) + list(refurbish_dir.glob("*.exc.dsl")) + list(refurbish_dir.glob("*.svg.dsl"))
    assert len(artifacts) >= 1, "expected at least one refurbished DSL artifact"


def test_polish_refurbish_produces_real_polished_pptx(tmp_path):
    """End-to-end: rough PPTX with 1 diagram → polished PPTX containing
    the brand-perfect diagram on a feinschliff layout."""
    _require_rendering_backend()
    src = _build_rough_deck(tmp_path)  # existing helper, 1 slide with 3 boxes
    out = tmp_path / "polished.pptx"
    result = subprocess.run(
        [
            "uv", "run", "feinschliff", "deck", "polish",
            str(src), "-o", str(out),
            "--refurbish-all", "--brand", "feinschliff",
        ],
        capture_output=True, encoding="utf-8", cwd=REPO,
        env={
            **__import__("os").environ,
            "DYLD_FALLBACK_LIBRARY_PATH": "/opt/homebrew/opt/cairo/lib",
            "PYTHONIOENCODING": "utf-8",
            "PYTHONUTF8": "1",
        },
    )
    assert result.returncode == 0, result.stderr
    # The output PPTX should NOT be a byte-for-byte copy of the input
    # (it should have been actually rebuilt).
    assert out.exists()
    assert out.stat().st_size > 0
    assert out.read_bytes() != src.read_bytes(), "output is identical to input — stub still active"

    # Open output and confirm it has at least 1 slide
    from pptx import Presentation
    pres = Presentation(str(out))
    assert len(pres.slides) >= 1
