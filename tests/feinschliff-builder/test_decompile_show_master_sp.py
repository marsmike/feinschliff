"""`showMasterSp="0"` on a slide / layout must hide master decorative shapes.

PowerPoint renders a slide's master chrome (logos, footer plates, decorative
blocks) only when the slide's — or, absent a slide-level flag, its layout's —
`showMasterSp` attribute is not "0". Corporate templates use hide-master
layouts for covers and dividers; walking the master unconditionally lands a
phantom plate (e.g. a dark tx2 box) on every decompiled hide-master slide.

Master PLACEHOLDER shapes are unaffected: the flag governs plain master
shapes only, and placeholder inheritance still flows through the layout.
"""
from __future__ import annotations

from lxml import etree
from pptx import Presentation

from feinschliff_builder.decompile.pptx_svg_decompile import CanvasMap, walk_slide

_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

_MASTER_RECT = f"""
  <p:sp xmlns:p="{_P}" xmlns:a="{_A}">
    <p:nvSpPr><p:cNvPr id="99" name="MasterPlate"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
    <p:spPr>
      <a:xfrm><a:off x="7670800" y="5238875"/><a:ext cx="4064000" cy="520700"/></a:xfrm>
      <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
      <a:solidFill><a:srgbClr val="262626"/></a:solidFill>
    </p:spPr>
  </p:sp>
"""


def _slide_with_master_rect():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    master = slide.slide_layout.slide_master
    spTree = master.element.find(f".//{{{_P}}}cSld/{{{_P}}}spTree")
    spTree.append(etree.fromstring(_MASTER_RECT.encode()))
    return slide


def _master_plates(slide):
    cmap = CanvasMap(12192000, 6858000, 1920, 1080)
    shapes = walk_slide(slide, cmap, {}, {})
    return [s for s in shapes if s.fill == "#262626"]


def test_master_shape_inherited_by_default():
    slide = _slide_with_master_rect()
    assert len(_master_plates(slide)) == 1


def test_layout_show_master_sp_0_hides_master_shape():
    slide = _slide_with_master_rect()
    slide.slide_layout.element.set("showMasterSp", "0")
    assert _master_plates(slide) == []


def test_slide_show_master_sp_0_hides_master_shape():
    slide = _slide_with_master_rect()
    slide.element.set("showMasterSp", "0")
    assert _master_plates(slide) == []


def test_slide_flag_overrides_hiding_layout():
    slide = _slide_with_master_rect()
    slide.slide_layout.element.set("showMasterSp", "0")
    slide.element.set("showMasterSp", "1")
    assert len(_master_plates(slide)) == 1
