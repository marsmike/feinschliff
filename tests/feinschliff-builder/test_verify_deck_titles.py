"""Unit tests for lib/verify/deck/titles — title extractors."""
from __future__ import annotations

import json
from pathlib import Path

import pytest

from feinschliff_builder.verify.deck.titles import (
    extract_titles_from_plan,
    extract_titles_from_pptx,
)


def test_extract_titles_from_deck_plan(tmp_path: Path):
    """deck_plan.json: list of {layout_id, slot_values: {title: "..."}}."""
    plan = {
        "slides": [
            {"layout_id": "title-orange", "slot_values": {"title": "Q3 in one slide"}},
            {"layout_id": "key-takeaways", "slot_values": {"title": "Three things to remember"}},
            {"layout_id": "end", "slot_values": {}},  # no title
        ],
    }
    p = tmp_path / "deck_plan.json"
    p.write_text(json.dumps(plan))
    titles = extract_titles_from_plan(p)
    assert titles == ["Q3 in one slide", "Three things to remember", ""]


def test_extract_titles_from_content_plan(tmp_path: Path):
    """content_plan.json: list of slides with raw `title` field."""
    plan = {
        "slides": [
            {"title": "Why now", "concept": "context"},
            {"title": "What changed", "concept": "complication"},
        ],
    }
    p = tmp_path / "content_plan.json"
    p.write_text(json.dumps(plan))
    titles = extract_titles_from_plan(p)
    assert titles == ["Why now", "What changed"]


def test_extract_titles_from_plan_missing_file(tmp_path: Path):
    """Missing file raises FileNotFoundError with a clear message."""
    with pytest.raises(FileNotFoundError):
        extract_titles_from_plan(tmp_path / "nope.json")


def test_extract_titles_from_plan_invalid_shape(tmp_path: Path):
    """A JSON file without a `slides` list raises ValueError."""
    p = tmp_path / "bad.json"
    p.write_text(json.dumps({"not_slides": []}))
    with pytest.raises(ValueError, match="slides"):
        extract_titles_from_plan(p)


def test_extract_titles_from_pptx(tmp_path: Path):
    """Build a 2-slide .pptx with python-pptx, extract titles via the helper."""
    from pptx import Presentation
    prs = Presentation()
    blank_layout = prs.slide_layouts[5]  # title-only
    for title_text in ["First slide claim", "Second slide claim"]:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.title.text = title_text
    out = tmp_path / "deck.pptx"
    prs.save(str(out))
    assert extract_titles_from_pptx(out) == ["First slide claim", "Second slide claim"]


def test_extract_titles_from_pptx_slide_without_title(tmp_path: Path):
    """A slide whose layout has no title placeholder → empty string."""
    from pptx import Presentation
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # truly blank
    prs.slides.add_slide(blank_layout)
    out = tmp_path / "deck.pptx"
    prs.save(str(out))
    assert extract_titles_from_pptx(out) == [""]


def test_extract_titles_from_deck_yaml(tmp_path):
    """Phase 6 /deck plan must work end-to-end on YAML deck plans
    (which is what `feinschliff deck build` consumes). Regression: prior
    to this fix, extract_titles_from_plan was JSON-only and crashed
    with 'Expecting value: line 1 column 1' on any YAML input.

    Two slides — first uses action_title, second uses content.title.
    """
    plan = """
brand: feinschliff
slides:
  - layout: layouts/recommendation.slide.dsl
    content:
      action_title: "Close the gap with three moves."
  - layout: layouts/full-bleed-cover.slide.dsl
    content:
      title: "Cover headline."
"""
    p = tmp_path / "deck.yaml"
    p.write_text(plan)
    assert extract_titles_from_plan(p) == [
        "Close the gap with three moves.",
        "Cover headline.",
    ]


def test_extract_titles_from_deck_yml_extension(tmp_path):
    """`.yml` extension also auto-detected as YAML."""
    plan = "slides:\n  - content:\n      action_title: \"Y\"\n"
    p = tmp_path / "deck.yml"
    p.write_text(plan)
    assert extract_titles_from_plan(p) == ["Y"]
