"""Subprocess-level CLI tests for `feinschliff deck apply-fixes` and
`feinschliff deck build --autofix`.

These tests exercise the full CLI binary-I/O surface (exit codes, file I/O,
stdout/stderr content) to complement the unit-level tests in
test_verify_autofix_apply.py.

Choice of approach for test_cli_build_autofix_inner_loop:
  - `deck build --autofix` writes the fixed plan back to disk before calling
    compile_slide.  compile_slide uses pure python-pptx (no soffice), so the
    build can run end-to-end in the test environment.  We therefore let the
    full build run, confirm the plan was mutated on disk, and check that exit
    code is 0.  The output .pptx is written to tmp_path to keep things clean.
  - No monkeypatching needed; the only external binary is soffice for the
    pptx-to-PNG render path, which is NOT part of `deck build`.

Critical-bug regression (PPTX compiled from pre-fix content):
  - The test verifies that the produced PPTX reflects the FIXED plan, not the
    pre-fix snapshot.  It opens the output .pptx with python-pptx and asserts
    that the slide text matches the shortened content, not the original 90-char
    string that exceeded the budget.  This directly catches the bug where
    `slides_spec` was captured before `apply_fixes` ran, causing the compile
    loop to iterate the pre-fix content.

Round-trip regression (engine-severity round-trip):
  - test_cli_round_trip_verify_static_to_apply_fixes exercises the full
    documented round-trip: deck verify-static --json → deck apply-fixes --defects.
    verify-static emits ENGINE severity values ("error"/"warning"/"info") because
    it calls validate() which maps legacy FATAL→ERROR, WARN→WARNING.
    apply-fixes must accept those engine values, not silently drop them as
    "malformed defect entry" (the bug that existed before this fix).
"""
from __future__ import annotations

import json
import subprocess
import sys
from pathlib import Path

import yaml
from pptx import Presentation

REPO = Path(__file__).resolve().parents[2]
FEINSCHLIFF = REPO / "feinschliff"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_plan_yaml(tmp_path: Path, content: dict) -> Path:
    """Write a minimal single-slide deck plan to disk and return its path."""
    plan_path = tmp_path / "plan.yaml"
    plan_path.write_text(
        yaml.safe_dump(content, sort_keys=False, allow_unicode=True),
        encoding="utf-8",
    )
    return plan_path


def _make_defects_json(tmp_path: Path, defects: list[dict]) -> Path:
    """Write a flat-list defects JSON to disk and return its path."""
    defects_path = tmp_path / "defects.json"
    defects_path.write_text(
        json.dumps(defects, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )
    return defects_path


def _run(cmd: list[str]) -> subprocess.CompletedProcess:
    import os
    env = os.environ.copy()
    env["PYTHONIOENCODING"] = "utf-8"
    env["PYTHONUTF8"] = "1"
    return subprocess.run(
        cmd,
        capture_output=True,
        text=True,
        encoding="utf-8",
        env=env,
        cwd=FEINSCHLIFF,
    )


# ---------------------------------------------------------------------------
# 1. apply-fixes — end-to-end binary I/O + exit code
# ---------------------------------------------------------------------------

def test_cli_apply_fixes_end_to_end(tmp_path: Path):
    """apply-fixes reads a plan + defects JSON, mutates the plan, writes -o.

    Defect shape matches Defect.to_dict() exactly (the flat-list shape that
    `deck verify-static --json` emits).  We plant a SLOT_OVERFLOW against
    action_title with a small budget so the shorten_slot patch fires, then
    confirm the output plan has the slot shortened.

    Fixture is sized to be BELOW the 20% swap threshold so shorten_slot fires
    (not swap_layout_larger).  budget=40, content=45 chars → 45 < 40*1.20=48.
    """
    budget = 40
    long_text = "This is a title too long for"  # 28 chars — pad to 45
    long_text = long_text.ljust(45, ".")       # 45 chars: above budget, below threshold
    assert len(long_text) == 45, f"Fixture length changed: {len(long_text)}"
    assert budget < len(long_text) < budget * 1.20, (
        f"Fixture must be above budget ({budget}) and below swap threshold "
        f"({budget * 1.20}); got {len(long_text)}"
    )

    plan = {
        "brand": "feinschliff",
        "out": str(tmp_path / "deck.pptx"),
        "slides": [
            {
                "layout": "layouts/action-title.slide.dsl",
                "content": {
                    "action_title": long_text,
                    "footer_left": "Corp",
                    "footer_right": "2026",
                },
            }
        ],
    }
    plan_path = _make_plan_yaml(tmp_path, plan)

    # Defect uses Defect.to_dict() shape (kind is the string value, not enum name)
    defects = [
        {
            "slide_index": 1,
            "kind": "slot-overflow",
            "severity": "warn",
            "message": f"slot 'action_title' overflows by {len(long_text) - budget} chars",
            "meta": {
                "slot": "action_title",
                "budget_chars": budget,
                "over_by": len(long_text) - budget,
            },
        }
    ]
    defects_path = _make_defects_json(tmp_path, defects)
    out_path = tmp_path / "fixed_plan.yaml"

    result = _run([
        sys.executable, "-m", "feinschliff.cli",
        "deck", "apply-fixes",
        str(plan_path),
        "--defects", str(defects_path),
        "-o", str(out_path),
    ])

    assert result.returncode == 0, (
        f"Expected exit code 0 (patches applied), got {result.returncode}.\n"
        f"stdout: {result.stdout}\nstderr: {result.stderr}"
    )
    assert out_path.is_file(), "Expected fixed plan YAML to be written"

    fixed = yaml.safe_load(out_path.read_text(encoding="utf-8"))
    fixed_text = fixed["slides"][0]["content"]["action_title"]
    assert len(fixed_text) <= budget, (
        f"Expected action_title shortened to ≤{budget} chars, "
        f"got {len(fixed_text)!r}: {fixed_text!r}"
    )


# ---------------------------------------------------------------------------
# 2. apply-fixes — exits 1 when no patches applicable
# ---------------------------------------------------------------------------

def test_cli_apply_fixes_exits_one_when_no_patches(tmp_path: Path):
    """apply-fixes exits 1 (informational) when defects have no mechanical fix.

    CLAIM_TITLE is not handled by plan_fixes (ambiguous, left for LLM revise),
    so a defects file containing only CLAIM_TITLE defects should yield exit 1.
    """
    plan = {
        "brand": "feinschliff",
        "out": str(tmp_path / "deck.pptx"),
        "slides": [
            {
                "layout": "layouts/action-title.slide.dsl",
                "content": {
                    "action_title": "Revenue grew",
                    "footer_left": "Corp",
                    "footer_right": "2026",
                },
            }
        ],
    }
    plan_path = _make_plan_yaml(tmp_path, plan)

    # CLAIM_TITLE is not patchable by plan_fixes — class is left for LLM.
    defects = [
        {
            "slide_index": 1,
            "kind": "claim-title",
            "severity": "warn",
            "message": "Title is not an action-oriented claim",
            "meta": {},
        }
    ]
    defects_path = _make_defects_json(tmp_path, defects)

    result = _run([
        sys.executable, "-m", "feinschliff.cli",
        "deck", "apply-fixes",
        str(plan_path),
        "--defects", str(defects_path),
    ])

    assert result.returncode == 1, (
        f"Expected exit code 1 (no patches applicable), got {result.returncode}.\n"
        f"stdout: {result.stdout}\nstderr: {result.stderr}"
    )


# ---------------------------------------------------------------------------
# 3. apply-fixes — exits 1 on empty defects file
# ---------------------------------------------------------------------------

def test_cli_apply_fixes_exits_one_on_empty_defects(tmp_path: Path):
    """apply-fixes exits 1 (nothing to do) when the defects list is empty."""
    plan = {
        "brand": "feinschliff",
        "out": str(tmp_path / "deck.pptx"),
        "slides": [
            {
                "layout": "layouts/action-title.slide.dsl",
                "content": {"action_title": "Clean plan"},
            }
        ],
    }
    plan_path = _make_plan_yaml(tmp_path, plan)
    defects_path = _make_defects_json(tmp_path, [])  # empty list

    result = _run([
        sys.executable, "-m", "feinschliff.cli",
        "deck", "apply-fixes",
        str(plan_path),
        "--defects", str(defects_path),
    ])

    assert result.returncode == 1, (
        f"Expected exit code 1 (no defects → nothing to do), got {result.returncode}.\n"
        f"stdout: {result.stdout}\nstderr: {result.stderr}"
    )


# ---------------------------------------------------------------------------
# 4. build --autofix — inner fix loop mutates plan on disk before compile
# ---------------------------------------------------------------------------

def test_cli_build_autofix_inner_loop(tmp_path: Path, monkeypatch):
    """deck build --autofix runs up to 3 autofix cycles, writes the mutated
    plan back to disk before compile, AND compiles the PPTX from the FIXED
    content — not the pre-fix snapshot.

    Fixture: executive-summary layout, action_title at 180 chars (budget = 168,
    overflow = 12 chars, BELOW the 20% swap threshold of 201.6 chars).  This
    forces shorten_slot (not swap_layout_larger), making it straightforward to
    assert that the PPTX text frame contains the shortened string.

    The fixture pins HEURISTIC textfit numbers (budget=168 from the Noto Sans
    table ratio at 2 bleed-model lines) and asserts the fixed string verbatim — with real metrics,
    prevent_orphan may legitimately NBSP-glue the last two words.  Force the
    heuristic path so the pins stay machine-independent; the autofix loop
    under test is metrics-agnostic.  The env var propagates to the `_run`
    subprocess via the os.environ copy.

    Regression guard for the critical bug: prior to the fix, `slides_spec` was
    captured before `apply_fixes` ran.  The plan on disk received the fix but
    the compile loop iterated the stale pre-fix list, so the PPTX was built
    from the original 90-char string.  After the fix, `slides_spec` is
    re-captured immediately after the autofix block, so the PPTX reflects the
    shortened content.

    We pass --skip-content-lint so the build proceeds to the PPTX write step
    regardless of any residual defects — our interest is in confirming the
    inner loop ran (plan-on-disk mutated) and that the PPTX was compiled from
    the fixed plan (pure python-pptx, no soffice needed).
    """
    monkeypatch.setenv("FEINSCHMIEDE_NO_REAL_METRICS", "1")
    # action_title budget for executive-summary = 168 chars (2 bleed-model lines).
    # Use 180 chars: above budget (triggers SLOT_OVERFLOW) but below the 20%
    # swap threshold (168 * 1.20 = 201.6), so shorten_slot fires.
    # The text must be a real string (not "X"*90) so sentence-boundary trimming
    # produces a stable shortened value we can assert against in the PPTX.
    long_text = "Action required: revenue declined three quarters in a row this year."
    assert len(long_text) == 68, f"Fixture length changed: {len(long_text)}"
    # Pad to exactly 90 chars so it's above budget=84 but below threshold=100.8
    long_text = long_text + " " + "A" * (180 - len(long_text) - 1)
    assert len(long_text) == 180, f"Padded fixture is {len(long_text)} chars, expected 180"
    assert 168 < len(long_text) <= 201, "Fixture must be above budget and below swap threshold"

    budget = 168
    out_pptx = tmp_path / "deck.pptx"
    plan = {
        "brand": "feinschliff",
        "out": str(out_pptx),
        "slides": [
            {
                "layout": "layouts/executive-summary.slide.dsl",
                "content": {
                    "action_title": long_text,
                    "footer_left": "Corp",
                    "footer_right": "2026",
                },
            }
        ],
    }
    plan_path = _make_plan_yaml(tmp_path, plan)

    result = _run([
        sys.executable, "-m", "feinschliff.cli",
        "deck", "build",
        "--autofix",
        "--skip-content-lint",   # let compile proceed; we test the fix loop, not lint
        str(plan_path),
    ])

    assert result.returncode == 0, (
        f"Expected exit code 0 from deck build --autofix --skip-content-lint, "
        f"got {result.returncode}.\n"
        f"stdout: {result.stdout}\nstderr: {result.stderr}"
    )

    # The autofix cycle count must appear in stdout (confirms the loop ran).
    assert "autofix cycle" in result.stdout, (
        f"Expected 'autofix cycle' in stdout — fix loop did not fire.\n"
        f"stdout: {result.stdout}"
    )

    # ── Plan-on-disk assertion ────────────────────────────────────────────────
    fixed_plan = yaml.safe_load(plan_path.read_text(encoding="utf-8"))
    fixed_text = fixed_plan["slides"][0]["content"]["action_title"]
    assert fixed_text != long_text, (
        "Expected autofix to shorten action_title on disk, but it was unchanged."
    )
    assert len(fixed_text) <= budget, (
        f"Expected action_title ≤{budget} chars in plan on disk after autofix, "
        f"got {len(fixed_text)}: {fixed_text!r}"
    )

    # ── PPTX content assertion (critical-bug regression) ─────────────────────
    # Open the produced PPTX and collect all text frame content from slide 0.
    # The compile loop must have used the FIXED plan (slides_spec re-captured
    # after apply_fixes), so the PPTX must NOT contain the original long_text
    # and the action_title frame must reflect the shortened content.
    assert out_pptx.is_file(), f"Expected PPTX output at {out_pptx}"
    prs = Presentation(str(out_pptx))
    assert len(prs.slides) == 1, f"Expected 1 slide in PPTX, got {len(prs.slides)}"
    slide_0 = prs.slides[0]

    all_text_in_pptx = []
    for shape in slide_0.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs)
                if line.strip():
                    all_text_in_pptx.append(line)

    full_pptx_text = " ".join(all_text_in_pptx)

    # The original (unfixed) long_text must NOT appear verbatim in the PPTX.
    assert long_text not in full_pptx_text, (
        f"PPTX was compiled from pre-fix content (the critical bug). "
        f"The original long_text ({len(long_text)} chars) appears verbatim in "
        f"the PPTX output, meaning slides_spec was not re-captured after autofix.\n"
        f"PPTX text: {full_pptx_text!r}\nOriginal: {long_text!r}"
    )

    # The fixed_text (shortened to ≤budget) must appear in the PPTX.
    assert fixed_text in full_pptx_text, (
        f"PPTX does not contain the fixed action_title text.\n"
        f"Expected to find: {fixed_text!r}\n"
        f"PPTX text found: {full_pptx_text!r}"
    )


# ---------------------------------------------------------------------------
# 5. Round-trip: verify-static --json → apply-fixes (engine severity values)
# ---------------------------------------------------------------------------

def test_cli_round_trip_verify_static_to_apply_fixes(tmp_path: Path):
    """True round-trip: verify-static --json output feeds directly into apply-fixes.

    verify-static --json emits ENGINE severity values ("error", "warning",
    "info") because it calls validate() which maps:
      legacy FATAL → engine "error"
      legacy WARN  → engine "warning"
      legacy INFO  → engine "info"

    apply-fixes must accept these engine values (not drop them as "malformed
    defect entry").  This test is the regression guard for the bug where
    apply-fixes only accepted legacy vocabulary ("fatal", "warn", "info").

    Steps:
    1. Build a plan with a genuinely overflowing slot (action_title 180 chars,
       budget 168) and write it to disk.
    2. Run `deck verify-static --json` — expect exit 1 and ENGINE severity
       values ("error") in the output JSON.
    3. Pipe that JSON file directly into `deck apply-fixes --defects`.
    4. Assert apply-fixes exits 0 (patch applied, NOT exit 1 "no defects").
    5. Assert the fixed plan has the slot shortened — confirming the defect
       was actually processed, not silently dropped as malformed.
    """
    import os

    # Use heuristic metrics so the budget is deterministic (machine-independent).
    env = os.environ.copy()
    env["PYTHONIOENCODING"] = "utf-8"
    env["PYTHONUTF8"] = "1"
    env["FEINSCHMIEDE_NO_REAL_METRICS"] = "1"

    # action_title budget for executive-summary = 168 chars (heuristic path,
    # 2 bleed-model lines). 180 chars: above budget (triggers SLOT_OVERFLOW),
    # below swap threshold (168 * 1.20 = 201.6), so shorten_slot fires.
    budget = 168
    long_text = "Action required: revenue declined three quarters in a row this year."
    assert len(long_text) == 68
    long_text = long_text + " " + "A" * (180 - len(long_text) - 1)
    assert len(long_text) == 180
    assert budget < len(long_text) < budget * 1.20, (
        f"Fixture must be above budget ({budget}) and below swap threshold "
        f"({budget * 1.20}); got {len(long_text)}"
    )

    plan = {
        "brand": "feinschliff",
        "out": str(tmp_path / "deck.pptx"),
        "slides": [
            {
                "layout": "layouts/executive-summary.slide.dsl",
                "content": {
                    "action_title": long_text,
                    "footer_left": "Corp",
                    "footer_right": "2026",
                },
            }
        ],
    }
    plan_path = tmp_path / "plan.yaml"
    plan_path.write_text(
        yaml.safe_dump(plan, sort_keys=False, allow_unicode=True), encoding="utf-8"
    )

    defects_path = tmp_path / "defects.json"
    fixed_plan_path = tmp_path / "fixed_plan.yaml"

    # Step 2: run verify-static --json; expect exit 1 (defects found)
    vs_result = subprocess.run(
        [
            sys.executable, "-m", "feinschliff.cli",
            "deck", "verify-static", str(plan_path), "--json",
        ],
        capture_output=True,
        text=True,
        encoding="utf-8",
        env=env,
        cwd=str(FEINSCHLIFF),
    )
    assert vs_result.returncode == 1, (
        f"Expected exit 1 from verify-static (defects present), "
        f"got {vs_result.returncode}.\nstdout: {vs_result.stdout}\nstderr: {vs_result.stderr}"
    )

    # Parse the JSON output — must be a non-empty list of defects.
    defects_json = json.loads(vs_result.stdout)
    assert isinstance(defects_json, list) and len(defects_json) >= 1, (
        f"Expected a non-empty list of defects in JSON output; got: {defects_json!r}"
    )

    # At least one defect must carry an ENGINE severity value ("error" or "warning").
    # This is the key assertion: if it carries "fatal"/"warn" the bug isn't triggered.
    engine_severities = {"error", "warning", "info"}
    found_engine_sev = any(d.get("severity") in engine_severities for d in defects_json)
    assert found_engine_sev, (
        f"Expected at least one defect with engine severity (error/warning/info) "
        f"in verify-static --json output; got severities: "
        f"{[d.get('severity') for d in defects_json]}"
    )

    # Write the verify-static JSON output directly to disk (no transformation).
    defects_path.write_text(vs_result.stdout, encoding="utf-8")

    # Step 3+4: run apply-fixes with the raw verify-static output.
    af_result = subprocess.run(
        [
            sys.executable, "-m", "feinschliff.cli",
            "deck", "apply-fixes", str(plan_path),
            "--defects", str(defects_path),
            "-o", str(fixed_plan_path),
        ],
        capture_output=True,
        text=True,
        encoding="utf-8",
        env=env,
        cwd=str(FEINSCHLIFF),
    )

    # Step 4: apply-fixes must exit 0 (patch applied).
    # Before the fix, it exits 1 ("no defects to process") because every defect
    # was dropped as "malformed" due to the legacy-only severity parser.
    assert af_result.returncode == 0, (
        f"Expected exit 0 from apply-fixes (patch applied). "
        f"If you see exit 1 / 'no defects to process', the engine-severity "
        f"round-trip bug is present.\n"
        f"stdout: {af_result.stdout}\nstderr: {af_result.stderr}"
    )

    # Step 5: the fixed plan must exist and have the slot shortened.
    assert fixed_plan_path.is_file(), "Expected fixed plan YAML to be written."
    fixed = yaml.safe_load(fixed_plan_path.read_text(encoding="utf-8"))
    fixed_text = fixed["slides"][0]["content"]["action_title"]
    assert len(fixed_text) <= budget, (
        f"Expected action_title shortened to ≤{budget} chars after round-trip fix; "
        f"got {len(fixed_text)}: {fixed_text!r}"
    )
