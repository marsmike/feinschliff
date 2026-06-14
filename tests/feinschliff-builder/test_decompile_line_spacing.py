"""Paragraph line spacing must round-trip by provenance: an explicit
`<a:lnSpc><a:spcPct>` becomes `linespacing:<multiplier>`, while a paragraph
with NO explicit spacing becomes `linespacing:native` (the emitter then
writes no lnSpc so the renderer's single spacing applies, like PowerPoint).

Regression for the Bosch header shift: the emitter's 1.2 default leading on
source-default paragraphs pushed every decompiled headline ~25 design-px
down the slide on all 99 layouts.
"""
from __future__ import annotations

from lxml import etree

from feinschliff_builder.decompile.pptx_svg_decompile import (
    CanvasMap,
    _walk,
    emit_dsl,
)

_EMU_W, _EMU_H = 12192000, 6858000

_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _text_sp(lnspc_val: int | None) -> str:
    lnspc = (
        f"<a:pPr><a:lnSpc><a:spcPct val=\"{lnspc_val}\"/></a:lnSpc></a:pPr>"
        if lnspc_val is not None else ""
    )
    return f"""
      <p:sp xmlns:p="{_P}" xmlns:a="{_A}">
        <p:nvSpPr><p:cNvPr id="2" name="Title 1"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
        <p:spPr>
          <a:xfrm><a:off x="205740" y="257175"/><a:ext cx="10556240" cy="388620"/></a:xfrm>
          <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
        </p:spPr>
        <p:txBody><a:bodyPr/><a:p>{lnspc}<a:r><a:rPr lang="en-US" sz="2800"/>
        <a:t>Headline</a:t></a:r></a:p></p:txBody>
      </p:sp>
    """


def _decompile_text(lnspc_val: int | None):
    tree = etree.fromstring(
        f'<p:spTree xmlns:p="{_P}" xmlns:a="{_A}">{_text_sp(lnspc_val)}</p:spTree>'.encode()
    )
    shapes: list = []
    cmap = CanvasMap(_EMU_W, _EMU_H, 1920, 1080)
    _walk(tree, (0, 0), shapes, None, cmap, {}, {})
    texts = [s for s in shapes if s.kind == "text"]
    assert len(texts) == 1
    return shapes, texts[0], cmap


def test_explicit_spcpct_is_captured_and_emitted():
    shapes, t, cmap = _decompile_text(90000)
    assert t.line_spacing == 0.9
    dsl = emit_dsl(shapes, cmap, "t")
    assert "linespacing:0.9" in dsl


def test_absent_lnspc_round_trips_as_native():
    shapes, t, cmap = _decompile_text(None)
    assert t.line_spacing is None
    dsl = emit_dsl(shapes, cmap, "t")
    assert "linespacing:native" in dsl


def test_master_txstyles_line_spacing_inherits():
    """A placeholder paragraph with NO own lnSpc must inherit the master
    txStyles value (Bosch sets bodyStyle 107% / titleStyle 89% ONLY on the
    master) — not fall through to `linespacing:native`."""
    from lxml import etree as _et
    from pptx import Presentation
    from feinschliff_builder.decompile.pptx_svg_decompile import (
        _layout_placeholder_line_spacing,
    )

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    master_el = slide.slide_layout.slide_master.element
    a = _A
    for style, val in (("titleStyle", 89000), ("bodyStyle", 107000)):
        st = master_el.find(f"{{{_P}}}txStyles/{{{_P}}}{style}")
        lvl1 = st.find(f"{{{a}}}lvl1pPr")
        if lvl1 is None:
            lvl1 = _et.SubElement(st, f"{{{a}}}lvl1pPr")
            st.insert(0, lvl1)
        lnspc = _et.fromstring(
            f'<a:lnSpc xmlns:a="{a}"><a:spcPct val="{val}"/></a:lnSpc>')
        lvl1.insert(0, lnspc)

    assert _layout_placeholder_line_spacing(slide, "title", None) == 0.89
    assert _layout_placeholder_line_spacing(slide, "body", "1") == 1.07
