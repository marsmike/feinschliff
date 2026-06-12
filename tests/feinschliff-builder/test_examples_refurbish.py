"""Verify the refurbish before/after demo produces expected artifacts."""
from __future__ import annotations

import os
import subprocess
from pathlib import Path

import pytest

REPO = Path(__file__).resolve().parents[2] / "feinschliff-builder"
BEFORE = REPO / "examples" / "refurbish" / "before" / "sewp-vi-overview.pptx"


def test_before_fixture_is_three_slides():
    if not BEFORE.exists():
        pytest.skip("sewp-vi-overview.pptx not yet built")
    from pptx import Presentation
    pres = Presentation(str(BEFORE))
    assert len(pres.slides) == 3, f"expected 3 slides, got {len(pres.slides)}"


def test_refurbish_produces_polished_pptx(tmp_path):
    if not BEFORE.exists():
        pytest.skip("sewp-vi-overview.pptx not yet built")
    try:
        import cairosvg
        cairosvg.svg2png(bytestring=b'<svg xmlns="http://www.w3.org/2000/svg" width="1" height="1"/>')
    except (ImportError, OSError):
        pytest.skip("render backend (cairo) unavailable")
    out = tmp_path / "polished.pptx"
    env = {**os.environ, "DYLD_FALLBACK_LIBRARY_PATH": "/opt/homebrew/opt/cairo/lib"}
    result = subprocess.run(
        ["uv", "run", "feinschliff", "deck", "polish",
         str(BEFORE), "-o", str(out),
         "--refurbish-all", "--brand", "feinschliff"],
        capture_output=True, text=True, cwd=REPO, env=env,
    )
    assert result.returncode == 0, f"polish failed: {result.stderr[-500:]}"
    assert out.exists()
    from pptx import Presentation
    pres = Presentation(str(out))
    assert len(pres.slides) >= 2


def test_refurbish_report_committed():
    # The report lives next to the output pptx in the after/ dir.
    candidates = [
        REPO / "examples" / "refurbish" / "refurbish_report.md",
        REPO / "examples" / "refurbish" / "after" / "refurbish_report.md",
    ]
    found = [c for c in candidates if c.exists()]
    if not found:
        pytest.skip("refurbish_report.md not yet committed")
    text = found[0].read_text()
    # Should mention at least one of the two diagram slides.
    assert ("slide 2" in text.lower() or "slide 3" in text.lower()), \
        "report should list at least one detected diagram slide"
