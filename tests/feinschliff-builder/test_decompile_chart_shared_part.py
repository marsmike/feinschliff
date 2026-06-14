"""Chart-frame part capture must record a part SHARED by two parents.

The chart branch of `_emit_graphic_frame` walks every `<c:chart r:id>` in the
frame and carries the part-graph rooted at each chart part. When two charts
share a part (here: chart2's rels also point at chart1's embedded workbook),
the part is materialised once — but the second parent's src_rid→part mapping
must still be folded into a "ref" entry, or chart2's reference dangles at
splice time. The group/diagram capture paths already did this; the chart
branch silently dropped the second reference.
"""
from __future__ import annotations

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Emu

from feinschliff_builder.decompile.pptx_svg_decompile import (
    CanvasMap,
    _emit_graphic_frame,
)

_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_RT_PACKAGE = f"{_R}/package"


def _add_chart(slide):
    data = CategoryChartData()
    data.categories = ["a", "b"]
    data.add_series("s", (1, 2))
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Emu(0), Emu(0), Emu(914400), Emu(914400), data,
    )
    return gframe.chart.part


def test_shared_xlsx_across_two_charts_gets_ref_entry():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    c1 = _add_chart(slide)
    c2 = _add_chart(slide)
    # chart1's embedded workbook part.
    [xlsx] = [r.target_part for r in c1.rels.values()
              if r.reltype == _RT_PACKAGE]
    # chart2 ALSO references chart1's workbook (the shared part).
    shared_rid = c2.relate_to(xlsx, _RT_PACKAGE)
    # slide-rel rIds for the two chart parts.
    rid_of = {r.target_part.partname: r.rId for r in slide.part.rels.values()
              if r.reltype.endswith("/chart")}
    frame = etree.fromstring(f"""
      <p:graphicFrame xmlns:p="{_P}" xmlns:a="{_A}" xmlns:c="{_C}" xmlns:r="{_R}">
        <p:nvGraphicFramePr><p:cNvPr id="9" name="Charts"/><p:cNvGraphicFramePr/><p:nvPr/></p:nvGraphicFramePr>
        <p:xfrm><a:off x="914400" y="914400"/><a:ext cx="6096000" cy="3048000"/></p:xfrm>
        <a:graphic><a:graphicData uri="{_C}">
          <c:chart r:id="{rid_of[c1.partname]}"/>
          <c:chart r:id="{rid_of[c2.partname]}"/>
        </a:graphicData></a:graphic>
      </p:graphicFrame>
    """.encode())

    shapes: list = []
    cmap = CanvasMap(12192000, 6858000, 1920, 1080)
    _emit_graphic_frame(frame, (0, 0), shapes, slide, cmap, {}, {})
    [shape] = shapes
    assert shape.native_xml is not None and shape.native_parts
    # The shared workbook is materialised exactly once …
    full = [e for e in shape.native_parts if "ref" not in e]
    assert len([e for e in full if e["partname"] == str(xlsx.partname)]) == 1
    # … and chart2's second reference is folded into a ref entry so the
    # emitter wires its relationship + rewrites its r:id.
    refs = [e for e in shape.native_parts if e.get("ref") == str(xlsx.partname)
            and e["parent"] == str(c2.partname)]
    assert refs and refs[0]["src_rid"] == shared_rid, (
        "second parent's reference to the shared part was dropped"
    )
