"""Integration test: `feinschliff deck build` threads top-level `notes:`
from the plan.yaml into the PPTX speaker-notes pane."""
from __future__ import annotations

import subprocess
import sys
from pathlib import Path

from pptx import Presentation

REPO = Path(__file__).resolve().parents[2]
FEINSCHLIFF = REPO / "feinschliff"
LAYOUT = FEINSCHLIFF / "layouts" / "action-title.slide.dsl"


def test_deck_build_threads_speaker_notes(tmp_path):
    """A plan with `notes:` on each slide produces a deck whose PPTX
    notes panes carry the authored strings. Slides without `notes:`
    don't materialise a notes slide.

    The hook slide (index 0) carries the deck-level storyline; the
    second slide carries listener bullets; the third slide is left
    bare to exercise the no-notes path."""
    storyline = (
        "Storyline: pain → demo → results.\n"
        "• Open with the time-collapse stat.\n"
        "• Hand off to the live demo at 0:45."
    )
    bullets = "• Polish loop now 15 min.\n• Slot budget enforces brevity."
    plan_yaml = f"""
brand: feinschliff
slides:
  - layout: {LAYOUT}
    notes: |
{chr(10).join("      " + line for line in storyline.splitlines())}
    content:
      title: "Polish time collapsed"
  - layout: {LAYOUT}
    notes: |
{chr(10).join("      " + line for line in bullets.splitlines())}
    content:
      title: "Q3 revenue rose 12 percent"
  - layout: {LAYOUT}
    content:
      title: "Customer churn dropped 8 points"
"""
    plan_file = tmp_path / "plan.yaml"
    plan_file.write_text(plan_yaml, encoding="utf-8")
    out_file = tmp_path / "out.pptx"
    import os
    env = os.environ.copy()
    env["PYTHONIOENCODING"] = "utf-8"
    env["PYTHONUTF8"] = "1"
    result = subprocess.run(
        [
            sys.executable, "-m", "feinschliff.cli", "deck", "build",
            str(plan_file),
            "-o", str(out_file),
        ],
        capture_output=True, text=True, encoding="utf-8", env=env, cwd=FEINSCHLIFF,
    )
    assert result.returncode == 0, result.stdout + result.stderr
    assert out_file.is_file()

    prs = Presentation(str(out_file))
    s0, s1, s2 = prs.slides

    assert s0.has_notes_slide
    assert "time-collapse stat" in s0.notes_slide.notes_text_frame.text

    assert s1.has_notes_slide
    assert "Slot budget" in s1.notes_slide.notes_text_frame.text

    assert not s2.has_notes_slide
