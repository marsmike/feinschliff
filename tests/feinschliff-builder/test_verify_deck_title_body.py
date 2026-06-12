"""Unit tests for lib/verify/deck/title_body — slide title + body extraction."""
from __future__ import annotations

from pathlib import Path

import pytest

from feinschliff_builder.verify.deck.title_body import extract_slide_title_and_body


def test_extract_title_and_body_from_pptx(tmp_path: Path):
    """Build a 2-slide .pptx, extract title + body from slide 1."""
    from pptx import Presentation
    prs = Presentation()
    title_layout = prs.slide_layouts[1]  # title + content
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = "Revenue dropped 12% in Q3"
    # Add a content placeholder body
    body = slide.placeholders[1]
    body.text = "Enterprise churn drove most of the loss."

    blank = prs.slide_layouts[5]
    prs.slides.add_slide(blank)

    out = tmp_path / "deck.pptx"
    prs.save(str(out))

    title, body_text = extract_slide_title_and_body(out, slide_index=1)
    assert title == "Revenue dropped 12% in Q3"
    assert "Enterprise churn" in body_text


def test_extract_returns_empty_for_no_title_shape(tmp_path: Path):
    """Slide with no title placeholder returns empty title."""
    from pptx import Presentation
    prs = Presentation()
    blank = prs.slide_layouts[6]  # truly blank
    prs.slides.add_slide(blank)
    out = tmp_path / "deck.pptx"
    prs.save(str(out))

    title, body = extract_slide_title_and_body(out, slide_index=1)
    assert title == ""
    assert body == ""


def test_extract_slide_index_out_of_range(tmp_path: Path):
    """slide_index beyond deck length raises IndexError."""
    from pptx import Presentation
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])
    out = tmp_path / "deck.pptx"
    prs.save(str(out))

    with pytest.raises(IndexError, match="out of range"):
        extract_slide_title_and_body(out, slide_index=2)
    with pytest.raises(IndexError, match="out of range"):
        extract_slide_title_and_body(out, slide_index=0)


def test_extract_multi_paragraph_body(tmp_path: Path):
    """Multiple body shapes get joined with double newline."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Title"
    # Two text boxes
    tb1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.5))
    tb1.text_frame.text = "First insight"
    tb2 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(0.5))
    tb2.text_frame.text = "Second insight"
    out = tmp_path / "deck.pptx"
    prs.save(str(out))

    title, body = extract_slide_title_and_body(out, slide_index=1)
    assert title == "Title"
    assert "First insight" in body
    assert "Second insight" in body
    assert body.count("\n\n") >= 1  # separator between bodies
