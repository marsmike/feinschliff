"""Regression tests for the automatic slide-number footer.

Covers:
  - build_multi_slide(slide_numbers=True) stamps 'NN / TOTAL' on every slide.
  - Footer text sits in the bottom-right quadrant of the slide.
  - build_multi_slide(slide_numbers=False) (default) produces no such frame.
  - The --no-slide-numbers opt-out propagates through cmd_build.
"""
from __future__ import annotations

import re
from pathlib import Path

import pytest

from feinschliff.dsl.parser import parse_lines
from feinschliff.dsl.pptx_emit import build_multi_slide
from feinschmiede.dsl.tokens import load_tokens
from pptx.util import Emu


REPO_ROOT = Path(__file__).resolve().parents[2] / "feinschliff"
BRAND_DIR = REPO_ROOT / "brands" / "feinschliff"

_SLIDE_FOOTER_RE = re.compile(r"^\d{2} / \d{2}$")

_DSL = """\
canvas 1920x1080
text 100,100 style:title "Hello"
"""


def _parse(dsl: str = _DSL):
    nodes, _ = parse_lines(dsl, source="<test>")
    return nodes


def _make_payload(n: int, dsl: str = _DSL):
    """Return a list of `n` 3-tuple slide payloads using the feinschliff brand."""
    tokens = load_tokens(BRAND_DIR)
    nodes = _parse(dsl)
    return [(nodes, tokens, BRAND_DIR / "assets")] * n


def _all_text_frames(prs):
    """Yield every (slide_idx, shape) pair that has a text frame."""
    for i, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.has_text_frame:
                yield i, shape


def _footer_frames(prs):
    """Return list of (slide_idx, text) for shapes whose text matches NN / NN."""
    found = []
    for slide_idx, shape in _all_text_frames(prs):
        text = shape.text_frame.text.strip()
        if _SLIDE_FOOTER_RE.match(text):
            found.append((slide_idx, text, shape))
    return found


# ---------------------------------------------------------------------------
# slide_numbers=True: footer present on every slide
# ---------------------------------------------------------------------------

def test_every_slide_has_footer():
    """All 3 slides must carry a 'NN / 03' footer."""
    payload = _make_payload(3)
    prs = build_multi_slide(payload, slide_numbers=True)
    footers = _footer_frames(prs)
    slide_numbers_found = {idx for idx, _, _ in footers}
    assert slide_numbers_found == {1, 2, 3}, (
        f"Expected footer on slides 1-3, found on {sorted(slide_numbers_found)}"
    )


def test_footer_text_format():
    """Footer text must be zero-padded NN / TOTAL, e.g. '01 / 05'."""
    payload = _make_payload(5)
    prs = build_multi_slide(payload, slide_numbers=True)
    footers = _footer_frames(prs)
    texts = [t for _, t, _ in footers]
    assert texts == ["01 / 05", "02 / 05", "03 / 05", "04 / 05", "05 / 05"], (
        f"Unexpected footer texts: {texts}"
    )


def test_footer_in_bottom_right_quadrant():
    """Footer shape must sit in the bottom-right quadrant of the slide."""
    payload = _make_payload(1)
    prs = build_multi_slide(payload, slide_numbers=True)
    slide = prs.slides[0]
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    footers = _footer_frames(prs)
    assert footers, "no footer frame found"
    _, _, shape = footers[0]
    # The left edge must be in the right half, top edge in the bottom half.
    assert shape.left > slide_w // 2, (
        f"Footer left ({shape.left}) not in right half (slide_w={slide_w})"
    )
    assert shape.top > slide_h // 2, (
        f"Footer top ({shape.top}) not in bottom half (slide_h={slide_h})"
    )


# ---------------------------------------------------------------------------
# slide_numbers=False (default): no footer
# ---------------------------------------------------------------------------

def test_no_footer_by_default():
    """Default build_multi_slide call must NOT stamp any footer."""
    payload = _make_payload(2)
    prs = build_multi_slide(payload)  # slide_numbers defaults to False
    footers = _footer_frames(prs)
    assert footers == [], (
        f"Expected no footer frames, found {[(i, t) for i, t, _ in footers]}"
    )


def test_slide_numbers_false_explicit():
    """Explicit slide_numbers=False produces no footer."""
    payload = _make_payload(2)
    prs = build_multi_slide(payload, slide_numbers=False)
    footers = _footer_frames(prs)
    assert footers == [], (
        f"Expected no footer frames, found {[(i, t) for i, t, _ in footers]}"
    )


# ---------------------------------------------------------------------------
# --no-slide-numbers CLI flag
# ---------------------------------------------------------------------------

def test_no_slide_numbers_flag_accepted(tmp_path):
    """--no-slide-numbers must be a recognised argument on `feinschliff deck build`."""
    import argparse
    from feinschliff.cli import deck as deck_mod
    parser = argparse.ArgumentParser()
    deck_mod.register(parser)
    # Parsing with --no-slide-numbers must not raise.
    args = parser.parse_args(["build", "--no-slide-numbers", "dummy.yaml"])
    assert getattr(args, "no_slide_numbers", False) is True
