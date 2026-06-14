"""Regression guards for the combined gem + wordmark logo asset.

Two concerns:
1. logo.png / logo-light.png exist at the expected path in the source tree
   (and therefore ship in the installed wheel via the brands/feinschliff/assets/*
   package-data glob in feinschliff/pyproject.toml).

2. A slide built from the feinschliff header compound contains exactly one
   Picture shape sourced from logo.png — not a gem.png Picture plus a Text Box.

These tests do NOT do pixel comparison — they're structural guards only.
"""
from __future__ import annotations

from pathlib import Path

import pytest
from pptx.enum.shapes import MSO_SHAPE_TYPE

from feinschliff.dsl.expander import (
    expand_compounds,
    interpolate_nodes,
    load_compounds_for_brand,
)
from feinschliff.dsl.parser import parse_lines
from feinschliff.dsl.pptx_emit import build_presentation
from feinschmiede.dsl.tokens import load_tokens

REPO_ROOT = Path(__file__).resolve().parents[2] / "feinschliff"
BRANDS_DIR = REPO_ROOT / "brands"
FEINSCHLIFF_ASSETS = BRANDS_DIR / "feinschliff" / "assets"
STD_COMPOUNDS = REPO_ROOT / "compounds"

# ---------------------------------------------------------------------------
# 1. Asset-existence regression guard (package-data glob coverage)
# ---------------------------------------------------------------------------


@pytest.mark.parametrize("filename", ["logo.png", "logo-light.png"])
def test_logo_png_exists_in_source_tree(filename: str) -> None:
    """logo.png and logo-light.png must be present in the feinschliff brand
    assets directory.  They are committed artefacts (rasterised from logo.svg
    by gen_logo.py) and ship in the installed wheel via the
    brands/feinschliff/assets/* package-data glob.
    """
    asset = FEINSCHLIFF_ASSETS / filename
    assert asset.is_file(), (
        f"{filename} not found at {asset}. "
        "Run: uv run --directory feinschliff python brands/feinschliff/assets/gen_logo.py"
    )
    assert asset.stat().st_size > 1000, f"{filename} is suspiciously small — regenerate it."


# ---------------------------------------------------------------------------
# 2. Header compound emits a single logo.png Picture, not gem + text
# ---------------------------------------------------------------------------


def _build_header_slide() -> "Presentation":
    """Parse a minimal DSL snippet that calls the feinschliff header compound
    and return the resulting python-pptx Presentation object."""
    brand_dir = BRANDS_DIR / "feinschliff"
    tokens = load_tokens(brand_dir, brands_dir=BRANDS_DIR)
    compounds = load_compounds_for_brand(
        brand_dir, std_dir=STD_COMPOUNDS, brands_dir=BRANDS_DIR
    )

    # Minimal DSL: canvas directive + single header compound call.
    dsl = 'canvas 1920x1080\nheader pgmeta:"Q2 2026 · 1 / 1"\n'
    nodes, extra_compounds = parse_lines(dsl, source="<test>")
    for cd in extra_compounds:
        compounds[cd.name] = cd

    interp = interpolate_nodes(nodes, {})
    primitives, _diagnostics = expand_compounds(interp, compounds)

    return build_presentation(
        primitives,
        tokens,
        asset_root=brand_dir / "assets",
    )


def test_header_compound_renders_logo_picture() -> None:
    """The header compound must produce a Picture shape (i.e. logo.png was
    resolved and embedded), not fall back to a text-only or rect-only render."""
    prs = _build_header_slide()
    slide = prs.slides[0]
    kinds = [s.shape_type for s in slide.shapes]
    assert MSO_SHAPE_TYPE.PICTURE in kinds, (
        f"Expected a Picture shape from the header compound (logo.png), "
        f"but got only: {kinds}. "
        "Check that logo.png exists and the header.dsl picture path is correct."
    )


def test_header_compound_no_gem_png_or_wordmark_textbox() -> None:
    """After the logo consolidation, the header compound must NOT produce a
    separate wordmark text box alongside the gem picture.  The old two-shape
    composition (gem.png Picture + 'FEINSCHLIFF.' Text Box) is replaced by a
    single logo.png Picture.

    We detect the old pattern by checking that the slide has at most one
    Picture shape (the combined logo) and that no Text Box contains the
    bare wordmark string 'FEINSCHLIFF.' without any template slots.
    """
    prs = _build_header_slide()
    slide = prs.slides[0]

    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(pictures) == 1, (
        f"Expected exactly 1 Picture shape (logo.png) in the header compound, "
        f"got {len(pictures)}.  The old gem.png + logo.png double-picture pattern "
        "would trigger this — check header.dsl."
    )

    # The wordmark should now be baked into the PNG, not a separate text run.
    wordmark_textboxes = [
        s for s in slide.shapes
        if hasattr(s, "text") and "FEINSCHLIFF." in s.text
    ]
    assert len(wordmark_textboxes) == 0, (
        f"Found {len(wordmark_textboxes)} shape(s) with bare 'FEINSCHLIFF.' text — "
        "the wordmark should be baked into logo.png, not emitted as a separate text box. "
        "Check that the style:wordmark text line was removed from header.dsl."
    )
