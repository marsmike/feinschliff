"""Real-world end-to-end deck-build test.

Exercises the full office build pipeline (compile_slide → build_presentation)
for a known-good layout+content+brand triple, then validates the produced .pptx
with python-pptx.  No subprocess round-trip — the test calls the same in-process
entrypoints the CLI uses so it stays under ~3 s and remains deterministic.

Case 1  — single-slide: quote layout + tests/fixtures/layouts/quote.yaml + feinschliff brand.
Case 2  — multi-slide deck plan: skipped because no deck-plan fixture exists under
          tests/fixtures/.  (The fixtures/layouts/*.yaml files are per-slide content
          files, not full deck plans accepted by `deck build`.)
"""
from __future__ import annotations

from pathlib import Path

import pytest
import yaml
from pptx import Presentation

from feinschmiede.brand_discovery import find_brand
from feinschliff.dsl.pptx_emit import build_presentation
from feinschliff.pipeline import compile_slide


PLUGIN_ROOT = Path(__file__).resolve().parents[2] / "feinschliff"
QUOTE_LAYOUT = PLUGIN_ROOT / "layouts" / "quote.slide.dsl"
QUOTE_CONTENT = Path(__file__).resolve().parent / "fixtures" / "layouts" / "quote.yaml"
BRAND_ID = "feinschliff"
BUNDLED_ASSETS = PLUGIN_ROOT / "feinschliff" / "assets"


def _build_single_slide(tmp_path: Path) -> Path:
    """Run the full compile → emit pipeline, return the written .pptx path."""
    brand = find_brand(BRAND_ID)
    brand_dir = brand.root

    ctx = yaml.safe_load(QUOTE_CONTENT.read_text()) or {}

    result = compile_slide(
        layout_path=QUOTE_LAYOUT,
        ctx=ctx,
        brand_dir=brand_dir,
        slide_index=1,
        diagrams_out_dir=tmp_path / "diagrams",
    )

    asset_root = brand_dir / "assets"
    asset_root_fallback = BUNDLED_ASSETS if BUNDLED_ASSETS.is_dir() else None

    prs = build_presentation(
        result.primitives,
        result.tokens,
        asset_root=asset_root,
        asset_root_fallback=asset_root_fallback,
    )

    out = tmp_path / "quote.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


# ── Case 1: single-slide quote ────────────────────────────────────────────────

def test_single_slide_produces_valid_pptx(tmp_path):
    assert QUOTE_LAYOUT.is_file(), f"layout not found: {QUOTE_LAYOUT}"
    assert QUOTE_CONTENT.is_file(), f"content fixture not found: {QUOTE_CONTENT}"

    out = _build_single_slide(tmp_path)

    assert out.is_file(), "build did not produce a .pptx file"
    assert out.stat().st_size > 1_000, "produced .pptx is suspiciously small"

    prs = Presentation(str(out))
    assert len(prs.slides) >= 1, "presentation has no slides"

    shapes_on_first = list(prs.slides[0].shapes)
    assert shapes_on_first, "first slide has no shapes — may be empty or corrupt"


# ── Case 2: multi-slide deck plan ─────────────────────────────────────────────

_DECK_PLAN_FIXTURE = next(
    (
        p for p in (Path(__file__).resolve().parent / "fixtures").rglob("*.yaml")
        if (yaml.safe_load(p.read_text()) or {}).get("slides")
    ),
    None,
)

@pytest.mark.skipif(
    _DECK_PLAN_FIXTURE is None,
    reason=(
        "no deck-plan fixture found under tests/fixtures/ "
        "(fixtures/layouts/*.yaml are per-slide content files, not deck plans)"
    ),
)
def test_multi_slide_plan_produces_valid_pptx(tmp_path):
    from feinschliff.dsl.pptx_emit import build_multi_slide

    plan = yaml.safe_load(_DECK_PLAN_FIXTURE.read_text()) or {}  # type: ignore[union-attr]
    default_brand_id = plan.get("brand", BRAND_ID)
    plan_dir = _DECK_PLAN_FIXTURE.parent  # type: ignore[union-attr]

    slides_payload = []
    for i, spec in enumerate(plan.get("slides", [])):
        layout_path = (plan_dir / spec["layout"]).resolve()
        brand_id = spec.get("brand", default_brand_id)
        brand_dir = find_brand(brand_id).root

        ctx: dict = spec.get("content") or {}
        if not ctx and "content_file" in spec:
            cf = (plan_dir / spec["content_file"]).resolve()
            ctx = yaml.safe_load(cf.read_text()) or {}

        result = compile_slide(
            layout_path=layout_path,
            ctx=ctx,
            brand_dir=brand_dir,
            slide_index=i + 1,
            diagrams_out_dir=tmp_path / "diagrams",
        )
        slides_payload.append((result.primitives, result.tokens, brand_dir / "assets"))

    asset_fallback = BUNDLED_ASSETS if BUNDLED_ASSETS.is_dir() else None
    out = tmp_path / "deck.pptx"
    prs = build_multi_slide(slides_payload, asset_root_fallback=asset_fallback)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))

    assert out.is_file()
    loaded = Presentation(str(out))
    assert len(loaded.slides) >= 1
    assert list(loaded.slides[0].shapes), "first slide has no shapes"
