"""Tests for Fix A (image_query slot routing) and Fix B (bundled placeholder).

Fix A: Image-bearing layouts now carry ``query:{{ image_query }}`` on every
``picture path:{{ image }}`` line so that a plan.yaml ``image_query: "..."``
field reaches the provider even when the ``image:`` slot is absent.

Fix B: ``assets/illustrations/placeholder.jpg`` is now listed under
``[tool.setuptools.package-data]`` so the wheel includes it and
``_placeholder_image_path`` finds it in an installed package.

Test 1 — image_query routes to provider
  Build a one-slide DSL with ``path:`` set to a non-existent path and
  ``query:"test scene"`` set.  Stub provider returns a real tiny PNG.
  Assert: PICTURE shape in output, provider called with "test scene".

Test 2 — no provider + bare query: raises DSLError (intentional fail-loud)
  Build a one-slide DSL with ``query:"test scene"`` and no ``path:``.
  Without a provider, the emitter must raise DSLError immediately.

Test 3 — _placeholder_image_path finds the gem in the installed package
  Unit test for Fix B: import the module, call the helper with a minimal
  context, and verify it returns a real existing Path.
"""
from __future__ import annotations

import io
from pathlib import Path
from unittest.mock import MagicMock

import pytest
from PIL import Image

from feinschliff.dsl.parser import parse_lines
from feinschliff.dsl.pptx_emit import (
    DSLError,
    EmitContext,
    build_presentation,
    _placeholder_image_path,
)
from feinschliff.io import image_provider
from feinschliff.io.image_provider import ImageHit, ImageProvider
from test_emitter_restraint import _minimal_tokens


# ---------------------------------------------------------------------------
# Registry isolation — prevents provider leaks between tests
# ---------------------------------------------------------------------------

@pytest.fixture(autouse=True)
def _isolate_registry(monkeypatch):
    monkeypatch.setattr(image_provider, "_REGISTRY", {})
    monkeypatch.setattr(image_provider, "_DISCOVERED", False)
    yield


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture
def tiny_jpg(tmp_path) -> Path:
    """Minimal 4×4 JPEG returned by the stub provider."""
    p = tmp_path / "stub.jpg"
    buf = io.BytesIO()
    Image.new("RGB", (4, 4), color=(100, 149, 237)).save(buf, "JPEG")
    p.write_bytes(buf.getvalue())
    return p


def _stub_provider(hit_path: Path) -> MagicMock:
    hit = ImageHit(
        url=f"file://{hit_path}",
        license="Test License",
        attribution="Stub Author",
        width=4,
        height=4,
        mime="image/jpeg",
    )
    prov = MagicMock(spec=ImageProvider)
    prov.name = "stub"
    prov.search.return_value = [hit]
    return prov


# ---------------------------------------------------------------------------
# Test 1 — image_query routes to provider when image: is absent/unresolved
# ---------------------------------------------------------------------------

def test_image_query_routes_to_provider(tmp_path, tiny_jpg):
    """Fix A regression guard.

    A layout line ``picture … path:{{ image }} query:{{ image_query }}``
    must pass the ``image_query`` value to the provider when ``image:`` is
    empty (or resolves to a non-existent path), i.e. the provider is called
    with "test scene", not skipped.
    """
    provider = _stub_provider(tiny_jpg)

    # DSL that mirrors what the layout DSL files now emit after Fix A:
    # path is missing (empty), query carries the image_query value.
    dsl = (
        "canvas 1920x1080\n"
        "theme test\n"
        'picture 1040,200 780x720 path:"" query:"test scene" cover:true'
    )
    nodes, _ = parse_lines(dsl)
    prs = build_presentation(
        nodes,
        _minimal_tokens(),
        image_provider=provider,
        deck_dir=tmp_path,
    )

    # Provider was called with the image_query string, not a synthesised path.
    provider.search.assert_called_once_with("test scene", count=1)

    # A real PICTURE shape was emitted (not a grey rect placeholder).
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shapes = list(prs.slides[0].shapes)
    assert any(s.shape_type == MSO_SHAPE_TYPE.PICTURE for s in shapes), (
        "expected a PICTURE shape from provider, but none was emitted. "
        f"Got shapes: {[s.shape_type for s in shapes]}"
    )


# ---------------------------------------------------------------------------
# Test 2 — bare query: with no provider raises DSLError (intentional)
# ---------------------------------------------------------------------------

def test_bare_query_without_provider_raises(tmp_path):
    """Fix A: brand packs that set image_query without $image_provider
    should see an explicit DSLError, not a silent grey placeholder.

    This is the *intentional* fail-loud behavior preserved by the fix —
    documented here so future changes don't accidentally soften it.
    """
    dsl = (
        "canvas 1920x1080\n"
        "theme test\n"
        'picture 1040,200 780x720 query:"test scene" cover:true'
    )
    nodes, _ = parse_lines(dsl)
    with pytest.raises(DSLError, match="image_provider"):
        build_presentation(
            nodes,
            _minimal_tokens(),
            image_provider=None,
            deck_dir=tmp_path,
        )


# ---------------------------------------------------------------------------
# Test 3 — _placeholder_image_path resolves the gem from the installed package
# ---------------------------------------------------------------------------

def test_placeholder_image_path_resolves(tmp_path):
    """Fix B regression guard.

    ``_placeholder_image_path`` must return a real existing Path so that
    missing-image slots render the gem illustration instead of a grey rect.
    This also verifies that ``assets/illustrations/placeholder.jpg`` is
    reachable from the installed package (via the ``_PLACEHOLDER_ILLUSTRATION``
    fallback that uses ``Path(__file__).resolve()``).
    """
    # Build a minimal context with no asset_root overrides so the function
    # falls through to the source-tree / wheel path (_PLACEHOLDER_ILLUSTRATION).
    ctx = EmitContext(
        tokens=_minimal_tokens(),
        asset_root=None,
        asset_root_fallback=None,
        image_provider=None,
        deck_dir=tmp_path,
    )
    result = _placeholder_image_path(ctx)
    assert result is not None, (
        "_placeholder_image_path returned None — assets/illustrations/placeholder.jpg "
        "is not visible from the package. Check [tool.setuptools.package-data]."
    )
    assert result.is_file(), (
        f"_placeholder_image_path returned {result!r} but that path does not exist."
    )
