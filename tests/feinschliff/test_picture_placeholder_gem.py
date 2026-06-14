"""A missing/unresolvable picture should fall back to the shared gem
illustration (`assets/illustrations/placeholder.jpg`), not a bare grey box.

Before: `_emit_picture_placeholder` drew a `paper-2`/`fog` rect, so image-
bearing layouts with missing source assets shipped grey boxes. Now it places
the gem placeholder image (grey rect only if the asset can't be located).
"""
from pathlib import Path

from pptx.enum.shapes import MSO_SHAPE_TYPE

from feinschliff.dsl.parser import DSLNode
from feinschliff.dsl.pptx_emit import build_presentation
from feinschmiede.dsl.tokens import load_tokens

REPO = Path(__file__).resolve().parents[2] / "feinschliff"
BRANDS = REPO / "brands"


def _missing_picture_slide():
    tokens = load_tokens(BRANDS / "feinschliff", brands_dir=BRANDS)
    pic = DSLNode(
        kind="picture",
        pos_args=["100,100", "800x600"],
        kw_args={"path": "does-not-exist-xyz.jpg"},  # unresolvable, not optional
        label=None, line_no=1, source="t",
    )
    return build_presentation([pic], tokens, asset_root=BRANDS / "feinschliff" / "assets")


def test_missing_picture_renders_gem_placeholder_picture():
    prs = _missing_picture_slide()
    kinds = [s.shape_type for s in prs.slides[0].shapes]
    assert MSO_SHAPE_TYPE.PICTURE in kinds, (
        f"expected the gem placeholder image, got shapes: {kinds}"
    )
