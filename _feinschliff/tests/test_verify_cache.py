"""Unit tests for lib/verify/cache — per-slide content-hash cache."""
from __future__ import annotations

from pathlib import Path
from unittest.mock import patch


from lib.verify.cache import CachedVerdict, VerifyCache, slide_hash


# ---------------------------------------------------------------------------
# slide_hash — stability and invalidation tests
# ---------------------------------------------------------------------------

def _make_slide(layout: str = "layouts/title-body.slide.dsl", content: dict | None = None) -> dict:
    return {
        "layout": layout,
        "content": content or {"title": "Hello world", "body": "Some body text"},
        "_meta": {"title": "Should be excluded"},
        "slot_budgets": {"title": 60},
    }


def test_slide_hash_same_content_same_hash():
    slide = _make_slide()
    h1 = slide_hash(slide, "feinschliff")
    h2 = slide_hash(slide, "feinschliff")
    assert h1 == h2


def test_slide_hash_meta_change_does_not_change_hash():
    slide_a = _make_slide()
    slide_b = {**slide_a, "_meta": {"title": "Completely different meta"}}
    assert slide_hash(slide_a, "feinschliff") == slide_hash(slide_b, "feinschliff")


def test_slide_hash_slot_budgets_change_does_not_change_hash():
    slide_a = _make_slide()
    slide_b = {**slide_a, "slot_budgets": {"title": 999}}
    assert slide_hash(slide_a, "feinschliff") == slide_hash(slide_b, "feinschliff")


def test_slide_hash_content_change_changes_hash():
    slide_a = _make_slide(content={"title": "Hello", "body": "Original body"})
    slide_b = _make_slide(content={"title": "Hello", "body": "Changed body"})
    assert slide_hash(slide_a, "feinschliff") != slide_hash(slide_b, "feinschliff")


def test_slide_hash_layout_change_changes_hash():
    slide_a = _make_slide(layout="layouts/title-body.slide.dsl")
    slide_b = _make_slide(layout="layouts/claim-evidence.slide.dsl")
    assert slide_hash(slide_a, "feinschliff") != slide_hash(slide_b, "feinschliff")


def test_slide_hash_brand_change_changes_hash():
    slide = _make_slide()
    assert slide_hash(slide, "feinschliff") != slide_hash(slide, "other-brand")


def test_slide_hash_content_key_order_stable():
    """Inserting keys in different orders produces the same hash (sort_keys)."""
    content_a = {"body": "Body text", "title": "Title text"}
    content_b = {"title": "Title text", "body": "Body text"}
    slide_a = _make_slide(content=content_a)
    slide_b = _make_slide(content=content_b)
    assert slide_hash(slide_a, "feinschliff") == slide_hash(slide_b, "feinschliff")


def test_slide_hash_returns_hex_string():
    h = slide_hash(_make_slide(), "feinschliff")
    assert isinstance(h, str)
    assert len(h) == 64  # sha256 hex digest


# ---------------------------------------------------------------------------
# VerifyCache — get / put / save / load round-trip
# ---------------------------------------------------------------------------

def test_verify_cache_get_miss(tmp_path: Path):
    cache = VerifyCache(tmp_path)
    assert cache.get("nonexistent_hash", "squint") is None


def test_verify_cache_put_then_get(tmp_path: Path):
    cache = VerifyCache(tmp_path)
    verdict = CachedVerdict(
        slide_hash="abc123",
        rubric="squint",
        status="pass",
        findings={"status": "pass", "reason": ""},
        cached_at="2026-01-01T00:00:00+00:00",
    )
    cache.put(verdict)
    result = cache.get("abc123", "squint")
    assert result is not None
    assert result.status == "pass"
    assert result.rubric == "squint"
    assert result.findings == {"status": "pass", "reason": ""}


def test_verify_cache_save_and_reload(tmp_path: Path):
    cache = VerifyCache(tmp_path)
    verdict = CachedVerdict(
        slide_hash="abc123",
        rubric="squint",
        status="fail",
        findings={"status": "fail", "reason": "unclear headline"},
        cached_at="2026-01-01T00:00:00+00:00",
    )
    cache.put(verdict)
    cache.save()

    # Load fresh instance from the same dir
    cache2 = VerifyCache(tmp_path)
    result = cache2.get("abc123", "squint")
    assert result is not None
    assert result.status == "fail"
    assert result.findings["reason"] == "unclear headline"


def test_verify_cache_file_path(tmp_path: Path):
    cache = VerifyCache(tmp_path)
    cache.put(CachedVerdict("h", "squint", "pass", {}, "2026-01-01T00:00:00+00:00"))
    cache.save()
    assert (tmp_path / ".verify_cache.json").exists()


def test_verify_cache_different_rubric_separate_keys(tmp_path: Path):
    cache = VerifyCache(tmp_path)
    cache.put(CachedVerdict("h1", "squint", "pass", {"r": "squint"}, "2026-01-01T00:00:00+00:00"))
    cache.put(CachedVerdict("h1", "title-body", "fail", {"r": "title-body"}, "2026-01-01T00:00:00+00:00"))
    assert cache.get("h1", "squint").findings["r"] == "squint"
    assert cache.get("h1", "title-body").findings["r"] == "title-body"


def test_verify_cache_get_returns_none_for_wrong_rubric(tmp_path: Path):
    cache = VerifyCache(tmp_path)
    cache.put(CachedVerdict("h1", "squint", "pass", {}, "2026-01-01T00:00:00+00:00"))
    assert cache.get("h1", "title-body") is None


def test_verify_cache_put_overwrites(tmp_path: Path):
    cache = VerifyCache(tmp_path)
    cache.put(CachedVerdict("h1", "squint", "pass", {"v": 1}, "2026-01-01T00:00:00+00:00"))
    cache.put(CachedVerdict("h1", "squint", "fail", {"v": 2}, "2026-01-01T00:00:00+00:00"))
    result = cache.get("h1", "squint")
    assert result.status == "fail"
    assert result.findings["v"] == 2


# ---------------------------------------------------------------------------
# Cache integration with run_squint — second call must not invoke _judge
# ---------------------------------------------------------------------------

def _make_plan(n_slides: int = 2) -> dict:
    return {
        "slides": [
            {
                "layout": "layouts/title-body.slide.dsl",
                "content": {"title": f"Slide {i}", "body": f"Body {i}"},
                "_meta": {"title": f"Slide {i}"},
            }
            for i in range(1, n_slides + 1)
        ]
    }


def _make_rendered_pngs(tmp_path: Path, n: int) -> dict[int, Path]:
    """Create dummy 1x1 PNG files to satisfy make_squint_thumbnail."""
    from PIL import Image
    pngs = {}
    for i in range(1, n + 1):
        p = tmp_path / f"slide-{i:03d}.png"
        Image.new("RGB", (10, 10), color=(128, 128, 128)).save(p, format="PNG")
        pngs[i] = p
    return pngs


def test_run_squint_second_call_uses_cache(tmp_path: Path):
    """Running run_squint twice with same plan+cache → _judge called only on first run."""
    from lib.verify.cache import VerifyCache
    from lib.verify.llm import rubric as rubric_mod

    plan = _make_plan(2)
    pngs = _make_rendered_pngs(tmp_path, 2)
    cache = VerifyCache(tmp_path)

    judge_return = {"status": "pass", "reason": "looks good"}

    with patch.object(rubric_mod, "_judge", return_value=judge_return) as mock_judge:
        # First run — should call _judge twice (one per slide)
        r1 = rubric_mod.run_squint(
            tmp_path / "deck.pptx",
            pngs,
            offline=False,
            cache=cache,
            plan=plan,
            brand="feinschliff",
        )
        assert mock_judge.call_count == 2

        mock_judge.reset_mock()

        # Second run — same plan, same cache → 0 LLM calls
        r2 = rubric_mod.run_squint(
            tmp_path / "deck.pptx",
            pngs,
            offline=False,
            cache=cache,
            plan=plan,
            brand="feinschliff",
        )
        assert mock_judge.call_count == 0, "Cache should prevent _judge calls on second run"

    # Both results should have per-slide entries
    assert len(r1.per_slide) == 2
    assert len(r2.per_slide) == 2
    # Second run results should be marked cached
    assert all(p.get("cached") is True for p in r2.per_slide)


def test_run_squint_cache_hit_only_for_unchanged_slides(tmp_path: Path):
    """After first run, changing one slide's content forces re-judge for that slide only."""
    from lib.verify.cache import VerifyCache
    from lib.verify.llm import rubric as rubric_mod

    plan = _make_plan(2)
    pngs = _make_rendered_pngs(tmp_path, 2)
    cache = VerifyCache(tmp_path)

    judge_return = {"status": "pass", "reason": "ok"}

    with patch.object(rubric_mod, "_judge", return_value=judge_return) as mock_judge:
        # First run — populate cache
        rubric_mod.run_squint(
            tmp_path / "deck.pptx", pngs,
            offline=False, cache=cache, plan=plan, brand="feinschliff",
        )
        assert mock_judge.call_count == 2
        mock_judge.reset_mock()

        # Modify slide 2's content
        plan2 = {
            "slides": [
                plan["slides"][0],  # slide 1 unchanged
                {
                    **plan["slides"][1],
                    "content": {"title": "Changed title", "body": "Changed body"},
                },
            ]
        }

        rubric_mod.run_squint(
            tmp_path / "deck.pptx", pngs,
            offline=False, cache=cache, plan=plan2, brand="feinschliff",
        )
        # Only slide 2 should re-invoke _judge
        assert mock_judge.call_count == 1


def test_run_squint_no_cache_arg_unchanged_behavior(tmp_path: Path):
    """Without cache args, run_squint behaves exactly as before (always calls _judge)."""
    from lib.verify.llm import rubric as rubric_mod

    pngs = _make_rendered_pngs(tmp_path, 2)
    judge_return = {"status": "pass", "reason": "ok"}

    with patch.object(rubric_mod, "_judge", return_value=judge_return) as mock_judge:
        rubric_mod.run_squint(tmp_path / "deck.pptx", pngs, offline=False)
        rubric_mod.run_squint(tmp_path / "deck.pptx", pngs, offline=False)
        # No cache → _judge called each time for each slide
        assert mock_judge.call_count == 4  # 2 slides × 2 runs


def test_no_cache_flag_in_verify_quality_cli(tmp_path: Path):
    """--no-cache flag is accepted without error (offline mode for simplicity)."""
    import subprocess
    from pathlib import Path as _P
    from pptx import Presentation

    # Build a minimal .pptx fixture
    deck = tmp_path / "deck.pptx"
    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Test Slide"
    prs.save(str(deck))

    proc = subprocess.run(
        ["uv", "run", "feinschliff", "verify-quality",
         str(deck), "--offline", "--no-cache", "--json",
         "--out", str(tmp_path / "report.json")],
        cwd=_P(__file__).resolve().parents[1],
        capture_output=True, text=True,
    )
    assert proc.returncode == 0, proc.stderr


# ---------------------------------------------------------------------------
# I1 — Data-loss path: missing --plan must not touch existing cache file
# ---------------------------------------------------------------------------

def test_no_plan_arg_does_not_overwrite_existing_cache(tmp_path: Path):
    """I1 fix: when --plan is absent, an existing .verify_cache.json is never touched."""
    import json as _json
    from unittest.mock import patch
    from pptx import Presentation

    # Pre-populate a cache file with real content
    cache_file = tmp_path / ".verify_cache.json"
    original_content = {"abc123::squint": {"slide_hash": "abc123", "rubric": "squint",
                                            "status": "pass", "findings": {}, "cached_at": "x"}}
    cache_file.write_text(_json.dumps(original_content))
    original_mtime = cache_file.stat().st_mtime

    # Build a minimal .pptx
    deck = tmp_path / "deck.pptx"
    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Test"
    prs.save(str(deck))

    # Run cmd_verify_quality WITHOUT --plan (plan arg defaults to None)
    import argparse
    from cli.verify_quality import cmd_verify_quality

    args = argparse.Namespace(
        deck=deck,
        offline=True,
        no_cache=False,     # caching NOT disabled by flag …
        plan=None,          # … but --plan is absent
        brand=None,
        rubric="squint",
        out=tmp_path / "report.md",
        json_out=False,
    )

    # Patch _render_slide_pngs to avoid needing a real PPTX renderer
    with patch("cli.verify_quality._render_slide_pngs", return_value={}):
        cmd_verify_quality(args)

    # The cache file must be byte-for-byte identical — mtime unchanged is sufficient
    assert cache_file.stat().st_mtime == original_mtime, (
        "cache file was written even though --plan was absent (data-loss bug)"
    )
    assert _json.loads(cache_file.read_text()) == original_content


# ---------------------------------------------------------------------------
# I2 — Stale slide_index in cached findings must be overwritten by loop idx
# ---------------------------------------------------------------------------

def test_run_squint_cached_stale_slide_index_is_overwritten(tmp_path: Path):
    """I2 fix: a stale slide_index=99 in stored findings is replaced by the loop's idx."""
    from lib.verify.cache import VerifyCache, CachedVerdict
    from lib.verify.llm import rubric as rubric_mod
    from unittest.mock import patch

    plan = _make_plan(2)
    pngs = _make_rendered_pngs(tmp_path, 2)

    # Pre-populate the cache with a verdict whose findings carry a bogus slide_index
    cache = VerifyCache(tmp_path)
    for i, slide in enumerate(plan["slides"], start=1):
        from lib.verify.cache import slide_hash
        h = slide_hash(slide, "feinschliff")
        cache.put(CachedVerdict(
            slide_hash=h,
            rubric="squint",
            status="pass",
            findings={"status": "pass", "reason": "ok", "slide_index": 99},  # stale!
            cached_at="2026-01-01T00:00:00+00:00",
        ))
    cache.save()

    # Re-load the cache (simulates a fresh run)
    fresh_cache = VerifyCache(tmp_path)
    with patch.object(rubric_mod, "_judge") as mock_judge:
        result = rubric_mod.run_squint(
            tmp_path / "deck.pptx", pngs,
            offline=False, cache=fresh_cache, plan=plan, brand="feinschliff",
        )
        # All slides served from cache — no LLM calls
        assert mock_judge.call_count == 0

    # Every per_slide entry must have the correct loop index, NOT the stale 99
    for entry in result.per_slide:
        assert entry["slide_index"] != 99, (
            f"slide_index was not re-asserted; got stale value 99 in {entry}"
        )
    assert {e["slide_index"] for e in result.per_slide} == {1, 2}


def test_run_text_rubric_cached_stale_slide_index_is_overwritten(tmp_path: Path):
    """I2 fix (text path): stale slide_index=99 in findings is replaced by loop idx."""
    from lib.verify.cache import VerifyCache, CachedVerdict, slide_hash
    from lib.verify.llm import rubric as rubric_mod
    from unittest.mock import patch
    from pptx import Presentation

    # Build a minimal 2-slide PPTX with titles
    deck = tmp_path / "deck.pptx"
    prs = Presentation()
    for i in range(1, 3):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Title {i}"
        slide.placeholders[1].text = f"Body {i}"
    prs.save(str(deck))

    plan = _make_plan(2)
    cache = VerifyCache(tmp_path)

    # Pre-populate cache with stale slide_index=99
    for i, slide_data in enumerate(plan["slides"], start=1):
        h = slide_hash(slide_data, "feinschliff")
        cache.put(CachedVerdict(
            slide_hash=h,
            rubric="title-body",
            status="pass",
            findings={"status": "pass", "reason": "ok", "slide_index": 99},  # stale!
            cached_at="2026-01-01T00:00:00+00:00",
        ))
    cache.save()

    fresh_cache = VerifyCache(tmp_path)
    with patch.object(rubric_mod, "_judge") as mock_judge:
        result = rubric_mod.run_title_body(
            deck, offline=False, cache=fresh_cache, plan=plan, brand="feinschliff",
        )
        assert mock_judge.call_count == 0

    for entry in result.per_slide:
        assert entry["slide_index"] != 99, (
            f"stale slide_index 99 was not overwritten in {entry}"
        )
    assert {e["slide_index"] for e in result.per_slide} == {1, 2}
