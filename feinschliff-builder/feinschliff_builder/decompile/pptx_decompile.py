"""Inverse of pptx_emit — read a .pptx, emit per-slide .slide.dsl files.

Walks each slide's shapes, classifies them (text / rect / native MSO_SHAPE /
line / picture), and emits the matching DSL primitive. Where possible:

  - Solid-fill and stroke hex colors are reverse-mapped to brand color
    token names (`accent`, `paper-2`, `fog`, …) by exact hex match.
  - Text run (font, size, weight, color) tuples are reverse-mapped to a
    `STYLE_BUNDLES` name (canonical + brand-defined). Falls back to a
    `body` style with inline `color:` override when no bundle matches.
  - Shape rotation is captured into `rotate:` kwarg.
  - Native MSO_SHAPE.PARALLELOGRAM's adj1 is read back from
    shape.adjustments[0] and emitted as `adj1:N`.
  - Picture shapes are extracted to `<assets_dir>/source-slide-NN-K.<ext>`
    and referenced via `picture X,Y WxH path:…`.

Day-one scope is *primitive* level: one DSL primitive per source shape.
Compound recognition (e.g. detecting a 5-shape group as a `footer(…)` call)
and slot extraction (replacing literal strings with `{{ title }}` etc.)
are explicitly out of scope and belong to a separate pass.
"""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_FILL_TYPE, MSO_THEME_COLOR_INDEX

from feinschliff.dsl.tokens import STYLE_BUNDLES, Tokens, load_tokens


# 1 design-px = 6350 EMU on the canonical 1920×1080 canvas. Mirrors the
# emitter side.
_EMU_PER_PX = 6350
_PX_TO_PT = 0.5  # design-px → typographic pt (matches pptx_emit)


# Reverse of `_SHAPE_KIND` in pptx_emit. The forward map collapses several
# kinds (chevron / right-arrow / ellipse / circle) onto the same MSO_SHAPE;
# the reverse picks the most-specific name we author with.
_MSO_TO_KIND: dict[MSO_SHAPE, str] = {
    MSO_SHAPE.RECTANGLE:           "rect",
    MSO_SHAPE.OVAL:                "oval",
    MSO_SHAPE.ISOSCELES_TRIANGLE:  "triangle",
    MSO_SHAPE.RIGHT_TRIANGLE:      "triangle-right",
    MSO_SHAPE.CHEVRON:             "chevron",
    MSO_SHAPE.RIGHT_ARROW:         "right-arrow",
    MSO_SHAPE.LEFT_ARROW:          "left-arrow",
    MSO_SHAPE.DIAMOND:             "diamond",
    MSO_SHAPE.TRAPEZOID:           "trapezoid",
    MSO_SHAPE.PARALLELOGRAM:       "parallelogram",
    MSO_SHAPE.PIE:                 "pie",
    MSO_SHAPE.PIE_WEDGE:           "pie-wedge",
    MSO_SHAPE.ARC:                 "arc",
    MSO_SHAPE.BLOCK_ARC:           "block-arc",
}


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _px(emu: int | None) -> int:
    """EMU → design-px, rounded to int (matches the emitter's grid)."""
    if emu is None:
        return 0
    return int(round(emu / _EMU_PER_PX))


def _scaled(value: int, scale: float) -> int:
    """Rescale a design-px coord by `scale` (used when source canvas isn't
    the toolkit's canonical 1920×1080)."""
    return int(round(value * scale))


def _escape_label(text: str) -> str:
    return text.replace("\\", "\\\\").replace('"', '\\"').replace("\n", " ")


def _rgb_to_hex(rgb: RGBColor | None) -> str | None:
    if rgb is None:
        return None
    return f"{rgb}".lower()


@dataclass
class _ColorIndex:
    """Maps hex strings (lowercase, no `#`) back to brand color token names."""
    hex_to_name: dict[str, str]

    @classmethod
    def from_tokens(cls, tokens: Tokens) -> _ColorIndex:
        idx: dict[str, str] = {}
        for name, value in tokens.raw.get("color", {}).items():
            if name.startswith("$"):
                continue
            try:
                hex_val = tokens.color(name)
            except KeyError:
                continue
            hex_val = hex_val.lower().lstrip("#")
            # First color name wins for a given hex; brand packs use
            # ordered keys (accent before accent-aliases).
            idx.setdefault(hex_val, name)
        return cls(idx)

    def resolve(self, hex_color: str | None) -> str | None:
        if not hex_color:
            return None
        return self.hex_to_name.get(hex_color.lower().lstrip("#"))


@dataclass
class _StyleIndex:
    """Reverse-style lookup keyed by (family, size_px, weight, color_role)."""
    entries: list[tuple[tuple[str, float, int, str], str]]

    @classmethod
    def from_tokens(cls, tokens: Tokens) -> _StyleIndex:
        brand_style_names = {k for k in tokens.raw.get("style", {})
                             if not k.startswith("$")}
        names = set(STYLE_BUNDLES.keys()) | brand_style_names
        entries: list[tuple[tuple[str, float, int, str], str]] = []
        for name in sorted(names):
            try:
                resolved = tokens.resolve_style(name)
            except KeyError:
                continue
            key = (
                (resolved.font_family[0] if resolved.font_family else "").lower(),
                round(resolved.size_px, 1),
                int(resolved.weight),
                resolved.color_role,
            )
            entries.append((key, name))
        return cls(entries)

    def resolve(self, family: str, size_px: float, weight: int,
                color_role: str | None) -> str | None:
        if not family or color_role is None:
            return None
        target = (family.lower(), round(size_px, 1), int(weight), color_role)
        for key, name in self.entries:
            if (key[0] == target[0]
                    and abs(key[1] - target[1]) <= 0.5
                    and key[2] == target[2]
                    and key[3] == target[3]):
                return name
        return None


# ---------------------------------------------------------------------------
# Per-shape emitters
# ---------------------------------------------------------------------------

def _shape_fill_hex(shape) -> str | None:
    try:
        fill = shape.fill
        if fill.type is None or fill.type == MSO_FILL_TYPE.BACKGROUND:
            return None
        if fill.type == MSO_FILL_TYPE.SOLID:
            return _rgb_to_hex(fill.fore_color.rgb)
    except (AttributeError, ValueError):
        pass
    return None


def _shape_stroke(shape) -> tuple[str | None, float | None]:
    """Return (stroke_hex, stroke_width_pt). Either may be None."""
    try:
        line = shape.line
        hex_color = _rgb_to_hex(line.color.rgb) if line.color and line.color.type else None
        width_pt = line.width.pt if line.width else None
        return hex_color, width_pt
    except (AttributeError, ValueError):
        return None, None


def _font_color_hex(font) -> str | None:
    color = font.color
    try:
        if color.type is None:
            return None
        if color.type == MSO_THEME_COLOR_INDEX:
            return None  # theme color — out of scope; fall back to no role
        return _rgb_to_hex(color.rgb)
    except (AttributeError, ValueError):
        return None


def _dominant_run_style(shape, color_index: _ColorIndex, style_index: _StyleIndex
                       ) -> tuple[str, str | None, str | None]:
    """Walk the shape's text frame, return (style_name, color_role, color_hex).

    Style picked from the first non-empty run; color from same run. Style
    falls back to `body` when no bundle matches.
    """
    style_name = "body"
    color_role: str | None = None
    color_hex: str | None = None
    if not shape.has_text_frame:
        return style_name, color_role, color_hex
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if not run.text.strip():
                continue
            font = run.font
            family = font.name or ""
            size_pt = float(font.size.pt) if font.size else 16.0
            size_px = size_pt / _PX_TO_PT       # 1 pt = 2 design-px
            weight = 700 if font.bold else 400
            color_hex = _font_color_hex(font)
            color_role = color_index.resolve(color_hex)
            matched = style_index.resolve(family, size_px, weight, color_role)
            if matched:
                style_name = matched
            return style_name, color_role, color_hex
    return style_name, color_role, color_hex


def _emit_text(shape, x: int, y: int, w: int, h: int,
               color_index: _ColorIndex, style_index: _StyleIndex) -> str | None:
    if not shape.has_text_frame or not shape.text_frame.text.strip():
        return None
    style_name, color_role, color_hex = _dominant_run_style(shape, color_index, style_index)
    label = _escape_label(shape.text_frame.text.strip())
    color_attr = ""
    if color_role:
        # Only emit color: when it differs from the style's default (so we
        # don't bloat the DSL with redundant overrides). Cheapest check:
        # always emit; the emitter no-ops on identical overrides.
        color_attr = f" color:{color_role}"
    elif color_hex:
        color_attr = f' color:"#{color_hex}"'  # inline literal hex fallback
    return f'text {x},{y} style:{style_name}{color_attr} maxwidth:{w} maxheight:{h} "{label}"'


def _emit_rect(shape, x: int, y: int, w: int, h: int,
               color_index: _ColorIndex) -> str:
    fill_hex = _shape_fill_hex(shape)
    fill_role = color_index.resolve(fill_hex)
    parts = [f"rect {x},{y} {w}x{h}"]
    if fill_role:
        parts.append(f"fill:{fill_role}")
    elif fill_hex:
        parts.append(f'fill:"#{fill_hex}"')
    stroke_hex, stroke_pt = _shape_stroke(shape)
    if stroke_hex:
        stroke_role = color_index.resolve(stroke_hex)
        parts.append(f"stroke:{stroke_role}" if stroke_role else f'stroke:"#{stroke_hex}"')
        if stroke_pt:
            # pt → design-px: stroke-width in DSL is design-px units
            parts.append(f"stroke-width:{int(round(stroke_pt / 0.75))}")
    return " ".join(parts)


def _emit_shape(shape, x: int, y: int, w: int, h: int,
                color_index: _ColorIndex) -> str:
    try:
        kind = _MSO_TO_KIND.get(shape.auto_shape_type, "rect")
    except (AttributeError, ValueError):
        kind = "rect"
    if kind == "rect":
        return _emit_rect(shape, x, y, w, h, color_index)

    parts = [f"shape {x},{y} {w}x{h} kind:{kind}"]
    fill_hex = _shape_fill_hex(shape)
    fill_role = color_index.resolve(fill_hex)
    if fill_role:
        parts.append(f"fill:{fill_role}")
    elif fill_hex:
        parts.append(f'fill:"#{fill_hex}"')
    stroke_hex, stroke_pt = _shape_stroke(shape)
    if stroke_hex:
        stroke_role = color_index.resolve(stroke_hex)
        parts.append(f"stroke:{stroke_role}" if stroke_role else f'stroke:"#{stroke_hex}"')
        if stroke_pt:
            parts.append(f"stroke-width:{int(round(stroke_pt / 0.75))}")
    if shape.rotation:
        parts.append(f"rotate:{int(round(shape.rotation))}")
    # Adjustment handle for parallelogram skew, pie angle, etc.
    try:
        adj = shape.adjustments[0]
        if adj is not None and abs(adj - _default_adj_for(kind)) > 0.001:
            parts.append(f"adj1:{adj:.3f}")
    except (IndexError, AttributeError, ValueError):
        pass
    return " ".join(parts)


def _default_adj_for(kind: str) -> float:
    """Per-kind default adjustment values from python-pptx — used to decide
    whether to emit an `adj1:` override or leave the kind's default."""
    return {
        "parallelogram": 0.25,
        "pie": 0.0,
        "pie-wedge": 0.0,
        "trapezoid": 0.20,
        "chevron": 0.50,
    }.get(kind, 0.0)


def _emit_line(shape, color_index: _ColorIndex) -> str:
    """Connector → `line X,Y X2,Y2 stroke:role stroke-width:N`."""
    x1, y1 = _px(shape.left), _px(shape.top)
    x2 = x1 + _px(shape.width)
    y2 = y1 + _px(shape.height)
    parts = [f"line {x1},{y1} {x2},{y2}"]
    stroke_hex, stroke_pt = _shape_stroke(shape)
    if stroke_hex:
        stroke_role = color_index.resolve(stroke_hex)
        parts.append(f"stroke:{stroke_role}" if stroke_role else f'stroke:"#{stroke_hex}"')
        if stroke_pt:
            parts.append(f"stroke-width:{int(round(stroke_pt / 0.75))}")
    return " ".join(parts)


def _emit_picture(shape, x: int, y: int, w: int, h: int,
                  assets_dir: Path, slide_idx: int, pic_idx: int) -> str:
    """Extract the embedded image to assets_dir; emit `picture X,Y WxH path:NAME`."""
    try:
        img = shape.image
        ext = img.ext or "png"
        asset_name = f"source-slide-{slide_idx:02d}-{pic_idx}.{ext}"
        out = assets_dir / asset_name
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_bytes(img.blob)
        return f"picture {x},{y} {w}x{h} path:{asset_name} cover:false"
    except Exception as exc:
        return f"# picture {x},{y} {w}x{h} — extract failed: {exc}"


# ---------------------------------------------------------------------------
# Top-level decompile
# ---------------------------------------------------------------------------

def decompile_pptx(pptx_path: Path, brand_pack_dir: Path,
                   output_dir: Path, assets_dir: Path | None = None,
                   brands_dir: Path | None = None) -> int:
    """Decompile every slide in `pptx_path` into per-slide .slide.dsl files
    under `output_dir`. Pictures are written to `assets_dir` (defaults to
    `brand_pack_dir / 'assets'`). `brands_dir` resolves the `extends:`
    parent chain — defaults to `brand_pack_dir.parent`, but pass the
    feinschliff repo's `brands/` for out-of-tree brand packs.
    Returns the slide count.
    """
    tokens = load_tokens(brand_pack_dir, brands_dir=brands_dir)
    color_index = _ColorIndex.from_tokens(tokens)
    style_index = _StyleIndex.from_tokens(tokens)
    if assets_dir is None:
        assets_dir = brand_pack_dir / "assets"

    prs = Presentation(str(pptx_path))
    output_dir.mkdir(parents=True, exist_ok=True)

    # Source canvas in design-px (1 px = 6350 EMU). Rescale every coord into
    # the toolkit's canonical 1920×1080 frame so emitted DSL composes with
    # existing layouts.
    src_w = _px(prs.slide_width)
    src_h = _px(prs.slide_height)
    target_w, target_h = 1920, 1080
    scale_x = target_w / src_w if src_w else 1.0
    scale_y = target_h / src_h if src_h else 1.0
    # Use a single uniform scale (x and y are tied to slide aspect ratio).
    scale = (scale_x + scale_y) / 2

    for slide_idx, slide in enumerate(prs.slides, start=1):
        lines: list[str] = [
            "# Decompiled by feinschliff decompile from "
            f"{pptx_path.name} slide {slide_idx}",
            f"# Source canvas {src_w}x{src_h} → rescaled to 1920x1080 (scale {scale:.3f})",
            "canvas 1920x1080",
            f"theme {tokens.brand_name}",
            "",
        ]
        _pic_counter = [0]
        for shape in slide.shapes:
            x = _scaled(_px(shape.left), scale)
            y = _scaled(_px(shape.top), scale)
            w = _scaled(_px(shape.width), scale)
            h = _scaled(_px(shape.height), scale)
            line = _emit_one(shape, x, y, w, h,
                             color_index=color_index, style_index=style_index,
                             assets_dir=assets_dir, slide_idx=slide_idx,
                             pic_idx_ref=lambda inc=False, _c=_pic_counter:
                                 (_c.__setitem__(0, _c[0] + 1) if inc else _c[0]))
            if line:
                lines.append(line)
        out = output_dir / f"slide-{slide_idx:02d}.slide.dsl"
        out.write_text("\n".join(lines) + "\n")
    return len(prs.slides)


def _emit_one(shape, x, y, w, h, *, color_index, style_index,
              assets_dir, slide_idx, pic_idx_ref) -> str | None:
    """Dispatch a single shape to the right primitive emitter."""
    stype = shape.shape_type
    if stype == MSO_SHAPE_TYPE.PICTURE:
        pic_idx_ref(True)
        return _emit_picture(shape, x, y, w, h, assets_dir, slide_idx,
                             pic_idx_ref())
    if stype == MSO_SHAPE_TYPE.LINE:
        return _emit_line(shape, color_index)
    if stype == MSO_SHAPE_TYPE.GROUP:
        # Groups are flattened: emit each member as its own primitive.
        out: list[str] = []
        for member in shape.shapes:
            mx = _px(member.left) + x
            my = _px(member.top) + y
            mw = _px(member.width)
            mh = _px(member.height)
            line = _emit_one(member, mx, my, mw, mh,
                             color_index=color_index, style_index=style_index,
                             assets_dir=assets_dir, slide_idx=slide_idx,
                             pic_idx_ref=pic_idx_ref)
            if line:
                out.append(line)
        return "\n".join(out) if out else None
    # AUTO_SHAPE or TEXT_BOX.
    has_text = shape.has_text_frame and shape.text_frame.text.strip()
    has_visual = (_shape_fill_hex(shape) is not None
                  or _shape_stroke(shape)[0] is not None)
    if has_text and not has_visual:
        return _emit_text(shape, x, y, w, h, color_index, style_index)
    if has_text and has_visual:
        # Both: emit the shape AND the text as two adjacent lines.
        return (_emit_shape(shape, x, y, w, h, color_index)
                + "\n" + (_emit_text(shape, x, y, w, h, color_index, style_index) or ""))
    return _emit_shape(shape, x, y, w, h, color_index)
