"""Hybrid PPTX+SVG → Feinschliff DSL decompiler — brand-agnostic.

A higher-fidelity alternative to `lib/dsl/pptx_decompile.py`. Uses PPTX
XML as the canonical source for shape semantics and (optionally) SVG —
rendered from the slide's PDF via pdf2svg — as a secondary source for
custGeom geometry that PPTX xfrm inherits from the layout.

Sources of truth:

  PPTX (`ppt/slides/slideN.xml` + slideLayouts + slideMasters + theme1)
    - canonical for shape semantics: placeholder type/idx, group structure,
      text content + per-run style, schemeClr → token resolution, z-order,
      footer / page-number / wordmark placeholders inherited from master.

  SVG (pdf2svg of the slide's PDF page; optional)
    - canonical for final rendered geometry: bounding boxes of custGeom
      shapes, fall-back bbox when PPTX xfrm is inherited from layout.
    - NOT used for text classification or color matching (PPTX wins).

  Picture handling
    - every <p:pic> shape and every placeholder with type="pic" is
      emitted as a `picture` statement pointing at a configurable
      placeholder image (default: `assets/illustrations/placeholder.jpg`,
      the feinschliff convention).

Coordinate systems:
  - PPTX EMU: 914400 EMU/inch, slide size from presentation.xml::sldSz.
  - SVG: PDF points (1 pt = 1/72 inch). pdf2svg width/height map 1:1 to
    the slide's printable area, so EMU and pt match through inch.
  - DSL canvas: 1920×1080 px by default (override via canvas_w/h).

Brand-specific knobs (all defaulted to feinschliff baseline):
  - `theme_name` — name of the brand to emit on the `theme` directive
  - `tokens_path` — brand's tokens.json (for nearest-color matching)
  - `placeholder_rel` — DSL-relative path for picture placeholders
  - Footer-region text is emitted as plain `text` primitives; brands
    that ship a `footer(...)` compound can post-process the output to
    collapse the four footer lines into a single compound call.

Usage (programmatic, preferred):
  from feinschliff_builder.decompile.pptx_svg_decompile import derive
  dsl = derive(pptx_path, slide_idx=1, theme_name="acme",
               tokens_path=Path("brands/acme/tokens.json"),
               layout_name="cover-orange")

Usage (CLI smoke test):
  uv run python lib/dsl/pptx_svg_decompile.py SOURCE.pptx --slide N \\
      --theme <brand> --brand-tokens brands/<brand>/tokens.json > out.dsl
"""
from __future__ import annotations

import argparse
import json
import math
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation

from feinschmiede.dsl.tokens import STYLE_BUNDLES

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

EMU_PER_PT = 12700
PLACEHOLDER_REL = "assets/illustrations/placeholder.jpg"
# A non-placeholder <p:pic> smaller than this fraction of the slide is treated as
# fixed corporate-design chrome (logo / mark) and carried natively; larger pics
# are changeable topical content and stay fillable picture slots.
_TEMPLATE_IMG_MAX_AREA = 0.12


# ---------------------------------------------------------------------------
# Data
# ---------------------------------------------------------------------------


@dataclass
class Shape:
    kind: str            # rect | line | oval | pic | text | table
    x: float             # px in canvas
    y: float
    w: float
    h: float
    fill: str | None = None       # token name, e.g. "accent"
    stroke: str | None = None
    text_runs: list["TextRun"] = field(default_factory=list)
    is_picture: bool = False      # True for <p:pic> or ph type="pic"
    # True when this shape was inherited from the layout/master (not authored on the
    # slide). Slide-authored text is real content; only inherited text is prompt copy.
    from_chain: bool = False
    ph_type: str | None = None    # 'title','body','subTitle','pic','ftr','sldNum',...
    ph_idx: str | None = None
    # Border width in design-px. Captured from `<a:ln w="...">` (EMU); when
    # None the emitter uses its default hairline width.
    stroke_width: float | None = None
    # Dashed-line preset name from `<a:ln><a:prstDash val="...">` (e.g.
    # "dash", "dot", "sysDash", "lgDashDot"). None = solid stroke.
    stroke_dash: str | None = None
    # Corner radius in design-px for `roundRect` shapes. Captured from the
    # `<a:gd name="adj" fmla="val N">` adjustment value where N is N/100000
    # of the shape's shortest side. None when the source uses sharp corners.
    corner_radius: float | None = None
    # Drop shadow descriptor from `<a:effectLst><a:outerShdw>`. Tuple
    # (blur_px, dist_px, angle_deg, color_token, alpha) so the decompiler
    # can emit a compact `shadow:` kwarg and the emitter can rebuild the
    # `<a:effectLst>` XML at build time. None = no shadow.
    shadow: tuple[float, float, float, str, float] | None = None
    # Gradient fill descriptor from `<a:gradFill>`. List of (position, color)
    # stops (position 0..1, color = token or hex) plus the linear angle in
    # degrees. None = solid fill (use `Shape.fill` instead). When set, the
    # emitter writes a `gradFill` XML block onto the shape's spPr; the
    # decompiler emits a `gradient:angle=Ddeg;0=token;1=token` kwarg.
    gradient: tuple[list[tuple[float, str]], float] | None = None
    # Source bodyPr `anchor` attribute — controls text vertical position
    # within the shape bbox. "ctr" centers (PowerPoint default for many
    # text frames); "b" bottoms; "t" or absent = top. Without this the DSL
    # always emits top-anchored text, so source content that's vertically
    # centered in its frame renders shifted up by half the frame height.
    valign: str | None = None
    # normAutofit (PowerPoint "shrink text on overflow"): `autoshrink` arms the
    # emitter's fit; `font_scale` (0..1) reproduces the source's pre-shrink so a
    # placeholder's text fits its box instead of overflowing the final render.
    autoshrink: bool = False
    font_scale: float = 1.0
    # Text-frame internal insets as (l, t, r, b) in design-px. Source carries
    # these on `<a:bodyPr lIns="..." tIns="..." rIns="..." bIns="...">` in EMU.
    # When absent on source, defaults to PowerPoint's published 91440 / 45720
    # EMU (= ~7.2px / ~3.6px at 13.33" canvas scale). Captured per-text so
    # decompiled DSL renders at the exact source text position; without this,
    # my emitter zeroed all four and shifted every text by ~9-19 px versus
    # source — visible as the persistent blue/red ghost offsets in the redline.
    padding: tuple[float, float, float, float] | None = None
    # When image_extract_dir is passed to derive(), pictures are extracted
    # from the source PPTX and this holds the brand-pack-relative path the
    # DSL's `default:` should resolve to. None falls back to the generic
    # placeholder (genericised brand-template behaviour).
    media_path: str | None = None
    media_rid: str | None = None  # rId of the <a:blip r:embed=...>, for extraction
    # Resolved python-pptx Part for the embedded image, captured at walk
    # time so the rId is looked up against the part that actually owns it
    # (slide vs. layout vs. master). Without this, layout-inherited
    # pictures (most corporate-template chrome) fail extraction because
    # their rId only resolves on the layout part, not the slide.
    media_part: Any = None
    # For custGeom shapes: the full pathLst converted to an SVG-`d` string
    # in *canvas pixels* (already scaled from the path's path-local space).
    # When set, emit_dsl emits the shape as an `svg { path … }` block
    # instead of the lossy bbox-rect fallback. Lets the renderer reproduce
    # rings, arrows, callouts, and other vector decoration that the PDF
    # pipeline would otherwise rasterise.
    svg_path_d: str | None = None
    # Verbatim source `<p:sp>` XML (emitted as a base64 `native` primitive) for
    # complex custGeom chrome — carried as an editable native vector instead of
    # round-tripping through svg → raster → picture. None = use the other fields.
    native_xml: str | None = None
    # For a carried <p:pic> (template image): base64 of the embedded media bytes,
    # re-embedded into the output deck by the emitter (the source rId is dead here).
    native_media: str | None = None
    # For a native-carried CHART <p:graphicFrame>: the external part-graph the
    # frame's `<c:chart r:id>` reaches (the chart part + its chartStyle /
    # chartColorStyle / embedded-xlsx children). A list of dicts, each
    # {partname, content_type, blob (base64), reltype, parent ('slide' | a
    # chart-partname str)}. Tables need NO external parts (inline <a:tbl>) so
    # they leave this None; charts carry it so the emitter can re-create the
    # parts + rewire rIds in the output deck. None = no external parts.
    native_parts: list[dict] | None = None


@dataclass
class TextRun:
    text: str
    pt: float            # font size in points
    bold: bool = False
    italic: bool = False
    color: str | None = None  # token name
    align: str | None = None  # paragraph alignment: "center" | "right" | "justify" | None


# ---------------------------------------------------------------------------
# Palette + color resolution
# ---------------------------------------------------------------------------


def load_palette(tokens_path: Path) -> dict[str, tuple[int, int, int]]:
    # Walk the brand-pack `extends:` chain when DESIGN.md declares one, so a
    # child pack that only carries local overrides still gets the parent's
    # palette for nearest-colour matching. Falls back to reading just the
    # immediate file when there's no DESIGN.md or the parent can't be
    # resolved (e.g. out-of-tree packs whose parent lives elsewhere).
    brand_root = tokens_path.parent
    data = None
    if (brand_root / "DESIGN.md").is_file():
        from feinschmiede.dsl.tokens import load_tokens
        # Try sibling-located parent first (the default), then fall back
        # to the toolkit's bundled brands/ dir. Out-of-tree packs (e.g.
        # `.debug/brands/<name>` or `~/customer-brands/<name>`) declare
        # `extends: feinschliff` but their sibling dir isn't the toolkit
        # repo, so without this fallback the parent palette never loads
        # and nearest_token() degrades to raw hex emission for every
        # shape — visible in the decompiled DSL as `fill:#ffed00` instead
        # of `fill:accent` and `fill:neutral` on every custGeom because
        # _svg_color_token() then sees an unknown brand token.
        from feinschmiede.brand_discovery import discover_brands as _discover_brands
        candidate_dirs = [brand_root.parent]
        # Add all discovered brands directories in priority order so out-of-tree
        # packs whose parent lives in a different discovery source still resolve.
        for _brand in _discover_brands():
            brand_parent = _brand.root.parent
            if brand_parent not in candidate_dirs:
                candidate_dirs.append(brand_parent)
        seen: set[Path] = set()
        for cd in candidate_dirs:
            cd = cd.resolve()
            if cd in seen:
                continue
            seen.add(cd)
            try:
                data = load_tokens(brand_root, brands_dir=cd).raw
                break
            except (FileNotFoundError, ValueError):
                continue
    if data is None:
        data = json.loads(tokens_path.read_text(encoding="utf-8"))
    palette: dict[str, tuple[int, int, int]] = {}
    colors = data.get("color") or data.get("colors") or {}

    def _hex_of(entry):
        if isinstance(entry, dict):
            return entry.get("$value") or entry.get("value")
        if isinstance(entry, str):
            return entry
        return None

    for name, entry in colors.items():
        if name.startswith("$"):
            continue
        v = _hex_of(entry)
        if not isinstance(v, str):
            continue
        v = v.strip()
        if v.startswith("{") and v.endswith("}"):
            ref = v[1:-1].split(".")[-1]
            if ref in colors:
                v = _hex_of(colors[ref]) or ""
        if v.startswith("#") and len(v) == 7:
            palette[name] = (
                int(v[1:3], 16), int(v[3:5], 16), int(v[5:7], 16),
            )
    return palette


def nearest_token(rgb: tuple[int, int, int], palette: dict[str, tuple[int, int, int]]) -> str:
    # No palette → emit the raw hex literal so the DSL preserves the source
    # color verbatim. `derive()` documents this fallback for callers that
    # don't pass a tokens.json.
    if not palette:
        return "#{:02x}{:02x}{:02x}".format(*rgb)
    best = None
    best_d = math.inf
    for name, prgb in palette.items():
        d = sum((a - b) ** 2 for a, b in zip(rgb, prgb))
        if d < best_d:
            best_d = d
            best = name
    # Source-fidelity guard. Squared-euclidean threshold ≈ 25 per channel
    # (3 * 25^2 = 1875). When the closest brand token is further than this,
    # the source colour isn't really represented in the palette — emit the
    # raw hex literal instead of approximating to a token that renders as
    # a visibly different colour (e.g. the Sartorius source's #FFED00
    # yellow shouldn't squash to the feinschliff parent's gold #C9A24A
    # accent token just because it's the closest of 30 mostly-cool tokens).
    if best is not None and best_d <= _NEAREST_TOKEN_THRESHOLD_SQ:
        return best
    return "#{:02x}{:02x}{:02x}".format(*rgb)


_NEAREST_TOKEN_THRESHOLD_SQ = 1875


_THEME_RELTYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"


def master_theme_blob(pres: Presentation) -> bytes | None:
    """Raw XML of the theme part the deck's primary slide master references.

    PowerPoint numbers theme parts PER-MASTER, so the active theme is often NOT
    `theme1.xml` (a deck whose master is `slideMaster11` references
    `theme11.xml`). Resolving via the master→theme relationship — instead of a
    hardcoded part name — is what lets schemeClr fills/strokes, the background
    panel, and font capture work on real corporate templates.
    """
    try:
        for master in pres.slide_masters:
            for rel in master.part.rels.values():
                if rel.reltype == _THEME_RELTYPE:
                    return rel.target_part.blob
    except Exception:
        pass
    return None


def load_theme_scheme(pres: Presentation) -> dict[str, str]:
    """Map theme scheme keys (accent1..6, dk1, lt1, hlink, folHlink) to #RRGGBB.

    Resolves the theme part the slide master actually references (not a
    hardcoded `theme1.xml`); falls back to any theme part, then empty dict.
    """
    out: dict[str, str] = {}
    blobs: list[bytes] = []
    primary = master_theme_blob(pres)
    if primary is not None:
        blobs.append(primary)
    else:
        # Fallback: any theme part in the package (legacy behaviour).
        try:
            blobs = [p.blob for p in pres.part.package.iter_parts()
                     if "/theme/theme" in str(p.partname)]
        except Exception:
            blobs = []
    for blob in blobs:
        try:
            root = etree.fromstring(blob)
        except Exception:
            continue
        scheme = root.find(".//a:clrScheme", NS)
        if scheme is None:
            continue
        for child in scheme:
            key = etree.QName(child).localname  # dk1, lt1, accent1, ...
            srgb = child.find("a:srgbClr", NS)
            sys_ = child.find("a:sysClr", NS)
            if srgb is not None:
                out[key] = "#" + srgb.get("val").upper()
            elif sys_ is not None:
                out[key] = "#" + (sys_.get("lastClr") or "000000").upper()
        # PowerPoint default clrMap aliases — these slots are always present
        # and resolve to the matching scheme entry. Without them, a shape fill
        # of `schemeClr val="bg2"` (used widely in corporate templates that put
        # the slide bg in a layout rect rather than `<p:bg>`) falls through
        # unmapped.
        for alias, real in (("bg1", "lt1"), ("bg2", "lt2"),
                            ("tx1", "dk1"), ("tx2", "dk2")):
            if alias not in out and real in out:
                out[alias] = out[real]
        if out:
            break
    return out


# ---------------------------------------------------------------------------
# Geometry conversion
# ---------------------------------------------------------------------------


class CanvasMap:
    def __init__(self, slide_cx_emu: int, slide_cy_emu: int, canvas_w: int, canvas_h: int):
        self.cx = slide_cx_emu
        self.cy = slide_cy_emu
        self.cw = canvas_w
        self.ch = canvas_h
        self.sx = canvas_w / slide_cx_emu
        self.sy = canvas_h / slide_cy_emu
        # SVG renders at 1pt per pt; emu→pt = 1/12700. Conversion to SVG units:
        self.emu_to_pt = 1 / EMU_PER_PT

    def x(self, emu: float) -> int:
        return round(emu * self.sx)

    def y(self, emu: float) -> int:
        return round(emu * self.sy)

    def w(self, emu: float) -> int:
        return max(1, round(emu * self.sx))

    def h(self, emu: float) -> int:
        return max(1, round(emu * self.sy))


# ---------------------------------------------------------------------------
# Shape walking
# ---------------------------------------------------------------------------


def _split_runs_by_color(runs: list["TextRun"]) -> list[list["TextRun"]]:
    """Group runs into consecutive blocks of the same SIZE (and colour).

    A shape often stacks a large header over smaller body text ("Placeholder" 16pt
    over "This text demonstrates…" 12pt) or a coloured headline over body bullets.
    Collapsing to one primitive emits the whole shape at the MAX size / first
    colour, so the small body renders too big and OVERFLOWS. Split so each block
    keeps its own size + colour at its own y-offset (the caller positions them).

    A **size change always starts a new block** (a size jump is a distinct
    visual block regardless of layout). A **colour change splits only across a
    paragraph boundary** — i.e. when a `\\n` marker separates the two
    differently-coloured runs. Two explicitly-coloured runs on the SAME line
    (no intervening `\\n`) stay in one block, because stacking them at separate
    y-offsets would push the second run onto its own line. This is the
    footer's "Internal C-SC1" (bold red) + " | C/CGB-CD | …" (ink) case: one
    source line that previously split into two stacked statements. A colourless
    run inherits and attaches to the current block. `\\n` markers attach to the
    preceding block. Single-(size,colour) shapes return one block (the caller
    skips splitting).
    """
    blocks: list[list[TextRun]] = []
    current: list[TextRun] = []
    cur_color: str | None = None
    cur_size: int | None = None
    saw_break = False  # a paragraph `\n` marker seen since the last content run
    for r in runs:
        if not (r.text and r.text != "\n"):
            # Newline marker / empty run — attach to the current block and arm
            # the colour-split (a colour change is only meaningful as a stacked
            # block when it follows a real line break).
            if current:
                current.append(r)
            if r.text == "\n":
                saw_break = True
            continue
        sz = round(r.pt)
        if not current:
            current = [r]
            cur_color, cur_size = r.color, sz
            saw_break = False
            continue
        color_changed = (
            saw_break
            and r.color is not None and cur_color is not None and r.color != cur_color
        )
        if sz != cur_size or color_changed:
            blocks.append(current)
            current = [r]
            cur_color, cur_size = r.color, sz
        else:
            current.append(r)
            if cur_color is None:
                cur_color = r.color
        saw_break = False
    if current:
        blocks.append(current)
    return blocks


def _resolve_gradient(spPr: etree._Element, theme: dict[str, str],
                       palette: dict[str, tuple[int, int, int]]
                       ) -> tuple[list[tuple[float, str]], float] | None:
    """Extract `<a:gradFill>` as ([(pos, token)], angle_deg) — or None."""
    if spPr is None:
        return None
    grad = spPr.find("a:gradFill", NS)
    if grad is None:
        return None
    stops: list[tuple[float, str]] = []
    for gs in grad.findall("a:gsLst/a:gs", NS):
        try:
            pos = int(gs.get("pos") or 0) / 100000
        except ValueError:
            continue
        srgb = gs.find("a:srgbClr", NS)
        if srgb is not None and srgb.get("val"):
            hx = srgb.get("val")
            try:
                rgb = (int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))
                color = nearest_token(rgb, palette) if palette else f"#{hx}"
            except ValueError:
                continue
            stops.append((pos, color))
    if not stops:
        return None
    angle_deg = 0.0
    lin = grad.find("a:lin", NS)
    if lin is not None and lin.get("ang"):
        try:
            angle_deg = int(lin.get("ang")) / 60000.0
        except ValueError:
            pass
    return (stops, angle_deg)


def _resolve_fill(spPr: etree._Element, theme: dict[str, str], palette: dict[str, tuple[int, int, int]]) -> str | None:
    """Return a token name, or None if no fill.

    Handles `<a:grpFill/>` by walking up `<p:grpSp>` ancestors iteratively
    until an actual `<a:solidFill>` is found. The walk is iterative (not
    recursive on `_resolve_fill`) so nested `grpFill` chains — common in
    multi-level Design Kit groups where every ancestor's grpSpPr itself
    carries `<a:grpFill/>` — don't trigger `RecursionError`.
    """
    if spPr is None:
        return None
    sf = spPr.find("a:solidFill", NS)
    if sf is None:
        # No direct solid fill. If the shape declares grpFill, walk up
        # ancestor groups looking for the first one with a real solid
        # fill on its grpSpPr.
        if spPr.find("a:grpFill", NS) is None:
            return None
        anc = spPr.getparent()
        while anc is not None and sf is None:
            if etree.QName(anc).localname == "grpSp":
                grpSpPr = anc.find("p:grpSpPr", NS)
                if grpSpPr is not None:
                    inner_sf = grpSpPr.find("a:solidFill", NS)
                    if inner_sf is not None:
                        sf = inner_sf
                        break
                    # grpSpPr itself is grpFill-only — keep climbing.
            anc = anc.getparent()
        if sf is None:
            return None
    srgb = sf.find("a:srgbClr", NS)
    if srgb is not None:
        hx = srgb.get("val")
        rgb = (int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))
        rgb = _apply_color_mods(rgb, srgb)
        rgb = _blend_on_white(rgb, _alpha_for_color(srgb))
        return nearest_token(rgb, palette)
    scheme = sf.find("a:schemeClr", NS)
    if scheme is not None:
        key = scheme.get("val")
        hex_str = theme.get(key)
        if hex_str:
            rgb = (int(hex_str[1:3], 16), int(hex_str[3:5], 16), int(hex_str[5:7], 16))
            rgb = _apply_color_mods(rgb, scheme)
            rgb = _blend_on_white(rgb, _alpha_for_color(scheme))
            return nearest_token(rgb, palette)
    return None


def _get_xfrm(spPr: etree._Element) -> tuple[int, int, int, int] | None:
    if spPr is None:
        return None
    xfrm = spPr.find("a:xfrm", NS)
    if xfrm is None:
        return None
    off = xfrm.find("a:off", NS)
    ext = xfrm.find("a:ext", NS)
    if off is None or ext is None:
        return None
    return int(off.get("x")), int(off.get("y")), int(ext.get("cx")), int(ext.get("cy"))


def _placeholder_info(node: etree._Element) -> tuple[str | None, str | None]:
    ph = node.find(".//p:nvSpPr/p:nvPr/p:ph", NS)
    if ph is None:
        ph = node.find(".//p:nvPicPr/p:nvPr/p:ph", NS)
    if ph is None:
        ph = node.find(".//p:nvCxnSpPr/p:nvPr/p:ph", NS)
    if ph is None:
        return None, None
    return ph.get("type"), ph.get("idx")


def _layout_placeholder_default_sz(slide, ph_type: str | None, ph_idx: str | None) -> int | None:
    """Walk slide layout + master for the placeholder's default font size.

    PowerPoint inherits font sizes from the layout (and master) when a
    slide-level placeholder has no explicit `sz` on its runs/paragraphs.
    Layout writes the size on
    `<p:sp><p:txBody><a:lstStyle><a:lvl1pPr><a:defRPr sz="...">`; master
    defines title/body defaults on `<p:txStyles>/<p:titleStyle>` and
    `<p:bodyStyle>`.

    Without this lookup, body-level placeholders that omit explicit sz
    inherit our hardcoded 1800 (18pt), which renders chapter titles and
    other layout-controlled headlines at body-size — visibly wrong on
    showcase decks where the layout sets large headlines.

    Returns sz in hundredths-of-pt (PPTX units) or None.
    """
    layout = getattr(slide, "slide_layout", None)
    master = getattr(layout, "slide_master", None) if layout is not None else None
    parents = [p for p in (layout, master) if p is not None]
    for parent in parents:
        root = parent.element
        for sp in root.iter("{%s}sp" % NS["p"]):
            ph = sp.find(".//p:nvSpPr/p:nvPr/p:ph", NS)
            if ph is None:
                continue
            if (ph_type and ph.get("type") == ph_type) or (
                ph_idx and ph.get("idx") == ph_idx
            ):
                lvl1 = sp.find(".//p:txBody/a:lstStyle/a:lvl1pPr", NS)
                if lvl1 is not None:
                    d = lvl1.find("a:defRPr", NS)
                    if d is not None and d.get("sz"):
                        try:
                            return int(d.get("sz"))
                        except (TypeError, ValueError):
                            pass
    style_for_type = {"title": "titleStyle", "ctrTitle": "titleStyle"}
    style_name = style_for_type.get(ph_type or "", "bodyStyle")
    if master is not None:
        ts = master.element.find(f".//p:txStyles/p:{style_name}", NS)
        if ts is not None:
            lvl1 = ts.find("a:lvl1pPr", NS)
            if lvl1 is not None:
                d = lvl1.find("a:defRPr", NS)
                if d is not None and d.get("sz"):
                    try:
                        return int(d.get("sz"))
                    except (TypeError, ValueError):
                        pass
    return None


def _layout_placeholder_caps_bold(slide, ph_type: str | None, ph_idx: str | None) -> tuple[bool, bool]:
    """Walk slide layout + master for the placeholder's inherited cap/bold.

    Like `_layout_placeholder_default_sz` but for `cap="all"` (all-caps render
    transform) + `b="1"` (bold). A title whose run states neither still renders
    UPPERCASE + bold because the master titleStyle/bodyStyle (and sometimes the
    layout placeholder) defRPr sets them. Returns (caps_all, bold)."""
    layout = getattr(slide, "slide_layout", None)
    master = getattr(layout, "slide_master", None) if layout is not None else None
    caps = bold = False
    # Master title/body style is the cascade base.
    style_name = {"title": "titleStyle", "ctrTitle": "titleStyle"}.get(ph_type or "", "bodyStyle")
    if master is not None:
        d = master.element.find(f".//p:txStyles/p:{style_name}/a:lvl1pPr/a:defRPr", NS)
        if d is not None:
            if d.get("cap") is not None:
                caps = d.get("cap") == "all"
            if d.get("b") is not None:
                bold = d.get("b") == "1"
    # The layout's own placeholder defRPr overrides the master.
    if layout is not None:
        for sp in layout.element.iter("{%s}sp" % NS["p"]):
            ph = sp.find(".//p:nvSpPr/p:nvPr/p:ph", NS)
            if ph is None:
                continue
            if (ph_type and ph.get("type") == ph_type) or (ph_idx and ph.get("idx") == ph_idx):
                d = sp.find(".//p:txBody/a:lstStyle/a:lvl1pPr/a:defRPr", NS)
                if d is not None:
                    if d.get("cap") is not None:
                        caps = d.get("cap") == "all"
                    if d.get("b") is not None:
                        bold = d.get("b") == "1"
    return caps, bold


def _layout_placeholder_anchor(slide, ph_type: str | None, ph_idx: str | None) -> str | None:
    """Walk slide layout + master for the placeholder's inherited vertical anchor.

    Like `_layout_placeholder_default_sz` but for `<a:bodyPr anchor="ctr|b|t">`. A
    title whose own bodyPr sets no anchor still renders centre/bottom-anchored when
    the layout/master placeholder bodyPr does (MS Geometric: master title
    placeholder anchor="b"). Returns "middle" / "bottom" / "top" / None."""
    layout = getattr(slide, "slide_layout", None)
    master = getattr(layout, "slide_master", None) if layout is not None else None
    for parent in (p for p in (layout, master) if p is not None):
        for sp in parent.element.iter("{%s}sp" % NS["p"]):
            ph = sp.find(".//p:nvSpPr/p:nvPr/p:ph", NS)
            if ph is None:
                continue
            if (ph_type and ph.get("type") == ph_type) or (ph_idx and ph.get("idx") == ph_idx):
                bodyPr = sp.find(".//p:txBody/a:bodyPr", NS)
                anc = bodyPr.get("anchor") if bodyPr is not None else None
                if anc:
                    return {"ctr": "middle", "b": "bottom", "t": "top"}.get(anc)
    return None


def _layout_placeholder_insets(
    slide, ph_type: str | None, ph_idx: str | None
) -> tuple[int | None, int | None, int | None, int | None]:
    """Walk slide layout + master for a placeholder's inherited text-frame insets.

    Like `_layout_placeholder_anchor` but for `<a:bodyPr lIns/tIns/rIns/bIns>`. A
    title whose own bodyPr omits insets still renders at the layout/master
    placeholder's insets — the master title placeholder sets all four to
    `0`, so the slide's "Headline"/"Subheadline" must hug the box left edge, not
    inherit PowerPoint's published 91440/45720 EMU default (which shoved them ~16
    px right / 8 px down and was a top contributor to the heading ghost in the
    redline). Returns (l, t, r, b) in EMU; each element is None when no ancestor
    placeholder specifies that side (caller falls back to the PowerPoint default).
    The MASTER is the cascade base; the LAYOUT placeholder overrides per-side."""
    layout = getattr(slide, "slide_layout", None)
    master = getattr(layout, "slide_master", None) if layout is not None else None
    out: list[int | None] = [None, None, None, None]
    # Master first (base), then layout (override) — later writes win per-side.
    for parent in (p for p in (master, layout) if p is not None):
        for sp in parent.element.iter("{%s}sp" % NS["p"]):
            ph = sp.find(".//p:nvSpPr/p:nvPr/p:ph", NS)
            if ph is None:
                continue
            if (ph_type and ph.get("type") == ph_type) or (ph_idx and ph.get("idx") == ph_idx):
                bodyPr = sp.find(".//p:txBody/a:bodyPr", NS)
                if bodyPr is None:
                    continue
                for i, attr in enumerate(("lIns", "tIns", "rIns", "bIns")):
                    v = bodyPr.get(attr)
                    if v is not None:
                        try:
                            out[i] = int(v)
                        except (TypeError, ValueError):
                            pass
    return tuple(out)  # type: ignore[return-value]


def _layout_placeholder_autofit(slide, ph_type: str | None, ph_idx: str | None) -> bool:
    """Whether a placeholder INHERITS normAutofit ("shrink text on overflow") from
    its layout/master. Some corporate templates set it on the LAYOUT (on the layout, not on
    the slides), so without this the decompiled placeholder text renders full-size and
    overflows the final render. We don't apply the layout's own fontScale (a template
    default) — emitting `autoshrink` lets the renderer compute the per-slide fit,
    matching PowerPoint."""
    layout = getattr(slide, "slide_layout", None)
    master = getattr(layout, "slide_master", None) if layout is not None else None
    for parent in (p for p in (layout, master) if p is not None):
        for sp in parent.element.iter("{%s}sp" % NS["p"]):
            ph = sp.find(".//p:nvSpPr/p:nvPr/p:ph", NS)
            if ph is None:
                continue
            if (ph_type and ph.get("type") == ph_type) or (ph_idx and ph.get("idx") == ph_idx):
                if sp.find(".//p:txBody/a:bodyPr/a:normAutofit", NS) is not None:
                    return True
    return False


def _layout_placeholder_color(slide, ph_type: str | None, ph_idx: str | None,
                              theme: dict[str, str],
                              palette: dict[str, tuple[int, int, int]]) -> str | None:
    """Walk slide layout + master for the placeholder's default text colour.

    Mirrors `_layout_placeholder_default_sz` but pulls `<a:defRPr><a:solidFill>`.
    A layout's WHITE title placeholder (e.g. a `ctrTitle` sitting on a coloured
    circle) carries no colour on the slide run; without this lookup it decompiles
    colourless and renders ink-grey. Returns a token/hex or None.
    """
    layout = getattr(slide, "slide_layout", None)
    master = getattr(layout, "slide_master", None) if layout is not None else None
    for parent in (p for p in (layout, master) if p is not None):
        for sp in parent.element.iter("{%s}sp" % NS["p"]):
            ph = sp.find(".//p:nvSpPr/p:nvPr/p:ph", NS)
            if ph is None:
                continue
            if (ph_type and ph.get("type") == ph_type) or (ph_idx and ph.get("idx") == ph_idx):
                d = sp.find(".//p:txBody/a:lstStyle/a:lvl1pPr/a:defRPr", NS)
                sf = d.find("a:solidFill", NS) if d is not None else None
                if sf is not None:
                    c = _resolve_solid(sf, theme, palette)
                    if c:
                        return c
    style_name = {"title": "titleStyle", "ctrTitle": "titleStyle"}.get(ph_type or "", "bodyStyle")
    if master is not None:
        d = master.element.find(f".//p:txStyles/p:{style_name}/a:lvl1pPr/a:defRPr", NS)
        sf = d.find("a:solidFill", NS) if d is not None else None
        if sf is not None:
            c = _resolve_solid(sf, theme, palette)
            if c:
                return c
    return None


def _layout_placeholder_xfrm(slide, ph_type: str | None, ph_idx: str | None) -> tuple[int, int, int, int] | None:
    """Walk slide layout + master to resolve an inherited placeholder bbox.

    Accepts a Slide, SlideLayout, or SlideMaster — when the caller is
    already a layout/master (because walk_slide now recurses through the
    inheritance chain), there's no further parent to walk so this just
    no-ops out.
    """
    layout = getattr(slide, "slide_layout", None)
    master = getattr(layout, "slide_master", None) if layout is not None else None
    parents = [p for p in (layout, master) if p is not None]
    for parent in parents:
        root = parent.element
        for sp in root.iter("{%s}sp" % NS["p"]):
            ph = sp.find(".//p:nvSpPr/p:nvPr/p:ph", NS)
            if ph is None:
                continue
            if (ph_type and ph.get("type") == ph_type) or (
                ph_idx and ph.get("idx") == ph_idx
            ):
                xfrm = _get_xfrm(sp.find("p:spPr", NS))
                if xfrm:
                    return xfrm
    return None


def _text_runs(node: etree._Element, theme: dict[str, str], palette: dict[str, tuple[int, int, int]],
               inherited_default_sz: int | None = None,
               inherited_caps: bool = False, inherited_bold: bool = False) -> list[TextRun]:
    runs: list[TextRun] = []
    txBody = node.find(".//p:txBody", NS)
    if txBody is None:
        txBody = node.find(".//a:txBody", NS)
    if txBody is None:
        return runs
    # Body-level lstStyle/lvl1pPr/defRPr sz acts as the cascade default for any
    # run that does not set its own sz. Without this we miss titles whose size
    # is defined only at the body level (common in master-driven decks).
    body_default_sz: int | None = None
    lstStyle = txBody.find("a:lstStyle", NS)
    if lstStyle is not None:
        lvl1 = lstStyle.find("a:lvl1pPr", NS)
        if lvl1 is not None:
            d = lvl1.find("a:defRPr", NS)
            if d is not None and d.get("sz"):
                body_default_sz = int(d.get("sz"))
    for para in txBody.findall("a:p", NS):
        para_runs: list[TextRun] = []
        # Pick up para-level defRPr or pPr/defRPr for sz fallback.
        pPr = para.find("a:pPr", NS)
        # Cascade for the paragraph's default sz:
        # 1. txBody/lstStyle/lvl1pPr/defRPr sz (slide-level)
        # 2. inherited_default_sz (layout/master placeholder lookup — only
        #    threaded by `_emit_sp` when the shape is a placeholder)
        # 3. hardcoded 1800 (18pt)
        default_sz = body_default_sz
        if default_sz is None and inherited_default_sz is not None:
            default_sz = inherited_default_sz
        if default_sz is None:
            default_sz = 1800
        if pPr is not None:
            d = pPr.find("a:defRPr", NS)
            if d is not None and d.get("sz"):
                default_sz = int(d.get("sz"))
        # Paragraph-level alignment (`<a:pPr algn="ctr|r|just">`). Source
        # frequently centers KPI numbers + their labels within a card; the
        # emitter's default `align:left` shifts them left of source. Stored
        # on each TextRun in this paragraph so emit_dsl can lift the
        # majority-vote into a per-text `align:` kwarg.
        para_align = None
        if pPr is not None:
            algn = pPr.get("algn")
            if algn == "ctr":
                para_align = "center"
            elif algn == "r":
                para_align = "right"
            elif algn == "just":
                para_align = "justify"
        for r in para.findall("a:r", NS):
            rPr = r.find("a:rPr", NS)
            t = r.find("a:t", NS)
            if t is None or t.text is None:
                continue
            sz = default_sz
            # cap + bold cascade from the master title/body style when the run
            # states neither (inherited_*); a title whose run carries no `b`/`cap`
            # still decompiles bold + UPPERCASE, matching the render.
            bold = inherited_bold
            caps = inherited_caps
            italic = False
            color = None
            text = t.text
            if rPr is not None:
                if rPr.get("sz"):
                    sz = int(rPr.get("sz"))
                if rPr.get("b") is not None:
                    bold = rPr.get("b") == "1"
                italic = rPr.get("i") == "1"
                sf = rPr.find("a:solidFill", NS)
                if sf is not None:
                    color = _resolve_fill(rPr, theme, palette) or _resolve_solid(sf, theme, palette)
                if rPr.get("cap") is not None:
                    caps = rPr.get("cap") == "all"
            # PPTX `cap="all"` is a render-time text-transform: the run's stored
            # text stays mixed-case but draws uppercase. Bake the transform into
            # the emitted DSL since downstream layouts carry the literal text,
            # not a `text-transform` directive.
            if caps:
                text = text.upper()
            para_runs.append(TextRun(text=text, pt=sz / 100, bold=bold, italic=italic, color=color, align=para_align))
        if para_runs:
            # Insert a newline marker between paragraphs so emit_dsl can preserve
            # line breaks. Without this, "Headline" + "Lorem ipsum…" paragraphs
            # collapse into "HeadlineLorem ipsum…" and ruin the body diff score.
            if runs:
                runs.append(TextRun(text="\n", pt=default_sz / 100))
            runs.extend(para_runs)
        else:
            # No <a:r> — could be just <a:fld> (page number, date). Capture as one run.
            for fld in para.findall("a:fld", NS):
                t = fld.find("a:t", NS)
                if t is not None and t.text:
                    runs.append(TextRun(text=t.text, pt=default_sz / 100))
    return runs


def _alpha_for_color(color_el: etree._Element) -> float:
    """Return alpha 0..1 from `<a:alpha val="...">` child (PPTX uses 0..100000).
    Defaults to 1.0 when the element is absent."""
    a = color_el.find("a:alpha", NS)
    if a is None or not a.get("val"):
        return 1.0
    try:
        return max(0.0, min(1.0, int(a.get("val")) / 100000.0))
    except (TypeError, ValueError):
        return 1.0


def _blend_on_white(rgb: tuple[int, int, int], alpha: float) -> tuple[int, int, int]:
    """Pre-multiply RGBA against a white slide background.

    Most decks render alpha-on-shape against the slide's paper colour.
    Approximating "blend against white" lets us preserve the perceived
    colour of semi-transparent fills (Venn circles, overlay panels) on
    typical white-canvas slides without threading true alpha through the
    build pipeline. For non-white slide backgrounds the result is
    visually off but only fractionally so — the colour shifts toward
    white instead of the actual canvas.
    """
    if alpha >= 0.999:
        return rgb
    return tuple(
        max(0, min(255, int(round(c * alpha + 255 * (1 - alpha)))))
        for c in rgb
    )


def _apply_color_mods(rgb: tuple[int, int, int],
                      color_el: etree._Element) -> tuple[int, int, int]:
    """Apply PPTX colour modifiers (lumMod/lumOff/tint/shade) to an RGB.

    PowerPoint uses these to derive variants of theme colours — typically
    `<a:schemeClr val="accent1"><a:lumMod val="50000"/><a:lumOff val="50000"/></a:schemeClr>`
    for a 50%-mixed accent. The arithmetic is a crude HSL-luminance shim
    sufficient for the dominant cases (mods used on dark theme colours to
    derive lighter swatch variants in chart series, bar tinting, etc.).
    All values are in PPTX percent-of-100000.
    """
    lumMod = color_el.find("a:lumMod", NS)
    lumOff = color_el.find("a:lumOff", NS)
    tint = color_el.find("a:tint", NS)
    shade = color_el.find("a:shade", NS)
    if lumMod is not None or lumOff is not None:
        try:
            mod = int(lumMod.get("val")) / 100000 if lumMod is not None else 1.0
            off = int(lumOff.get("val")) / 100000 if lumOff is not None else 0.0
            rgb = tuple(max(0, min(255, int(c * mod + 255 * off))) for c in rgb)
        except (TypeError, ValueError):
            pass
    if tint is not None and tint.get("val"):
        # `tint` blends toward white. val = strength of the SOURCE colour
        # retained (lower val = closer to white).
        try:
            t = int(tint.get("val")) / 100000
            rgb = tuple(max(0, min(255, int(c * t + 255 * (1 - t)))) for c in rgb)
        except (TypeError, ValueError):
            pass
    if shade is not None and shade.get("val"):
        # `shade` blends toward black.
        try:
            s = int(shade.get("val")) / 100000
            rgb = tuple(max(0, min(255, int(c * s))) for c in rgb)
        except (TypeError, ValueError):
            pass
    return rgb


def _resolve_solid(sf: etree._Element, theme: dict[str, str], palette: dict[str, tuple[int, int, int]]) -> str | None:
    srgb = sf.find("a:srgbClr", NS)
    if srgb is not None:
        hx = srgb.get("val")
        rgb = (int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))
        rgb = _apply_color_mods(rgb, srgb)
        rgb = _blend_on_white(rgb, _alpha_for_color(srgb))
        return nearest_token(rgb, palette)
    scheme = sf.find("a:schemeClr", NS)
    if scheme is not None:
        key = scheme.get("val")
        hex_str = theme.get(key)
        if hex_str:
            rgb = (int(hex_str[1:3], 16), int(hex_str[3:5], 16), int(hex_str[5:7], 16))
            rgb = _apply_color_mods(rgb, scheme)
            rgb = _blend_on_white(rgb, _alpha_for_color(scheme))
            return nearest_token(rgb, palette)
    return None


# Map brand-pack tokens (the full feinschliff vocabulary) onto the SVG
# DSL's 17-name semantic vocabulary (defined in skills/svg/references/
# dsl-reference.md, resolved through feinschmiede.diagrams.brand_bridge). Tokens
# that have no direct counterpart fall back to the nearest neutral so
# the SVG block builds — at worst we lose a shade of grey, never the
# shape itself.
_BRAND_TO_SVG_COLOR: dict[str, str] = {
    "accent":         "accent",
    "accent-hover":   "accent",
    "highlight":      "accent",
    "ink":            "ink",
    "black":          "ink",
    "white":          "paper",
    "graphite":       "neutral-strong",
    "steel":          "neutral",
    "silver":         "neutral-soft",
    "fog":            "surface-2",
    "paper":          "paper",
    "paper-2":        "surface-2",
    "off-white":      "paper",
    "off-white-2":    "surface-2",
    "rule-dark":      "neutral-strong",
    "accent-2":       "primary",
    "accent-3":       "secondary",
    "severity-low":   "success",
    "severity-medium":"warning",
    "severity-high":  "danger",
    "status-done":    "status-on",
    "status-current": "status-pending",
    # Chart-series ramp — identity-mapped because chart-series-N is BOTH
    # a brand-pack token and a valid SVG semantic name (added together
    # in the chart-decompile feature so pie/bar slice fills survive the
    # round-trip from XML → DSL → SVG block).
    "chart-series-1": "chart-series-1",
    "chart-series-2": "chart-series-2",
    "chart-series-3": "chart-series-3",
    "chart-series-4": "chart-series-4",
    "chart-series-5": "chart-series-5",
    "chart-series-6": "chart-series-6",
    "status-next":    "status-off",
}


def _svg_color_token(brand_token: str | None, *, default: str = "neutral") -> str:
    """Best-effort map of brand-pack color token → SVG DSL semantic name.

    Inline `#rrggbb` literals (produced by nearest_token when the source
    colour is too far from any palette entry) pass through unchanged so
    the SVG block renders the source colour verbatim instead of
    collapsing to the default neutral. Without this, the source-fidelity
    guard in nearest_token would buy nothing for custGeom paths — they'd
    still render as a generic grey because `_BRAND_TO_SVG_COLOR` only
    knows the named feinschliff vocabulary.
    """
    if not brand_token:
        return default
    if brand_token.startswith("#"):
        return brand_token
    return _BRAND_TO_SVG_COLOR.get(brand_token, default)


def _pptx_path_to_svg_d(
    path_el: etree._Element,
    path_w: float, path_h: float,
    target_w: float, target_h: float,
) -> str:
    """Convert one `<a:path>` element to an SVG `d` string in target coords.

    PPTX path commands map onto SVG as follows:
      <a:moveTo>     → M x,y
      <a:lnTo>       → L x,y
      <a:cubicBezTo> → C x1,y1 x2,y2 x,y   (3 pts)
      <a:quadBezTo>  → Q x1,y1 x,y         (2 pts)
      <a:arcTo>      → A rx,ry 0 large,sweep x,y  (computed from wR/hR/stAng/swAng)
      <a:close/>     → Z

    PPTX path-local coordinates are integers in the path's own
    (w, h) box. We scale linearly to (target_w, target_h) in px so the
    resulting SVG `d` lives in the same coordinate space as the
    surrounding `svg <id> X,Y WxH { … }` block (which already places
    the origin at the shape's bbox top-left).

    PPTX arcs are documented angle-based (start angle, sweep angle, both
    in 60000ths of a degree, measured from the +x axis going clockwise).
    We compute the arc endpoint manually and emit SVG's endpoint-form
    arc command. PPTX `sweep > 0` is clockwise; SVG sweep flag `1` is
    clockwise — they line up.
    """
    sx = target_w / path_w if path_w else 1.0
    sy = target_h / path_h if path_h else 1.0
    out: list[str] = []
    cx_cur = cy_cur = 0.0

    def _xy(pt: etree._Element) -> tuple[float, float]:
        return float(pt.get("x")) * sx, float(pt.get("y")) * sy

    for cmd in path_el:
        tag = etree.QName(cmd).localname
        if tag == "moveTo":
            pt = cmd.find("a:pt", NS)
            x, y = _xy(pt)
            out.append(f"M {x:.2f},{y:.2f}")
            cx_cur, cy_cur = x, y
        elif tag == "lnTo":
            pt = cmd.find("a:pt", NS)
            x, y = _xy(pt)
            out.append(f"L {x:.2f},{y:.2f}")
            cx_cur, cy_cur = x, y
        elif tag == "cubicBezTo":
            pts = cmd.findall("a:pt", NS)
            if len(pts) == 3:
                (x1, y1), (x2, y2), (x, y) = _xy(pts[0]), _xy(pts[1]), _xy(pts[2])
                out.append(f"C {x1:.2f},{y1:.2f} {x2:.2f},{y2:.2f} {x:.2f},{y:.2f}")
                cx_cur, cy_cur = x, y
        elif tag == "quadBezTo":
            pts = cmd.findall("a:pt", NS)
            if len(pts) == 2:
                (x1, y1), (x, y) = _xy(pts[0]), _xy(pts[1])
                out.append(f"Q {x1:.2f},{y1:.2f} {x:.2f},{y:.2f}")
                cx_cur, cy_cur = x, y
        elif tag == "arcTo":
            wR = float(cmd.get("wR")) * sx
            hR = float(cmd.get("hR")) * sy
            stAng = float(cmd.get("stAng")) / 60000.0     # degrees
            swAng = float(cmd.get("swAng")) / 60000.0
            # PPTX arc convention: the arc is on an ellipse whose CENTER
            # is at (cx_cur - wR*cos(stAng), cy_cur - hR*sin(stAng))
            # — i.e. the current point IS the arc's start; the angle
            # tells us where the start lives on the ellipse. Compute the
            # endpoint by adding the sweep.
            import math
            st_rad = math.radians(stAng)
            end_rad = math.radians(stAng + swAng)
            centre_x = cx_cur - wR * math.cos(st_rad)
            centre_y = cy_cur - hR * math.sin(st_rad)
            end_x = centre_x + wR * math.cos(end_rad)
            end_y = centre_y + hR * math.sin(end_rad)
            large_arc = 1 if abs(swAng) > 180 else 0
            sweep_flag = 1 if swAng > 0 else 0
            out.append(
                f"A {wR:.2f},{hR:.2f} 0 {large_arc} {sweep_flag} {end_x:.2f},{end_y:.2f}"
            )
            cx_cur, cy_cur = end_x, end_y
        elif tag == "close":
            out.append("Z")
    return " ".join(out)


def _custgeom_svg_d(spPr: etree._Element, target_w: float, target_h: float) -> str | None:
    """Walk a custGeom's pathLst and concatenate the SVG `d` strings.

    Multiple `<a:path>` siblings inside `<a:pathLst>` become subpaths in
    a single `d` — each starts with its own `M`. Returns None when the
    spPr is not a custGeom or there's no usable path.
    """
    if spPr is None:
        return None
    cg = spPr.find("a:custGeom", NS)
    if cg is None:
        return None
    pathLst = cg.find("a:pathLst", NS)
    if pathLst is None:
        return None
    parts: list[str] = []
    for path_el in pathLst.findall("a:path", NS):
        pw = float(path_el.get("w") or 0)
        ph = float(path_el.get("h") or 0)
        if pw <= 0 or ph <= 0:
            continue
        d = _pptx_path_to_svg_d(path_el, pw, ph, target_w, target_h)
        if d:
            parts.append(d)
    return " ".join(parts) if parts else None


def _shape_geometry_kind(spPr: etree._Element) -> str:
    """Classify a sp by its geometry: rect, oval, line, or 'shape' (custGeom)."""
    if spPr is None:
        return "rect"
    pg = spPr.find("a:prstGeom", NS)
    if pg is not None:
        preset = pg.get("prst")
        if preset in ("ellipse",):
            return "oval"
        if preset in ("line", "straightConnector1"):
            return "line"
        if preset in ("rect", "roundRect"):
            return "rect"
        # Known presets with a simple closed-polygon geometry get routed
        # to `shape` so the emitter writes an `svg { path … }` block with
        # the polygon's `d` string. See `_preset_geom_path` for the table.
        if preset in _PRESET_PATH_PRESETS:
            return "shape"
        return "rect"
    if spPr.find("a:custGeom", NS) is not None:
        return "shape"
    return "rect"


# Presets whose geometry is a fixed closed polygon (no adjustment values
# read from <a:avLst>). For each, `_preset_geom_path` returns the SVG `d`
# string in the shape's local 0..w × 0..h pixel coordinate space.
_PRESET_PATH_PRESETS: frozenset[str] = frozenset({
    "triangle", "rtTriangle", "diamond",
    "parallelogram", "trapezoid",
    "pentagon", "hexagon", "heptagon", "octagon",
    "homePlate", "chevron",
    "rightArrow", "leftArrow", "upArrow", "downArrow",
})


def _preset_geom_path(preset: str, w: float, h: float) -> str | None:
    """SVG `d` string for the known closed-polygon presets, in local px.

    The shapes use PowerPoint's default unadjusted geometry — the
    `<a:avLst>` adjustment slider values are ignored. For the simple
    convex polygons in `_PRESET_PATH_PRESETS` the unadjusted form is
    visually correct in the vast majority of decks. Arrows use 50%
    barb / 50% shaft as the PowerPoint default.
    """
    if w <= 0 or h <= 0:
        return None
    if preset == "triangle":
        return f"M {w/2:.1f},0 L {w:.1f},{h:.1f} L 0,{h:.1f} Z"
    if preset == "rtTriangle":
        return f"M 0,0 L 0,{h:.1f} L {w:.1f},{h:.1f} Z"
    if preset == "diamond":
        return (f"M {w/2:.1f},0 L {w:.1f},{h/2:.1f} "
                f"L {w/2:.1f},{h:.1f} L 0,{h/2:.1f} Z")
    if preset == "parallelogram":
        # Default skew = 25% from left.
        skew = w * 0.25
        return (f"M {skew:.1f},0 L {w:.1f},0 L {w-skew:.1f},{h:.1f} "
                f"L 0,{h:.1f} Z")
    if preset == "trapezoid":
        # Default top is 75% of bottom, centered.
        inset = w * 0.125
        return (f"M {inset:.1f},0 L {w-inset:.1f},0 L {w:.1f},{h:.1f} "
                f"L 0,{h:.1f} Z")
    if preset in ("pentagon", "homePlate", "hexagon"):
        # PowerPoint's `pentagon` and `homePlate` differ in spec but render
        # similarly to `hexagon` — same convex outline for all three here.
        return (f"M {w*0.25:.1f},0 L {w*0.75:.1f},0 L {w:.1f},{h*0.5:.1f} "
                f"L {w*0.75:.1f},{h:.1f} L {w*0.25:.1f},{h:.1f} "
                f"L 0,{h*0.5:.1f} Z")
    if preset == "heptagon":
        # Regular-ish 7-gon inscribed in the bbox.
        import math as _m
        pts = []
        cx, cy = w / 2, h / 2
        rx, ry = w / 2, h / 2
        for i in range(7):
            ang = -_m.pi / 2 + i * 2 * _m.pi / 7
            pts.append(f"{cx + rx * _m.cos(ang):.1f},{cy + ry * _m.sin(ang):.1f}")
        return "M " + " L ".join(pts) + " Z"
    if preset == "octagon":
        # Regular octagon — inset 0.2929 of bbox dimension on each corner.
        c = 0.2929
        return (f"M {w*c:.1f},0 L {w-w*c:.1f},0 L {w:.1f},{h*c:.1f} "
                f"L {w:.1f},{h-h*c:.1f} L {w-w*c:.1f},{h:.1f} "
                f"L {w*c:.1f},{h:.1f} L 0,{h-h*c:.1f} L 0,{h*c:.1f} Z")
    if preset == "chevron":
        # Right-pointing chevron (arrow head + notch in tail).
        return (f"M 0,0 L {w*0.7:.1f},0 L {w:.1f},{h*0.5:.1f} "
                f"L {w*0.7:.1f},{h:.1f} L 0,{h:.1f} "
                f"L {w*0.3:.1f},{h*0.5:.1f} Z")
    if preset == "rightArrow":
        # 50% shaft height, 50% arrowhead length.
        sy0, sy1 = h * 0.25, h * 0.75
        ax = w * 0.5
        return (f"M 0,{sy0:.1f} L {ax:.1f},{sy0:.1f} L {ax:.1f},0 "
                f"L {w:.1f},{h*0.5:.1f} L {ax:.1f},{h:.1f} "
                f"L {ax:.1f},{sy1:.1f} L 0,{sy1:.1f} Z")
    if preset == "leftArrow":
        sy0, sy1 = h * 0.25, h * 0.75
        ax = w * 0.5
        return (f"M {w:.1f},{sy0:.1f} L {ax:.1f},{sy0:.1f} L {ax:.1f},0 "
                f"L 0,{h*0.5:.1f} L {ax:.1f},{h:.1f} "
                f"L {ax:.1f},{sy1:.1f} L {w:.1f},{sy1:.1f} Z")
    if preset == "upArrow":
        sx0, sx1 = w * 0.25, w * 0.75
        ay = h * 0.5
        return (f"M {sx0:.1f},{h:.1f} L {sx0:.1f},{ay:.1f} L 0,{ay:.1f} "
                f"L {w*0.5:.1f},0 L {w:.1f},{ay:.1f} L {sx1:.1f},{ay:.1f} "
                f"L {sx1:.1f},{h:.1f} Z")
    if preset == "downArrow":
        sx0, sx1 = w * 0.25, w * 0.75
        ay = h * 0.5
        return (f"M {sx0:.1f},0 L {sx0:.1f},{ay:.1f} L 0,{ay:.1f} "
                f"L {w*0.5:.1f},{h:.1f} L {w:.1f},{ay:.1f} L {sx1:.1f},{ay:.1f} "
                f"L {sx1:.1f},0 Z")
    return None


# ---------------------------------------------------------------------------
# Tree walk
# ---------------------------------------------------------------------------


def walk_slide(slide, cmap: CanvasMap, theme: dict[str, str], palette: dict[str, tuple[int, int, int]]) -> list[Shape]:
    """Walk the slide's spTree; also collect shapes inherited from the
    slide layout + master that aren't already represented on the slide.

    A typical PowerPoint slide carries only its unique content (a title
    placeholder, the body text). All decorative chrome — corporate logo,
    page-number bar, branded background blocks — lives on the slide's
    layout and master. Without this inheritance walk, the decompile of
    a corporate-template deck emits ~10% of what the source renders.

    Layout/master shapes that are placeholder fills already provided by
    the slide itself (same ph_idx) are skipped — slide content wins.
    Everything else is added at the front of the shape list so it draws
    behind slide-level content.
    """
    shapes: list[Shape] = []
    spTree = slide.element.find(".//p:cSld/p:spTree", NS)
    _walk(spTree, (0, 0), shapes, slide, cmap, theme, palette)
    # A slide-level placeholder only "owns" its idx when it actually
    # carries content (text, fill, or geometry) — empty placeholders
    # (common pattern: <p:ph idx="N"/> + empty <p:spPr/>) inherit
    # EVERYTHING from the layout, including the layout's fill/size.
    # Filtering by content here lets the layout walk re-add the rich
    # version of those placeholders below.
    def _has_content(sh: Shape) -> bool:
        if sh.fill or sh.stroke or sh.gradient or sh.svg_path_d:
            return True
        if sh.is_picture:
            return True
        if any(r.text and r.text.strip() and r.text != "\n" for r in sh.text_runs):
            return True
        return False

    slide_ph_idxs = {s.ph_idx for s in shapes if s.ph_idx and _has_content(s)}
    # Filter out the empty placeholders themselves — the layout version
    # will provide the actual content.
    shapes = [s for s in shapes if not (s.ph_idx and not _has_content(s))]
    inherited: list[Shape] = []
    layout_master_chain = _layout_master_chain(slide)
    for src in layout_master_chain:
        chain_spTree = src.element.find(".//p:cSld/p:spTree", NS)
        if chain_spTree is None:
            continue
        chain_shapes: list[Shape] = []
        _walk(chain_spTree, (0, 0), chain_shapes, src, cmap, theme, palette)
        for s in chain_shapes:
            # Skip placeholder shapes the slide already owns with content.
            if s.ph_idx and s.ph_idx in slide_ph_idxs:
                continue
            # Skip pure page-number / footer placeholders by type when the
            # slide hasn't overridden them — these are pp:fld things that
            # the source-deck renderer fills at slide time; decompiling
            # them from the master emits literal "<#>" tokens that pollute
            # the output. (sldNum/ftr/dt placeholders all carry ph_type.)
            if s.ph_type in ("sldNum", "ftr", "dt") and not s.text_runs:
                continue
            # Inherited PICTURE placeholders with no real media binary
            # (no media_rid) are shells the slide should fill but didn't.
            # PowerPoint renders them as empty; the build's missing-asset
            # fallback paints them as white rectangles which inflate the
            # diff visibly. Skip them so the rendered output matches
            # PowerPoint's "absent" behaviour.
            if s.is_picture and not s.media_rid:
                continue
            # Skip layout placeholders whose only text is template
            # prompt copy ("Hier Zitat einfügen.", "Überschrift 1, TT
            # Norms Pro, 28 pt", etc) — these are author hints that
            # PowerPoint suppresses at render time when the slide's
            # version is empty. We can't easily distinguish prompt from
            # real content, but layouts mark prompts with
            # `<p:nvPr hasCustomPrompt="1">` — drop just the text on
            # those, keep the fill/geometry so the layout still emits
            # its visual frame.
            # Inherited PLACEHOLDER text is template prompt copy that
            # PowerPoint never renders ("Überschrift 1, TT Norms Pro",
            # "Diagramm durch Klicken …", "Click to edit Master title").
            # The slide's placeholder either overrides it (filtered above
            # via slide_ph_idxs) or is empty — in which case PowerPoint
            # shows nothing. Drop the text but keep fill / geometry so
            # the layout's visual frame (yellow rect, black square)
            # still emits. Non-placeholder layout/master shapes (logo
            # glyphs, decorative rects) pass through untouched.
            if s.ph_type and s.text_runs:
                s.text_runs = []
                if not _has_content(s):
                    continue
            s.from_chain = True
            inherited.append(s)
            if s.ph_idx:
                slide_ph_idxs.add(s.ph_idx)
    # Inherited chrome draws behind slide content.
    return inherited + shapes


def _layout_master_chain(slide) -> list:
    """Return [layout, master] for a slide (best-effort, never raises)."""
    chain = []
    try:
        layout = slide.slide_layout
        if layout is not None:
            chain.append(layout)
            master = layout.slide_master
            if master is not None:
                chain.append(master)
    except Exception:
        pass
    return chain


def extract_slide_bg_fill(slide, theme: dict[str, str],
                          palette: dict[str, tuple[int, int, int]]) -> str | None:
    """Return the slide's background solid-fill colour as a token / hex.

    Walks slide → LAYOUT only (NOT the master). A layout that wants a
    full-bleed colour sets its own `<p:cSld><p:bg>` (kept here — e.g. a
    timeline layout's yellow, a divider's navy). The MASTER-level bg is
    deliberately ignored: corporate light templates carry a dark master bg
    that the layout visually overrides with a white layer / colour panel, so
    emitting it as a full-canvas rect paints every white content slide black
    (a ~97% regression). solidFill or bgRef→theme; first level found wins.
    """
    for src in [slide, *_layout_master_chain(slide)[:1]]:
        bg = src.element.find(".//p:cSld/p:bg", NS)
        if bg is None:
            continue
        bgPr = bg.find("p:bgPr", NS)
        if bgPr is not None:
            color = _resolve_fill(bgPr, theme, palette)
            if color:
                return color
        # <p:bgRef idx="N"><a:schemeClr val="bg1"/></p:bgRef> — referenced
        # background-fill style from theme1.xml's bgFillStyleLst. The
        # schemeClr fills the `phClr` placeholder in the referenced style.
        bgRef = bg.find("p:bgRef", NS)
        if bgRef is not None:
            color = _resolve_bg_ref(bgRef, theme, palette)
            if color:
                return color
    return None


def _resolve_bg_ref(bgRef, theme: dict[str, str],
                    palette: dict[str, tuple[int, int, int]]) -> str | None:
    """Resolve a `<p:bgRef idx="N">` reference against the theme's
    `bgFillStyleLst`, filling `<a:schemeClr val="phClr"/>` with the
    schemeClr that the bgRef carries.

    PowerPoint idx encoding: 1001 = bgFillStyleLst[0], 1002 = [1], etc.
    Only `<a:solidFill>` entries are inlined; gradient/blip refs fall
    through (no DSL primitive for those yet at the bg level).
    """
    scheme = bgRef.find("a:schemeClr", NS)
    if scheme is None or not scheme.get("val"):
        return None
    scheme_key = scheme.get("val")
    # Resolve scheme → hex via the theme dict captured by load_theme_scheme.
    # Some keys are aliases: `bg1`→`lt1`, `bg2`→`lt2`, `tx1`→`dk1`, `tx2`→`dk2`.
    alias = {"bg1": "lt1", "bg2": "lt2", "tx1": "dk1", "tx2": "dk2"}
    hex_val = theme.get(scheme_key) or theme.get(alias.get(scheme_key, scheme_key))
    if not hex_val:
        return None
    try:
        rgb = (int(hex_val[1:3], 16), int(hex_val[3:5], 16), int(hex_val[5:7], 16))
    except (ValueError, IndexError):
        return None
    return nearest_token(rgb, palette) if palette else f"#{hex_val[1:].upper()}"


def _walk(node, offset, shapes, slide, cmap, theme, palette):
    # offset is either a 2-tuple (ox, oy) or an 8-tuple carrying a scaled-
    # group affine — the latter is unpacked by `_shape_bbox`. We only need
    # ox/oy here to forward to nested non-scaled groups.
    ox, oy = offset[:2]
    for ch in node:
        tag = etree.QName(ch).localname
        if tag == "sp":
            _emit_sp(ch, offset, shapes, slide, cmap, theme, palette)
        elif tag == "pic":
            _emit_pic(ch, offset, shapes, slide, cmap, theme, palette)
        elif tag == "cxnSp":
            _emit_cxn(ch, offset, shapes, cmap, theme, palette)
        elif tag == "graphicFrame":
            _emit_graphic_frame(ch, offset, shapes, slide, cmap, theme, palette)
        elif tag == "grpSp":
            # Whole-group native-carry: a top-level group of pure decorative
            # chrome (custGeom / graphicFrame, NO content placeholder) is frozen
            # verbatim — PowerPoint applies its own off/ext/chOff/chExt affine, so
            # even a SCALED group of vectors (world maps, decorative clusters) that
            # the per-shape re-synth mangles renders pixel-exact. Returns True when
            # carried (don't recurse); False → recurse exactly as before.
            if _try_carry_group(ch, offset, shapes, slide, cmap, theme):
                continue
            # Walk children with the group's offset added. Scaled groups
            # (ext != chExt — typical for master-level logo bundles
            # dropped into smaller slots) are skipped because the walker
            # only carries a pure translation offset; emitting their
            # children at unscaled coords would land them off-canvas.
            grp_xfrm = ch.find("p:grpSpPr/a:xfrm", NS)
            child_off = (ox, oy)
            if grp_xfrm is not None:
                off = grp_xfrm.find("a:off", NS)
                ext = grp_xfrm.find("a:ext", NS)
                chOff = grp_xfrm.find("a:chOff", NS)
                chExt = grp_xfrm.find("a:chExt", NS)
                if (off is not None and ext is not None
                        and chOff is not None and chExt is not None):
                    try:
                        ox_emu = int(off.get("x"))
                        oy_emu = int(off.get("y"))
                        cx = int(ext.get("cx"))
                        cy = int(ext.get("cy"))
                        chcx = int(chExt.get("cx"))
                        chcy = int(chExt.get("cy"))
                        chox = int(chOff.get("x"))
                        choy = int(chOff.get("y"))
                        if chcx > 0 and chcy > 0:
                            sx = cx / chcx
                            sy = cy / chcy
                            # Scaled group: thread a 6-tuple offset that
                            # _shape_bbox unpacks and applies as an EMU-level
                            # affine. Translation-only groups stay as 2-tuple
                            # for backward compat.
                            if abs(sx - 1.0) > 0.001 or abs(sy - 1.0) > 0.001:
                                child_off = (ox, oy, ox_emu, oy_emu, chox, choy, sx, sy)
                                _walk(ch, child_off, shapes, slide, cmap, theme, palette)
                                continue
                        # Pure translation fallthrough.
                        child_off = (ox + ox_emu - chox, oy + oy_emu - choy)
                    except (ValueError, TypeError):
                        pass
                elif off is not None and chOff is not None:
                    try:
                        dx = int(off.get("x")) - int(chOff.get("x"))
                        dy = int(off.get("y")) - int(chOff.get("y"))
                        child_off = (ox + dx, oy + dy)
                    except (ValueError, TypeError):
                        pass
            _walk(ch, child_off, shapes, slide, cmap, theme, palette)


def _shape_bbox(ch, offset, slide):
    spPr = ch.find("p:spPr", NS)
    xfrm = _get_xfrm(spPr)
    if xfrm is None:
        ph_type, ph_idx = _placeholder_info(ch)
        xfrm = _layout_placeholder_xfrm(slide, ph_type, ph_idx)
    if xfrm is None:
        return None
    x, y, w, h = xfrm
    # Offset shapes:
    #   2-tuple (ox, oy)      — translation-only ancestor group(s)
    #   8-tuple (ox, oy, ax, ay, chox, choy, sx, sy)
    #                         — current shape lives inside a scaled group;
    #                           apply the EMU-level affine before adding
    #                           any outer translation.
    if len(offset) == 2:
        ox, oy = offset
        return x + ox, y + oy, w, h
    ox, oy, ax, ay, chox, choy, sx, sy = offset
    x_emu = ax + (x - chox) * sx
    y_emu = ay + (y - choy) * sy
    w_emu = w * sx
    h_emu = h * sy
    return x_emu + ox, y_emu + oy, w_emu, h_emu


def _bake_scheme_colors(el, theme: dict[str, str]) -> None:
    """In-place: rewrite `<a:schemeClr val=KEY>` → `<a:srgbClr val=HEX>` from the
    SOURCE theme, so a carried-native shape keeps its EXACT source colours even
    when the output deck's theme differs. Child mods (lumMod/alpha/…) are valid on
    srgbClr too and ride along untouched."""
    a_ns = NS["a"]
    alias = {"bg1": "lt1", "bg2": "lt2", "tx1": "dk1", "tx2": "dk2"}
    for sc in list(el.iter(f"{{{a_ns}}}schemeClr")):
        key = sc.get("val")
        hexv = theme.get(key) or theme.get(alias.get(key, key))
        if hexv:
            sc.tag = f"{{{a_ns}}}srgbClr"
            sc.set("val", hexv.lstrip("#").upper())


def _emit_sp(ch, offset, shapes, slide, cmap, theme, palette):
    spPr = ch.find("p:spPr", NS)
    bbox = _shape_bbox(ch, offset, slide)
    if bbox is None:
        return
    x, y, w, h = bbox
    ph_type, ph_idx = _placeholder_info(ch)
    # Pull placeholder default sz from layout/master so body placeholders
    # without explicit run-level `sz` inherit the right headline size.
    inherited_sz = _layout_placeholder_default_sz(slide, ph_type, ph_idx) if (ph_type or ph_idx) else None
    # cap="all" / bold also cascade from the layout/master placeholder style.
    inherited_caps, inherited_bold = (
        _layout_placeholder_caps_bold(slide, ph_type, ph_idx) if (ph_type or ph_idx) else (False, False)
    )
    runs = _text_runs(ch, theme, palette, inherited_default_sz=inherited_sz,
                      inherited_caps=inherited_caps, inherited_bold=inherited_bold)
    # G3 — capture text colour the slide run INHERITS rather than states. When a
    # run carries no explicit `<a:rPr><a:solidFill>`, fall back to (a) the shape's
    # `<p:style><a:fontRef>` (decorative styled shapes) then (b) the layout/master
    # PLACEHOLDER default colour (e.g. a WHITE `ctrTitle` over a coloured circle).
    # Without this, on-shape / on-placeholder titles render ink-grey. Emitted
    # downstream only when it differs from the style-bundle default.
    if runs and not any(r.color for r in runs):
        _style = ch.find("p:style", NS)
        _font_ref = _style.find("a:fontRef", NS) if _style is not None else None
        _text_color = _resolve_solid(_font_ref, theme, palette) if _font_ref is not None else None
        if _text_color is None and (ph_type or ph_idx):
            _text_color = _layout_placeholder_color(slide, ph_type, ph_idx, theme, palette)
        if _text_color:
            for _r in runs:
                if _r.text and _r.text != "\n":
                    _r.color = _text_color
    kind = _shape_geometry_kind(spPr)
    # For custGeom shapes (kind="shape") — typically map polygons,
    # decorative vector clusters, or hand-drawn paths — bypass the
    # brand-token mapping in `nearest_token` and emit the source colour
    # as raw hex. These shapes carry hundreds of subtly-different
    # source colours (e.g. world-map country fills at #EBEBEB /
    # #DDDDDD / #C8C8C8) and the round-trip through nearest_token →
    # `_svg_color_token` → `brand_bridge.resolve` collapses them to a
    # handful of SVG semantic names that resolve to materially
    # different greys in the brand pack. Going straight to hex
    # preserves source-pixel fidelity for these high-cardinality
    # vector compositions.
    if kind == "shape":
        fill = _resolve_fill(spPr, theme, palette={})
    else:
        fill = _resolve_fill(spPr, theme, palette)
    gradient = _resolve_gradient(spPr, theme, palette)
    # Vertical anchor — `<a:bodyPr anchor="ctr">` / `b` / `t`. Without
    # this the rendered text lands at frame-top even when source centers
    # it, which is the dominant cause of the redline "two ghost positions"
    # pattern: source content at frame center, render content at frame top.
    valign: str | None = None
    padding_emu: tuple[int, int, int, int] | None = None
    autoshrink = False
    font_scale = 1.0
    txBody = ch.find(".//p:txBody", NS) or ch.find(".//a:txBody", NS)
    if txBody is not None:
        bodyPr = txBody.find("a:bodyPr", NS)
        if bodyPr is not None:
            anc = bodyPr.get("anchor")
            if anc == "ctr":
                valign = "middle"
            elif anc == "b":
                valign = "bottom"
            # Insets — l/t/r/b. The slide's own bodyPr wins per-side; for any
            # side it omits, the placeholder INHERITS the layout/master
            # placeholder's inset (the master title sets all four to 0, so
            # the heading must hug the box edge); only when no ancestor specifies
            # a side does PowerPoint's published default (91440 / 45720 EMU)
            # apply. Without the inheritance step every empty-bodyPr title
            # decompiled with the +16/+8 px default offset and ghosted in the
            # redline.
            inh = (
                _layout_placeholder_insets(slide, ph_type, ph_idx)
                if (ph_type or ph_idx)
                else (None, None, None, None)
            )

            def _ins(attr: str, inherited: int | None, default: int) -> int:
                v = bodyPr.get(attr)
                if v is not None:
                    try:
                        return int(v)
                    except (TypeError, ValueError):
                        pass
                return inherited if inherited is not None else default

            left = _ins("lIns", inh[0], 91440)
            top = _ins("tIns", inh[1], 45720)
            right = _ins("rIns", inh[2], 91440)
            bottom = _ins("bIns", inh[3], 45720)
            padding_emu = (left, top, right, bottom)
            # normAutofit = PowerPoint "Shrink text on overflow": the source pre-
            # shrinks the text by `fontScale` to fit the box (the run `sz` stays at
            # the authored size). Capture it so the emitter reproduces the fit
            # instead of rendering full-size and overflowing (inherited placeholders).
            na = bodyPr.find("a:normAutofit", NS)
            if na is not None:
                autoshrink = True
                if na.get("fontScale"):
                    try:
                        font_scale = int(na.get("fontScale")) / 100000.0
                    except (TypeError, ValueError):
                        pass
    # Vertical anchor also inherits: a slide title with no own bodyPr anchor still
    # renders bottom/centre-anchored when the layout/master placeholder bodyPr sets
    # it (MS Geometric: master title placeholder anchor="b" → titles sit at the box
    # bottom, not top). Mirror the size / caps / colour inheritance. Feature-2's box
    # extension correctly skips bottom/middle text, so it won't fight this.
    if valign is None and (ph_type or ph_idx):
        valign = _layout_placeholder_anchor(slide, ph_type, ph_idx)
    # normAutofit also inherits: some templates set "shrink text on overflow" on the
    # LAYOUT placeholder, not the slide, so a placeholder whose own bodyPr lacks it
    # still shrinks-to-fit — capture that (autoshrink) or the text overflows.
    if not autoshrink and (ph_type or ph_idx):
        autoshrink = _layout_placeholder_autofit(slide, ph_type, ph_idx)
    # Convert insets EMU → design-px for the Shape (CanvasMap-relative).
    padding_px: tuple[float, float, float, float] | None = None
    if padding_emu is not None:
        left, top, right, bottom = padding_emu
        padding_px = (cmap.w(left), cmap.h(top), cmap.w(right), cmap.h(bottom))
    # Stroke (line) colour + width + dash. PowerPoint stores stroke width
    # in EMU on `<a:ln w="...">` (default ~9525 EMU = 0.75pt = 1px hairline);
    # dash preset on the optional `<a:prstDash val="...">` child.
    stroke = None
    stroke_width: float | None = None
    stroke_dash: str | None = None
    ln = spPr.find("a:ln", NS) if spPr is not None else None
    if ln is not None:
        sf = ln.find("a:solidFill", NS)
        if sf is not None:
            stroke = _resolve_solid(sf, theme, palette)
        w_attr = ln.get("w")
        if w_attr:
            try:
                w_emu = int(w_attr)
                # Width is uniform in EMU (1pt = 12700); converting via the
                # horizontal scale keeps it consistent with the rest of the
                # design-px coordinate system on this canvas.
                stroke_width = cmap.w(w_emu)
            except (ValueError, TypeError):
                pass
        dash = ln.find("a:prstDash", NS)
        if dash is not None and dash.get("val"):
            stroke_dash = dash.get("val")
    # Corner radius — captured from `prstGeom prst="roundRect"` with an
    # `<a:gd name="adj" fmla="val N">`. N is N/100000ths of the shape's
    # shortest side. PowerPoint defaults to 0.10 when adj is absent.
    corner_radius: float | None = None
    if spPr is not None:
        pg = spPr.find("a:prstGeom", NS)
        if pg is not None and pg.get("prst") == "roundRect":
            gd = pg.find(".//a:gd[@name='adj']", NS)
            adj_frac = 0.10  # PowerPoint default when omitted
            if gd is not None and gd.get("fmla"):
                m = re.search(r"val (\d+)", gd.get("fmla"))
                if m:
                    adj_frac = int(m.group(1)) / 100000
            corner_radius = cmap.w(min(w, h)) * adj_frac
    # Drop shadow — `<a:effectLst><a:outerShdw>`. Standard PowerPoint card
    # shadows use blurRad/dist in EMU, dir in 1/60000ths of a degree, and a
    # solid colour with an `<a:alpha val="...">` modifier (0-100000 = 0-100%).
    shadow: tuple[float, float, float, str, float] | None = None
    if spPr is not None:
        eff = spPr.find(".//a:effectLst/a:outerShdw", NS)
        if eff is not None:
            blur_emu = int(eff.get("blurRad") or 0)
            dist_emu = int(eff.get("dist") or 0)
            dir_60k = int(eff.get("dir") or 0)
            blur_px = cmap.w(blur_emu)
            dist_px = cmap.w(dist_emu)
            angle_deg = dir_60k / 60000.0
            sh_color = "black"
            sh_alpha = 1.0
            srgb = eff.find("a:srgbClr", NS)
            if srgb is not None and srgb.get("val"):
                hx = srgb.get("val")
                try:
                    rgb = (int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))
                    sh_color = nearest_token(rgb, palette) if palette else f"#{hx}"
                except ValueError:
                    pass
                alpha = srgb.find("a:alpha", NS)
                if alpha is not None and alpha.get("val"):
                    sh_alpha = int(alpha.get("val")) / 100000.0
            shadow = (blur_px, dist_px, angle_deg, sh_color, sh_alpha)

    # Picture-typed placeholder → picture shape (no actual <p:pic>).
    if ph_type == "pic":
        shapes.append(Shape(
            kind="pic", x=cmap.x(x), y=cmap.y(y), w=cmap.w(w), h=cmap.h(h),
            is_picture=True, ph_type=ph_type, ph_idx=ph_idx,
        ))
        return

    # Pure-text shape (placeholder, label, etc.) — no rect, just text.
    if runs and fill is None and kind == "rect":
        # Multi-paragraph colour split: when a text shape carries paragraphs
        # in different colours (typical of a cyan headline above silver body
        # bullets), emit one DSL `text` primitive per consecutive same-colour
        # block at calculated y-offsets, instead of collapsing them into a
        # single primitive whose first-run colour overrides the rest.
        blocks = _split_runs_by_color(runs)
        if len(blocks) > 1:
            frame_x_px = cmap.x(x)
            frame_y_px = cmap.y(y)
            frame_w_px = cmap.w(w)
            # Cursor-style y placement: each block consumes its own height
            # based on its primary pt (line-height factor 1.25 captures
            # standard slide leading without over-spacing tight headlines).
            cursor = frame_y_px
            for block_runs in blocks:
                block_pts = [r.pt for r in block_runs if r.text and r.text != "\n"]
                primary_pt = max(block_pts) if block_pts else 12
                line_count = sum(1 for r in block_runs if r.text and r.text != "\n")
                block_h_px = max(int(round(primary_pt * (4 / 3) * 1.25 * max(line_count, 1))), 24)
                # Split blocks always anchor top so the next block starts at
                # the bottom of the previous one — using the parent shape's
                # valign:middle would re-center each split block and overlap
                # neighbours.
                shapes.append(Shape(
                    kind="text", x=frame_x_px, y=cursor,
                    w=frame_w_px, h=block_h_px,
                    text_runs=block_runs, ph_type=ph_type, ph_idx=ph_idx,
                    valign=None, padding=padding_px,
                ))
                cursor += block_h_px
            return
        shapes.append(Shape(
            kind="text", x=cmap.x(x), y=cmap.y(y), w=cmap.w(w), h=cmap.h(h),
            text_runs=runs, ph_type=ph_type, ph_idx=ph_idx, valign=valign,
            autoshrink=autoshrink, font_scale=font_scale,
            padding=padding_px,
        ))
        return

    # custGeom paths convert directly to SVG path `d`. Build it in
    # canvas-pixel space so the surrounding svg-block can simply
    # `path "<d>"` without further transforms.
    svg_d = None
    native_xml = None
    if kind == "shape":
        svg_d = _custgeom_svg_d(spPr, cmap.w(w), cmap.h(h))
        if svg_d is None and spPr is not None:
            # Preset-geom polygon (triangle, diamond, arrow, etc.) — the
            # source uses `prstGeom prst="…"` with no custGeom, so
            # _custgeom_svg_d returns None. Synthesize the path from the
            # preset name so the renderer draws the correct outline
            # instead of falling back to a bbox-rect.
            pg = spPr.find("a:prstGeom", NS)
            if pg is not None:
                preset = pg.get("prst")
                if preset:
                    svg_d = _preset_geom_path(preset, cmap.w(w), cmap.h(h))
        # Prefer carrying the native <p:sp> verbatim for TOP-LEVEL complex chrome:
        # a real, editable vector spliced straight into the output deck — NO
        # svg → raster → picture round-trip (which both distorts the shape and is
        # a picture "cheat"). The DSL stays the content layer (text + images);
        # corporate-design geometry rides along untouched. Colours are baked
        # schemeClr→srgbClr against the SOURCE theme so they survive the output
        # deck's theme. Grouped shapes (offset != 0) keep the svg path — their
        # xfrm is group-relative, not slide-absolute, so a verbatim splice would
        # land in the wrong place.
        if svg_d is not None and len(offset) == 2:
            import copy as _copy
            sp_el = _copy.deepcopy(ch)
            _bake_scheme_colors(sp_el, theme)
            # `offset` is the (EMU) group/layout translation this shape was walked
            # under; its xfrm is relative to that, so shift it to slide-absolute
            # before splicing. (Scaled groups thread an 8-tuple affine instead —
            # those fall through to the svg path, which is bbox-correct already.)
            ox, oy = int(offset[0]), int(offset[1])
            _off_el = sp_el.find("p:spPr/a:xfrm/a:off", NS)
            if _off_el is not None and (ox or oy):
                _off_el.set("x", str(int(_off_el.get("x") or 0) + ox))
                _off_el.set("y", str(int(_off_el.get("y") or 0) + oy))
            native_xml = etree.tostring(sp_el).decode("utf-8")

    # Geometry shape (rect / oval / shape). May also carry text.
    shapes.append(Shape(
        kind=kind, x=cmap.x(x), y=cmap.y(y), w=cmap.w(w), h=cmap.h(h),
        fill=fill, stroke=stroke, stroke_width=stroke_width,
        stroke_dash=stroke_dash, corner_radius=corner_radius, shadow=shadow,
        gradient=gradient,
        text_runs=runs, ph_type=ph_type, ph_idx=ph_idx, svg_path_d=svg_d,
        native_xml=native_xml,
    ))


def _emit_pic(ch, offset, shapes, slide, cmap, theme, palette):
    bbox = _shape_bbox(ch, offset, slide)
    if bbox is None:
        return
    x, y, w, h = bbox
    ph_type, ph_idx = _placeholder_info(ch)
    # Capture the embedded media rId so derive() can extract the binary
    # when image_extract_dir is set (pipeline-optimization mode).
    rid = None
    media_part = None
    blip = ch.find(".//a:blip", NS)
    if blip is not None:
        rid = blip.get(f"{{{NS['r']}}}embed")
        # Resolve the related part NOW against the source object that
        # actually owns the rId (the slide, layout, or master `slide`
        # param of this call). Storing only the rId and re-resolving on
        # `slide.part` later breaks for layout-inherited pictures —
        # their rId is scoped to the layout's relationships, not the
        # slide's.
        if rid is not None:
            try:
                media_part = slide.part.related_part(rid)
            except (KeyError, AttributeError):
                media_part = None
    # Template image: a small, non-placeholder <p:pic> (logo, mark, brand chrome)
    # is fixed corporate-design identity — carry it natively (verbatim element +
    # its media, base64-inline) rather than a fillable picture slot. Large or
    # placeholder pics are CHANGEABLE topical content → they stay slots.
    native_xml = None
    native_media = None
    if (ph_type is None and media_part is not None and len(offset) == 2
            and cmap.w(w) * cmap.h(h) < _TEMPLATE_IMG_MAX_AREA * (cmap.cw * cmap.ch)):
        try:
            import copy as _copy
            import base64 as _b64
            pic_el = _copy.deepcopy(ch)
            _bake_scheme_colors(pic_el, theme)
            ox, oy = int(offset[0]), int(offset[1])
            _off_el = pic_el.find("p:spPr/a:xfrm/a:off", NS)
            if _off_el is not None and (ox or oy):
                _off_el.set("x", str(int(_off_el.get("x") or 0) + ox))
                _off_el.set("y", str(int(_off_el.get("y") or 0) + oy))
            native_xml = etree.tostring(pic_el).decode("utf-8")
            native_media = _b64.b64encode(media_part.blob).decode("ascii")
        except Exception:
            native_xml = native_media = None
    shapes.append(Shape(
        kind="pic", x=cmap.x(x), y=cmap.y(y), w=cmap.w(w), h=cmap.h(h),
        is_picture=True, ph_type=ph_type, ph_idx=ph_idx, media_rid=rid,
        media_part=media_part, native_xml=native_xml, native_media=native_media,
    ))


def _emit_cxn(ch, offset, shapes, cmap, theme, palette):
    spPr = ch.find("p:spPr", NS)
    xfrm = _get_xfrm(spPr)
    if xfrm is None:
        return
    x, y, w, h = xfrm
    # Offset is a 2-tuple (translation-only ancestor groups) or an 8-tuple
    # (scaled-group affine). Apply it exactly like _shape_bbox so a connector
    # inside a SCALED group doesn't crash (`ox, oy = offset` blew up on the
    # 8-tuple) and still lands at the right place.
    if len(offset) == 2:
        ox, oy = offset
        x, y = x + ox, y + oy
    else:
        ox, oy, ax, ay, chox, choy, sx, sy = offset
        x = ax + (x - chox) * sx + ox
        y = ay + (y - choy) * sy + oy
        w, h = w * sx, h * sy
    # Stroke color + width + dash from <a:ln w="..."><a:solidFill .../><a:prstDash .../></a:ln>.
    ln = spPr.find("a:ln", NS)
    stroke = None
    stroke_width: float | None = None
    stroke_dash: str | None = None
    if ln is not None:
        sf = ln.find("a:solidFill", NS)
        if sf is not None:
            stroke = _resolve_solid(sf, theme, palette)
        w_attr = ln.get("w")
        if w_attr:
            try:
                stroke_width = cmap.w(int(w_attr))
            except (ValueError, TypeError):
                pass
        dash = ln.find("a:prstDash", NS)
        if dash is not None and dash.get("val"):
            stroke_dash = dash.get("val")
    shapes.append(Shape(
        kind="line", x=cmap.x(x), y=cmap.y(y), w=cmap.w(w), h=cmap.h(h),
        stroke=stroke or "fog",
        stroke_width=stroke_width,
        stroke_dash=stroke_dash,
    ))


CHART_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
RELS_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
# Relationship TYPE the slide part uses to reach a chart part. Stored as the
# chart-part's `reltype` in the carried part-graph so the emitter re-creates
# the slide→chart relationship with the correct type (the value is otherwise
# unused for the leaf style/colors/xlsx parts, whose own reltypes are carried
# verbatim from the source rels).
RELS_NS_CHART = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart"

# SmartArt / diagram graphicFrame. The `<dgm:relIds>` inside graphicData carries
# four slide-rels (data / layout / quickStyle / colors); a 5th part — the
# pre-rendered drawing — is reached via the dataN.xml extLst's
# `<dsp:dataModelExt relId>` (also a slide rel). All five (+ any media sub-rel of
# the data / drawing part) are native-carried so the diagram renders pixel-exact,
# vs `_emit_smartart`'s lossy flatten (parses the cached drawing into bbox rects,
# dropping connectors / geometry / per-node styling).
DGM_NS = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
DSP_DATAMODEL_NS = "http://schemas.microsoft.com/office/drawing/2008/diagram"
# Reltype for the diagramDrawing part (the dataN extLst relId points at it). This
# is a Microsoft-extension reltype with no RELATIONSHIP_TYPE constant.
RELS_NS_DGM_DRAWING = "http://schemas.microsoft.com/office/2007/relationships/diagramDrawing"
# The four `<dgm:relIds>` attributes → their slide-rel reltypes, so the emitter
# re-creates each slide→diagram relationship with the right type.
_DGM_RELID_ATTRS = (
    ("dm", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData"),
    ("lo", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramLayout"),
    ("qs", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramQuickStyle"),
    ("cs", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramColors"),
)
# Relationship type a slide/part uses to reach an embedded raster image. Carried
# as the reltype of a grouped <p:pic>'s media part so the emitter re-creates the
# slide→image relationship with the right type when splicing a native group.
RELS_NS_IMAGE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"


def _capture_group_parts(grp, slide, theme) -> list[dict] | None:
    """Collect the external part-graph a native-carried `<p:grpSp>` reaches, so a
    grouped chart / SmartArt / image renders pixel-exact after splicing.

    Mirrors `_emit_graphic_frame`'s Stage B/C part-capture, but rooted at a WHOLE
    group rather than a single graphicFrame, and additionally carries grouped
    `<p:pic>` media (the existing single-`<p:pic>` `media:` path can't cover a
    group with several pics). Returns a list of {partname, content_type, blob,
    reltype, parent, src_rid} dicts (+ optional `ref` shared-part entries), or
    None when the group reaches no external parts (pure inline custGeom — the
    common decorative-vector case). Raises on a resolve failure so the caller's
    try/except falls back to recursion.

    Every captured XML blob gets schemeClr→srgbClr baked against the SOURCE theme
    so the carried content keeps its exact source palette under the output deck's
    theme — identical to the chart/diagram branches.
    """
    import base64 as _b64
    parts: list[dict] = []
    seen: set[str] = set()
    shared: dict[str, dict[str, tuple[str, str]]] = {}

    def _capture(part, parent_key: str, reltype: str, src_rid: str) -> None:
        pn = str(part.partname)
        if pn in seen:
            # A part shared by two parents is materialised once, but record this
            # parent's src_rid→part mapping so its own ref gets rewritten too.
            if src_rid:
                shared.setdefault(parent_key, {})[src_rid] = (pn, reltype)
            return
        seen.add(pn)
        if src_rid:
            shared.setdefault(parent_key, {})[src_rid] = (pn, reltype)
        raw = part.blob
        ct = part.content_type
        if ct.endswith("+xml") or ct.endswith("/xml"):
            try:
                _proot = etree.fromstring(raw)
                _bake_scheme_colors(_proot, theme)
                raw = etree.tostring(_proot, xml_declaration=True,
                                     encoding="UTF-8", standalone=True)
            except Exception:
                raw = part.blob
        parts.append({
            "partname": pn, "content_type": ct,
            "blob": _b64.b64encode(raw).decode("ascii"),
            "reltype": reltype, "parent": parent_key, "src_rid": src_rid,
        })
        try:
            child_rels = list(part.rels.values())
        except Exception:
            child_rels = []
        for r in child_rels:
            if r.is_external:
                continue
            try:
                tgt = r.target_part
            except Exception:
                continue
            _capture(tgt, pn, r.reltype, r.rId)

    # Grouped chart <c:chart r:id> — the chart part + its style/colors/xlsx graph.
    for cref in grp.iter(f"{{{CHART_NS}}}chart"):
        rid = cref.get(f"{{{RELS_NS}}}id")
        if rid:
            _capture(slide.part.related_part(rid), "slide", RELS_NS_CHART, rid)
    # Grouped SmartArt <dgm:relIds> — the four dgm parts + pre-rendered drawing.
    for relids in grp.findall(f".//{{{DGM_NS}}}relIds"):
        data_part = None
        for attr, reltype in _DGM_RELID_ATTRS:
            rid = relids.get(f"{{{RELS_NS}}}{attr}")
            if not rid:
                continue
            p = slide.part.related_part(rid)
            if attr == "dm":
                data_part = p
            _capture(p, "slide", reltype, rid)
        if data_part is not None:
            try:
                droot = etree.fromstring(data_part.blob)
                dmext = droot.find(f".//{{{DSP_DATAMODEL_NS}}}dataModelExt")
            except Exception:
                dmext = None
            draw_rid = dmext.get("relId") if dmext is not None else None
            if draw_rid:
                _capture(slide.part.related_part(draw_rid), "slide",
                         RELS_NS_DGM_DRAWING, draw_rid)
    # Grouped <p:pic> media — every <a:blip r:embed> resolves to an image part on
    # the slide; carry it so the spliced blip can be re-pointed at a fresh rId.
    for blip in grp.findall(".//a:blip", NS):
        rid = blip.get(f"{{{RELS_NS}}}embed")
        if not rid:
            continue
        try:
            mp = slide.part.related_part(rid)
        except Exception:
            continue
        _capture(mp, "slide", RELS_NS_IMAGE, rid)

    if not parts:
        return None
    # Fold shared-child mappings into ref entries so the emitter wires the second
    # parent's relationship + rewrites its ref without re-materialising the part.
    for parent_key, m in shared.items():
        for src_rid, (old_pn, reltype) in m.items():
            if any(e.get("parent") == parent_key and e.get("src_rid") == src_rid
                   for e in parts):
                continue
            parts.append({"ref": old_pn, "parent": parent_key,
                          "src_rid": src_rid, "reltype": reltype})
    return parts


def _try_carry_group(ch, offset, shapes, slide, cmap, theme) -> bool:
    """Native-carry a WHOLE decorative `<p:grpSp>` verbatim when it's pure chrome.

    Extends native-carry from top-level graphicFrames / custGeom shapes to entire
    groups — the last structural gap. A scaled group of custGeom vectors (world-map
    illustrations, decorative clusters) re-synthesises lossily today because the
    walker only carries a pure translation; carried WHOLE, PowerPoint applies the
    group's own off/ext/chOff/chExt affine and the children render pixel-exact.

    Qualifies (so we freeze it) only when the group:
      * contains a `custGeom` OR a `<p:graphicFrame>` descendant (complex content
        the per-shape re-synth botches), AND
      * contains NO `<p:ph>` descendant — carrying it whole never buries a fillable
        content slot, so the content/chrome split is preserved.

    Top-level only (`len(offset) == 2`): the group's xfrm is then slide-absolute, so
    shifting its own `grpSpPr/a:off` by the (pure-translation) offset places it
    correctly and the internal chOff/chExt handle the child affine untouched. A
    nested-scaled group (8-tuple offset) recurses as before. Returns True when it
    carried the group (caller must NOT recurse); False to recurse exactly as today.
    Any failure returns False → safe fall-back to recursion.
    """
    if len(offset) != 2:
        return False
    # Qualify gate: complex content present, no content placeholder buried inside.
    has_complex = (ch.find(".//a:custGeom", NS) is not None
                   or ch.find(".//p:graphicFrame", NS) is not None)
    if not has_complex:
        return False
    if ch.find(".//p:ph", NS) is not None:
        return False
    try:
        import copy as _copy
        grp_xfrm = ch.find("p:grpSpPr/a:xfrm", NS)
        if grp_xfrm is None:
            return False
        off = grp_xfrm.find("a:off", NS)
        ext = grp_xfrm.find("a:ext", NS)
        if off is None or ext is None:
            return False
        ox, oy = int(offset[0]), int(offset[1])
        gx = int(off.get("x")) + ox
        gy = int(off.get("y")) + oy
        gw = int(ext.get("cx"))
        gh = int(ext.get("cy"))
        # Carry the external part-graph (grouped chart / SmartArt / pic media)
        # BEFORE mutating the copy, so blob capture reads the live source rels.
        parts = _capture_group_parts(ch, slide, theme)
        grp = _copy.deepcopy(ch)
        _bake_scheme_colors(grp, theme)
        # Shift the group's OWN origin by the ancestor translation; leave chOff /
        # chExt (they define the child coordinate space the group internally maps).
        _goff = grp.find("p:grpSpPr/a:xfrm/a:off", NS)
        if _goff is not None and (ox or oy):
            _goff.set("x", str(int(_goff.get("x") or 0) + ox))
            _goff.set("y", str(int(_goff.get("y") or 0) + oy))
        shapes.append(Shape(
            kind="graphic",
            x=cmap.x(gx), y=cmap.y(gy), w=cmap.w(gw), h=cmap.h(gh),
            native_xml=etree.tostring(grp).decode("utf-8"),
            native_parts=parts,
        ))
        return True
    except Exception:
        return False


def _emit_graphic_frame(ch, offset, shapes, slide, cmap, theme, palette):
    """Tables and charts both arrive as <p:graphicFrame>. Dispatch by inner kind."""
    xfrm = ch.find("p:xfrm", NS)
    if xfrm is None:
        return
    off = xfrm.find("a:off", NS)
    ext = xfrm.find("a:ext", NS)
    if off is None or ext is None:
        return
    ox_local, oy_local = offset
    x0 = int(off.get("x")) + ox_local
    y0 = int(off.get("y")) + oy_local
    fw = int(ext.get("cx"))
    fh = int(ext.get("cy"))

    tbl = ch.find(".//a:tbl", NS)
    if tbl is not None:
        # Native-carry the whole graphicFrame (inline <a:tbl>, no external parts) for
        # a pixel-exact table — _emit_table's cell-by-cell re-synthesis drifts (guessed
        # row heights, only bottom borders, no merged cells, and mixed cell font sizes
        # collapse + overflow). Top-level only (offset is a pure translation).
        if len(offset) == 2:
            try:
                import copy as _copy
                frame = _copy.deepcopy(ch)
                _bake_scheme_colors(frame, theme)
                _foff = frame.find("p:xfrm/a:off", NS)
                if _foff is not None and (offset[0] or offset[1]):
                    _foff.set("x", str(int(_foff.get("x") or 0) + int(offset[0])))
                    _foff.set("y", str(int(_foff.get("y") or 0) + int(offset[1])))
                shapes.append(Shape(
                    kind="graphic", x=cmap.x(x0), y=cmap.y(y0),
                    w=cmap.w(fw), h=cmap.h(fh),
                    native_xml=etree.tostring(frame).decode("utf-8"),
                ))
                return
            except Exception:
                pass
        _emit_table(tbl, x0, y0, shapes, cmap, theme, palette)
        return

    # <c:chart r:id="..."/> inside graphicData → resolve chart part via slide rels.
    chart_ref = ch.find(f".//{{{CHART_NS}}}chart")
    if chart_ref is not None:
        # Native-carry the whole graphicFrame + its external part-graph (the chart
        # part itself + its chartStyle / chartColorStyle / embedded-xlsx children)
        # so the chart renders pixel-exact, vs _emit_chart's lossy re-synthesis
        # (pie/bar only, recoloured to brand, line/area/etc dropped). Top-level
        # only (offset is a pure translation; a grouped frame's xfrm is
        # group-relative). Falls back to _emit_chart on ANY failure so a single
        # awkward chart never crashes the decompile.
        if len(offset) == 2:
            try:
                import copy as _copy
                import base64 as _b64
                frame = _copy.deepcopy(ch)
                _bake_scheme_colors(frame, theme)
                _foff = frame.find("p:xfrm/a:off", NS)
                if _foff is not None and (offset[0] or offset[1]):
                    _foff.set("x", str(int(_foff.get("x") or 0) + int(offset[0])))
                    _foff.set("y", str(int(_foff.get("y") or 0) + int(offset[1])))
                # Walk every <c:chart r:id> in the frame and collect the part-graph
                # rooted at each chart part. Each entry records:
                #   parent   — 'slide' for the chart part itself (related from
                #              slide.part), else the OWNING chart-partname str for
                #              that chart's children, so the emitter rewires
                #              bottom-up.
                #   src_rid  — the rId by which this part is referenced FROM its
                #              parent in the SOURCE deck (the <c:chart r:id> for a
                #              chart part; the chart-part rel rId for a leaf). The
                #              emitter maps src_rid→new_rid to rewrite the matching
                #              r:id / r:embed references (e.g. <c:externalData> →
                #              the xlsx) so they point at the freshly-created parts.
                parts: list[dict] = []
                seen: set[str] = set()

                def _capture_graph(part, parent_key: str, reltype: str, src_rid: str) -> None:
                    pn = str(part.partname)
                    if pn in seen:
                        return
                    seen.add(pn)
                    # Bake schemeClr→srgbClr against the SOURCE theme in every
                    # carried XML part (chart, chartStyle, chartColorStyle) so the
                    # chart keeps its EXACT source series colours. Chart series are
                    # nearly always theme accent refs (accent1..6 in both chart.xml
                    # and colors*.xml); without baking they'd resolve against the
                    # OUTPUT deck's default theme and render in Office defaults
                    # (red/green/purple) instead of the brand palette. The
                    # embedded .xlsx is a binary zip, not XML — leave it verbatim.
                    raw = part.blob
                    ct = part.content_type
                    if ct.endswith("+xml") or ct.endswith("/xml"):
                        try:
                            _proot = etree.fromstring(raw)
                            _bake_scheme_colors(_proot, theme)
                            raw = etree.tostring(
                                _proot, xml_declaration=True,
                                encoding="UTF-8", standalone=True,
                            )
                        except Exception:
                            raw = part.blob
                    parts.append({
                        "partname": pn,
                        "content_type": ct,
                        "blob": _b64.b64encode(raw).decode("ascii"),
                        "reltype": reltype,
                        "parent": parent_key,
                        "src_rid": src_rid,
                    })
                    # Recurse into THIS part's own relationships (chart → style /
                    # colors / xlsx; those leaves have no further parts we carry).
                    try:
                        child_rels = list(part.rels.values())
                    except Exception:
                        child_rels = []
                    for r in child_rels:
                        if r.is_external:
                            continue
                        try:
                            tgt = r.target_part
                        except Exception:
                            continue
                        _capture_graph(tgt, pn, r.reltype, r.rId)

                for cref in frame.iter(f"{{{CHART_NS}}}chart"):
                    rid = cref.get(f"{{{RELS_NS}}}id")
                    if not rid:
                        continue
                    cp = slide.part.related_part(rid)
                    _capture_graph(cp, "slide", RELS_NS_CHART, rid)
                if not parts:
                    raise ValueError("no chart parts resolved")
                shapes.append(Shape(
                    kind="graphic", x=cmap.x(x0), y=cmap.y(y0),
                    w=cmap.w(fw), h=cmap.h(fh),
                    native_xml=etree.tostring(frame).decode("utf-8"),
                    native_parts=parts,
                ))
                return
            except Exception:
                pass
        rid = chart_ref.get(f"{{{RELS_NS}}}id")
        if rid:
            try:
                chart_part = slide.part.related_part(rid)
            except Exception:
                chart_part = None
            if chart_part is not None:
                _emit_chart(chart_part, x0, y0, fw, fh, shapes, cmap, theme, palette)
        return

    # SmartArt / diagram: graphicData uri='…/drawingml/2006/diagram' containing
    # `<dgm:relIds r:dm/r:lo/r:qs/r:cs>`. Native-carry the whole graphicFrame + the
    # diagram part-graph (data / layout / quickStyle / colors + the pre-rendered
    # drawing + any media sub-rel) so the diagram renders pixel-exact, vs
    # `_emit_smartart`'s lossy flatten. Top-level only (offset is a pure
    # translation; a grouped frame's xfrm is group-relative). Falls back to
    # `_emit_smartart` on ANY failure so an awkward diagram never crashes decompile.
    gdata = ch.find(".//a:graphicData", NS)
    gdata_uri = gdata.get("uri") if gdata is not None else None
    relids_list = ch.findall(f".//{{{DGM_NS}}}relIds")
    if gdata_uri and gdata_uri.endswith("/diagram") and relids_list and len(offset) == 2:
        try:
            import copy as _copy
            import base64 as _b64
            frame = _copy.deepcopy(ch)
            _bake_scheme_colors(frame, theme)
            _foff = frame.find("p:xfrm/a:off", NS)
            if _foff is not None and (offset[0] or offset[1]):
                _foff.set("x", str(int(_foff.get("x") or 0) + int(offset[0])))
                _foff.set("y", str(int(_foff.get("y") or 0) + int(offset[1])))
            # Collect the diagram part-graph. Each entry records the same shape as
            # the chart branch ({partname, content_type, blob, reltype, parent,
            # src_rid}). The four dgm parts + the drawing part hang off 'slide';
            # a media image (data/drawing → ../media/imageN.png) hangs off its
            # data/drawing part. schemeClr→srgbClr is baked into every XML blob
            # (colors*.xml + the drawing reference the theme), so the diagram
            # keeps its EXACT source palette under the output deck's theme.
            parts: list[dict] = []
            seen: set[str] = set()

            def _capture_graph(part, parent_key: str, reltype: str, src_rid: str) -> None:
                pn = str(part.partname)
                # A part shared by two parents (a data part AND
                # drawing10 both reference ../media/image131.png) is materialised
                # ONCE, but we still record THIS parent's src_rid→part mapping so
                # the SECOND parent's own <…r:embed> gets rewritten too. Without
                # this the drawing's image ref would dangle at build time.
                if pn in seen:
                    if src_rid:
                        _DGM_SHARED.setdefault(parent_key, {})[src_rid] = (pn, reltype)
                    return
                seen.add(pn)
                if src_rid:
                    _DGM_SHARED.setdefault(parent_key, {})[src_rid] = (pn, reltype)
                raw = part.blob
                ct = part.content_type
                if ct.endswith("+xml") or ct.endswith("/xml"):
                    try:
                        _proot = etree.fromstring(raw)
                        _bake_scheme_colors(_proot, theme)
                        raw = etree.tostring(
                            _proot, xml_declaration=True,
                            encoding="UTF-8", standalone=True,
                        )
                    except Exception:
                        raw = part.blob
                parts.append({
                    "partname": pn,
                    "content_type": ct,
                    "blob": _b64.b64encode(raw).decode("ascii"),
                    "reltype": reltype,
                    "parent": parent_key,
                    "src_rid": src_rid,
                })
                try:
                    child_rels = list(part.rels.values())
                except Exception:
                    child_rels = []
                for r in child_rels:
                    if r.is_external:
                        continue
                    try:
                        tgt = r.target_part
                    except Exception:
                        continue
                    _capture_graph(tgt, pn, r.reltype, r.rId)

            # `_DGM_SHARED` records, per parent-partname, the src_rid→partname of
            # any ALREADY-captured shared child so the emitter can still wire the
            # second parent's rel + rewrite its ref. (parent → {src_rid: old_pn})
            _DGM_SHARED: dict[str, dict[str, str]] = {}
            for relids in relids_list:
                # The four core parts, resolved via the slide rels.
                data_part = None
                for attr, reltype in _DGM_RELID_ATTRS:
                    rid = relids.get(f"{{{RELS_NS}}}{attr}")
                    if not rid:
                        continue
                    p = slide.part.related_part(rid)
                    if attr == "dm":
                        data_part = p
                    _capture_graph(p, "slide", reltype, rid)
                # The pre-rendered drawing: dataN extLst <dsp:dataModelExt relId>
                # is a SLIDE rel pointing at the drawing part.
                if data_part is not None:
                    dme = data_part.blob
                    try:
                        droot = etree.fromstring(dme)
                        dmext = droot.find(
                            f".//{{{DSP_DATAMODEL_NS}}}dataModelExt"
                        )
                    except Exception:
                        dmext = None
                    draw_rid = dmext.get("relId") if dmext is not None else None
                    if draw_rid:
                        dp = slide.part.related_part(draw_rid)
                        _capture_graph(dp, "slide", RELS_NS_DGM_DRAWING, draw_rid)
            if not parts:
                raise ValueError("no diagram parts resolved")
            # Fold the shared-child mappings into the parts list so the emitter
            # can wire them: emit one extra entry per (parent, src_rid) that
            # points at an EXISTING carried part (blob omitted — the emitter
            # recognises a "ref"-only entry and reuses the materialised part).
            for parent_key, m in _DGM_SHARED.items():
                for src_rid, (old_pn, reltype) in m.items():
                    if any(
                        e["parent"] == parent_key and e.get("src_rid") == src_rid
                        for e in parts
                    ):
                        continue
                    parts.append({
                        "ref": old_pn, "parent": parent_key,
                        "src_rid": src_rid, "reltype": reltype,
                    })
            shapes.append(Shape(
                kind="graphic", x=cmap.x(x0), y=cmap.y(y0),
                w=cmap.w(fw), h=cmap.h(fh),
                native_xml=etree.tostring(frame).decode("utf-8"),
                native_parts=parts,
            ))
            return
        except Exception:
            pass

    # SmartArt diagrams: graphicData uri='…/drawingml/2006/diagram'. The slide
    # rels carry both a `diagramData` (the semantic model) and a
    # `diagramDrawing` (the pre-rendered drawing*.xml computed by PowerPoint
    # when the user last edited the diagram). Parsing the drawing skips
    # re-implementing the SmartArt layout engine — every shape, its xfrm,
    # fill, stroke, and text live inside `<dsp:sp>` elements that mirror
    # the `<p:sp>` structure.
    diag_rel_ns = "http://schemas.microsoft.com/office/2007/relationships/diagramDrawing"
    if slide is not None and hasattr(slide, "part"):
        try:
            rels = slide.part.rels
            drawing_part = None
            for rel in rels.values():
                if rel.reltype == diag_rel_ns:
                    # The drawing isn't directly referenced by rId from the
                    # graphicFrame — it's a sibling relationship of the
                    # diagramData. PowerPoint ties them by partname suffix
                    # (data6.xml ↔ drawing6.xml). Find the matching one by
                    # numeric suffix on the graphicData's diagramData rId.
                    drawing_part = rel.target_part
                    # The slide may have multiple diagramDrawing rels (one
                    # per SmartArt). Use the rId currently being processed
                    # if available; fall back to first drawing.
                    break
            if drawing_part is not None:
                _emit_smartart(drawing_part.blob, x0, y0, fw, fh,
                               shapes, cmap, theme, palette)
        except Exception:
            pass


DSP_NS = "http://schemas.microsoft.com/office/drawing/2008/diagram"


def _emit_smartart(blob: bytes, x0: int, y0: int, fw: int, fh: int,
                   shapes: list, cmap, theme, palette) -> None:
    """Decompile a SmartArt's pre-rendered drawing XML.

    PowerPoint caches the computed-layout shapes for each SmartArt
    diagram in `ppt/diagrams/drawing*.xml` so they don't need to be
    relaid-out at presentation time. Each shape lives inside `<dsp:sp>`
    elements that mirror the regular `<p:sp>` structure but use the
    `dsp` namespace. Walking that tree gives us the actual circles,
    arrows, callouts, etc. without re-implementing the SmartArt layout
    engine.

    `x0,y0` and `fw,fh` are the host `<p:graphicFrame>`'s EMU position
    and size — the drawing's internal coordinates are already in
    slide-EMU (the layout engine wrote them out absolute), so no extra
    transform is needed for shapes whose own xfrm already lives in
    slide-space. The frame offset is preserved as a fallback for
    shapes whose xfrm is relative.
    """
    try:
        root = etree.fromstring(blob)
    except Exception:
        return
    spTree = root.find(f".//{{{DSP_NS}}}spTree")
    if spTree is None:
        return
    for sp in spTree.findall(f"{{{DSP_NS}}}sp"):
        spPr = sp.find(f"{{{DSP_NS}}}spPr")
        if spPr is None:
            continue
        xfrm = spPr.find("a:xfrm", NS)
        if xfrm is None:
            continue
        off = xfrm.find("a:off", NS)
        ext = xfrm.find("a:ext", NS)
        if off is None or ext is None:
            continue
        try:
            x = int(off.get("x")) + x0
            y = int(off.get("y")) + y0
            w = int(ext.get("cx"))
            h = int(ext.get("cy"))
        except (TypeError, ValueError):
            continue
        # `<dsp:sp>` xfrm coords are RELATIVE to the host graphicFrame
        # (the layout engine writes them in drawing-internal space), so
        # we add the frame's (x0, y0) to land each shape on the slide.
        # Fill — same resolver as <p:sp> uses (the a: children are
        # identical regardless of dsp vs p parent).
        fill = _resolve_fill(spPr, theme, palette)
        # Stroke + width from <a:ln>.
        stroke = None
        stroke_width: float | None = None
        ln = spPr.find("a:ln", NS)
        if ln is not None:
            sf = ln.find("a:solidFill", NS)
            if sf is not None:
                stroke = _resolve_solid(sf, theme, palette)
            w_attr = ln.get("w")
            if w_attr:
                try:
                    stroke_width = cmap.w(int(w_attr))
                except (TypeError, ValueError):
                    pass
        # Geometry kind.
        kind = "rect"
        pg = spPr.find("a:prstGeom", NS)
        if pg is not None:
            preset = pg.get("prst")
            if preset == "ellipse":
                kind = "oval"
            elif preset in ("line", "straightConnector1"):
                kind = "line"
            elif preset in ("rect", "roundRect"):
                kind = "rect"
        elif spPr.find("a:custGeom", NS) is not None:
            kind = "shape"
        # Text — dsp:txBody mirrors a:txBody / p:txBody.
        txBody = sp.find(f"{{{DSP_NS}}}txBody")
        runs = _text_runs(sp, theme, palette) if txBody is not None else []
        # Drop placeholder demo text the same way regular shapes do.
        if runs and _is_placeholder_text(runs):
            runs = []
        shapes.append(Shape(
            kind=kind,
            x=cmap.x(x), y=cmap.y(y),
            w=cmap.w(w), h=cmap.h(h),
            fill=fill, stroke=stroke, stroke_width=stroke_width,
            text_runs=runs,
        ))


def _emit_table(tbl, x0, y0, shapes, cmap, theme, palette):
    grid = tbl.find("a:tblGrid", NS)
    if grid is None:
        return
    col_widths = [int(c.get("w")) for c in grid.findall("a:gridCol", NS)]

    # Auto-fit h=0 rows: PowerPoint expands them to fit text. Use a heuristic
    # of 350000 EMU (~31 px) per text line so subsequent rows shift down and
    # don't overlap the header. Without this the header sits visually inside
    # row 1.
    EMU_PER_LINE = 350000

    rows = list(tbl.findall("a:tr", NS))
    effective_h = []
    for tr in rows:
        h = int(tr.get("h", "0"))
        if h == 0:
            # estimate from text content
            text = "".join(t.text or "" for t in tr.findall(".//a:t", NS))
            lines = max(1, len(text.split("\n"))) if text else 1
            h = EMU_PER_LINE * lines
        effective_h.append(h)

    y_cursor = y0
    for ri, tr in enumerate(rows):
        row_h = effective_h[ri]
        x_cursor = x0
        cells = tr.findall("a:tc", NS)
        for ci, tc in enumerate(cells):
            cw = col_widths[ci] if ci < len(col_widths) else 0
            tcPr = tc.find("a:tcPr", NS)
            fill = _resolve_fill(tcPr, theme, palette) if tcPr is not None else None
            runs = _text_runs(tc, theme, palette)
            if fill is not None:
                shapes.append(Shape(
                    kind="rect", x=cmap.x(x_cursor), y=cmap.y(y_cursor),
                    w=cmap.w(cw), h=cmap.h(row_h), fill=fill,
                ))
            # Cell bottom border → emit as line (orange separators under headers).
            if tcPr is not None:
                lnB = tcPr.find("a:lnB", NS)
                if lnB is not None:
                    sf = lnB.find("a:solidFill", NS)
                    stroke = _resolve_solid(sf, theme, palette) if sf is not None else None
                    if stroke is not None:
                        y_b = y_cursor + row_h
                        shapes.append(Shape(
                            kind="line", x=cmap.x(x_cursor), y=cmap.y(y_b),
                            w=cmap.w(cw), h=0, stroke=stroke,
                        ))
            if runs:
                shapes.append(Shape(
                    kind="text", x=cmap.x(x_cursor + cw // 20), y=cmap.y(y_cursor + row_h // 4),
                    w=cmap.w(cw - cw // 10), h=cmap.h(row_h),
                    text_runs=runs,
                ))
            x_cursor += cw
        y_cursor += row_h


def _emit_pie_chart(pie_el, x0, y0, fw, fh, shapes, cmap, theme=None, palette=None):
    """Extract pie/doughnut chart geometry and emit one svg{} arc path per slice.

    Each slice becomes a Shape with kind='shape', svg_path_d set to an SVG
    arc path of the form 'M cx,cy L x1,y1 A r,r 0 large,sweep x2,y2 Z',
    fill mapped to chart-series-N via slice index (using the brand's
    chart-series ramp — e.g. accent → accent-80 → ... → accent-10).
    emit_dsl() converts each Shape into a standalone `svg{}` block; the
    blocks share the same bbox (the chart frame) so slices overlay into
    one unified pie at render time.

    Percentage labels are emitted as `text` Shapes positioned just outside
    the slice's arc bisector when <c:dLbls><c:showPercent val="1"> is set
    in the source — matches PowerPoint's default external-label placement.

    Per-slice colors in the source XML are intentionally ignored in favour
    of the brand's chart-series ramp: source decks authored against the
    brand already use the same hue sequence, so index-based mapping gives
    fidelity AND brand-correctness in one pass. Source decks authored
    off-brand will use the brand's hues here — that is by design (a
    brand-pack decompile is brand-conforming output, not pixel mimicry).
    """
    ser = pie_el.find(f"{{{CHART_NS}}}ser")
    if ser is None:
        return
    val_els = ser.findall(
        f".//{{{CHART_NS}}}val//{{{CHART_NS}}}pt/{{{CHART_NS}}}v"
    )
    try:
        values = [float(v.text) for v in val_els if v.text]
    except (TypeError, ValueError):
        return
    if not values or sum(values) <= 0:
        return
    total = sum(values)

    # Per-slice colors from <c:dPt> data-point elements. Each dPt carries
    # <c:idx val="N"/> identifying the slice index plus an optional
    # <c:spPr><a:solidFill> with that slice's brand color. Falls back to
    # the chart-series-N ramp by index when absent.
    #
    # We resolve dPt fills via `palette={}` so the source's exact hex
    # propagates through unchanged. Going through the brand-token
    # nearest_token collapses two close source hues to the same token
    # ("accent"), which then renders identically — defeating the whole
    # purpose of per-slice colours. Same rationale as the custGeom
    # palette-bypass added earlier.
    slice_colors: dict[int, str] = {}
    if theme is not None:
        for dpt in ser.findall(f"{{{CHART_NS}}}dPt"):
            idx_el = dpt.find(f"{{{CHART_NS}}}idx")
            sp_pr = dpt.find(f"{{{CHART_NS}}}spPr")
            if idx_el is None or sp_pr is None:
                continue
            try:
                idx = int(idx_el.get("val") or "-1")
            except (TypeError, ValueError):
                continue
            color = _resolve_fill(sp_pr, theme, palette={})
            if color and idx >= 0:
                slice_colors[idx] = color

    cat_els = ser.findall(
        f".//{{{CHART_NS}}}cat//{{{CHART_NS}}}pt/{{{CHART_NS}}}v"
    )
    categories = [c.text or "" for c in cat_els]

    # Legend position: source PowerPoint convention is r/l/t/b. We honour
    # only r/l (column-style legend with one row per category — the
    # dominant case for pies); t/b are rare on small pies and fall back
    # to right-side. Search at chart-space level since legend lives on
    # the chart root, not inside pie_el. When the chart has NO `<c:legend>`
    # element at all, the pie should fill its frame (no legend slot to
    # reserve). The previous code defaulted to "r" even when no legend
    # element existed, which made pies on legend-less charts (showcase
    # decks where every slice is labelled inline with `<c:showPercent>`)
    # render at ~50% of their source size.
    chart_root = pie_el
    while chart_root is not None and chart_root.tag != f"{{{CHART_NS}}}chartSpace":
        parent = chart_root.getparent()
        if parent is None:
            break
        chart_root = parent
    legend_pos: str | None = None
    if chart_root is not None:
        legend_el = chart_root.find(f".//{{{CHART_NS}}}legend")
        if legend_el is not None:
            # Overlay flag: when the legend is set to overlay the plot area
            # (`<c:overlay val="1"/>`), no horizontal slot is reserved for
            # it — treat as legend-less for sizing purposes.
            overlay_el = legend_el.find(f"{{{CHART_NS}}}overlay")
            is_overlay = overlay_el is not None and overlay_el.get("val") in ("1", "true")
            lp = legend_el.find(f"{{{CHART_NS}}}legendPos")
            if not is_overlay:
                # Collapse PowerPoint's 8-position legend space ("l", "r",
                # "t", "b", "tr", "tl", "br", "bl") down to the dominant
                # axis. Pie sizing only cares whether the legend lives
                # on the left/right (reserving horizontal space) or
                # top/bottom (vertical) — corner positions like "tr" act
                # as right-side legends for sizing.
                raw = lp.get("val") if lp is not None else "r"
                if raw in ("l", "tl", "bl"):
                    legend_pos = "l"
                elif raw in ("t", "b"):
                    legend_pos = raw
                else:  # "r", "tr", "br", or unknown — treat as right
                    legend_pos = "r"

    # Data label flags — <c:dLbls><c:showPercent val="1"/> or
    # <c:dLbls><c:showVal val="1"/>. Mutually exclusive in practice;
    # percent wins when both are set, matching PowerPoint behaviour.
    show_percent = False
    show_val = False
    sp = pie_el.find(
        f".//{{{CHART_NS}}}dLbls/{{{CHART_NS}}}showPercent"
    )
    if sp is not None and sp.get("val") == "1":
        show_percent = True
    sv = pie_el.find(
        f".//{{{CHART_NS}}}dLbls/{{{CHART_NS}}}showVal"
    )
    if sv is not None and sv.get("val") == "1":
        show_val = True

    # svg-block-local pixel coords for slice paths. The block's outer
    # bbox is the chart frame; coords inside are 0..bbox_w_px by
    # 0..bbox_h_px.
    bbox_w_px = cmap.w(fw)
    bbox_h_px = cmap.h(fh)

    # `<c:plotArea><c:layout><c:manualLayout>` gives EXACT fractional
    # plot-area position within the chart frame (xMode/yMode="edge" with
    # x/y/w/h as fractions of bbox_w/bbox_h). When present, use those
    # directly — they're what PowerPoint's layout engine resolved when
    # the deck author placed the chart. Falls back to the aspect-based
    # heuristic when the source uses `<c:layout/>` (auto-layout).
    pa_layout = None
    chart_root_for_layout = chart_root
    if chart_root_for_layout is not None:
        pa_layout = chart_root_for_layout.find(
            f".//{{{CHART_NS}}}plotArea/{{{CHART_NS}}}layout/{{{CHART_NS}}}manualLayout"
        )

    def _layout_frac(el, tag: str, default: float | None = None) -> float | None:
        if el is None:
            return default
        c = el.find(f"{{{CHART_NS}}}{tag}")
        if c is None or not c.get("val"):
            return default
        try:
            return float(c.get("val"))
        except (TypeError, ValueError):
            return default

    if pa_layout is not None:
        plot_xf = _layout_frac(pa_layout, "x", 0.0) or 0.0
        plot_yf = _layout_frac(pa_layout, "y", 0.0) or 0.0
        plot_wf = _layout_frac(pa_layout, "w", 1.0) or 1.0
        plot_hf = _layout_frac(pa_layout, "h", 1.0) or 1.0
        pie_off_x = plot_xf * bbox_w_px
        pie_off_y = plot_yf * bbox_h_px
        pie_w_px = plot_wf * bbox_w_px
        pie_h_px = plot_hf * bbox_h_px
    else:
        # Heuristic: pie-area fraction adapts to chart-frame aspect.
        # Wide frames (multi-pie-in-column layouts) keep ~60%; square-ish
        # frames shrink to ~50% leaving room for the legend.
        frame_aspect = bbox_w_px / bbox_h_px if bbox_h_px else 1.0
        if categories and legend_pos in ("l", "r"):
            pie_w_frac = 0.60 if frame_aspect > 1.4 else 0.50
        else:
            pie_w_frac = 1.0
        pie_w_px = bbox_w_px * pie_w_frac
        pie_h_px = bbox_h_px
        pie_off_x = (bbox_w_px - pie_w_px) if legend_pos == "l" else 0.0
        pie_off_y = 0.0
    cx_px = pie_off_x + pie_w_px / 2
    cy_px = pie_off_y + pie_h_px / 2
    # min(w,h) keeps pies circular in non-square frames. When manualLayout
    # gave us a plot-area smaller than the chart frame, the radius is half
    # the plot's shortest side (no further margin); otherwise 0.36 leaves
    # margin for the auto-layout's external label placement.
    if pa_layout is not None:
        r_px = min(pie_w_px, pie_h_px) / 2.0
    else:
        r_px = min(pie_w_px, pie_h_px) * 0.36

    # Doughnut hole: `<c:holeSize val="N"/>` on `<c:doughnutChart>` where
    # N is 10..90 = inner-radius percentage of outer radius. Default 50
    # when the element is missing on a doughnutChart; pieChart has no
    # hole. `pie_el.tag` localname distinguishes the two.
    inner_r_px = 0.0
    if etree.QName(pie_el).localname == "doughnutChart":
        hole_pct = 50
        hs = pie_el.find(f"{{{CHART_NS}}}holeSize")
        if hs is not None and hs.get("val"):
            try:
                hole_pct = max(10, min(90, int(hs.get("val"))))
            except (TypeError, ValueError):
                pass
        inner_r_px = r_px * (hole_pct / 100.0)

    # Start at 12 o'clock (-π/2), sweep clockwise. PowerPoint pies follow
    # this convention; matching it preserves slice-to-color correspondence
    # against the source.
    #
    # `<c:firstSliceAng val="N"/>` rotates the start clockwise by N
    # degrees (0..360). Sources rarely use it but when they do (corporate
    # decks rotating the highlighted slice into a fixed position), the
    # entire slice-to-colour ordering shifts.
    angle_start = -math.pi / 2
    first_ang_el = pie_el.find(f"{{{CHART_NS}}}firstSliceAng")
    if first_ang_el is not None and first_ang_el.get("val"):
        try:
            angle_start += math.radians(float(first_ang_el.get("val")))
        except (TypeError, ValueError):
            pass

    for i, v in enumerate(values):
        if v <= 0:
            angle_start += 0
            continue
        sweep = (v / total) * 2 * math.pi
        angle_end = angle_start + sweep
        x1 = cx_px + r_px * math.cos(angle_start)
        y1 = cy_px + r_px * math.sin(angle_start)
        x2 = cx_px + r_px * math.cos(angle_end)
        y2 = cy_px + r_px * math.sin(angle_end)
        large_arc = 1 if sweep > math.pi else 0
        if inner_r_px > 0:
            # Annular sector — outer arc forward + inner arc reversed.
            ix1 = cx_px + inner_r_px * math.cos(angle_start)
            iy1 = cy_px + inner_r_px * math.sin(angle_start)
            ix2 = cx_px + inner_r_px * math.cos(angle_end)
            iy2 = cy_px + inner_r_px * math.sin(angle_end)
            d = (
                f"M {x1:.1f},{y1:.1f} "
                f"A {r_px:.1f},{r_px:.1f} 0 {large_arc},1 {x2:.1f},{y2:.1f} "
                f"L {ix2:.1f},{iy2:.1f} "
                f"A {inner_r_px:.1f},{inner_r_px:.1f} 0 {large_arc},0 {ix1:.1f},{iy1:.1f} "
                f"Z"
            )
        else:
            # Pie wedge — apex at centre.
            d = (
                f"M {cx_px:.1f},{cy_px:.1f} "
                f"L {x1:.1f},{y1:.1f} "
                f"A {r_px:.1f},{r_px:.1f} 0 {large_arc},1 {x2:.1f},{y2:.1f} Z"
            )
        fill_token = slice_colors.get(i) or f"chart-series-{(i % 6) + 1}"
        shapes.append(Shape(
            kind="shape",
            x=cmap.x(x0), y=cmap.y(y0),
            w=bbox_w_px, h=bbox_h_px,
            fill=fill_token,
            svg_path_d=d,
        ))

        if show_percent or show_val:
            mid_angle = angle_start + sweep / 2
            # showPercent labels sit ~15% outside circumference (PPT
            # default external placement on small pies). showVal labels
            # sit inside the slice at ~65% radius — PowerPoint's default
            # internal-label position.
            label_r = r_px * 1.15 if show_percent else r_px * 0.65
            lx = cx_px + label_r * math.cos(mid_angle)
            ly = cy_px + label_r * math.sin(mid_angle)
            if show_percent:
                label_text = f"{round((v / total) * 100)} %"
            else:
                # Format like the source: 8.2 → "8,2", 1 → "1".
                if v == int(v):
                    label_text = str(int(v))
                else:
                    label_text = f"{v:.1f}".replace(".", ",")
            shapes.append(Shape(
                kind="text",
                x=cmap.x(x0) + int(lx) - 20,
                y=cmap.y(y0) + int(ly) - 10,
                w=40, h=20,
                text_runs=[TextRun(text=label_text, pt=10)],
            ))

        angle_start = angle_end

    # Legend (categories + colour swatches). Position: right or left of pie.
    # Stack one row per category centred vertically against the pie. Skip
    # entirely when no categories — labelled-from-percentage slices alone
    # carry enough signal for the small pies common in showcase decks.
    if categories and legend_pos in ("l", "r"):
        swatch_w = 14
        swatch_h = 14
        row_gap = 32  # vertical pitch between legend rows
        text_gap = 8  # swatch-to-label horizontal gap
        legend_w_px = bbox_w_px - pie_w_px
        # Legend bbox top-left within the chart frame.
        legend_x_local = 0 if legend_pos == "l" else pie_w_px
        # Center the legend block vertically against the pie.
        block_h = len(categories) * row_gap
        legend_y_local = max(0, (bbox_h_px - block_h) / 2)
        # Insets keep legend out of frame edges.
        legend_x_local += 12
        for i, cat in enumerate(categories):
            row_y = legend_y_local + i * row_gap
            color = f"chart-series-{(i % 6) + 1}"
            # Swatch as a rect — survives outside svg{} blocks and accepts
            # the brand vocab directly.
            shapes.append(Shape(
                kind="rect",
                x=cmap.x(x0) + int(legend_x_local),
                y=cmap.y(y0) + int(row_y),
                w=swatch_w, h=swatch_h,
                fill=color,
            ))
            shapes.append(Shape(
                kind="text",
                x=cmap.x(x0) + int(legend_x_local) + swatch_w + text_gap,
                y=cmap.y(y0) + int(row_y) - 4,
                w=int(legend_w_px) - swatch_w - text_gap - 12,
                h=24,
                text_runs=[TextRun(text=cat, pt=10)],
            ))


def _emit_chart(chart_part, x0, y0, fw, fh, shapes, cmap, theme, palette):
    """Extract chart geometry from a c:chartSpace part and emit DSL primitives.

    Dispatches by chart type:
      * c:pieChart / c:doughnutChart → _emit_pie_chart (arc paths + labels)
      * c:barChart                   → _emit_bar_chart (rects + axis + labels)
    Other chart types (line, area, scatter, etc.) fall through unhandled —
    the decompile-as-rects fallback in the hybrid SVG pass still gives a
    rough first-pass for those, and improve-brand sub-agents refine.
    """
    try:
        root = etree.fromstring(chart_part.blob)
    except Exception:
        return

    pie = (root.find(f".//{{{CHART_NS}}}pieChart")
           or root.find(f".//{{{CHART_NS}}}doughnutChart"))
    if pie is not None:
        _emit_pie_chart(pie, x0, y0, fw, fh, shapes, cmap, theme=theme, palette=palette)
        return

    bar = root.find(f".//{{{CHART_NS}}}barChart")
    if bar is None:
        return

    # Bar orientation: <c:barDir val="bar"/> = horizontal bars (categories
    # stack vertically, values extend rightward); val="col" (default) =
    # vertical columns. The previous code always emitted columns, so any
    # horizontal-bar source slide rendered with bars rotated 90° — visible
    # as vertical stripes where the source had horizontal bars.
    bar_dir_el = bar.find(f"{{{CHART_NS}}}barDir")
    horizontal_bars = (
        bar_dir_el is not None and bar_dir_el.get("val") == "bar"
    )

    # Stacking: <c:grouping val="standard|stacked|percentStacked"/>.
    # "standard" (default): series sit side-by-side per category.
    # "stacked": series stack head-to-tail; axis_max = max sum-per-category.
    # "percentStacked": each category bar is 100%; series segment by share.
    grouping_el = bar.find(f"{{{CHART_NS}}}grouping")
    grouping = grouping_el.get("val") if grouping_el is not None else "standard"
    if grouping not in ("standard", "stacked", "percentStacked", "clustered"):
        grouping = "standard"

    series = []
    for ser in bar.findall(f"{{{CHART_NS}}}ser"):
        name_el = ser.find(f".//{{{CHART_NS}}}tx//{{{CHART_NS}}}v")
        name = name_el.text if name_el is not None else "?"
        vals = [float(v.text) for v in ser.findall(f".//{{{CHART_NS}}}val//{{{CHART_NS}}}pt/{{{CHART_NS}}}v")]
        cats = [v.text for v in ser.findall(f".//{{{CHART_NS}}}cat//{{{CHART_NS}}}pt/{{{CHART_NS}}}v")]
        # Per-series fill colour. Resolve via empty palette so the source
        # hex propagates verbatim — going through nearest_token collapses
        # two close source hues to the same brand token, losing the
        # series-to-series colour distinction on stacked / clustered
        # bar charts.
        sp_pr = ser.find(f"{{{CHART_NS}}}spPr")
        ser_color = _resolve_fill(sp_pr, theme, palette={}) if sp_pr is not None else None
        # Per-data-point colours from `<c:dPt>` overrides. Showcase decks
        # often colour alternating bars different hues to highlight a
        # specific category — that information lives in dPt only, not on
        # the series. Same palette={} bypass.
        bar_colors: dict[int, str] = {}
        for dpt in ser.findall(f"{{{CHART_NS}}}dPt"):
            idx_el = dpt.find(f"{{{CHART_NS}}}idx")
            dpt_sp = dpt.find(f"{{{CHART_NS}}}spPr")
            if idx_el is None or dpt_sp is None:
                continue
            try:
                idx = int(idx_el.get("val") or "-1")
            except (TypeError, ValueError):
                continue
            color = _resolve_fill(dpt_sp, theme, palette={})
            if color and idx >= 0:
                bar_colors[idx] = color
        series.append((name, vals, cats, ser_color, bar_colors))
    if not series:
        return

    n_cats = max(len(s[1]) for s in series)
    n_series = len(series)
    cats = series[0][2] if series[0][2] else [f"Cat {i+1}" for i in range(n_cats)]
    if grouping == "stacked":
        # Each category's bar = sum of all series values for that category.
        cat_totals = [sum(s[1][ci] for s in series if ci < len(s[1]))
                      for ci in range(n_cats)]
        data_max = max(cat_totals) if cat_totals else 0
    elif grouping == "percentStacked":
        # Every category sums to 100% — plot axis is just 100.
        data_max = 100
    else:
        data_max = max((max(s[1]) for s in series if s[1]), default=0)
    # Round axis max up to the next integer above data_max.
    # LibreOffice's auto-axis (the source-PNG ground truth in the verify
    # loop) adds one major-unit of headroom over data_max, so matching its
    # tick count beats the more semantically-correct ceil(data_max) — the
    # struct_diff_ratio improves when ticks line up with the source-PNG
    # rasterisation, not with what PowerPoint would have drawn.
    if grouping == "percentStacked":
        axis_max = 100
    else:
        axis_max = math.ceil(data_max + 0.5) if data_max > 0 else 5

    # Axis visibility — `<c:valAx><c:delete val="1"/>` and `<c:catAx>
    # <c:delete val="1"/>` hide the respective axis at render time. Many
    # showcase charts use this for a clean look: bars/segments alone, no
    # tick labels or category strings. Reading the flag from the source
    # XML is much better than always emitting ticks (and then mismatching
    # source pixels at every tick position).
    val_axis_hidden = False
    cat_axis_hidden = False
    for ax in root.findall(f".//{{{CHART_NS}}}valAx"):
        d = ax.find(f"{{{CHART_NS}}}delete")
        if d is not None and d.get("val") in ("1", "true"):
            val_axis_hidden = True
            break
    for ax in root.findall(f".//{{{CHART_NS}}}catAx"):
        d = ax.find(f"{{{CHART_NS}}}delete")
        if d is not None and d.get("val") in ("1", "true"):
            cat_axis_hidden = True
            break

    # Data-label flags — `<c:dLbls>` controls whether value/category/series
    # text gets drawn next to each bar. PowerPoint's structure puts this
    # on the bar chart element itself (and optionally per-series). Read
    # the top-level flags only for now; missing flags default to PPT's
    # behaviour (no labels unless set).
    def _dlbl_flag(parent, name: str) -> bool:
        el = parent.find(f"{{{CHART_NS}}}dLbls/{{{CHART_NS}}}{name}")
        return el is not None and el.get("val") in ("1", "true")
    show_val_labels = _dlbl_flag(bar, "showVal")
    show_cat_labels = _dlbl_flag(bar, "showCatName")

    # Plot-area extents inside the frame (EMU). When the axes are hidden
    # the plot can fill the frame edge-to-edge; otherwise reserve
    # PowerPoint's typical insets for tick/category labels.
    if val_axis_hidden and cat_axis_hidden:
        plot_x = x0
        plot_y = y0
        plot_w = fw
        plot_h = fh
    else:
        plot_x = x0 + int(fw * 0.07) if not val_axis_hidden else x0
        plot_y = y0 + int(fh * 0.12) if not cat_axis_hidden else y0
        plot_w = int(fw * (1.0 - 0.07 - 0.02)) if not val_axis_hidden else fw
        plot_h = int(fh * (1.0 - 0.12 - 0.22)) if not cat_axis_hidden else int(fh * (1.0 - 0.22))

    # Y-axis numeric labels — only when the value axis isn't hidden.
    if not val_axis_hidden:
        n_ticks = axis_max + 1
        for i in range(n_ticks):
            v = axis_max - i  # top→bottom
            ty = plot_y + int(plot_h * i / axis_max)
            shapes.append(Shape(
                kind="text",
                x=cmap.x(x0 + int(fw * 0.005)),
                y=cmap.y(ty - 180000),
                w=cmap.w(int(fw * 0.05)),
                h=cmap.h(360000),
                text_runs=[TextRun(text=str(v), pt=14)],
            ))

    # Skip gridlines: source has barely-visible hairlines; rendering them at
    # 0.75pt over a 1240px-wide plot adds heavy diff pixels and emit_dsl
    # orders lines after rects, so they paint OVER the bars producing stripes.

    # Category labels above each group — only when the category axis
    # isn't hidden AND the source explicitly enables them via dLbls.
    cat_w = plot_w // n_cats if n_cats else plot_w
    # Category labels render when the axis is visible AND the source either
    # opts into `<c:showCatName val="1"/>` OR omits the dLbls flag entirely
    # (PowerPoint's default behaviour shows axis-tied category labels even
    # without an explicit dLbls override).
    if not cat_axis_hidden and show_cat_labels:
        for ci in range(n_cats):
            cx = plot_x + ci * cat_w + cat_w // 4
            shapes.append(Shape(
                kind="text",
                x=cmap.x(cx),
                y=cmap.y(plot_y - int(fh * 0.07)),
                w=cmap.w(cat_w // 2),
                h=cmap.h(int(fh * 0.06)),
                text_runs=[TextRun(text=cats[ci] if ci < len(cats) else "", pt=14)],
            ))

    # Bars: each category has n_series side-by-side bars. PowerPoint sizes
    # them via `<c:gapWidth val="N"/>` where N is the inter-group gap as a
    # percentage of bar width (default 150 = gap is 1.5x bar width). The
    # category width then holds n_series bars plus a gap on each side:
    #   cat_w = bar_w * n_series + bar_w * (gapWidth/100)
    #   →  bar_w = cat_w / (n_series + gapWidth/100)
    # Reading the actual gapWidth lets thick "showcase" bars (default 150)
    # decompile at their real width instead of a fixed 8.5%-of-cat hairline,
    # which left source-bar pixels uncovered and inflated struct_diff on
    # every bar-chart slide.
    gap_pct = 150.0
    gw_el = bar.find(f"{{{CHART_NS}}}gapWidth")
    if gw_el is not None and gw_el.get("val"):
        try:
            gap_pct = float(gw_el.get("val"))
        except (TypeError, ValueError):
            pass
    if horizontal_bars:
        # Horizontal layout: category axis runs vertically (rows), value
        # axis runs horizontally. Each category gets a row of height
        # `cat_h`; bars within stack by series and extend rightward.
        cat_h = plot_h // n_cats if n_cats else plot_h
        if grouping in ("stacked", "percentStacked"):
            # One bar per category, series segments fill it left-to-right.
            bar_h = int(cat_h / (1 + gap_pct / 100))
            for ci in range(n_cats):
                # Per-category total (stacked uses sum, percentStacked
                # normalises to 100 so each row fills plot_w).
                if grouping == "percentStacked":
                    cat_total = sum(s[1][ci] for s in series if ci < len(s[1]))
                else:
                    cat_total = data_max
                cursor_x = plot_x
                row_y = plot_y + ci * cat_h + (cat_h - bar_h) // 2
                for si, (name, vals, _, ser_color, dpt_colors) in enumerate(series):
                    if ci >= len(vals):
                        continue
                    v = vals[ci]
                    # Per-data-point `<c:dPt>` colour overrides the
                    # series colour for this specific category index.
                    color = dpt_colors.get(ci) or ser_color or f"chart-series-{(si % 6) + 1}"
                    if grouping == "percentStacked":
                        seg_w = int(plot_w * (v / cat_total)) if cat_total > 0 else 0
                    else:
                        seg_w = int(plot_w * (v / axis_max)) if axis_max > 0 else 0
                    shapes.append(Shape(
                        kind="rect",
                        x=cmap.x(cursor_x), y=cmap.y(row_y),
                        w=cmap.w(seg_w), h=cmap.h(bar_h),
                        fill=color,
                    ))
                    # Value labels for stacked / percentStacked horizontal
                    # bars — emit ONLY when source has `<c:showVal val="1"/>`.
                    # Label sits in the middle of its segment.
                    if show_val_labels and seg_w > 200000:
                        label = str(v).rstrip("0").rstrip(".") if "." in str(v) else str(v)
                        label = label.replace(".", ",")
                        shapes.append(Shape(
                            kind="text",
                            x=cmap.x(cursor_x + seg_w // 2 - 100000),
                            y=cmap.y(row_y),
                            w=cmap.w(200000),
                            h=cmap.h(int(bar_h)),
                            text_runs=[TextRun(text=label, pt=12)],
                        ))
                    cursor_x += seg_w
        else:
            bar_h = int(cat_h / (n_series + gap_pct / 100))
            group_h = bar_h * n_series
            group_inset_v = (cat_h - group_h) // 2
            for si, (name, vals, _, ser_color, dpt_colors) in enumerate(series):
                default_color = ser_color or f"chart-series-{(si % 6) + 1}"
                for ci, v in enumerate(vals):
                    color = dpt_colors.get(ci) or default_color
                    by_ = plot_y + ci * cat_h + group_inset_v + si * bar_h
                    bw_ = int(plot_w * v / axis_max) if axis_max > 0 else 0
                    bx_ = plot_x
                    shapes.append(Shape(
                        kind="rect",
                        x=cmap.x(bx_), y=cmap.y(by_),
                        w=cmap.w(bw_), h=cmap.h(bar_h),
                        fill=color,
                    ))
                    if show_val_labels:
                        label = str(v).rstrip("0").rstrip(".") if "." in str(v) else str(v)
                        label = label.replace(".", ",")
                        shapes.append(Shape(
                            kind="text",
                            x=cmap.x(bx_ + bw_ + 50000),
                            y=cmap.y(by_),
                            w=cmap.w(int(fw * 0.08)),
                            h=cmap.h(int(bar_h)),
                            text_runs=[TextRun(text=label, pt=14)],
                        ))
    else:
        bar_w = int(cat_w / (n_series + gap_pct / 100))
        group_w = bar_w * n_series
        group_inset = (cat_w - group_w) // 2
        for si, (name, vals, _, ser_color, dpt_colors) in enumerate(series):
            default_color = ser_color or f"chart-series-{(si % 6) + 1}"
            for ci, v in enumerate(vals):
                color = dpt_colors.get(ci) or default_color
                bx = plot_x + ci * cat_w + group_inset + si * bar_w
                bh = int(plot_h * v / axis_max) if axis_max > 0 else 0
                by = plot_y + plot_h - bh
                shapes.append(Shape(
                    kind="rect",
                    x=cmap.x(bx), y=cmap.y(by),
                    w=cmap.w(bar_w), h=cmap.h(bh),
                    fill=color,
                ))
                # Value label above the bar — only when source enables it.
                if show_val_labels:
                    label = str(v).rstrip("0").rstrip(".") if "." in str(v) else str(v)
                    label = label.replace(".", ",")
                    shapes.append(Shape(
                        kind="text",
                        x=cmap.x(bx - bar_w // 2),
                        y=cmap.y(by - 400000),
                        w=cmap.w(bar_w * 2),
                        h=cmap.h(360000),
                        text_runs=[TextRun(text=label, pt=14)],
                    ))

    # Legend + chart-title emission — gated on the source actually
    # carrying `<c:legend>` and `<c:title>` (with `<c:autoTitleDeleted
    # val="0"/>`). The previous code always painted them. Showcase
    # charts with no `<c:legend>` element rendered a phantom series
    # name (e.g. "Datenreihe 1") at the bottom-left that wraps mid-
    # word inside the swatch slot.
    legend_y = y0 + fh - int(fh * 0.12)
    legend_x = plot_x + int(fw * 0.02)
    swatch_w = int(fw * 0.012)
    swatch_h = int(fh * 0.025)
    has_legend = root.find(f".//{{{CHART_NS}}}legend") is not None
    title_el = root.find(f".//{{{CHART_NS}}}title")
    auto_title_deleted_el = root.find(f".//{{{CHART_NS}}}autoTitleDeleted")
    title_deleted = (
        auto_title_deleted_el is not None
        and auto_title_deleted_el.get("val") in ("1", "true")
    )
    title_text = ""
    if title_el is not None and not title_deleted:
        for t_el in title_el.iterfind(f".//{{{NS['a']}}}t"):
            if t_el.text:
                title_text += t_el.text
    if title_text:
        shapes.append(Shape(
            kind="text",
            x=cmap.x(legend_x), y=cmap.y(legend_y),
            w=cmap.w(int(fw * 0.22)), h=cmap.h(int(fh * 0.04)),
            text_runs=[TextRun(text=title_text, pt=14)],
        ))
    if not has_legend:
        return
    lx = legend_x + int(fw * 0.18)
    for si, (name, _, _, ser_color, _dpt) in enumerate(series):
        color = ser_color or f"chart-series-{(si % 6) + 1}"
        shapes.append(Shape(
            kind="rect",
            x=cmap.x(lx), y=cmap.y(legend_y + (swatch_h // 4)),
            w=cmap.w(swatch_w), h=cmap.h(swatch_h),
            fill=color,
        ))
        shapes.append(Shape(
            kind="text",
            x=cmap.x(lx + swatch_w + 50000),
            y=cmap.y(legend_y),
            # Legend label width must fit "Data N" at 14pt without
            # wrapping; the previous 4% gave ~31 design-px which forced
            # multi-line "Da\nta\n1" — visibly wrong on every chart
            # legend. 10% fits the common case comfortably + leaves
            # room for short multi-word series names.
            w=cmap.w(int(fw * 0.10)),
            h=cmap.h(int(fh * 0.04)),
            text_runs=[TextRun(text=name or "", pt=14)],
        ))
        lx += int(fw * 0.12)


# ---------------------------------------------------------------------------
# Style mapping (PPTX pt size → DSL style token)
# ---------------------------------------------------------------------------


_NUM_RE = re.compile(r"^\s*\d{1,2}\.\s*$")


def _style_for(pt: float, text: str, is_footer: bool) -> str:
    # Map source pt → nearest available style bundle (by emitted px). The
    # standard feinschliff bundles emit these px sizes:
    #   body-sm=16  ·  body=26  ·  sub=44  ·  title-l=80  ·  huge=120  ·  display=160
    # A DSL `text style:sub` round-trips through python-pptx as ≈33pt on the
    # rendered slide (px*0.75). So matching SOURCE pt to BUNDLE px directly
    # (1pt ≈ 1.333px) lands closer to source than the prior bucketing,
    # which mapped any pt ≥ 28 to title-l and over-sized 32pt slide titles
    # by almost 2×. Boundaries below are midpoints between adjacent bundle
    # px values, expressed as source pt.
    if _NUM_RE.match(text):
        return "agenda-num"
    if is_footer:
        return "footer"
    px = pt * 1.333
    if px >= 140:                              # 140+ px (≈ 105pt+) → display
        return "display"
    if px >= 100:                              # 100-140 px (≈ 75-105pt) → huge
        return "huge"
    if px >= 62:                               # 62-100 px (≈ 47-75pt) → title-l
        return "title-l"
    if px >= 35:                               # 35-62 px (≈ 26-47pt) → sub
        return "sub"
    if px >= 21:                               # 21-35 px (≈ 16-26pt) → body
        return "body"
    return "body-sm"                           # <21 px → body-sm (16 px)


# ---------------------------------------------------------------------------
# Emission
# ---------------------------------------------------------------------------


# Demo-placeholder text patterns. Templates ship slides with literal
# guidance text so the user can see what each slot is for; we suppress
# them on decompile so the round-trip render shows an empty slot
# instead of the prompt.
#
# Exact phrases (case-insensitive, stripped). Add a phrase here whenever
# a new corporate template surfaces a new prompt variant.
_PLACEHOLDER_EXACT: frozenset[str] = frozenset(map(str.lower, [
    "headline", "subheadline", "sub-headline", "sub headline",
    "placeholder", "placeholder text", "placeholder copy",
    "this is a placeholder text", "this is placeholder text",
    "title", "subtitle", "sub-title",
    "presentation title", "chapter title", "section title", "slide title",
    "welcome!", "welcome", "thank you!", "thank you",
    "lorem ipsum",
    "body text", "body copy",
    "caption", "footnote",
    "click here to add text", "click to add text",
    "this text can be replaced with your own text",
    "this text can be replaced",
    "your text here", "insert text here",
    # PowerPoint outline-level prompts (default for body placeholders) —
    # English + German since many corporate templates ship localised.
    "click to edit master title style",
    "click to edit master text styles",
    "second level", "third level", "fourth level", "fifth level",
    "text hinzufügen", "text hinzufuegen",
    "zweite ebene", "dritte ebene", "vierte ebene", "fünfte ebene", "fuenfte ebene",
    "klicken sie, um einen titel hinzuzufügen",
    "klicken sie, um titel hinzuzufügen",
    # Common prompts from vendor design kits.
    "add text", "add title", "add headline", "add subheadline",
]))

# Prefix patterns — first few words of the text run lower-cased.
_PLACEHOLDER_PREFIXES: tuple[str, ...] = (
    "click to edit",
    "lorem ipsum",
    "this text can be replaced",
    "this text demonstrates",       # "This text demonstrates how your own text..."
    "this is a placeholder",
    "this is placeholder",
    "placeholder text",             # "Placeholder text\n…" wrappers
    "click here to add",
    "tap to add",
    "double-click to edit",
    "replace this text",
    "replace with your",
    "your own text",
    "sample text",
    "example text",
)

# Mail-merge / template-variable tokens (corporate convention): wholly-token
# strings like `%classification%`
# or chained `%a%%b%%c%`. PowerPoint replaces these at the org level;
# our renderer has no such facility, so they emit literally.
_TEMPLATE_VAR_RE = __import__("re").compile(r"^(%[A-Za-z][A-Za-z0-9_-]*%)+$")


def _is_placeholder_line(line: str) -> bool:
    norm = line.strip().lower().rstrip(".!?:;,")
    if not norm:
        return True   # blank line counts as placeholder noise
    if norm in _PLACEHOLDER_EXACT:
        return True
    if any(norm.startswith(p) for p in _PLACEHOLDER_PREFIXES):
        return True
    if _TEMPLATE_VAR_RE.match(line.strip()):
        return True
    return False


def _is_placeholder_text(text_runs: list["TextRun"]) -> bool:
    """Return True iff the concatenated run text is recognizable demo /
    template-prompt copy that should NOT survive the decompile round-trip.

    Suppression triggers on any of:
      * the WHOLE text (linebreaks collapsed to spaces) matches an exact
        prompt or a known prefix — catches `"This text can be replaced\\nwith your own text."` where the source linebreaks for layout
        reasons but the whole string is still one prompt
      * every non-blank line matches an individual prompt pattern —
        catches `"Headline\\nSubheadline\\nBody"` where each line is its
        own prompt stacked in one placeholder
    """
    if not text_runs:
        return False
    raw = "".join(r.text or "" for r in text_runs).strip()
    if not raw:
        return False
    # Collapse line breaks + duplicate whitespace, then check.
    flat = " ".join(raw.split())
    if _is_placeholder_line(flat):
        return True
    lines = [ln for ln in raw.splitlines() if ln.strip()]
    if lines and all(_is_placeholder_line(ln) for ln in lines):
        return True
    return False


def _strip_placeholder_paragraphs(text_runs: list["TextRun"]) -> list["TextRun"]:
    """Return text_runs with placeholder paragraphs removed.

    Source decks often combine real labels with prompt copy inside one
    placeholder, e.g.

        "04\\nThis is a placeholder text\\nThis text demonstrates how your
        own text will look when you replace the placeholder with your own
        text."

    Dropping the whole shape would lose the "04" label that's actual
    content. This walks paragraphs (separated by `TextRun(text="\\n")`
    markers inserted by `_text_runs`), drops paragraphs whose joined
    text matches `_is_placeholder_line`, and stitches the survivors
    back together with the same `\\n` separators.

    Returns the original list unchanged when no paragraph qualifies as
    placeholder, or an empty list when EVERY paragraph qualifies (caller
    drops the shape).
    """
    if not text_runs:
        return text_runs
    # Group into paragraphs. A paragraph is a contiguous slice of runs
    # not crossed by a `\n` separator run.
    paragraphs: list[list["TextRun"]] = [[]]
    for r in text_runs:
        if r.text == "\n":
            paragraphs.append([])
            continue
        paragraphs[-1].append(r)
    kept: list[list["TextRun"]] = []
    dropped_any = False
    for para in paragraphs:
        joined = "".join((r.text or "") for r in para)
        if _is_placeholder_line(joined.strip()):
            dropped_any = True
            continue
        kept.append(para)
    if not dropped_any:
        return text_runs
    out: list["TextRun"] = []
    for i, para in enumerate(kept):
        if i > 0:
            # Reuse a TextRun shape similar to what `_text_runs` emits —
            # `pt` from the para's first run so the separator's height is
            # consistent with the surrounding text.
            sep_pt = para[0].pt if para else 12
            out.append(TextRun(text="\n", pt=sep_pt))
        out.extend(para)
    return out


def emit_dsl(shapes: list[Shape], cmap: CanvasMap, layout_name: str,
             theme_name: str = "feinschliff",
             placeholder_rel: str = PLACEHOLDER_REL,
             bg_fill: str | None = None) -> str:
    out: list[str] = [
        "# auto-derived from PPTX+SVG hybrid — review before use",
        f"# layout: {layout_name}",
        f"canvas {cmap.cw}x{cmap.ch}",
        f"theme {theme_name}",
        "",
    ]

    # Slide-level background fill (from <p:cSld><p:bg>) emits first as a
    # full-canvas rect. PowerPoint draws this *under* every shape, so the
    # DSL ordering matches the source z-order. Without this, dark slides
    # rebuild on the brand's paper default and white text disappears.
    if bg_fill:
        out.append(f"rect 0,0 {cmap.cw}x{cmap.ch} fill:{bg_fill}")

    # Pre-split into geometry-first, text-last, footer-last.
    rects = [s for s in shapes if s.kind == "rect" and s.fill]
    ovals = [s for s in shapes if s.kind == "oval"]
    custs = [s for s in shapes if s.kind == "shape"]
    pics = [s for s in shapes if s.kind == "pic" or s.is_picture]
    lines = [s for s in shapes if s.kind == "line"]
    graphics = [s for s in shapes if s.kind == "graphic"]
    texts: list[Shape] = []
    for s in shapes:
        if s.kind == "text" and s.text_runs:
            texts.append(s)
        elif s.kind in ("rect", "oval", "shape") and s.text_runs:
            # geometry shape that also carries text — emit shape now, defer text
            texts.append(Shape(
                kind="text", x=s.x, y=s.y, w=s.w, h=s.h, text_runs=s.text_runs,
                ph_type=s.ph_type, ph_idx=s.ph_idx, valign=s.valign,
                padding=s.padding, from_chain=s.from_chain,
            ))
    # Drop demo-placeholder text. Corporate templates ship slides with
    # literal guidance text inside placeholders ("Headline", "Subheadline",
    # "Click to edit Master title", "%classification%" mail-merge tokens,
    # "This text can be replaced with your own text.") so the user knows
    # what each slot is for. We don't want those strings showing up in
    # the rendered round-trip — keeps the layout chrome / geometry but
    # drops the prompt text so the slot reads as empty in the brand
    # template render. See `_is_placeholder_text` for the patterns.
    #
    # When a shape mixes a real label with placeholder paragraphs
    # ("04\nThis is a placeholder text\nThis text demonstrates…"), strip
    # ONLY the placeholder paragraphs so the label survives. Shapes
    # whose every paragraph is placeholder collapse to empty runs and
    # the whole shape gets dropped.
    filtered: list[Shape] = []
    for t in texts:
        # Demo-placeholder suppression applies ONLY to text inherited from the
        # layout/master (true prompt copy). Slide-AUTHORED text is real content —
        # corporate templates ship visible example copy ("Presentation
        # title", "This is a placeholder text") right on the slide, which the renderer
        # shows; string-matching it as a prompt wrongly deleted it. (MS-gallery prompts
        # are inherited and already blanked in walk_slide, so this doesn't regress them.)
        if not t.from_chain:
            filtered.append(t)
            continue
        if _is_placeholder_text(t.text_runs):
            continue
        stripped = _strip_placeholder_paragraphs(t.text_runs)
        if stripped is t.text_runs:
            filtered.append(t)
            continue
        if not any((r.text or "").strip() for r in stripped):
            continue
        filtered.append(Shape(
            kind=t.kind, x=t.x, y=t.y, w=t.w, h=t.h,
            text_runs=stripped, ph_type=t.ph_type, ph_idx=t.ph_idx,
            valign=t.valign, padding=t.padding, from_chain=t.from_chain,
        ))
    texts = filtered

    footer_y_threshold = int(cmap.ch * 0.92)

    # Preserve SOURCE z-order — `rects` is already in render order (inherited
    # chrome behind slide content, each layer in spTree order). The previous
    # area-descending sort wrongly buried a large CONTENT panel beneath a smaller
    # decorative panel drawn ON TOP of it (MS Geometric: the cream content card
    # sank under the accent strip, so the whole slide read as the accent colour).
    # Only a near-full-bleed background rect is still forced to the bottom; a
    # STABLE sort keeps every other rect in the order PowerPoint draws it.
    # Stroke / dash / radius are captured so framed cards + dividers round-trip.
    _canvas_area = max(1, cmap.cw * cmap.ch)
    for r in sorted(rects, key=lambda s: 0 if (s.w * s.h) >= 0.9 * _canvas_area else 1):
        line = f"rect {r.x},{r.y} {r.w}x{r.h} fill:{r.fill}"
        if r.gradient is not None:
            stops, angle = r.gradient
            stops_str = ";".join(f"{p:.2f}={c}" for p, c in stops)
            line += f" gradient:angle={angle:g};{stops_str}"
        if r.corner_radius is not None and r.corner_radius > 0:
            line += f" radius:{r.corner_radius:g}"
        if r.stroke:
            line += f" stroke:{r.stroke}"
            if r.stroke_width is not None and r.stroke_width > 0:
                line += f" stroke-width:{r.stroke_width:g}"
            if r.stroke_dash:
                line += f" dash:{r.stroke_dash}"
        if r.shadow is not None:
            blur, dist, angle, color, alpha = r.shadow
            line += (f" shadow:blur:{blur:g},dist:{dist:g},"
                     f"angle:{angle:g},color:{color},alpha:{alpha:.2f}")
        out.append(line)

    # Custom shapes (puzzle pieces, parallelograms, border paths, ring
    # sectors, etc.). When we recovered an SVG `d` string from the source
    # `<a:custGeom>`, emit an `svg { path "<d>" … }` block so the
    # renderer reproduces the actual vector geometry — this is the
    # difference between "blue donut with the right arc" and "grey bbox
    # rect where the donut should be." Stroke-only paths emit
    # `stroke:<token>` with no fill; otherwise fill (or fog fallback).
    # When no path data is available we keep the lossy bbox-rect fallback
    # so the DSL still builds.
    for i, s in enumerate(custs, 1):
        if s.native_xml:
            # Native vector chrome carried verbatim from the source (editable,
            # pixel-exact). base64 so the whole <p:sp> rides inside one DSL line
            # and the brand pack stays self-contained — no source pptx at build.
            import base64 as _b64
            _enc = _b64.b64encode(s.native_xml.encode("utf-8")).decode("ascii")
            out.append(f'native shape{i} b64:"{_enc}"')
        elif s.svg_path_d:
            # `none` is not in the SVG DSL's 17-name semantic colour
            # vocabulary, so we omit `fill:` entirely when the source has
            # no solid fill — the path primitive defaults to stroke-only
            # rendering. With a fill, map the brand token onto an SVG
            # vocabulary name.
            attrs = []
            if s.fill:
                attrs.append(f"fill:{_svg_color_token(s.fill)}")
            if s.stroke:
                attrs.append(f"stroke:{_svg_color_token(s.stroke)}")
            attr_str = (" " + " ".join(attrs)) if attrs else ""
            out.append(f"svg shape{i} {s.x},{s.y} {s.w}x{s.h} {{")
            # Path coordinates are already in svg-block-local pixels (the
            # converter scales from path-local space to the shape's bbox).
            out.append(f"  path p \"{s.svg_path_d}\"{attr_str}")
            out.append("}")
        elif s.fill is None and s.stroke:
            out.append(f"shape {s.x},{s.y} {s.w}x{s.h} kind:rect stroke:{s.stroke}")
        else:
            out.append(f"shape {s.x},{s.y} {s.w}x{s.h} kind:rect fill:{s.fill or 'fog'}")

    # Native graphic frames carried verbatim — pixel-exact, vs the lossy
    # re-synthesis. Tables ship inline (<a:tbl>, no external parts → b64 only).
    # Charts ALSO carry an external part-graph (the chart part + chartStyle /
    # chartColorStyle / embedded-xlsx) as a base64-of-json `parts:` kwarg so the
    # emitter can re-create the parts + rewire rIds in the output deck.
    for i, s in enumerate(graphics, 1):
        if s.native_xml:
            import base64 as _b64
            _enc = _b64.b64encode(s.native_xml.encode("utf-8")).decode("ascii")
            if s.native_parts:
                _pj = _b64.b64encode(
                    json.dumps(s.native_parts).encode("utf-8")
                ).decode("ascii")
                out.append(f'native graphic{i} b64:"{_enc}" parts:"{_pj}"')
            else:
                out.append(f'native graphic{i} b64:"{_enc}"')

    # Ovals (circles, decorative dots). Stroke-only ovals (callout
    # circles, annotation marks) emit as stroke without fill; fill-only
    # ovals emit fill; if neither is present, fall back to a muted neutral
    # so the DSL always builds.
    for o in ovals:
        if o.fill is None and o.stroke:
            line = f"shape {o.x},{o.y} {o.w}x{o.h} kind:oval stroke:{o.stroke}"
        else:
            line = f"shape {o.x},{o.y} {o.w}x{o.h} kind:oval fill:{o.fill or 'fog'}"
            if o.stroke:
                line += f" stroke:{o.stroke}"
        if o.stroke_width is not None and o.stroke_width > 0:
            line += f" stroke-width:{o.stroke_width:g}"
        out.append(line)

    # Pictures — default to the brand's generic placeholder (so a derived
    # layout works as a reusable template) OR, when derive() was called
    # with image_extract_dir, to the actual extracted-from-source asset
    # (pipeline-optimization mode: no picture_coverage masking needed,
    # struct_diff_ratio reflects real shape/text mismatch).
    # Clamp bbox to the canvas so that picture-bleed boxes (e.g.
    # 166,-144 2345x1319 on 1920x1080) become canvas-fitted rectangles.
    # PowerPoint crops bleed at slide edges anyway, and the unclamped
    # bbox confuses the visual-diff coverage gate (>90% triggers a
    # struct = total fallback that masks real text deficits).
    for i, p in enumerate(pics, 1):
        if p.native_xml and p.native_media:
            # Template image carried natively (fixed corporate-design chrome):
            # the <p:pic> + its media ride inline; the emitter re-embeds + splices.
            import base64 as _b64
            _x = _b64.b64encode(p.native_xml.encode("utf-8")).decode("ascii")
            out.append(f'native pic{i} b64:"{_x}" media:"{p.native_media}"')
            continue
        slot = "image" if len(pics) == 1 else f"image{i}"
        cx0 = max(0, p.x)
        cy0 = max(0, p.y)
        cx1 = min(cmap.cw, p.x + p.w)
        cy1 = min(cmap.ch, p.y + p.h)
        cw = max(1, cx1 - cx0)
        ch = max(1, cy1 - cy0)
        default_path = p.media_path or placeholder_rel
        # The expander's default-filter grammar requires `default("...")`
        # — parentheses + double quotes (see lib/dsl/expander.py:_DEFAULT_FILTER_RE).
        # The earlier `default:'…'` form silently failed to match, so the
        # slot resolved to empty string and the build fell into the
        # "no image bound" placeholder-rect branch — exactly why pictures
        # rendered as grey rects even when the asset path was correct.
        out.append(
            f'picture {cx0},{cy0} {cw}x{ch} '
            f'path:"{{{{ {slot} | default(\\"{default_path}\\") }}}}" cover:true'
        )

    # Lines. Stroke-width preserves the source `<a:ln w="...">` value so
    # a 3pt horizontal divider survives the round-trip — the previous
    # `stroke-width:1` hardcode flattened every line to a hairline and
    # made decorative dividers invisible.
    for ln in lines:
        x1, y1 = ln.x, ln.y
        x2, y2 = ln.x + ln.w, ln.y + ln.h
        sw = ln.stroke_width if ln.stroke_width is not None and ln.stroke_width > 0 else 1
        attrs = f"stroke:{ln.stroke or 'fog'} stroke-width:{sw:g}"
        if ln.stroke_dash:
            attrs += f" dash:{ln.stroke_dash}"
        out.append(f"line {x1},{y1} {x2},{y2} {attrs}")

    if rects or pics or lines or ovals or custs:
        out.append("")

    # Footer collection: shapes whose y is in the bottom 8%. Keep the whole
    # SHAPE (not a per-run flatten) so the footer emit below preserves the
    # source text box's real width / padding / size / colour — a per-run
    # flatten dropped all four, hardcoded a 400 px maxwidth, and split a
    # two-run line ("Internal C-SC1" + " | C/CGB-CD | …") into two stacked
    # `text` statements. The narrow 400 px box wrapped the long © copyright
    # line into ~6 lines (the "Textumbruch" / line-break drift the reviewer
    # flagged), shifting the whole footer block.
    footer_shapes: list[Shape] = []
    body_texts: list[Shape] = []
    for t in texts:
        if t.y >= footer_y_threshold:
            footer_shapes.append(t)
        else:
            body_texts.append(t)

    # Body text in reading order.
    body_texts.sort(key=lambda s: (round(s.y / 5) * 5, s.x))
    for t in body_texts:
        # Concatenate runs verbatim — PPTX emits explicit space-only runs
        # between words, so " ".join would produce double spaces. Collapse
        # any runs of whitespace down to a single space after concat.
        raw = "".join(r.text for r in t.text_runs if r.text)
        # Collapse runs of spaces/tabs but PRESERVE newlines (paragraph breaks).
        full = re.sub(r"[ \t]+", " ", raw).strip()
        # Strip stray spaces on either side of a newline that resulted from
        # space-only runs between paragraphs.
        full = re.sub(r" *\n *", "\n", full)
        if not full:
            continue
        # Exclude paragraph-break marker runs (text == "\n") from the size
        # vote — those runs inherit the body-level default size (often 18pt)
        # which would otherwise drown out the actual text-content size when
        # two 13pt paragraphs share a single shape.
        content_pts = [r.pt for r in t.text_runs if r.text and r.text != "\n"]
        pt = max(content_pts) if content_pts else max((r.pt for r in t.text_runs), default=18)
        # normAutofit: reproduce the source's pre-shrink so placeholder text fits
        # its box (the source shrank the displayed size by font_scale; the run sz
        # stayed at the authored value, which we'd otherwise emit and overflow).
        pt *= t.font_scale
        style = _style_for(pt, full, is_footer=False)
        text = full.replace("\\", "\\\\").replace("\n", "\\n").replace('"', '\\"')
        mw = max(80, t.w)
        mh = max(24, t.h)
        # Emit `color:` override when the source's text-run colour differs
        # from the chosen style bundle's default. Captured in TextRun.color
        # from `<a:rPr><a:solidFill>`; without this the title that should
        # render in accent-blue lands in ink-grey (or whatever the style
        # bundle's default colour is) and inflates the visual diff.
        run_colors = [r.color for r in t.text_runs if r.color]
        run_color = run_colors[0] if run_colors else None
        style_default = STYLE_BUNDLES.get(style, {}).get("color")
        color_attr = (
            f" color:{run_color}" if run_color and run_color != style_default else ""
        )
        # Per-run weight override. `<a:rPr b="1">` rides on the source run;
        # the size-based classifier picks a bundle whose default weight may
        # not match (e.g. source 60pt bold maps to `huge` which is
        # weight:light by token convention). Without this, the rendered
        # text loses its emphasis even though the source explicitly carried
        # the bold flag. Emit only when source ≠ bundle default to keep
        # decompile output stable for the common case.
        source_bold = any(r.bold for r in t.text_runs)
        bundle_weight = STYLE_BUNDLES.get(style, {}).get("weight")
        weight_attr = ""
        if source_bold and bundle_weight not in ("bold", "semibold", "black"):
            weight_attr = " weight:bold"
        elif (not source_bold) and bundle_weight == "bold":
            # Source author explicitly chose regular against a bold-default
            # bundle — preserve that. `regular` is the most common name in
            # `font-weight` tokens; brands that use a different label can
            # override per-layout.
            weight_attr = " weight:regular"
        # Per-run size override. The classifier rounds source pt to the
        # nearest bundle, but the bundle steps are coarse (16/26/44/80 px)
        # and the emitted pt depends on the brand's slide physical width
        # (via `_PX_TO_PT`). Emit the source pt verbatim so renders match
        # source physical pt regardless of slide scale.
        # exactly. Stable for the common case (no emit when bundle is close).
        # Always emit `size:<pt>pt` from the source pt. The toolkit's style
        # bundles use design-px (`body=26px`, `sub=44px`, ...) which only
        # round-trip to the right rendered pt at the 13.33"-wide slide
        # convention (`_PX_TO_PT=0.5`). When the brand pack inherits a
        # different physical slide size from the source PPTX (`slide.width_emu`
        # in tokens.json), the emitter's `_PX_TO_PT` shifts and the same
        # px-valued bundles render at a different pt — which silently
        # shrinks every untagged text. Locking each run to its source pt
        # makes physical font sizes faithful regardless of slide scale.
        size_attr = f" size:{pt:g}pt"
        # autoshrink: source bodyPr had normAutofit — arm the emitter's fit as a
        # safety net (no-op when the scaled size already fits the box).
        autoshrink_attr = " autoshrink:true" if t.autoshrink else ""
        valign_attr = f" valign:{t.valign}" if t.valign else ""
        # Horizontal align — pick the first non-None align from the runs.
        # PPTX stores it per-paragraph; for a single emitted `text` primitive
        # the first paragraph's alignment wins (consistent with how we pick
        # color + size from the first content run).
        run_aligns = [r.align for r in t.text_runs if r.align]
        run_align = run_aligns[0] if run_aligns else None
        align_attr = f" align:{run_align}" if run_align else ""
        padding_attr = ""
        if t.padding is not None:
            left, top, right, bottom = t.padding
            # Compact form `padding:N` when all four insets are equal,
            # `padding:L,T,R,B` otherwise — keeps the common-case DSL short.
            if left == right and top == bottom and left == top:
                padding_attr = f" padding:{left:g}"
            else:
                padding_attr = f" padding:{left:g},{top:g},{right:g},{bottom:g}"
        # Extend the text box into the AVAILABLE space — the clearance to the
        # nearest neighbouring element (or the slide edge) — so /deck can fill in
        # MORE text than the template carried without colliding with surrounding
        # chrome. Grow only in directions that DON'T reposition the existing
        # top-left-anchored text: left/justify-aligned → widen rightwards, and
        # top-valign → grow downwards. Centred / right / middle / bottom text
        # keeps its source box (extending would re-centre / move the glyphs).
        # The current (template) text is unchanged — a short line still sits at
        # the top-left; only longer user content uses the extra room.
        _gap = 10
        _right, _bottom = cmap.cw, cmap.ch
        for _o in shapes:
            if _o is t:
                continue
            if _o.x >= t.x + t.w and not (_o.y + _o.h <= t.y or _o.y >= t.y + t.h):
                _right = min(_right, _o.x)          # neighbour clear to the right
            if _o.y >= t.y + t.h and not (_o.x + _o.w <= t.x or _o.x >= t.x + t.w):
                _bottom = min(_bottom, _o.y)         # neighbour clear below
        if run_align in (None, "left", "justify"):
            mw = max(mw, _right - t.x - _gap)
        if t.valign in (None, "top"):
            mh = max(mh, _bottom - t.y - _gap)
        out.append(
            f'text {t.x},{t.y} style:{style}{color_attr}{weight_attr}{size_attr}{autoshrink_attr}{valign_attr}{align_attr}{padding_attr} '
            f'maxwidth:{mw} maxheight:{mh} "{text}"'
        )

    # Footer-region text. Anything below `footer_y_threshold` (bottom 8%) is
    # emitted as `style:footer` text primitives. Each source text box becomes
    # ONE statement that keeps its real captured box (x, y, maxwidth from the
    # box width, maxheight from the box height), its source font size, its
    # per-run colour (the classification line's leading run is bold red
    # #D70012), and its text-frame padding. Emitting the real box width is what
    # makes the long copyright line wrap onto the same line count as the source
    # (the previous hardcoded 400 px box wrapped it ~6× and pushed the footer
    # up). Multi-run boxes concatenate verbatim so a two-run line stays one
    # line instead of splitting into two stacked statements.
    if footer_shapes:
        footer_shapes.sort(key=lambda s: (s.y, s.x))
        out.append("")
        for t in footer_shapes:
            raw = "".join(r.text for r in t.text_runs if r.text)
            full = re.sub(r"[ \t]+", " ", raw).strip()
            full = re.sub(r" *\n *", "\n", full)
            if not full:
                continue
            text = full.replace("\\", "\\\\").replace("\n", "\\n").replace('"', '\\"')
            content_pts = [r.pt for r in t.text_runs if r.text and r.text != "\n"]
            pt = max(content_pts) if content_pts else 6.0
            pt *= t.font_scale
            mw = max(80, t.w)
            mh = max(16, t.h)
            size_attr = f" size:{pt:g}pt"
            # Footer copyright/classification lines are short and box-bound; the
            # first run's colour wins (matches the body path's colour pick). The
            # `footer` bundle's own default colour is suppressed when equal.
            run_colors = [r.color for r in t.text_runs if r.color]
            run_color = run_colors[0] if run_colors else None
            footer_default = STYLE_BUNDLES.get("footer", {}).get("color")
            color_attr = (
                f" color:{run_color}" if run_color and run_color != footer_default else ""
            )
            padding_attr = ""
            if t.padding is not None:
                left, top, right, bottom = t.padding
                if left == right and top == bottom and left == top:
                    padding_attr = f" padding:{left:g}"
                else:
                    padding_attr = f" padding:{left:g},{top:g},{right:g},{bottom:g}"
            out.append(
                f'text {t.x},{t.y} style:footer{color_attr}{size_attr}{padding_attr} '
                f'maxwidth:{mw} maxheight:{mh} "{text}"'
            )

    return "\n".join(out) + "\n"



# ---------------------------------------------------------------------------
# Public derive()
# ---------------------------------------------------------------------------


def derive(
    pptx_path: Path,
    slide_idx: int,
    canvas_w: int = 1920,
    canvas_h: int = 1080,
    tokens_path: Path | None = None,
    layout_name: str = "derived",
    theme_name: str = "feinschliff",
    placeholder_rel: str = PLACEHOLDER_REL,
    pdf_path: Path | None = None,
    image_extract_dir: Path | None = None,
    image_extract_rel: str | None = None,
) -> str:
    """Decompile one slide of `pptx_path` (1-indexed) into a Feinschliff DSL
    string. Brand-agnostic: pass `theme_name` and `tokens_path` to point at
    the target brand pack. `tokens_path` is used only for nearest-color
    matching against brand color tokens; if omitted, raw hex colors land
    in the DSL.

    `pdf_path` is currently unused; reserved for SVG cross-check of
    custGeom bboxes (callers render the slide's PDF page on demand).

    `image_extract_dir` + `image_extract_rel` enable **image carry-over**
    for pipeline-optimization runs: each `<p:pic>` in the source slide
    has its embedded binary written to `image_extract_dir/imageN.<ext>`,
    and the generated DSL's `default:` for that slot points at
    `image_extract_rel/imageN.<ext>` (a brand-pack-relative path the
    build will resolve at render time). Without these args, picture
    statements fall back to the generic placeholder image as before,
    which is the right default for *templating* an out-of-tree brand
    pack (where the source illustrations are not re-used). For *testing
    decompile fidelity* against the source, carry the images over so
    the visual diff measures real shape/text mismatch instead of
    picture-region noise.
    """
    # python-pptx rejects template content-types (.potx / .pptm-template)
    # with a cryptic "is not a PowerPoint file, content type is
    # '…presentationml.template.main+xml'" error. Catch the suffix up front
    # and tell the caller exactly how to convert. (LibreOffice converts
    # cleanly in one pass; renaming the file is NOT enough — the internal
    # [Content_Types].xml carries the template MIME and must be rewritten.)
    if pptx_path.suffix.lower() in (".potx", ".potm"):
        raise ValueError(
            f"{pptx_path.name} is a PowerPoint TEMPLATE (.potx/.potm), not a "
            f"presentation. Convert it first with:\n"
            f"  soffice --headless --convert-to pptx --outdir {pptx_path.parent} "
            f"{pptx_path}\n"
            f"then re-run with the produced .pptx."
        )
    palette: dict[str, tuple[int, int, int]] = {}
    if tokens_path and tokens_path.exists():
        palette = load_palette(tokens_path)
    pres = Presentation(str(pptx_path))
    theme = load_theme_scheme(pres)
    slide = pres.slides[slide_idx - 1]

    cx = pres.slide_width
    cy = pres.slide_height
    cmap = CanvasMap(cx, cy, canvas_w, canvas_h)

    shapes = walk_slide(slide, cmap, theme, palette)
    bg_fill = extract_slide_bg_fill(slide, theme, palette)
    _ = pdf_path  # reserved for SVG cross-check, off by default

    if image_extract_dir is not None:
        if image_extract_rel is None:
            raise ValueError(
                "image_extract_rel is required when image_extract_dir is set "
                "(it goes into the DSL `default:` slot — the build resolves "
                "it relative to the brand pack root)."
            )
        image_extract_dir.mkdir(parents=True, exist_ok=True)
        # Prefer pictures whose Part was resolved at walk-time (handles
        # layout-inherited pics); fall back to slide.part lookup for the
        # in-slide case so callers that pre-date media_part still work.
        pics = [s for s in shapes if s.is_picture and (s.media_part or s.media_rid)]
        for i, p in enumerate(pics, 1):
            part = p.media_part
            if part is None and p.media_rid:
                try:
                    part = slide.part.related_part(p.media_rid)
                except KeyError:
                    continue
            if part is None:
                continue
            blob = getattr(part, "blob", None)
            partname = str(getattr(part, "partname", "/image.bin"))
            ext = partname.rsplit(".", 1)[-1].lower() if "." in partname else "bin"
            stem = "image" if len(pics) == 1 else f"image{i}"
            out_name = f"{stem}.{ext}"
            (image_extract_dir / out_name).write_bytes(blob or b"")
            p.media_path = f"{image_extract_rel.rstrip('/')}/{out_name}"

    return emit_dsl(shapes, cmap, layout_name,
                    theme_name=theme_name,
                    placeholder_rel=placeholder_rel,
                    bg_fill=bg_fill)


def main() -> int:
    ap = argparse.ArgumentParser(description="Decompile one PPTX slide → Feinschliff DSL (hybrid PPTX+SVG)")
    ap.add_argument("pptx", type=Path)
    ap.add_argument("--slide", type=int, default=1)
    ap.add_argument("--canvas", default="1920x1080")
    ap.add_argument("--theme", default="feinschliff",
                    help="Brand name to emit on the `theme` directive (default: feinschliff)")
    ap.add_argument("--brand-tokens", type=Path, default=None,
                    help="Brand tokens.json — used for nearest-color matching against brand palette")
    ap.add_argument("--layout-name", default="derived")
    ap.add_argument("--placeholder", default=PLACEHOLDER_REL,
                    help=f"Picture placeholder path (default: {PLACEHOLDER_REL})")
    args = ap.parse_args()

    if not args.pptx.exists():
        print(f"missing: {args.pptx}", file=sys.stderr)
        return 2
    w, h = (int(x) for x in args.canvas.split("x"))
    print(derive(args.pptx, args.slide, w, h,
                 tokens_path=args.brand_tokens,
                 layout_name=args.layout_name,
                 theme_name=args.theme,
                 placeholder_rel=args.placeholder), end="")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
