"""Decompiler must resolve a placeholder's EFFECTIVE font size through the
inheritance chain (layout ph -> master ph -> master txStyles) instead of
falling back to a hardcoded default when the run has no explicit size.

This is the root cause of decompiled titles coming out at the wrong size
(BSH titles inherited 20pt from the master titleStyle, but the run carried no
explicit size, so the decompiler used its 16pt fallback).
"""
from pptx import Presentation
from pptx.util import Pt

from feinschliff_builder.decompile.pptx_decompile import _PlaceholderSizeIndex


def _title_slide():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    title = slide.shapes.title
    title.text = "Inherited Title"
    return slide, title


def test_inherited_title_size_resolves_from_master_txstyles():
    slide, title = _title_slide()
    run = title.text_frame.paragraphs[0].runs[0]
    assert run.font.size is None  # truly inherited, not explicit
    idx = _PlaceholderSizeIndex(slide)
    # default template master titleStyle lvl1 = 44pt; must NOT be the 16pt fallback
    assert idx.size_pt(title, run) == 44.0


def test_explicit_run_size_wins_over_inheritance():
    slide, title = _title_slide()
    run = title.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(28)
    idx = _PlaceholderSizeIndex(slide)
    assert idx.size_pt(title, run) == 28.0


def test_non_placeholder_falls_back_to_default():
    slide, _ = _title_slide()
    box = slide.shapes.add_textbox(Pt(10), Pt(10), Pt(100), Pt(30))
    box.text_frame.text = "free text"
    run = box.text_frame.paragraphs[0].runs[0]
    idx = _PlaceholderSizeIndex(slide)
    assert idx.size_pt(box, run) == 16.0  # no inheritance for a plain textbox
