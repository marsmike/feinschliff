from __future__ import annotations

from pathlib import Path

import pytest

# python-pptx may or may not be available; guard the import.
pptx = pytest.importorskip("pptx")
from pptx import Presentation
from pptx.util import Inches

from feinschliff_builder.diagrams.refurbish.extract_vector import extract_from_slide
from feinschliff_builder.diagrams.refurbish.ir import ExtractedDiagram


def _build_five_box_flow(tmp_path: Path) -> Path:
    """Build a fixture PPTX with 5 rectangles. MSO_SHAPE.RECTANGLE = 1."""
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])  # blank
    for i in range(5):
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(1 + i * 1.5), Inches(2),
            Inches(1.2), Inches(0.8),
        )
        shape.text_frame.text = f"Step {i+1}"
    out = tmp_path / "five-box-flow.pptx"
    pres.save(str(out))
    return out


def test_vector_extractor_finds_five_boxes(tmp_path):
    src = _build_five_box_flow(tmp_path)
    pres = Presentation(str(src))
    slide = pres.slides[0]

    ir = extract_from_slide(slide)
    assert isinstance(ir, ExtractedDiagram)
    assert len(ir.nodes) == 5
    assert all(n.type == "rect" for n in ir.nodes)
    assert all(n.label.startswith("Step") for n in ir.nodes)
    assert ir.confidence == 1.0
