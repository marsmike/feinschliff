"""Tests for lib/verify/chrome.py — deterministic pp-chrome scanner
and chrome-drift cross-slide comparator."""
from __future__ import annotations

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu


_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _new_blank_pptx() -> Presentation:
    prs = Presentation()
    prs.slide_width = 1920 * 6350
    prs.slide_height = 1080 * 6350
    prs.slides.add_slide(prs.slide_layouts[6])
    return prs


def test_pp_chrome_clean_slide_no_defects():
    from feinschliff_builder.verify.chrome import scan_pp_chrome
    prs = _new_blank_pptx()
    defects = scan_pp_chrome(prs)
    assert defects == []


def test_pp_chrome_detects_drop_shadow():
    """A shape with an outerShdw effect should produce one defect."""
    from lxml import etree
    from feinschliff_builder.verify.chrome import scan_pp_chrome

    prs = _new_blank_pptx()
    slide = prs.slides[0]
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(100000), Emu(100000))
    sp_pr = shape._element.find(f".//{{{_NS_P}}}spPr")
    effect_lst = etree.SubElement(sp_pr, f"{{{_NS_A}}}effectLst")
    etree.SubElement(effect_lst, f"{{{_NS_A}}}outerShdw")

    defects = scan_pp_chrome(prs)
    assert len(defects) == 1
    assert defects[0].kind == "drop-shadow"
    assert defects[0].slide_index == 1


def test_pp_chrome_detects_grad_fill():
    from lxml import etree
    from feinschliff_builder.verify.chrome import scan_pp_chrome

    prs = _new_blank_pptx()
    slide = prs.slides[0]
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(100000), Emu(100000))
    sp_pr = shape._element.find(f".//{{{_NS_P}}}spPr")
    etree.SubElement(sp_pr, f"{{{_NS_A}}}gradFill")

    defects = scan_pp_chrome(prs)
    assert any(d.kind == "gradient-fill" for d in defects)


def test_pp_chrome_detects_fat_outline():
    from lxml import etree
    from feinschliff_builder.verify.chrome import scan_pp_chrome

    prs = _new_blank_pptx()
    slide = prs.slides[0]
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(100000), Emu(100000))
    sp_pr = shape._element.find(f".//{{{_NS_P}}}spPr")
    ln = etree.SubElement(sp_pr, f"{{{_NS_A}}}ln")
    ln.set("w", "38100")   # 3pt — fat

    defects = scan_pp_chrome(prs)
    assert any(d.kind == "fat-outline" for d in defects)


def test_pp_chrome_respects_effect_allow_marker():
    """A <p:sp> with namespaced fs:effect-allow='1' keeps its effectLst (author opt-in)."""
    from lxml import etree
    from feinschliff_builder.verify.chrome import scan_pp_chrome, _FS_NS

    prs = _new_blank_pptx()
    slide = prs.slides[0]
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(100000), Emu(100000))
    shape._element.set(f"{{{_FS_NS}}}effect-allow", "1")
    sp_pr = shape._element.find(f".//{{{_NS_P}}}spPr")
    eff = etree.SubElement(sp_pr, f"{{{_NS_A}}}effectLst")
    etree.SubElement(eff, f"{{{_NS_A}}}outerShdw")

    defects = scan_pp_chrome(prs)
    assert not any(d.kind == "drop-shadow" for d in defects)


def test_pp_chrome_respects_legacy_bare_effect_allow_marker():
    """A <p:sp> with legacy bare effect-allow='1' still keeps its effectLst
    (backward-compat read path for old decks)."""
    from lxml import etree
    from feinschliff_builder.verify.chrome import scan_pp_chrome

    prs = _new_blank_pptx()
    slide = prs.slides[0]
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(100000), Emu(100000))
    shape._element.set("effect-allow", "1")
    sp_pr = shape._element.find(f".//{{{_NS_P}}}spPr")
    eff = etree.SubElement(sp_pr, f"{{{_NS_A}}}effectLst")
    etree.SubElement(eff, f"{{{_NS_A}}}outerShdw")

    defects = scan_pp_chrome(prs)
    assert not any(d.kind == "drop-shadow" for d in defects), (
        "legacy bare effect-allow='1' must suppress drop-shadow defect (backward compat)"
    )


# ---------------------------------------------------------------------------
# chrome-drift scanner
# ---------------------------------------------------------------------------

def _two_slide_pptx():
    prs = Presentation()
    prs.slide_width = 1920 * 6350
    prs.slide_height = 1080 * 6350
    for _ in range(2):
        prs.slides.add_slide(prs.slide_layouts[6])
    return prs


def test_chrome_drift_zero_when_positions_match():
    from feinschliff_builder.verify.chrome import scan_chrome_drift

    prs = _two_slide_pptx()
    for slide in prs.slides:
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(100), Emu(200), Emu(50), Emu(20))
        sh.name = "logo-chrome"
    defects = scan_chrome_drift(prs)
    assert defects == []


def test_chrome_drift_flags_position_mismatch():
    from feinschliff_builder.verify.chrome import scan_chrome_drift

    prs = _two_slide_pptx()
    sh1 = prs.slides[0].shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(100), Emu(200), Emu(50), Emu(20))
    sh1.name = "logo-chrome"
    # Move sh2 well past the 4-design-px tolerance (25400 EMU).
    sh2 = prs.slides[1].shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(200_000), Emu(200), Emu(50), Emu(20))
    sh2.name = "logo-chrome"

    defects = scan_chrome_drift(prs)
    assert any(d.kind == "chrome-drift" for d in defects)


def test_chrome_drift_skips_non_chrome_named_shapes():
    """Shapes whose name does not match the chrome-role prefixes are skipped."""
    from feinschliff_builder.verify.chrome import scan_chrome_drift

    prs = _two_slide_pptx()
    sh1 = prs.slides[0].shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(100), Emu(200), Emu(50), Emu(20))
    sh1.name = "body-content"
    sh2 = prs.slides[1].shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(999_000), Emu(200), Emu(50), Emu(20))
    sh2.name = "body-content"

    assert scan_chrome_drift(prs) == []


def test_chrome_drift_within_tolerance():
    """Drift smaller than the tolerance is not a defect."""
    from feinschliff_builder.verify.chrome import scan_chrome_drift

    prs = _two_slide_pptx()
    sh1 = prs.slides[0].shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(100), Emu(200), Emu(50), Emu(20))
    sh1.name = "footer-meta"
    # +2 design-px = 12700 EMU; under the 4-design-px (25400) tolerance.
    sh2 = prs.slides[1].shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(12800), Emu(200), Emu(50), Emu(20))
    sh2.name = "footer-meta"
    assert scan_chrome_drift(prs) == []
