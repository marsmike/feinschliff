"""A `<p:bg><p:bgPr><a:blipFill>` background image must be extracted.

Gallery templates paint full-bleed artwork as the LAYOUT background image
(Scientific's engraved petri dish). The solid-fill extractor can't see it,
so decompiled slides rendered bare white where the source shows artwork.
"""
from __future__ import annotations

from lxml import etree
from pptx import Presentation
from pptx.opc.package import Part
from pptx.opc.packuri import PackURI

from feinschliff_builder.decompile.pptx_svg_decompile import extract_slide_bg_image

_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

_PNG = bytes.fromhex(
    "89504e470d0a1a0a0000000d4948445200000001000000010806000000"
    "1f15c4890000000a49444154789c63000100000500010d0a2db400000000"
    "49454e44ae426082"
)


def _bg_xml(rid: str) -> str:
    return (
        f'<p:bg xmlns:p="{_P}" xmlns:a="{_A}" xmlns:r="{_R}"><p:bgPr>'
        f'<a:blipFill><a:blip r:embed="{rid}"/><a:stretch><a:fillRect/></a:stretch></a:blipFill>'
        f"<a:effectLst/></p:bgPr></p:bg>"
    )


def _add_bg_image(surface) -> None:
    pkg = surface.part.package
    img = Part(PackURI("/ppt/media/bgtest.png"), "image/png", pkg, blob=_PNG)
    rid = surface.part.relate_to(img, f"{_R}/image")
    cSld = surface.element.find(f"{{{_P}}}cSld")
    cSld.insert(0, etree.fromstring(_bg_xml(rid).encode()))


def test_layout_bg_image_is_extracted():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg_image(slide.slide_layout)
    got = extract_slide_bg_image(slide)
    assert got is not None
    blob, ext = got
    assert blob == _PNG
    assert ext == "png"


def test_no_bg_image_returns_none():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    assert extract_slide_bg_image(slide) is None
