"""ECMA-376 preset-geometry conformance (spec-audit quick wins).

- roundRect default adj is 16667/100000 of the shortest side, not 0.10.
- chevron/homePlate guides use ss (shortest side) × adj with default 50000;
  the old hardcoded 30%-of-width notch stretched wide process-step chevrons.
- pentagon is a regular pentagon, not a hexagon.
- explicit <a:avLst> adjustment values are honoured.
"""
import re

from feinschliff_builder.decompile.pptx_svg_decompile import _preset_geom_path


def _points(d: str) -> list[tuple[float, float]]:
    return [(float(x), float(y))
            for x, y in re.findall(r"([\d.]+),([\d.]+)", d)]


def test_chevron_default_notch_is_half_shortest_side():
    # 400 wide × 100 tall process-step chevron: notch = 0.5 * 100 = 50px,
    # NOT 0.3 * 400 = 120px.
    pts = _points(_preset_geom_path("chevron", 400, 100))
    assert len(pts) == 6
    xs = sorted({x for x, _ in pts})
    assert 350.0 in xs, f"point should start at w - ss/2 = 350: {pts}"
    assert 50.0 in xs, f"notch should sit at ss/2 = 50: {pts}"


def test_chevron_honours_explicit_adj():
    pts = _points(_preset_geom_path("chevron", 400, 100, {"adj": 1.0}))
    xs = sorted({x for x, _ in pts})
    assert 300.0 in xs and 100.0 in xs, pts


def test_homeplate_is_five_points_with_flat_left_edge():
    pts = _points(_preset_geom_path("homePlate", 400, 100))
    assert len(pts) == 5, f"homePlate must be a 5-gon, got {pts}"
    assert sum(1 for x, _ in pts if x == 0.0) == 2, "flat left edge expected"


def test_pentagon_is_five_points_not_hexagon():
    pts = _points(_preset_geom_path("pentagon", 200, 200))
    assert len(pts) == 5, f"pentagon must be a 5-gon, got {pts}"
    assert (100.0, 0.0) in pts, "apex must be top-centre"


def test_hexagon_inset_uses_shortest_side():
    pts = _points(_preset_geom_path("hexagon", 400, 100))
    xs = sorted({x for x, _ in pts})
    # x1 = 0.25 * ss = 25, not 0.25 * w = 100
    assert 25.0 in xs and 375.0 in xs, pts


def test_right_arrow_head_uses_shortest_side():
    pts = _points(_preset_geom_path("rightArrow", 400, 100))
    xs = sorted({x for x, _ in pts})
    # head length = 0.5 * ss = 50 → shaft ends at 350, not at w/2 = 200
    assert 350.0 in xs, pts


def test_roundrect_default_adj_is_16667(tmp_path):
    """Decompile a default roundRect (no avLst) — radius must be ~0.16667
    of the shortest side at canvas scale, not 0.10."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Emu
    from feinschliff_builder.decompile.pptx_svg_decompile import derive

    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Emu(1_000_000), Emu(1_000_000),
                                 Emu(3_000_000), Emu(1_500_000))
    from pptx.dml.color import RGBColor
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(0xD0, 0xD0, 0xD0)
    # python-pptx adds an explicit adj for ROUNDED_RECTANGLE on some
    # versions — strip avLst to test the spec DEFAULT.
    from pptx.oxml.ns import qn
    geom = shp._element.find(".//" + qn("a:prstGeom"))
    av = geom.find(qn("a:avLst"))
    if av is not None:
        geom.remove(av)

    p = tmp_path / "t.pptx"
    prs.save(str(p))
    dsl = derive(p, slide_idx=1, tokens_path=None,
                 layout_name="slide-01", theme_name="t")
    m = re.search(r"radius:([\d.]+)", dsl)
    assert m, f"no radius emitted:\n{dsl}"
    # shortest side = 1_500_000 EMU ≈ 236px at 1920-wide canvas;
    # spec default radius ≈ 236 * 0.16667 ≈ 39px (old bug: ≈ 24px).
    radius = float(m.group(1))
    assert 35 <= radius <= 44, f"radius {radius} not at the 16667 spec default"


def test_style_fillref_provides_fill_when_sppr_has_none(tmp_path):
    """A shape styled purely via <p:style><a:fillRef> (no spPr fill) must
    not lose its fill — the phClr resolution gap from the spec audit."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    from pptx.util import Emu
    from lxml import etree as _etree
    from feinschliff_builder.decompile.pptx_svg_decompile import derive

    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Emu(1_000_000), Emu(1_000_000),
                                 Emu(3_000_000), Emu(1_500_000))
    sp = shp._element
    # python-pptx leaves spPr fill-less by default; ensure that and attach
    # a style block referencing accent1.
    spPr = sp.find(qn("p:spPr"))
    for tag in ("a:solidFill", "a:noFill"):
        el = spPr.find(qn(tag))
        if el is not None:
            spPr.remove(el)
    a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    style = _etree.SubElement(sp, f"{{{p}}}style")
    for ref, idx in (("lnRef", "2"), ("fillRef", "1"),
                     ("effectRef", "0"), ("fontRef", "minor")):
        el = _etree.SubElement(style, f"{{{a}}}{ref}")
        el.set("idx" if ref != "fontRef" else "idx", idx)
        clr = _etree.SubElement(el, f"{{{a}}}schemeClr")
        clr.set("val", "accent1")

    path = tmp_path / "t.pptx"
    prs.save(str(path))
    dsl = derive(path, slide_idx=1, tokens_path=None,
                 layout_name="slide-01", theme_name="t")
    rect_lines = [ln for ln in dsl.splitlines() if ln.startswith("rect ")]
    assert any("fill:" in ln for ln in rect_lines), \
        f"style-fillRef shape lost its fill:\n{dsl}"
