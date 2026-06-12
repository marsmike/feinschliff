"""Bent connectors decompile into explicit axis-aligned `line` segments.

A bentConnectorN cxnSp routes H/V elbow segments under a rot/flip transform
(org-chart trees, hierarchy diagrams). A straight line between the bbox
corners misrepresents the route, and carrying the cxnSp verbatim breaks at
render time: LibreOffice routes a connector by its a:stCxn/a:endCxn shape
references when present (the IDs point at unrelated shapes in a rebuilt
deck), and its pure xfrm+rot+flip rendering of bent presets is wrong even
for the unmodified source once the references are stripped. The decompiler
therefore resolves the preset route + rot/flip into absolute segments.
"""
import re
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.util import Emu

from feinschliff_builder.decompile.pptx_svg_decompile import derive

_LINE_RE = re.compile(r"^line (\d+),(\d+) (\d+),(\d+)", re.M)


def _derive(prs) -> str:
    with tempfile.TemporaryDirectory() as d:
        p = Path(d) / "t.pptx"
        prs.save(str(p))
        return derive(p, slide_idx=1, tokens_path=None,
                      layout_name="slide-01", theme_name="t")


def _blank_16x9():
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _add_bent_connector(slide, x, y, cx, cy, *, rot=None, flip_v=False,
                        end_ref_id=None):
    """python-pptx only creates straight connectors — rewrite one."""
    cxn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Emu(x), Emu(y), Emu(x + cx), Emu(y + cy))
    el = cxn._element
    geom = el.find(".//" + qn("a:prstGeom"))
    geom.set("prst", "bentConnector3")
    av = geom.makeelement(qn("a:avLst"), {})
    gd = av.makeelement(qn("a:gd"), {"name": "adj1", "fmla": "val 50000"})
    av.append(gd)
    geom.append(av)
    xfrm = el.find(".//" + qn("a:xfrm"))
    if rot is not None:
        xfrm.set("rot", str(rot))
    if flip_v:
        xfrm.set("flipV", "1")
    if end_ref_id is not None:
        nv = el.find(qn("p:nvCxnSpPr")).find(qn("p:cNvCxnSpPr"))
        end = nv.makeelement(qn("a:endCxn"), {"id": str(end_ref_id), "idx": "2"})
        nv.append(end)
    return cxn


def test_bent_connector_becomes_axis_aligned_segments():
    """The BSH org-chart elbow shape: rot 270° + flipV. Expected route
    (EMU, computed by hand): start (5750000,2750000) → bend (5750000,2500000)
    → bend (2750000,2500000) → end (2750000,2250000)."""
    prs, slide = _blank_16x9()
    _add_bent_connector(slide, 4_000_000, 1_000_000, 500_000, 3_000_000,
                        rot=16_200_000, flip_v=True, end_ref_id=99)
    dsl = _derive(prs)

    segs = [tuple(map(int, m)) for m in _LINE_RE.findall(dsl)]
    assert len(segs) == 3, f"expected 3 elbow segments, got {segs}"
    for x1, y1, x2, y2 in segs:
        assert x1 == x2 or y1 == y2, f"segment not axis-aligned: {(x1, y1, x2, y2)}"

    def px(emu):
        return round(emu * 1920 / 12192000)

    def py(emu):
        return round(emu * 1080 / 6858000)
    expected = {
        (px(5_750_000), py(2_500_000), px(5_750_000), py(2_750_000)),  # trunk stub
        (px(2_750_000), py(2_500_000), px(5_750_000), py(2_500_000)),  # crossbar
        (px(2_750_000), py(2_250_000), px(2_750_000), py(2_500_000)),  # end stub
    }
    for exp in expected:
        assert any(all(abs(a - b) <= 2 for a, b in zip(s, exp)) for s in segs), \
            f"missing segment ~{exp} in {segs}"

    # The connector must NOT be carried as a native payload (LibreOffice
    # re-routes it by the stale endCxn reference) nor as one diagonal line.
    assert "bentConnector" not in dsl


def test_sibling_elbows_share_trunk_segment_once():
    """Two elbows fanning out of one parent share their trunk — emit it once."""
    prs, slide = _blank_16x9()
    _add_bent_connector(slide, 4_000_000, 1_000_000, 500_000, 3_000_000,
                        rot=16_200_000, flip_v=True)
    _add_bent_connector(slide, 4_000_000, 1_000_000, 500_000, 3_000_000,
                        rot=16_200_000, flip_v=True)
    dsl = _derive(prs)
    segs = _LINE_RE.findall(dsl)
    assert len(segs) == len(set(segs)) == 3, \
        f"duplicate elbow segments not deduped: {segs}"


def test_straight_connector_still_emits_single_line():
    prs, slide = _blank_16x9()
    slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Emu(1_000_000), Emu(1_000_000),
        Emu(3_000_000), Emu(2_000_000))
    dsl = _derive(prs)
    assert len(_LINE_RE.findall(dsl)) == 1
