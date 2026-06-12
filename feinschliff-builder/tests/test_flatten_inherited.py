"""pptx_flatten_inherited — the shared layout part must never be mutated."""
from __future__ import annotations

import importlib.util
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches

_SCRIPT = Path(__file__).resolve().parents[1] / "scripts" / "pptx_flatten_inherited.py"
spec = importlib.util.spec_from_file_location("flatten_inherited", _SCRIPT)
flat = importlib.util.module_from_spec(spec)
spec.loader.exec_module(flat)


def _layout_ph_count(layout) -> int:
    return sum(1 for sh in flat._shapes(layout.shapes._spTree)
               if flat._ph(sh) is not None)


def test_flatten_does_not_strip_layout_placeholders(tmp_path):
    """Slide A leaves a layout placeholder uninstantiated -> step 3 clones
    it. The CLONE loses its <p:ph>; the layout's own shape must keep it,
    else every later slide on that layout sees its placeholders as plain
    chrome and inherits 'Add Text' sample boxes over real content."""
    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content
    before = _layout_ph_count(layout)
    s1 = prs.slides.add_slide(layout)
    # remove the content placeholder from s1 so step 3 has work to do
    for sh in list(s1.shapes):
        ph = flat._ph(sh._element)
        if ph is not None and ph.get("idx") == "1":
            sh._element.getparent().remove(sh._element)
    s2 = prs.slides.add_slide(layout)

    flat.flatten_slide(s1, None)
    assert _layout_ph_count(layout) == before, \
        "flatten_slide mutated the shared layout part"
    flat.flatten_slide(s2, None)
    # s2 instantiates everything — no prompt clones may appear
    texts = [t.text for sh in s2.shapes
             for t in sh._element.iter(qn("a:t")) if t.text]
    assert not any(t.startswith("Add Text") for t in texts)


def test_real_picture_placeholder_content_survives(tmp_path):
    """A <p:pic> carrying a real image must NOT be replaced by the bundled
    placeholder."""
    from PIL import Image
    img = tmp_path / "real.png"
    Image.new("RGB", (40, 40), (200, 30, 30)).save(img)
    placeholder = tmp_path / "placeholder.jpg"
    Image.new("RGB", (40, 40), (240, 240, 240)).save(placeholder)

    prs = Presentation()
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    pic = slide.shapes.add_picture(str(img), Inches(1), Inches(1),
                                   Inches(2), Inches(2))
    # make it a picture placeholder (ph type=pic) like template-authored decks
    nv = pic._element.find(qn("p:nvPicPr")).find(qn("p:nvPr"))
    ph = nv.makeelement(qn("p:ph"), {"type": "pic", "idx": "7"})
    nv.append(ph)

    flat.flatten_slide(slide, placeholder)
    pics = [sh for sh in slide.shapes if sh.shape_type is not None
            and sh._element.tag == qn("p:pic")]
    assert len(pics) == 1, "real picture was duplicated or replaced"
    blip = pics[0]._element.find(".//" + qn("a:blip"))
    assert blip is not None
