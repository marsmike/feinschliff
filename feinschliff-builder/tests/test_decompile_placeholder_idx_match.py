"""Slide placeholders inherit position from the layout placeholder with the
SAME idx — not from the first same-type placeholder.

A team slide carries 4 picture + 8 text placeholders, all without their own
`<a:xfrm>`; each pairs with its layout counterpart by `idx`. Matching by type
first collapsed all of them onto the layout's first slot — four photos and
eight name/role labels overprinted at one position.
"""
from __future__ import annotations

from lxml import etree
from pptx import Presentation

from feinschliff_builder.decompile.pptx_svg_decompile import _layout_placeholder_xfrm

_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _layout_body_ph(idx: int, x: int) -> str:
    return f"""
      <p:sp xmlns:p="{_P}" xmlns:a="{_A}">
        <p:nvSpPr><p:cNvPr id="{40 + idx}" name="Text Placeholder {idx}"/>
          <p:cNvSpPr/><p:nvPr><p:ph type="body" idx="{idx}"/></p:nvPr></p:nvSpPr>
        <p:spPr><a:xfrm><a:off x="{x}" y="2000000"/><a:ext cx="1500000" cy="1000000"/></a:xfrm></p:spPr>
      </p:sp>
    """


def test_same_type_placeholders_resolve_distinct_positions():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    spTree = slide.slide_layout.element.find(f".//{{{_P}}}cSld/{{{_P}}}spTree")
    spTree.append(etree.fromstring(_layout_body_ph(10, 1000000).encode()))
    spTree.append(etree.fromstring(_layout_body_ph(11, 5000000).encode()))

    x10 = _layout_placeholder_xfrm(slide, "body", "10")
    x11 = _layout_placeholder_xfrm(slide, "body", "11")
    assert x10 is not None and x11 is not None
    assert x10[0] == 1000000
    assert x11[0] == 5000000, "idx=11 must pair with ITS layout slot, not idx=10's"


def test_type_fallback_still_works_without_idx_match():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    spTree = slide.slide_layout.element.find(f".//{{{_P}}}cSld/{{{_P}}}spTree")
    spTree.append(etree.fromstring(_layout_body_ph(10, 1000000).encode()))

    # Slide ph idx=99 has no layout counterpart — legacy type fallback applies.
    x = _layout_placeholder_xfrm(slide, "body", "99")
    assert x is not None and x[0] == 1000000
