"""Tests for `feinschliff verify` CLI."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from feinschliff_builder.cli.main import main


def _clean_deck(path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(0.5))
    tb.text_frame.text = "Top left"
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4), Inches(0.5))
    tb2.text_frame.text = "Far below"
    prs.save(path)
    return path


def _overlapping_deck(path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Two textboxes whose bboxes overlap.
    a = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(2))
    a.text_frame.text = "Title overlapping content"
    b = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    b.text_frame.text = "Body that overlaps the title"
    prs.save(path)
    return path


def _out_of_bounds_deck(path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(9), Inches(0.5), Inches(4), Inches(0.5))
    tb.text_frame.text = "Spills past right edge"
    prs.save(path)
    return path


def test_verify_clean_deck_exits_zero(tmp_path, capsys):
    deck = _clean_deck(tmp_path / "clean.pptx")
    rc = main(["verify", str(deck)])
    assert rc == 0
    out = capsys.readouterr().out
    assert "clean" in out


def test_verify_overlap_deck_exits_nonzero(tmp_path, capsys):
    deck = _overlapping_deck(tmp_path / "overlap.pptx")
    rc = main(["verify", str(deck)])
    assert rc == 1
    out = capsys.readouterr().out
    assert "text-overlap" in out


def test_verify_out_of_bounds_deck_exits_nonzero(tmp_path, capsys):
    deck = _out_of_bounds_deck(tmp_path / "oob.pptx")
    rc = main(["verify", str(deck)])
    assert rc == 1
    out = capsys.readouterr().out
    assert "out-of-bounds" in out


def test_verify_ignore_flag_suppresses_failure(tmp_path, capsys):
    deck = _overlapping_deck(tmp_path / "overlap.pptx")
    rc = main(["verify", str(deck), "--ignore-overlap"])
    assert rc == 0
    out = capsys.readouterr().out
    # Still reports the defect — just doesn't fail on it.
    assert "text-overlap" in out


def test_verify_missing_deck_exits_two(tmp_path, capsys):
    rc = main(["verify", str(tmp_path / "does-not-exist.pptx")])
    assert rc == 2
    err = capsys.readouterr().err
    assert "not found" in err


def _drop_shadow_deck(path: Path) -> Path:
    """Build a deck with a shape that carries a drop shadow — should
    trigger the new pp-chrome scanner."""
    from lxml import etree
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(4), Inches(2)
    )
    ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    sp_pr = shape._element.find(f".//{{{ns_p}}}spPr")
    eff = etree.SubElement(sp_pr, f"{{{ns_a}}}effectLst")
    etree.SubElement(eff, f"{{{ns_a}}}outerShdw")
    prs.save(path)
    return path


def test_verify_pp_chrome_drop_shadow_fails(tmp_path, capsys):
    deck = _drop_shadow_deck(tmp_path / "shadow.pptx")
    rc = main(["verify", str(deck)])
    assert rc == 1
    out = capsys.readouterr().out
    assert "drop-shadow" in out


def test_verify_json_surfaces_pp_chrome(tmp_path, capsys):
    """--json should serialize the new chrome defects."""
    import json
    deck = _drop_shadow_deck(tmp_path / "shadow.pptx")
    main(["verify", str(deck), "--json"])
    out = capsys.readouterr().out
    payload = json.loads(out)
    found = any(
        d.get("kind") == "drop-shadow"
        for ds in payload["defects"].values()
        for d in ds
    )
    assert found, payload
