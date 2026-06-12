"""The decompiler extends a text slot's box into the available clearance (to the
nearest neighbour, else the slide edge) so /deck can fill in MORE text than the
template carried — without repositioning the existing top-left-anchored text.
"""
import re
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

from feinschliff_builder.decompile.pptx_svg_decompile import derive


def _derive(prs) -> str:
    with tempfile.TemporaryDirectory() as d:
        p = Path(d) / "t.pptx"
        prs.save(str(p))
        return derive(p, slide_idx=1, tokens_path=None,
                      layout_name="slide-01", theme_name="t")


def _blank_16x9():
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    return prs, prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def _maxheight(dsl: str, needle: str) -> int:
    line = next(ln for ln in dsl.splitlines() if ln.startswith("text ") and needle in ln)
    return int(re.search(r"maxheight:(\d+)", line).group(1))


def test_text_box_extends_down_to_neighbour():
    """A neighbour below caps the downward growth at the gap, not the slide edge."""
    prs, slide = _blank_16x9()
    tb = slide.shapes.add_textbox(Emu(1_000_000), Emu(1_000_000), Emu(2_000_000), Emu(800_000))
    tb.text_frame.text = "Short title"
    ob = slide.shapes.add_textbox(Emu(1_000_000), Emu(4_000_000), Emu(2_000_000), Emu(500_000))
    ob.text_frame.text = "below"
    mh = _maxheight(_derive(prs), "Short title")
    # source box ~126px tall; clearance to the neighbour ~470px; slide edge ~920px.
    # Must land in the clearance band — extended past the source, short of the edge.
    assert 300 < mh < 700, f"box not extended to neighbour clearance: maxheight={mh}"


def test_text_inside_card_rect_clamps_to_card():
    """A filled card rect that contains the text box caps growth at the card's
    edges — text in an org-chart card must not grow across the slide even when
    nothing else shares its row (slide-45 root-body regression)."""
    import re as _re
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor

    prs, slide = _blank_16x9()
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(1_000_000), Emu(1_000_000),
        Emu(2_200_000), Emu(1_000_000))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    card.line.fill.background()
    tb = slide.shapes.add_textbox(Emu(1_100_000), Emu(1_100_000),
                                  Emu(2_000_000), Emu(800_000))
    tb.text_frame.text = "Card body"
    dsl = _derive(prs)
    line = next(ln for ln in dsl.splitlines()
                if ln.startswith("text ") and "Card body" in ln)
    mw = int(_re.search(r"maxwidth:(\d+)", line).group(1))
    mh = int(_re.search(r"maxheight:(\d+)", line).group(1))
    # card right edge = 3_200_000 EMU ≈ 504px; text starts at ≈173px.
    # growth must stop at the card (≤ ~331px wide), not run to the slide edge.
    assert mw <= 340, f"text grew past its containing card: maxwidth={mw}"
    assert mh <= 170, f"text grew past its containing card: maxheight={mh}"


def test_text_box_extends_to_slide_edge_when_alone():
    """With no neighbour below, the box grows toward the slide bottom — strictly
    more room than when a neighbour caps it."""
    prs, slide = _blank_16x9()
    tb = slide.shapes.add_textbox(Emu(1_000_000), Emu(1_000_000), Emu(2_000_000), Emu(800_000))
    tb.text_frame.text = "Lone title"
    mh = _maxheight(_derive(prs), "Lone title")
    assert mh > 700, f"box did not extend toward the slide edge: maxheight={mh}"
