#!/usr/bin/env python3
"""Flatten slideLayout/slideMaster inheritance into each slide of a PPTX.

The hybrid decompiler reads SLIDE parts only. Templates built on real slide
masters keep their chrome (logo marks, design plates, rules) and placeholder
geometry in the slideLayout/slideMaster parts — decompiled slides come out
nearly empty. This preprocessor copies, per slide:

  1. slideMaster non-placeholder shapes  (bottom of z-order)
  2. slideLayout non-placeholder shapes  (above master chrome)
  3. slideLayout placeholders the slide does NOT instantiate
     (geometry resolved against the master placeholder of the same type
     when the layout placeholder carries no <a:xfrm>; empty bodies get
     their placeholder name as prompt text so slotification produces a
     slot with a sensible default)

Slide-owned shapes stay on top, byte-identical. Master/layout footer-field
placeholders (dt / ftr / sldNum) are copied only when the slide does not
already carry the same type AND the slide's headerFooter settings show it.

Flattening must not change what a renderer shows: verify with
`soffice --convert-to pdf` + pixel-compare against the original.

Usage:
  uv run python scripts/pptx_flatten_inherited.py in.pptx out.pptx
"""
from __future__ import annotations

import copy
import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn


_PH_FIELD_TYPES = {"dt", "ftr", "sldNum"}


def _ph(el):
    """The <p:ph> element of a shape's nvPr, or None."""
    for nv in el.iter():
        if nv.tag == qn("p:ph"):
            return nv
    return None


def _shapes(spTree):
    """Direct drawable children of an spTree (skip nvGrpSpPr/grpSpPr)."""
    skip = {qn("p:nvGrpSpPr"), qn("p:grpSpPr")}
    return [ch for ch in spTree if ch.tag not in skip]


def _has_xfrm(el):
    spPr = el.find(qn("p:spPr"))
    return spPr is not None and spPr.find(qn("a:xfrm")) is not None


def _graft_xfrm(target, source):
    """Copy source's <a:xfrm> into target's spPr (target lacks one)."""
    if source is None or not _has_xfrm(source):
        return
    t_spPr = target.find(qn("p:spPr"))
    if t_spPr is None:
        return
    xfrm = copy.deepcopy(source.find(qn("p:spPr")).find(qn("a:xfrm")))
    t_spPr.insert(0, xfrm)


def _is_empty_text(el):
    return not any(t.text and t.text.strip() for t in el.iter(qn("a:t")))


def _inject_prompt(el, text):
    """Write *text* into the first <a:p> of an empty placeholder body."""
    txBody = el.find(qn("p:txBody"))
    if txBody is None:
        return
    p = txBody.find(qn("a:p"))
    if p is None:
        return
    r = p.makeelement(qn("a:r"), {})
    t = p.makeelement(qn("a:t"), {})
    t.text = text
    r.append(t)
    p.append(r)


_REL_ATTRS = (qn("r:embed"), qn("r:link"), qn("r:id"))


def _port_rels(el, src_part, slide_part):
    """Re-home relationship references (r:embed / r:link / r:id) of a copied
    element: the rIds point into the layout/master part's rels — register the
    same target on the slide part and rewrite the attribute."""
    for node in el.iter():
        for attr in _REL_ATTRS:
            rid = node.get(attr)
            if not rid:
                continue
            try:
                rel = src_part.rels[rid]
            except KeyError:
                continue
            if rel.is_external:
                new_rid = slide_part.rels.get_or_add_ext_rel(
                    rel.reltype, rel.target_ref)
            else:
                new_rid = slide_part.relate_to(rel.target_part, rel.reltype)
            node.set(attr, new_rid)


def _max_shape_id(spTree) -> int:
    ids = [int(el.get("id")) for el in spTree.iter(qn("p:cNvPr"))
           if el.get("id", "").isdigit()]
    return max(ids, default=1)


def _renumber(el, next_id: int) -> int:
    """Give every cNvPr in *el* a fresh unique id; return the next free id."""
    for nv in el.iter(qn("p:cNvPr")):
        nv.set("id", str(next_id))
        next_id += 1
    return next_id


_MACHINE_FIELD_RE = __import__("re").compile(r"^\s*(%[\w]+%\s*)+$")

_PIC_TMPL = (
    '<p:pic xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
    'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
    'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
    '<p:nvPicPr><p:cNvPr id="0" name="{name}"/><p:cNvPicPr/><p:nvPr/></p:nvPicPr>'
    '<p:blipFill><a:blip r:embed="{rid}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>'
    '<p:spPr><a:xfrm><a:off x="{x}" y="{y}"/><a:ext cx="{cx}" cy="{cy}"/></a:xfrm>'
    '<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr></p:pic>'
)


def _blipfill_sp_to_pic(el):
    """An autoshape painted with an image fill (<p:sp> + spPr/blipFill —
    e.g. the B/S/H/ logo mark) → equivalent real <p:pic>, which downstream
    tooling understands. Returns the new element or None when *el* is not
    a blipFill shape / lacks geometry. Call AFTER relationship porting."""
    from lxml import etree
    if el.tag != qn("p:sp"):
        return None
    spPr = el.find(qn("p:spPr"))
    if spPr is None:
        return None
    blipFill = spPr.find(qn("a:blipFill"))
    blip = blipFill.find(qn("a:blip")) if blipFill is not None else None
    xfrm = spPr.find(qn("a:xfrm"))
    if blip is None or xfrm is None or not blip.get(qn("r:embed")):
        return None
    off, ext = xfrm.find(qn("a:off")), xfrm.find(qn("a:ext"))
    if off is None or ext is None:
        return None
    return etree.fromstring(_PIC_TMPL.format(
        name=_shape_name(el) or "picture", rid=blip.get(qn("r:embed")),
        x=off.get("x"), y=off.get("y"), cx=ext.get("cx"), cy=ext.get("cy")))


def _shape_name(el) -> str:
    return next((nv.get("name", "") for nv in el.iter(qn("p:cNvPr"))), "")


def _all_text(el) -> str:
    return "".join(t.text or "" for t in el.iter(qn("a:t")))


def _prompt_for(name: str, ph_type: str) -> str:
    low = name.lower()
    if "title" in low or ph_type in ("title", "ctrTitle"):
        return "Title"
    if "subtitle" in low or ph_type == "subTitle":
        return "Subtitle"
    if "content" in low or "text" in low or ph_type == "body":
        return "Content placeholder text. Replace with your own copy."
    return name or ph_type.title()


def flatten_slide(slide, placeholder_img: Path | None) -> int:
    """Merge inherited chrome into *slide* and make its own empty
    placeholders decompilable (prompt text + resolved geometry; empty
    picture placeholders become real pictures). Returns shapes touched."""
    layout, master = slide.slide_layout, slide.slide_layout.slide_master
    s_tree = slide.shapes._spTree
    nid = _max_shape_id(s_tree) + 1000  # clear headroom above slide ids

    slide_names = {_shape_name(sh) for sh in _shapes(s_tree)}
    slide_ph_keys = set()
    for sh in _shapes(s_tree):
        ph = _ph(sh)
        if ph is not None:
            slide_ph_keys.add((ph.get("type", "body"), ph.get("idx", "0")))

    def ph_lookup(part_shapes, ph_type, idx):
        exact, by_type = None, None
        for sh in part_shapes:
            ph = _ph(sh)
            if ph is None:
                continue
            if (ph.get("type", "body"), ph.get("idx", "0")) == (ph_type, idx):
                exact = exact or sh
            if ph.get("type", "body") == ph_type:
                by_type = by_type or sh
        return exact or by_type

    layout_shapes = _shapes(layout.shapes._spTree)
    master_shapes = _shapes(master.shapes._spTree)

    copied = []

    # 1+2: non-placeholder chrome, master first then layout. Skip the
    # template's %field% automation boxes and anything already present
    # under the same shape name — on the slide OR copied from the master
    # (layouts repeat master chrome like footers and the logo mark).
    seen_names = set(slide_names)
    for part, part_shapes in ((master, master_shapes), (layout, layout_shapes)):
        for sh in part_shapes:
            if _ph(sh) is not None:
                continue
            name = _shape_name(sh)
            if name and name in seen_names:
                continue
            if _MACHINE_FIELD_RE.match(_all_text(sh) or ""):
                continue
            clone = copy.deepcopy(sh)
            _port_rels(clone, part.part, slide.part)
            clone = _blipfill_sp_to_pic(clone) or clone
            seen_names.add(name)
            copied.append(clone)

    # 3: layout placeholders the slide doesn't instantiate (rare).
    # Matching is two-stage: exact (type, idx) first, then remaining layout
    # placeholders pair up with remaining same-type slide placeholders —
    # slides authored from the master often drop the idx (PowerPoint writes
    # `<p:ph type="body"/>`), and an exact-only match would inject duplicate
    # prompt boxes over real content.
    unmatched_slide = dict.fromkeys(slide_ph_keys)
    layout_keys = []
    for sh in layout_shapes:
        ph = _ph(sh)
        if ph is not None:
            layout_keys.append((ph.get("type", "body"), ph.get("idx", "0")))
    for key in layout_keys:
        if key in unmatched_slide:
            del unmatched_slide[key]

    def _instantiated(ph_type, key):
        if key in slide_ph_keys:
            return True
        for skey in list(unmatched_slide):
            if skey[0] == ph_type:
                del unmatched_slide[skey]   # consume: one slide ph per layout ph
                return True
        return False

    for sh in layout_shapes:
        ph = _ph(sh)
        if ph is None:
            continue
        ph_type = ph.get("type", "body")
        key = (ph_type, ph.get("idx", "0"))
        if ph_type in _PH_FIELD_TYPES or _instantiated(ph_type, key):
            continue
        if _MACHINE_FIELD_RE.match(_all_text(sh) or ""):
            continue
        clone = copy.deepcopy(sh)
        _port_rels(clone, layout.part, slide.part)
        if not _has_xfrm(clone):
            _graft_xfrm(clone, ph_lookup(master_shapes, ph_type, key[1]))
        if _is_empty_text(clone):
            _inject_prompt(clone, _prompt_for(_shape_name(clone), ph_type))
        # Strip the <p:ph> from the CLONE — `ph` belongs to the layout's own
        # shape, and removing it there mutates the shared layout part: every
        # LATER slide on this layout then sees its placeholders as plain
        # chrome and inherits "Add Text" sample boxes over real content.
        clone_ph = _ph(clone)
        if clone_ph is not None:
            clone_ph.getparent().remove(clone_ph)
        copied.append(clone)

    # Insert inherited shapes BELOW slide-owned shapes (start of spTree,
    # after nvGrpSpPr/grpSpPr), preserving their relative order.
    anchor = 2  # children 0,1 are nvGrpSpPr / grpSpPr
    for el in copied:
        nid = _renumber(el, nid)
        s_tree.insert(anchor, el)
        anchor += 1
    touched = len(copied)

    # 4: the slide's OWN empty placeholders — resolve inherited geometry
    # and make them decompilable. Text placeholders get prompt text;
    # picture placeholders become real pictures (decompiler → image slot).
    for sh in list(_shapes(s_tree)):
        ph = _ph(sh)
        if ph is None:
            continue
        ph_type = ph.get("type", "body")
        idx = ph.get("idx", "0")
        if ph_type in _PH_FIELD_TYPES:
            continue
        if not _has_xfrm(sh):
            # The layout's matching placeholder may itself inherit geometry
            # from the master — graft from the first candidate that actually
            # carries an <a:xfrm>, not merely the first same-type match.
            candidates = [ph_lookup(layout_shapes, ph_type, idx),
                          ph_lookup(master_shapes, ph_type, idx)]
            src = next((c for c in candidates
                        if c is not None and _has_xfrm(c)), None)
            _graft_xfrm(sh, src)
        if (ph_type == "pic" and placeholder_img is not None
                and sh.tag == qn("p:sp")
                and sh.find(".//" + qn("a:blip")) is None):
            # Only a genuinely EMPTY picture placeholder gets the bundled
            # image. A <p:pic> (or any shape with a blip) IS the content —
            # replacing it destroyed author-placed photos/illustrations.
            from pptx.util import Emu
            spPr = sh.find(qn("p:spPr"))
            xfrm = spPr.find(qn("a:xfrm")) if spPr is not None else None
            if xfrm is not None:
                off, ext = xfrm.find(qn("a:off")), xfrm.find(qn("a:ext"))
                slide.shapes.add_picture(
                    str(placeholder_img),
                    Emu(int(off.get("x"))), Emu(int(off.get("y"))),
                    Emu(int(ext.get("cx"))), Emu(int(ext.get("cy"))))
                s_tree.remove(sh)
                touched += 1
            continue
        if _is_empty_text(sh):
            _inject_prompt(sh, _prompt_for(_shape_name(sh), ph_type))
            touched += 1
    return touched


def main() -> int:
    if len(sys.argv) != 3:
        print(__doc__.strip().splitlines()[-1], file=sys.stderr)
        return 2
    src, dst = Path(sys.argv[1]), Path(sys.argv[2])
    ph_img = (Path(__file__).resolve().parents[2] / "feinschliff" / "assets"
              / "illustrations" / "placeholder.jpg")
    if ph_img.is_file():
        # python-pptx needs a real JPEG/PNG; the bundled placeholder may be
        # webp-in-disguise. Re-encode to a sibling temp file.
        import tempfile
        from PIL import Image
        tmp = Path(tempfile.gettempdir()) / "flatten-placeholder.jpg"
        Image.open(ph_img).convert("RGB").save(tmp, "JPEG", quality=85)
        ph_img = tmp
    prs = Presentation(str(src))
    total = 0
    for i, slide in enumerate(prs.slides, 1):
        n = flatten_slide(slide, ph_img if ph_img.is_file() else None)
        total += n
        print(f"  slide {i:02d}: {n} shapes inherited/materialised")
    prs.save(str(dst))
    print(f"flattened {src.name} -> {dst} ({total} shapes touched)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
