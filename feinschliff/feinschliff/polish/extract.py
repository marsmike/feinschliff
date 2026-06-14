from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path


@dataclass
class SlideText:
    index: int          # 1-based
    title: str          # "" if missing
    body: str           # bullets joined by \n
    bullets: list[str]  # raw bullet texts in reading order
    notes: str          # speaker notes; "" if absent
    has_image: bool     # True if any non-placeholder picture shape
    has_chart: bool     # True if any chart shape
    has_table: bool     # True if any table shape


def _find_title_shape(slide):
    """Return (element, title_text) for the best title candidate, or (None, '').

    Returns the underlying XML *element* (not the wrapper object) so that
    callers can reliably skip the same shape in a second iteration (python-pptx
    creates a fresh wrapper on every ``slide.shapes`` access, so ``is`` checks
    across two iterations always fail).

    Priority:
    1. PPTX title placeholder (shapes.title or placeholder idx==0).
    2. Topmost free-floating textbox (smallest top coordinate).
    """
    # 1. Named title placeholder
    try:
        ts = slide.shapes.title
        if ts is not None and getattr(ts, "has_text_frame", False):
            t = ts.text_frame.text.strip()
            if t:
                return ts._element, t.split("\n")[0]
    except (AttributeError, ValueError):
        pass

    # 2. Placeholder idx == 0
    for sh in slide.shapes:
        try:
            pf = sh.placeholder_format
        except (AttributeError, ValueError):
            continue
        if pf is not None and pf.idx == 0 and getattr(sh, "has_text_frame", False):
            t = sh.text_frame.text.strip()
            if t:
                return sh._element, t.split("\n")[0]

    # 3. Topmost free-floating textbox (skip auto-shapes that have their own label)
    best_el = None
    best_top = None
    best_text = ""
    for sh in slide.shapes:
        try:
            _ = sh.auto_shape_type
            continue  # skip named auto-shapes (boxes, ellipses — diagram nodes)
        except (ValueError, AttributeError):
            pass
        if not getattr(sh, "has_text_frame", False):
            continue
        t = sh.text_frame.text.strip()
        if not t:
            continue
        top = getattr(sh, "top", None)
        if top is None:
            continue
        if best_top is None or top < best_top:
            best_el = sh._element
            best_top = top
            best_text = t.split("\n")[0]

    if best_el is not None:
        return best_el, best_text

    return None, ""


def extract_text_from_pptx(path: Path) -> list[SlideText]:
    """Extract text, notes and shape-type flags from every slide in *path*."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(path)
    results: list[SlideText] = []

    for idx, slide in enumerate(prs.slides, start=1):
        title_el, title_str = _find_title_shape(slide)

        bullets: list[str] = []
        for sh in slide.shapes:
            if title_el is not None and sh._element is title_el:
                continue
            if not getattr(sh, "has_text_frame", False):
                continue
            for para in sh.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    bullets.append(t)

        # Speaker notes
        notes_str = ""
        if slide.has_notes_slide:
            try:
                notes_str = slide.notes_slide.notes_text_frame.text.strip()
            except (AttributeError, ValueError):
                pass

        # Shape-type flags
        has_image = any(
            sh.shape_type == MSO_SHAPE_TYPE.PICTURE
            for sh in slide.shapes
            if not getattr(sh, "is_placeholder", False)
        )
        has_chart = any(getattr(sh, "has_chart", False) for sh in slide.shapes)
        has_table = any(getattr(sh, "has_table", False) for sh in slide.shapes)

        results.append(
            SlideText(
                index=idx,
                title=title_str,
                body="\n".join(bullets),
                bullets=bullets,
                notes=notes_str,
                has_image=has_image,
                has_chart=has_chart,
                has_table=has_table,
            )
        )

    return results
