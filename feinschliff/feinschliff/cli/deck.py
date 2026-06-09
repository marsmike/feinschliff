"""`feinschliff deck …` subcommands: multi-slide composer + layout picker.

Subcommands:

  feinschliff deck build <plan.yaml> [-o OUT.pptx]
      Build one .pptx from a plan that lists N slides, each pinning a
      layout and inline content (or a content file).

  feinschliff deck pick <signals.yaml>
      Print a recommended layout id based on structured signals.

  feinschliff deck wireframe <layout.slide.dsl> --brand <id> [-o out.svg]
      Render a DSL layout as an annotated SVG wireframe (bounding boxes for
      text slots, picture slots, rect backgrounds). No PPTX round-trip needed.
      Add --overlay-pptx <file.pptx> to embed the actual rendered slide behind
      the wireframe boxes for pixel-accurate deviation analysis.

  feinschliff deck wireframe-sheet <plan.yaml> [-o sheet.svg]
      Render every slide in a plan as a wireframe and compose them into a
      single SVG contact sheet. Useful as a fast layout-regression baseline:
      store the SVGs in git, regenerate after DSL changes, and diff.

Plan schema (build):

  brand: feinschliff                       # default brand for slides
  out:   deck.pptx                         # output path; --output overrides
  slides:
    - layout: layouts/title-orange.slide.dsl
      content:
        pgmeta: "Q1 2026"
        title:  "..."
    - layout: layouts/quote.slide.dsl
      content_file: examples/v2/quote.yaml  # alternative to inline
    - layout: brands/gs-ramspau/layouts/stundenplan.slide.dsl
      brand:  gs-ramspau                    # per-slide override
      content_file: brands/gs-ramspau/examples/v2/stundenplan.yaml
"""
from __future__ import annotations

import argparse
import os
import sys
import time
import tempfile
from pathlib import Path

import yaml

from feinschliff.deck.orchestrate import (
    patch_set_hash as _patch_set_hash_fn,
    build_primitives_for_layout as _build_primitives_for_layout_fn,
    build_refurbished_deck as _build_refurbished_deck_fn,
)

from feinschliff.dsl.parser import parse_file
from feinschmiede.dsl.tokens import load_tokens
from feinschliff.dsl.expander import (
    interpolate_nodes,
    expand_compounds,
    load_compounds_for_brand,
)
from feinschliff.dsl.pptx_emit import build_multi_slide
from feinschliff.content_validator import (
    emit_defects_and_abort_message, validate_content,
)
from feinschliff.slot_budget import compute_slot_budgets
from feinschliff.pipeline import compile_slide
from feinschliff.defects import fatal_kinds, format_defect
from feinschmiede.brand_discovery import find_brand
from feinschliff.io.image_provider import discover_providers, get_provider

# Builder-side imports (feinschliff-builder optional dependency).
# Fall back to no-ops when builder is not installed.
try:
    from feinschliff_builder.verify.deck.notes_budget import validate_notes  # type: ignore[import]
    from feinschliff_builder.verify.deck.titles import extract_titles_from_plan  # type: ignore[import]
    from feinschliff_builder.verify.deck.storyline import render_contact_sheet, write_storyline_report  # type: ignore[import]
    from feinschliff_builder.verify.deck.claim_evidence import judge_plan, write_report as write_claim_evidence_report  # type: ignore[import]
except ImportError:
    validate_notes = None  # type: ignore[assignment]
    extract_titles_from_plan = None  # type: ignore[assignment]
    render_contact_sheet = None  # type: ignore[assignment]
    write_storyline_report = None  # type: ignore[assignment]
    judge_plan = None  # type: ignore[assignment]
    write_claim_evidence_report = None  # type: ignore[assignment]
from feinschliff.pipeline_log import log_event


def _require_builder(feature: str) -> None:
    """Raise SystemExit with a helpful message if feinschliff-builder isn't installed."""
    try:
        import feinschliff_builder  # noqa: F401
    except ImportError as e:
        sys.stderr.write(
            f"error: '{feature}' requires the feinschliff-builder plugin.\n"
            f"  Install it with: uv add feinschliff-builder  (or via Claude Code marketplace)\n"
        )
        raise SystemExit(2) from e


def _bundled_assets() -> Path:
    """Return the assets/ directory shipped inside this plugin."""
    return Path(__file__).resolve().parents[1] / "assets"


def _bundled_compounds() -> Path:
    """Return the compounds/ directory shipped inside this plugin."""
    return Path(__file__).resolve().parents[1] / "compounds"


def _find_toolkit_file(rel: str) -> Path | None:
    """Resolve *rel* against each discovered layout dir's parent (i.e. the
    plugin root), returning the first match.  Used to replace hard-coded
    ``REPO_ROOT / rel`` fallback lookups in the deck CLI.
    """
    from feinschliff.layout_discovery import all_layout_dirs
    for layout_dir in all_layout_dirs():
        candidate = (layout_dir.parent / rel).resolve()
        if candidate.is_file():
            return candidate
    return None


def register(parser: argparse.ArgumentParser) -> None:
    sub = parser.add_subparsers(dest="deck_command", required=True)

    p_build = sub.add_parser("build", help="Build a multi-slide deck from a plan YAML")
    p_build.add_argument("plan", help="Path to the deck plan YAML")
    p_build.add_argument("-o", "--output", help="Output .pptx (overrides plan.out)")
    p_build.add_argument(
        "--skip-content-lint",
        action="store_true",
        help="Skip pre-render content lints (title-length, action-verb-leading). "
             "For emergency overrides only.",
    )
    p_build.add_argument(
        "--allow-diagram-warnings",
        action="store_true",
        help="Ship even when diagram-overflow or diagram-text-too-small "
             "defects surface. Otherwise these are fatal by default — same "
             "policy as single-slide `feinschliff build`.",
    )
    p_build.add_argument(
        "--allow-missing-assets",
        action="store_true",
        help="Ship even when a picture slot points at a missing file or is "
             "unset. Default: fatal. Mark intentionally-empty slots with "
             "`optional:true` to skip the abort without this flag.",
    )
    p_build.add_argument(
        "--strict-static",
        action="store_true",
        help="Run the pre-render static geometry verifier (feinschliff_builder.verify.static) "
             "before compile. Aborts with exit 1 and prints defects when any "
             "slot-overflow or empty-placeholder issues are detected. Off by "
             "default — opting in avoids surprising existing automation.",
    )
    p_build.add_argument(
        "--autofix",
        action="store_true",
        help="Run the static verifier before compile and automatically apply "
             "mechanical fixes (shorten_slot, delete_word, drop_bullet, "
             "swap_layout_*) for known defect classes.  Up to 3 inner fix "
             "cycles are attempted; residual defects are printed but do NOT "
             "block the compile.  The fixed plan is written back to disk "
             "before compile.",
    )
    p_build.set_defaults(func=cmd_build)

    p_pick = sub.add_parser("pick", help="Recommend a layout for the given signals")
    p_pick.add_argument("signals", help="Path to a signals YAML (or '-' for stdin)")
    p_pick.add_argument("--top-k", type=int, default=3,
                        help="Print the top-K candidates with scores (default 3)")
    p_pick.set_defaults(func=cmd_pick)

    p_storyline = sub.add_parser(
        "storyline",
        help="Emit a title-only contact sheet from a deck_plan.json / "
             "content_plan.json (Phase 2 storyline gate).",
    )
    p_storyline.add_argument("plan", help="Path to deck_plan.json or content_plan.json")
    p_storyline.add_argument(
        "-o", "--output", required=True,
        help="Output path for the storyline_report.md",
    )
    p_storyline.add_argument(
        "--brief-summary", default=None,
        help="Optional one-line summary shown above the contact sheet.",
    )
    p_storyline.set_defaults(func=cmd_storyline)

    p_ce = sub.add_parser(
        "claim-evidence",
        help="Mid-plan claim-evidence text gate (step 2b). "
             "Judges each claim-carrying slide for title-body coherence "
             "before render. Cheap Haiku pass, no PPTX round-trip.",
    )
    p_ce.add_argument("plan", help="Path to plan.yaml or plan.json")
    p_ce.add_argument(
        "--design-brief",
        default=None,
        help="Path to design_brief.json (optional — enables per-slide claim hints).",
    )
    p_ce.add_argument(
        "-o", "--output", required=True,
        help="Output path for claim_evidence_report.md",
    )
    p_ce.add_argument(
        "--offline", action="store_true",
        help="Skip all LLM calls; return clean verdicts (for testing).",
    )
    p_ce.add_argument(
        "--model", default="claude-haiku-4-5-20251001",
        help="Model to use for judgment (default: claude-haiku-4-5-20251001).",
    )
    p_ce.set_defaults(func=cmd_claim_evidence)

    p_wf = sub.add_parser(
        "wireframe",
        help="Render a DSL layout as an annotated SVG wireframe.",
    )
    p_wf.add_argument("layout", help="Path to a .slide.dsl file")
    p_wf.add_argument("--brand", required=True, help="Brand id (dir name under brands/)")
    p_wf.add_argument("--content", help="YAML file with slot values (optional)")
    p_wf.add_argument(
        "--show-slots",
        action="store_true",
        help="Preserve {{ slot_name }} labels even when --content is supplied. "
             "Forces skip_interpolation=True so the wireframe shows slot structure.",
    )
    p_wf.add_argument("-o", "--output", required=True, help="Output .svg path")
    p_wf.add_argument(
        "--overlay-pptx",
        help="Path to a .pptx file whose first slide is embedded as background. "
             "Requires LibreOffice on PATH.",
    )
    p_wf.add_argument(
        "--overlay-slide", type=int, default=0,
        help="0-based slide index to use from --overlay-pptx (default 0).",
    )
    p_wf.add_argument(
        "--overlay-opacity", type=float, default=0.55,
        help="Opacity of the background slide image (0.0–1.0, default 0.55).",
    )
    p_wf.set_defaults(func=cmd_wireframe)

    p_wfs = sub.add_parser(
        "wireframe-sheet",
        help="Render all slides in a plan as a SVG wireframe contact sheet.",
    )
    p_wfs.add_argument("plan", help="Path to the deck plan YAML")
    p_wfs.add_argument("-o", "--output", required=True, help="Output .svg path")
    p_wfs.add_argument(
        "--overlay-pptx",
        help="Path to a .pptx file; each slide is embedded behind its wireframe.",
    )
    p_wfs.add_argument(
        "--overlay-opacity", type=float, default=0.55,
        help="Opacity of the background slide images (0.0–1.0, default 0.55).",
    )
    p_wfs.add_argument(
        "--show-slots",
        action="store_true",
        help="Preserve {{ slot_name }} labels in every cell — skip interpolation "
             "even when plans provide inline content or content_file. The contact "
             "sheet is intended as a layout-regression baseline, so showing slot "
             "structure is often more useful than filled content.",
    )
    p_wfs.set_defaults(func=cmd_wireframe_sheet)

    p_polish = sub.add_parser(
        "polish",
        help="Refurbish old PPTX diagrams into brand-perfect DSL artifacts.",
    )
    p_polish.add_argument("input", help=".pptx file to refurbish")
    p_polish.add_argument("-o", "--output", required=True, help="output .pptx path")
    p_polish.add_argument("--brand", default="feinschliff")
    p_polish.add_argument(
        "--refurbish-all",
        action="store_true",
        help="Auto-accept all refurbish proposals (emit DSL for every diagram slide).",
    )
    p_polish.add_argument(
        "--no-refurbish",
        action="store_true",
        help="Skip refurbish; just copy the input to the output path unchanged.",
    )
    p_polish.add_argument(
        "--refurbish-default",
        choices=("excalidraw", "svg"),
        default=None,
        help="Force a specific emitter instead of letting kind_selector choose.",
    )
    p_polish.set_defaults(func=cmd_polish)

    p_book = sub.add_parser(
        "book",
        help="Render an annotated speaker-book PDF from a deck plan + "
             "design brief. Front matter (takeaway, audience, frame, "
             "red_line, hook) + one page per slide with the rendered "
             "thumbnail, claim, speaker notes, audience_fit, and role.",
    )
    p_book.add_argument("plan", help="Path to the deck plan YAML.")
    p_book.add_argument(
        "--design-brief", required=True,
        help="Path to design_brief.json (front matter + per-slide role / "
             "audience_fit come from here).",
    )
    p_book.add_argument(
        "--pptx", default=None,
        help="Path to the rendered .pptx (used to extract per-slide PNG "
             "thumbnails). If omitted, the book is rendered without "
             "thumbnails — useful for fast preview while authoring.",
    )
    p_book.add_argument(
        "-o", "--output", required=True,
        help="Output path for the speaker-book .pdf.",
    )
    p_book.set_defaults(func=cmd_book)

    # `deck plan` is a thin alias for `deck storyline` — same CLI surface,
    # same handler. The mode-level work (steps 0 → 1c only, no render) is
    # the skill orchestrator's job; the CLI part is just the storyline
    # report materialization helper. See
    # skills/deck/references/modes.md::plan for the mode semantics.
    p_plan = sub.add_parser(
        "plan",
        help="Emit a title-only storyline report (alias for `deck storyline`). "
             "Implements the CLI surface of the /deck plan mode — steps 0 → 1c "
             "only, no render. See skills/deck/references/modes.md::plan.",
    )
    p_plan.add_argument("plan", help="Path to deck_plan.json or content_plan.json")
    p_plan.add_argument(
        "-o", "--output", required=True,
        help="Output path for the storyline_report.md",
    )
    p_plan.add_argument(
        "--brief-summary", default=None,
        help="Optional one-line summary shown above the contact sheet.",
    )
    p_plan.set_defaults(func=cmd_storyline)

    # ── Pipeline timing logs + parallel plan authoring ──────────────────
    # (`log-event`, `timing`, `plan-skeleton`, `plan-merge`) — extracted
    # into deck_subcommands/plan_log.py to keep this file from growing.
    from feinschliff.cli.deck_subcommands import plan_log
    plan_log.register(sub)

    # ── Parallel-verify aspect helpers ──────────────────────────────────
    p_aspect = sub.add_parser(
        "verify-aspect",
        help="Run one focused verify aspect on a built deck. Each aspect "
             "is independently runnable; spawning multiple aspects in "
             "parallel + collating gives the parallel-verify path.",
    )
    p_aspect.add_argument(
        "aspect",
        choices=["bbox", "font", "narrative", "brand", "image", "content",
                 "notes-coherence"],
        help="bbox = bounding-box / overflow; font = legibility / role "
             "mismatches; narrative = SCQA / claim-title across titles; "
             "brand = color contrast / token discipline; image = picture "
             "slot fit / style consistency; content = title-body coherence; "
             "notes-coherence = per-slide speaker notes track the deck's "
             "red_line.",
    )
    p_aspect.add_argument("--plan", required=True,
                          help="Path to the deck plan.yaml.")
    p_aspect.add_argument("--pptx", default=None,
                          help="Path to the built .pptx (required for bbox/font/brand/image).")
    p_aspect.add_argument("--png-dir", default=None,
                          help="Directory with rendered slide-NN.png files.")
    p_aspect.add_argument("--design-brief", default=None,
                          help="Path to design_brief.json (used by narrative).")
    p_aspect.add_argument("-o", "--output", required=True,
                          help="Output path for verify-<aspect>.json.")
    p_aspect.set_defaults(func=cmd_verify_aspect)

    p_vs = sub.add_parser(
        "verify-static",
        help="Pre-render static geometry verify: detect slot-overflow and "
             "empty-placeholder defects from a plan.yaml without rendering. "
             "Exit 0 = clean, 1 = defects found, 2 = plumbing error.",
    )
    p_vs.add_argument("plan", help="Path to the deck plan YAML")
    p_vs.add_argument(
        "--json",
        action="store_true",
        help="Emit defects as a JSON array to stdout instead of the human "
             "readable format. Shape: [{slide_index, kind, severity, "
             "message, meta}, ...]",
    )
    p_vs.add_argument(
        "--brand", default=None,
        help="Override brand. Default: from plan.brand or 'feinschliff'.",
    )
    p_vs.set_defaults(func=cmd_verify_static)

    p_af = sub.add_parser(
        "apply-fixes",
        help="Apply mechanical fixes to a plan.yaml from a verify-static "
             "defects JSON.  Mutates plan in place (or writes -o out.yaml). "
             "Exit 0 = patches applied, 1 = no patches applied.",
    )
    p_af.add_argument("plan", help="Path to the deck plan YAML to fix.")
    p_af.add_argument(
        "--defects", required=True,
        help="Path to a defects JSON file: either a flat list of Defect "
             "dicts [{slide_index, kind, severity, message, meta}, ...] "
             "as emitted by `deck verify-static --json`, or a "
             '{"defects": {slide_idx: [...]}, ...} collated shape.',
    )
    p_af.add_argument(
        "-o", "--output",
        help="Output path for the fixed plan YAML.  When omitted the plan "
             "is updated in place.",
    )
    p_af.add_argument(
        "--brand", default=None,
        help="Override brand. Default: from plan.brand or 'feinschliff'.",
    )
    p_af.set_defaults(func=cmd_apply_fixes)

    p_collate = sub.add_parser(
        "verify-collate",
        help="Merge per-aspect verify outputs into a single verify_report.md.",
    )
    p_collate.add_argument(
        "--aspect", action="append", default=[], metavar="PATH",
        help="One per aspect: path to a verify-<aspect>.json file.",
    )
    p_collate.add_argument("--plan", required=True,
                           help="Path to the deck plan.yaml (for slide titles).")
    p_collate.add_argument("--iteration", type=int, default=1,
                           help="Iteration number (header field).")
    p_collate.add_argument("--budget", type=int, default=3,
                           help="Iteration budget (header field).")
    p_collate.add_argument("--png-dir", default=None,
                           help="PNG directory (header reference).")
    p_collate.add_argument("-o", "--output", required=True,
                           help="Output path for verify_report.md.")
    p_collate.set_defaults(func=cmd_verify_collate)


def _patch_set_hash(patches: list) -> str:
    """Delegate to feinschliff.deck.orchestrate.patch_set_hash."""
    return _patch_set_hash_fn(patches)


def cmd_build(args) -> int:
    plan_path = Path(args.plan).resolve()
    if not plan_path.is_file():
        print(f"deck: plan not found: {plan_path}", file=sys.stderr)
        return 2
    plan = yaml.safe_load(plan_path.read_text()) or {}

    # Surface the soft degradation: without feinschliff-builder, per-slide
    # speaker-notes validation is silently skipped. Print a one-line hint so
    # operators know what's missing instead of finding out via a downstream
    # surprise. Suppress with FEINSCHLIFF_QUIET_NOTES_BUDGET=1 for CI runs
    # that intentionally ship without builder.
    if validate_notes is None and not os.environ.get("FEINSCHLIFF_QUIET_NOTES_BUDGET"):
        print(
            "deck build: notes-budget validation skipped "
            "(feinschliff-builder not installed). "
            "Install feinschliff-builder to enable per-slide notes lint, "
            "or set FEINSCHLIFF_QUIET_NOTES_BUDGET=1 to silence this hint.",
            file=sys.stderr,
        )

    default_brand = plan.get("brand", "feinschliff")
    # Notes lint reads the deck-level verbosity (mirrors design_brief.verbosity)
    # to pick a per-slide word budget. Unset → budget check is skipped.
    plan_verbosity = plan.get("verbosity")
    slides_spec = plan.get("slides") or []
    if not slides_spec:
        print(f"deck: plan '{plan_path}' has no slides", file=sys.stderr)
        return 2

    # Resolve build-time image provider from the default brand's
    # `$image_provider` (extends-resolved). Per-slide `brand:` overrides
    # still drive tokens/compounds, but the provider used to resolve
    # `picture query:` nodes is deck-wide — `asset_lock.json` lives next
    # to the output deck. Absent → provider is None; brands using only
    # `picture path:` build as before. A typo'd kind raises KeyError with
    # a registry listing — surfaces as the normal CLI traceback.
    discover_providers()
    provider = None
    try:
        default_brand_obj = find_brand(default_brand)
    except ValueError as e:
        print(f"deck: {e}", file=sys.stderr)
        return 2
    if default_brand_obj.image_provider_config:
        cfg = default_brand_obj.image_provider_config
        provider = get_provider(cfg["kind"], cfg.get("config"))

    # ── Pre-render static geometry verify (--strict-static) ──────────────
    if getattr(args, "strict_static", False):
        _require_builder("deck build --strict-static")
        from feinschliff_builder.verify.static import validate as _validate_static
        _static_bag = _validate_static(
            plan, brand=default_brand_obj, plan_dir=plan_path.parent
        )
        if _static_bag:
            for _d in _static_bag:
                _loc = _d.location or "slide ?"
                print(
                    f"deck build: static: {_loc}: "
                    f"[{_d.severity.value.upper()}] {_d.kind.value} — {_d.message}",
                    file=sys.stderr,
                )
            print(
                f"deck build: --strict-static: {len(_static_bag)} static "
                f"defect(s) found. Fix them or remove --strict-static to skip "
                f"this gate.",
                file=sys.stderr,
            )
            return 1

    # ── Auto-fix loop (--autofix) ─────────────────────────────────────────
    if getattr(args, "autofix", False):
        _require_builder("deck build --autofix")
        from feinschliff_builder.verify.static import validate as _validate_static_af
        from feinschliff_builder.verify.autofix import plan_fixes, apply_fixes, diff_summary

        _MAX_AUTOFIX_CYCLES = 3
        _total_patches = 0
        _seen_hashes: set[str] = set()
        _oscillation_detected = False
        for _cycle in range(_MAX_AUTOFIX_CYCLES):
            _static_bag = _validate_static_af(
                plan, brand=default_brand_obj, plan_dir=plan_path.parent
            )
            if not _static_bag:
                break
            _patches = plan_fixes(_static_bag, plan, default_brand_obj.root)
            if not _patches:
                # No mechanical fix available; leave residuals for compile.
                break
            _h = _patch_set_hash(_patches)
            if _h in _seen_hashes:
                print(
                    f"deck build: autofix cycle {_cycle + 1}: identical patch set "
                    f"seen before; halting to avoid oscillation",
                    file=sys.stderr,
                )
                _oscillation_detected = True
                break
            _seen_hashes.add(_h)
            _before = plan
            plan = apply_fixes(plan, _patches)
            _total_patches += len(_patches)
            _summary = diff_summary(_before, plan)
            print(
                f"deck build: autofix cycle {_cycle + 1}: "
                f"{len(_patches)} patch(es) applied",
            )
            if _summary:
                for _line in _summary.splitlines():
                    print(f"  {_line}")
        else:
            # Exhausted cycles — check if residuals remain.
            _residuals = _validate_static_af(
                plan, brand=default_brand_obj, plan_dir=plan_path.parent
            )
            if _residuals:
                print(
                    f"deck build: --autofix: {len(_residuals)} residual static "
                    f"defect(s) after {_MAX_AUTOFIX_CYCLES} cycle(s) — proceeding "
                    f"to compile (orchestrator may revise).",
                    file=sys.stderr,
                )
        if _total_patches > 0:
            # Write the auto-fixed plan back to disk before compile.
            plan_path.write_text(
                yaml.safe_dump(plan, sort_keys=False, allow_unicode=True),
                encoding="utf-8",
            )
            print(
                f"deck build: auto-fix passes: {_total_patches} total patch(es); "
                f"plan written back to {plan_path}"
            )
        # Re-capture slides_spec from the (potentially mutated) plan so the
        # compile loop below uses the fixed content, not the pre-fix snapshot.
        slides_spec = plan.get("slides") or []

    # Compute the output deck path up front so it can serve as `deck_dir`
    # for `asset_lock.json` + `.cache/` during the build.
    out_path = Path(args.output or plan.get("out", "deck.pptx")).resolve()

    plan_dir = plan_path.parent
    slides_payload: list[tuple[list, object, Path]] = []
    content_defects_by_slide: dict[int, list] | None = (
        {} if not args.skip_content_lint else None
    )
    all_diagram_defects: list = []

    # Top-level timing phase for the build. Per-slide events would push
    # indentation past sanity here; instead we emit one `build:total`
    # event manually around the whole compile loop (see end of function)
    # and emit `build:slide` lightweight events inline via log_event.
    import time as _time
    _build_t0 = _time.perf_counter()
    log_event(plan_dir, "build:total", "start",
              slides=len(slides_spec), brand=default_brand)

    with tempfile.TemporaryDirectory() as _tmp:
        diagrams_out = Path(_tmp) / "diagrams"
        diagrams_out.mkdir()

        for i, spec in enumerate(slides_spec):
            layout_path = (plan_dir / spec["layout"]).resolve()
            if not layout_path.is_file():
                # also try toolkit-relative (plugin root / rel).
                alt = _find_toolkit_file(spec["layout"])
                if alt is not None:
                    layout_path = alt
                else:
                    print(f"deck: slide {i}: layout not found: {spec['layout']}", file=sys.stderr)
                    return 2

            brand = spec.get("brand", default_brand)
            try:
                brand_dir = find_brand(brand).root
            except ValueError as e:
                print(f"deck: slide {i}: {e}", file=sys.stderr)
                return 2

            tokens = load_tokens(brand_dir)
            compounds = load_compounds_for_brand(
                brand_dir, std_dir=_bundled_compounds()
            )

            layout_nodes, layout_compounds = parse_file(layout_path)
            for cd in layout_compounds:
                compounds[cd.name] = cd

            ctx = spec.get("content") or {}
            if not ctx and "content_file" in spec:
                content_path = (plan_dir / spec["content_file"]).resolve()
                if not content_path.is_file():
                    alt = _find_toolkit_file(spec["content_file"])
                    if alt is not None:
                        content_path = alt
                    else:
                        print(f"deck: slide {i}: content_file not found: {spec['content_file']}", file=sys.stderr)
                        return 2
                ctx = yaml.safe_load(content_path.read_text()) or {}

            if content_defects_by_slide is not None:
                slide_index = i + 1
                # Strip `.slide.dsl` (and any leading path) to get the bare
                # layout name — this is what the structural validators key on.
                layout_name = layout_path.name
                if layout_name.endswith(".slide.dsl"):
                    layout_name = layout_name[: -len(".slide.dsl")]
                slot_budgets = compute_slot_budgets(layout_nodes, tokens, compounds=compounds)
                slide_defects = validate_content(
                    ctx, slide_index=slide_index, layout=layout_name,
                    slot_budgets=slot_budgets,
                )
                if validate_notes is not None:
                    slide_defects.extend(validate_notes(
                        spec.get("notes"),
                        slide_index=slide_index,
                        is_hook=(i == 0),
                        verbosity=plan_verbosity,
                    ))
                if slide_defects:
                    content_defects_by_slide[slide_index] = slide_defects

            _slide_t0 = _time.perf_counter()
            log_event(plan_dir, "build:slide", "start", slide=i + 1,
                      layout=Path(spec["layout"]).name)
            slide_result = compile_slide(
                layout_path=layout_path,
                ctx=ctx,
                brand_dir=brand_dir,
                slide_index=i + 1,
                diagrams_out_dir=diagrams_out,
            )
            log_event(
                plan_dir, "build:slide", "end", slide=i + 1,
                layout=Path(spec["layout"]).name,
                elapsed_ms=int((_time.perf_counter() - _slide_t0) * 1000),
                defects=len(slide_result.defects),
            )
            for d in slide_result.defects:
                print(f"deck: slide {i + 1}: {format_defect(d)}", file=sys.stderr)
            all_diagram_defects.extend(slide_result.defects)
            notes = spec.get("notes")
            if notes is not None:
                slides_payload.append(
                    (slide_result.primitives, slide_result.tokens,
                     brand_dir / "assets", notes)
                )
            else:
                slides_payload.append(
                    (slide_result.primitives, slide_result.tokens,
                     brand_dir / "assets")
                )

        if content_defects_by_slide:
            emit_defects_and_abort_message(content_defects_by_slide, cli_name="deck build")
            return 1

        allowed_to_skip: set[str] = set()
        if getattr(args, "allow_diagram_warnings", False):
            allowed_to_skip |= {"diagram-overflow", "diagram-text-too-small"}
        _fatal_diagram = [
            d for d in all_diagram_defects
            if d.kind.value in fatal_kinds() and d.kind.value not in allowed_to_skip
        ]
        if _fatal_diagram:
            print(
                f"deck build: aborting — {len(_fatal_diagram)} fatal defect(s) "
                f"across {len(slides_spec)} slide(s). Pass "
                f"--allow-diagram-warnings to demote "
                f"diagram-overflow/diagram-text-too-small (if those are the only blockers).",
                file=sys.stderr,
            )
            return 1

        prs = build_multi_slide(
            slides_payload,
            asset_root_fallback=_bundled_assets(),
            image_provider=provider,
            deck_dir=out_path.parent,
        )
        missing = getattr(prs, "missing_assets", []) or []
        if missing and not getattr(args, "allow_missing_assets", False):
            for entry in missing:
                kind = entry.get("kind", "missing")
                path = entry.get("path") or "(unset)"
                slide_n = entry.get("slide_index", "?")
                line = entry.get("line_no", "?")
                print(
                    f"deck build: slide {slide_n}: missing asset ({kind}) "
                    f"at line {line}: {path}",
                    file=sys.stderr,
                )
            print(
                f"deck build: aborting — {len(missing)} missing required "
                f"asset(s) across {len(slides_spec)} slide(s). Mark optional "
                f"slots with `optional:true` or pass --allow-missing-assets.",
                file=sys.stderr,
            )
            return 1
        out_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out_path))
        print(f"wrote {out_path} ({len(prs.slides)} slides)")
        log_event(
            plan_dir, "build:total", "end",
            elapsed_ms=int((_time.perf_counter() - _build_t0) * 1000),
            slides=len(prs.slides), output=str(out_path),
        )
        return 0


def cmd_pick(args) -> int:
    from feinschliff.deck.picker import LayoutPicker

    if args.signals == "-":
        raw = sys.stdin.read()
    else:
        raw = Path(args.signals).read_text()
    signals = yaml.safe_load(raw) or {}

    picker = LayoutPicker(top_k=args.top_k)
    candidates = picker.candidates(signals)
    if not candidates:
        print("deck: no candidate layouts matched the signals", file=sys.stderr)
        return 1
    for c in candidates:
        print(f"{c.score:5.2f}  {c.layout_name:<24}  {c.reason}")
    return 0


def _build_primitives_for_layout(
    layout_path: Path, brand: str, content_path: Path | None,
    *, skip_interpolation: bool = False,
) -> tuple[list, object]:
    """Delegate to feinschliff.deck.orchestrate.build_primitives_for_layout."""
    return _build_primitives_for_layout_fn(
        layout_path, brand, content_path,
        skip_interpolation=skip_interpolation,
    )


def cmd_wireframe(args) -> int:
    _require_builder("deck wireframe")
    from feinschliff_builder.decompile.wireframe import render_wireframe

    layout_path = Path(args.layout).resolve()
    if not layout_path.is_file():
        alt = _find_toolkit_file(args.layout)
        if alt is not None:
            layout_path = alt
        else:
            print(f"deck wireframe: layout not found: {args.layout}", file=sys.stderr)
            return 2

    if args.content:
        content_path = Path(args.content).resolve()
        if not content_path.is_file():
            print(f"deck wireframe: content file not found: {args.content}", file=sys.stderr)
            return 2
    else:
        content_path = None
    # Skip interpolation when no content is given (preserves slot labels) or
    # when --show-slots forces slot-structure mode even with content provided.
    skip_interp = (content_path is None) or args.show_slots
    try:
        primitives, tokens = _build_primitives_for_layout(
            layout_path, args.brand, content_path,
            skip_interpolation=skip_interp,
        )
    except ValueError as exc:
        print(f"deck wireframe: {exc}", file=sys.stderr)
        return 2

    bg_b64: str | None = None
    if args.overlay_pptx:
        from feinschliff.io.pptx_to_png import slide_to_b64
        try:
            bg_b64 = slide_to_b64(
                Path(args.overlay_pptx).resolve(),
                slide_index=args.overlay_slide,
            )
        except (FileNotFoundError, RuntimeError, IndexError) as exc:
            print(f"deck wireframe: overlay failed — {exc}", file=sys.stderr)
            return 2

    title = layout_path.name.replace(".slide.dsl", "")
    svg = render_wireframe(
        primitives, tokens,
        title=title,
        background_png_b64=bg_b64,
        background_opacity=args.overlay_opacity,
    )
    out = Path(args.output).resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(svg, encoding="utf-8")
    mode = "overlay" if bg_b64 else "wireframe"
    print(f"wrote {out} ({mode}, {len(primitives)} primitives)")
    return 0


def cmd_wireframe_sheet(args) -> int:
    _require_builder("deck wireframe-sheet")
    from feinschliff_builder.decompile.wireframe import render_wireframe_sheet

    plan_path = Path(args.plan).resolve()
    if not plan_path.is_file():
        print(f"deck wireframe-sheet: plan not found: {plan_path}", file=sys.stderr)
        return 2

    plan = yaml.safe_load(plan_path.read_text()) or {}
    default_brand = plan.get("brand", "feinschliff")
    slides_spec = plan.get("slides") or []
    if not slides_spec:
        print(f"deck wireframe-sheet: plan '{plan_path}' has no slides", file=sys.stderr)
        return 2

    plan_dir = plan_path.parent
    slides_data: list[tuple[list, object, str]] = []

    for i, spec in enumerate(slides_spec):
        layout_rel = spec["layout"]
        layout_path = (plan_dir / layout_rel).resolve()
        if not layout_path.is_file():
            alt = _find_toolkit_file(layout_rel)
            if alt is not None:
                layout_path = alt
            else:
                print(f"deck wireframe-sheet: slide {i}: layout not found: {layout_rel}",
                      file=sys.stderr)
                return 2

        brand = spec.get("brand", default_brand)
        ctx_inline = spec.get("content") or {}
        content_path: Path | None = None
        if not ctx_inline and "content_file" in spec:
            cp = (plan_dir / spec["content_file"]).resolve()
            if not cp.is_file():
                cp = _find_toolkit_file(spec["content_file"]) or cp
            content_path = cp if cp.is_file() else None

        try:
            brand_dir = find_brand(brand).root
            tokens = load_tokens(brand_dir)
            compounds = load_compounds_for_brand(
                brand_dir, std_dir=_bundled_compounds()
            )
            layout_nodes, layout_compounds = parse_file(layout_path)
            for cd in layout_compounds:
                compounds[cd.name] = cd
            ctx: dict = ctx_inline.copy()
            if not ctx and content_path:
                ctx = yaml.safe_load(content_path.read_text()) or {}
            if args.show_slots:
                # Skip interpolation so {{ slot_name }} labels survive into the cell.
                primitives, _ = expand_compounds(layout_nodes, compounds)
            else:
                interp = interpolate_nodes(layout_nodes, ctx)
                primitives, _ = expand_compounds(interp, compounds)
        except (ValueError, OSError, yaml.YAMLError, KeyError) as exc:
            print(f"deck wireframe-sheet: slide {i}: {exc}", file=sys.stderr)
            return 1

        title = layout_path.name.replace(".slide.dsl", "")
        slides_data.append((primitives, tokens, title))

    bg_list: list[str | None] | None = None
    if args.overlay_pptx:
        from feinschliff.io.pptx_to_png import pptx_to_pngs_b64
        try:
            pngs = pptx_to_pngs_b64(Path(args.overlay_pptx).resolve())
        except (FileNotFoundError, RuntimeError) as exc:
            print(f"deck wireframe-sheet: overlay failed — {exc}", file=sys.stderr)
            return 2
        if len(pngs) != len(slides_data):
            if len(pngs) < len(slides_data):
                detail = (
                    f"only {len(pngs)} overlay(s) for {len(slides_data)} slide(s); "
                    "unmatched slides will have no overlay. "
                    "(Some LibreOffice versions only export the first slide.)"
                )
            else:
                detail = (
                    f"{len(pngs)} overlay(s) for {len(slides_data)} slide(s); "
                    "extra overlays will be ignored. Check that --overlay-pptx "
                    "matches the deck plan."
                )
            print(f"deck wireframe-sheet: warning — {detail}", file=sys.stderr)
        # Pad to match slide count; truncate is implicit via index guard in render.
        bg_list = list(pngs) + [None] * max(0, len(slides_data) - len(pngs))

    svg = render_wireframe_sheet(
        slides_data,
        background_pngs_b64=bg_list,
    )
    out = Path(args.output).resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(svg, encoding="utf-8")
    mode = "overlay sheet" if bg_list else "wireframe sheet"
    print(f"wrote {out} ({mode}, {len(slides_data)} slide(s))")
    return 0


def cmd_polish(args) -> int:
    """Walk a PPTX, extract diagram IR from each slide, emit DSL artifacts,
    and rebuild a brand-perfect polished deck.

    WIRED:
    - Reads input .pptx with python-pptx.
    - Calls extract_vector.extract_from_slide() for every slide.
    - Calls kind_selector.select_kind() (or uses --refurbish-default) to pick
      the emitter target.
    - Emits .exc.dsl (excalidraw) or .svg.dsl (svg) into <output-dir>/refurbished/.
    - Writes refurbish_report.md next to the output file.
    - Rebuilds a polished PPTX from the refurbished diagram slides using the
      excalidraw-diagram or svg-infographic layout with the brand token pack.

    NOTE: Non-diagram slides (no detectable nodes) are skipped in the rebuilt
    deck. The output PPTX contains only refurbished diagram slides. If no
    diagram slides are found, the input is copied to the output as a fallback.

    --refurbish-all is the only interactive-mode variant implemented here;
    interactive per-slide confirmation is not yet wired.
    """
    _require_builder("deck polish")
    import shutil

    from pathlib import Path as _Path
    from pptx import Presentation
    from feinschliff_builder.diagrams.refurbish.extract_vector import extract_from_slide
    from feinschliff_builder.diagrams.refurbish.kind_selector import select_kind
    from feinschliff_builder.diagrams.refurbish.emit_excalidraw import emit as emit_excalidraw_dsl
    from feinschliff_builder.diagrams.refurbish.emit_svg import emit as emit_svg_dsl

    def _extract_slide_title(slide) -> str:
        """Pick the best title candidate from a refurbished slide.

        Order: (1) PPTX title placeholder if it has text; (2) the largest
        free-floating textbox (i.e., one whose shape isn't an auto-shape with
        its own label). Never returns a shape-label as the slide title.
        """
        # 1. Slide title placeholder
        try:
            title_shape = slide.shapes.title
            if title_shape and title_shape.has_text_frame and title_shape.text_frame.text.strip():
                return title_shape.text_frame.text.strip().split("\n")[0]
        except (AttributeError, ValueError):
            pass
        # 2. Largest free-floating textbox (skip shapes with their own label)
        best = ""
        best_area = 0
        for sh in slide.shapes:
            try:
                _ = sh.auto_shape_type  # raises if not an auto-shape
                continue  # skip boxes/ellipses (they're nodes, not titles)
            except (ValueError, AttributeError):
                pass
            if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip():
                area = (sh.width or 0) * (sh.height or 0)
                if area > best_area:
                    best = sh.text_frame.text.strip().split("\n")[0]
                    best_area = area
        return best

    src_path = _Path(args.input).resolve()
    if not src_path.is_file():
        print(f"deck polish: input not found: {src_path}", file=sys.stderr)
        return 2

    out_path = _Path(args.output).resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)

    if args.no_refurbish:
        shutil.copy(src_path, out_path)
        print(f"deck polish: --no-refurbish — copied {src_path.name} → {out_path}")
        return 0

    refurbish_dir = out_path.parent / "refurbished"
    refurbish_dir.mkdir(exist_ok=True)
    report_lines = ["# Refurbish Report\n"]

    refurbished_slides: list[dict] = []  # each: {layout, content}

    in_pres = Presentation(str(src_path))
    for idx, slide in enumerate(in_pres.slides, start=1):
        ir = extract_from_slide(slide)
        if not ir.nodes:
            report_lines.append(f"- slide {idx}: no diagram nodes detected — skipped")
            continue

        kind = args.refurbish_default or select_kind(ir)

        if args.refurbish_all or args.refurbish_default:
            if kind == "excalidraw":
                dsl = emit_excalidraw_dsl(ir, 1720, 480)
                ext = ".exc.dsl"
                layout_name = "excalidraw-diagram.slide.dsl"
            else:
                dsl = emit_svg_dsl(ir, 1720, 480)
                ext = ".svg.dsl"
                layout_name = "svg-infographic.slide.dsl"
            artifact = refurbish_dir / f"slide-{idx}{ext}"
            artifact.write_text(dsl, encoding="utf-8")
            report_lines.append(
                f"- slide {idx}: detected {len(ir.nodes)}-node {kind} "
                f"(confidence {ir.confidence:.2f}) → {artifact.name}"
            )

            # Best-effort: pull a title from the source slide.
            # Prefer (1) the slide's title placeholder, (2) the largest free-
            # floating textbox that isn't a shape label. Never use a box's own
            # label as the slide title — that's noise.
            title = _extract_slide_title(slide)

            # Strip the leading "canvas WxH\n" line — the layout template
            # embeds diagram_dsl inside its own canvas block.
            dsl_body = dsl.partition("\n")[2]

            refurbished_slides.append({
                "layout": layout_name,
                "content": {
                    "pgmeta": f"Slide {idx}",
                    "tracker": "Refurbished",
                    "action_title": title or f"Diagram {idx}",
                    "so_what": "",
                    "source": f"Refurbished from input slide {idx}",
                    "diagram_dsl": dsl_body,
                    "footer_left": "Feinschliff",
                    "footer_right": f"Slide {idx}",
                },
            })
        else:
            report_lines.append(
                f"- slide {idx}: detected {len(ir.nodes)}-node {kind} "
                f"(confidence {ir.confidence:.2f}) — no flag set, skipped"
            )

    (out_path.parent / "refurbish_report.md").write_text("\n".join(report_lines), encoding="utf-8")

    if not refurbished_slides:
        # Nothing extracted — fall back to copying input
        shutil.copy(src_path, out_path)
        slide_count = len(in_pres.slides)
        print(
            f"deck polish: processed {slide_count} slide(s), "
            f"0 diagram slides found — input copied to {out_path.name}"
        )
        return 0

    _build_refurbished_deck(
        refurbished_slides,
        brand=args.brand,
        out_path=out_path,
    )

    slide_count = len(in_pres.slides)
    artifact_count = len(refurbished_slides)
    print(
        f"deck polish: processed {slide_count} slide(s), "
        f"{artifact_count} refurbished diagram slide(s) → {out_path.name}"
    )
    return 0


def _build_refurbished_deck(slides_plan: list[dict], brand: str, out_path: Path) -> None:
    """Delegate to feinschliff.deck.orchestrate.build_refurbished_deck."""
    _build_refurbished_deck_fn(slides_plan, brand, out_path)


def cmd_book(args) -> int:
    """`feinschliff deck book` — annotated speaker-book PDF.

    Loads the deck plan + design brief, optionally renders per-slide
    PNG thumbnails from the built .pptx (via the existing
    `feinschliff_builder.verify.render_pngs.render_slides_to_png` helper), and writes a
    multi-page PDF: front matter page + one page per slide.

    The brief is the source of truth for deck-level fields (takeaway,
    audience, frame, …) and per-slide role / audience_fit. Speaker
    notes come from the plan's per-slide `notes:` field (preferred);
    when absent, the brief's per-slide `notes` is used as fallback.
    """
    _require_builder("deck book")
    import json as _json
    import tempfile as _tempfile

    from feinschliff.book import (
        BookSlide, DeckFrontMatter, compose_book_pdf,
    )
    from feinschliff_builder.verify.render_pngs import render_slides_to_png

    plan_path = Path(args.plan).resolve()
    brief_path = Path(args.design_brief).resolve()
    out_path = Path(args.output).resolve()

    if not plan_path.is_file():
        print(f"deck book: plan not found: {plan_path}", file=sys.stderr)
        return 2
    if not brief_path.is_file():
        print(f"deck book: design_brief not found: {brief_path}",
              file=sys.stderr)
        return 2

    plan = yaml.safe_load(plan_path.read_text()) or {}
    brief = _json.loads(brief_path.read_text(encoding="utf-8"))

    front = DeckFrontMatter(
        takeaway=brief.get("takeaway", ""),
        audience=brief.get("audience", ""),
        audience_notes=brief.get("audience_notes", ""),
        frame=brief.get("frame", ""),
        frame_rationale=brief.get("frame_rationale", ""),
        red_line=brief.get("red_line", ""),
        hook_technique=(brief.get("hook") or {}).get("technique", ""),
        hook_opener=(brief.get("hook") or {}).get("opener", ""),
        deck_title=brief.get("takeaway") or plan.get("title"),
    )

    # Optional thumbnail render. The PDF is still useful without —
    # the speaker reads notes + claims even when the .pptx isn't
    # built yet (early-iteration mode).
    thumbnails: dict[int, Path] = {}
    tmp_ctx = None
    if args.pptx:
        pptx_path = Path(args.pptx).resolve()
        if pptx_path.is_file():
            tmp_ctx = _tempfile.TemporaryDirectory()
            thumbnails = render_slides_to_png(pptx_path, Path(tmp_ctx.name))
        else:
            print(f"deck book: --pptx not found, skipping thumbnails: "
                  f"{pptx_path}", file=sys.stderr)

    plan_slides = plan.get("slides") or []
    brief_slides = {int(s.get("index", 0)): s
                    for s in (brief.get("slides") or [])}

    book_slides: list[BookSlide] = []
    for i, spec in enumerate(plan_slides):
        brief_s = brief_slides.get(i, {})
        content = spec.get("content") or {}
        claim = (brief_s.get("claim")
                 or content.get("title")
                 or content.get("action_title")
                 or "")
        notes = spec.get("notes") or brief_s.get("notes") or ""
        thumb = thumbnails.get(i + 1)  # render_slides_to_png is 1-based
        book_slides.append(BookSlide(
            index=i,
            role=brief_s.get("role", ""),
            claim=claim,
            notes=notes,
            audience_fit=brief_s.get("audience_fit", ""),
            thumbnail_path=thumb,
            section_label=spec.get("section"),
        ))

    try:
        compose_book_pdf(front, book_slides, out_path)
    finally:
        if tmp_ctx is not None:
            tmp_ctx.cleanup()

    print(f"wrote {out_path} ({len(book_slides)} slide page(s) + 1 "
          f"front-matter page; thumbnails: "
          f"{'yes' if thumbnails else 'no'})")
    return 0


def cmd_storyline(args) -> int:
    _require_builder("deck storyline")
    plan_path = Path(args.plan).resolve()
    out_path = Path(args.output).resolve()
    try:
        titles = extract_titles_from_plan(plan_path)
    except FileNotFoundError as exc:
        print(f"deck storyline: {exc}", file=sys.stderr)
        return 2
    except ValueError as exc:
        print(f"deck storyline: {exc}", file=sys.stderr)
        return 2

    contact_sheet = render_contact_sheet(titles, brief_summary=args.brief_summary)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    write_storyline_report(out_path, contact_sheet=contact_sheet)
    print(f"wrote {out_path} ({len([t for t in titles if t])} non-empty title(s) "
          f"across {len(titles)} slide(s))")
    return 0


def cmd_claim_evidence(args) -> int:
    """``feinschliff deck claim-evidence`` — mid-plan claim-evidence gate.

    Exit codes:
    - 0: clean (all judged slides pass)
    - 1: dirty (at least one slide has a claim-evidence defect)
    - 2: plumbing error (plan not found, parse failure, etc.)
    """
    _require_builder("deck claim-evidence")
    plan_path = Path(args.plan).resolve()
    out_path = Path(args.output).resolve()

    # Load plan
    try:
        plan_text = plan_path.read_text()
    except FileNotFoundError as exc:
        print(f"deck claim-evidence: {exc}", file=sys.stderr)
        return 2

    try:
        plan: dict = yaml.safe_load(plan_text) or {}
    except Exception as exc:  # noqa: BLE001
        print(f"deck claim-evidence: failed to parse plan: {exc}", file=sys.stderr)
        return 2

    # Load optional design brief
    design_brief: dict | None = None
    if args.design_brief:
        brief_path = Path(args.design_brief).resolve()
        try:
            import json as _json
            brief_text = brief_path.read_text()
            # Accept both JSON and YAML
            try:
                design_brief = _json.loads(brief_text)
            except _json.JSONDecodeError:
                design_brief = yaml.safe_load(brief_text) or {}
        except FileNotFoundError as exc:
            print(f"deck claim-evidence: {exc}", file=sys.stderr)
            return 2
        except Exception as exc:  # noqa: BLE001
            print(f"deck claim-evidence: failed to parse design brief: {exc}", file=sys.stderr)
            return 2

    slide_count = len(plan.get("slides") or [])

    try:
        results = judge_plan(
            plan,
            design_brief=design_brief,
            offline=args.offline,
            model=args.model,
        )
    except SystemExit:
        raise  # propagate ANTHROPIC_API_KEY error
    except Exception as exc:  # noqa: BLE001
        print(f"deck claim-evidence: judgment failed: {exc}", file=sys.stderr)
        return 2

    # Token-cost estimate (AC6)
    judged_count = len(results)
    if not args.offline and results:
        from feinschliff_builder.verify.llm.prompts import claim_evidence_prompt
        # Rough estimate: average prompt length × slides judged / 4 chars/token
        sample_prompt = claim_evidence_prompt("Sample title", "Sample body text.")
        avg_tokens_per_slide = len(sample_prompt) // 4
        total_k = (avg_tokens_per_slide * judged_count) // 1000
        print(
            f"claim-evidence: {judged_count} slides judged, "
            f"~{max(total_k, 1)}k input tokens via {args.model}",
            file=sys.stderr,
        )
    else:
        print(
            f"claim-evidence: {judged_count} slides judged (--offline, 0 tokens)",
            file=sys.stderr,
        )

    overall = write_claim_evidence_report(out_path, results, slide_count=slide_count)
    print(f"wrote {out_path} (verdict: {overall}, {judged_count} judged / {slide_count} total)")

    return 0 if overall == "clean" else 1


# ──────────────────────────────────────────────────────────────────────────────
# Parallel verify — focused aspect checks
# ──────────────────────────────────────────────────────────────────────────────

def cmd_verify_aspect(args) -> int:
    """`feinschliff deck verify-aspect <aspect> --plan plan.yaml -o out.json`

    Each aspect runs an independent narrow check. Designed to be spawned
    as one subagent per aspect — N aspects run in parallel, then
    `verify-collate` merges them into a single verify_report.md.

    Implementation note: aspects that need LLM judgment (narrative,
    image-style, content-cohesion) currently emit a stub finding entry
    pointing the orchestrator at the relevant PNGs/files. The deterministic
    aspects (bbox, font, brand) emit actual findings. The orchestrator
    fills LLM-judged aspects via subagent dispatch.
    """
    import json as _json
    aspect = args.aspect
    plan_path = Path(args.plan).resolve()
    plan = yaml.safe_load(plan_path.read_text(encoding="utf-8")) or {}
    slides = plan.get("slides") or []

    out: dict = {
        "aspect": aspect,
        "plan": str(plan_path),
        "iteration_check_ms": None,
        "findings": [],  # list of {slide, kind, message, severity, hint}
        "summary": "",
        "needs_llm": False,
    }

    _t0 = time.perf_counter()

    if aspect == "bbox":
        # Deterministic: re-run feinschliff verify on the plan, surface
        # text-overflow / out-of-bounds / diagram-overflow defects.
        from feinschliff.pipeline import compile_slide
        import tempfile

        with tempfile.TemporaryDirectory() as _tmp:
            diagrams_out = Path(_tmp) / "diagrams"
            diagrams_out.mkdir()
            for i, spec in enumerate(slides):
                layout_path = (plan_path.parent / spec["layout"]).resolve()
                if not layout_path.is_file():
                    layout_path = _find_toolkit_file(spec["layout"]) or layout_path
                try:
                    brand_dir = find_brand(spec.get("brand")
                                           or plan.get("brand")
                                           or "feinschliff").root
                except ValueError:
                    continue
                try:
                    r = compile_slide(
                        layout_path=layout_path,
                        ctx=spec.get("content") or {},
                        brand_dir=brand_dir,
                        slide_index=i + 1,
                        diagrams_out_dir=diagrams_out,
                    )
                except Exception as e:
                    out["findings"].append({
                        "slide": i + 1, "kind": "compile-error",
                        "severity": "fatal",
                        "message": str(e)[:200], "hint": "see plan.yaml",
                    })
                    continue
                for d in r.defects:
                    out["findings"].append({
                        "slide": i + 1, "kind": d.kind.value,
                        "severity": d.severity.value,
                        "message": d.message[:240], "hint": "",
                    })

    elif aspect == "font":
        # Deterministic-ish: surface diagram-text-too-small from build,
        # plus a hint pointing the orchestrator at any slide where a text
        # primitive uses role=detail (most fragile per past runs).
        for i, spec in enumerate(slides):
            dsl = (spec.get("content") or {}).get("diagram_dsl") or ""
            if "size:detail" in dsl or " detail " in dsl:
                out["findings"].append({
                    "slide": i + 1, "kind": "detail-role-fragile",
                    "severity": "warn",
                    "message": "Diagram uses role=detail — close to the 10pt floor "
                               "after region/slide scaling; consider promoting to body.",
                    "hint": "search the diagram_dsl for `size:detail` or `detail `",
                })

    elif aspect == "narrative":
        # Reads titles + design_brief + storyline_report (if present).
        # The actual SCQA / claim-title judgment is LLM work; emit a
        # contact-sheet snapshot for the orchestrator subagent.
        titles = [
            (s.get("content") or {}).get("title")
            or (s.get("content") or {}).get("action_title")
            or s.get("_meta", {}).get("title", "?")
            for s in slides
        ]
        out["needs_llm"] = True
        out["contact_sheet"] = list(enumerate(titles, start=1))
        if args.design_brief:
            db = Path(args.design_brief).resolve()
            if db.is_file():
                try:
                    out["design_brief"] = _json.loads(db.read_text(encoding="utf-8"))
                except _json.JSONDecodeError:
                    pass

    elif aspect == "brand":
        # Deterministic-ish: scan diagram_dsl for fill tokens not in the
        # canonical 17-name vocabulary. Brand pack discipline check.
        canonical = {
            "primary", "accent", "secondary", "tertiary", "highlight",
            "ink", "neutral-strong", "neutral-soft", "graphite",
            "paper", "paper-2", "off-white", "surface-2",
            "severity-high", "severity-medium", "severity-low",
            "status-current", "status-next", "status-done",
            "success", "warning", "error", "code", "data", "inactive",
        }
        import re
        for i, spec in enumerate(slides):
            dsl = (spec.get("content") or {}).get("diagram_dsl") or ""
            for m in re.finditer(r"fill:([a-z0-9_-]+)", dsl):
                tok = m.group(1)
                if tok not in canonical:
                    out["findings"].append({
                        "slide": i + 1, "kind": "non-canonical-token",
                        "severity": "warn",
                        "message": f"Diagram uses fill:{tok}, not in the 17-name "
                                   f"semantic vocabulary.",
                        "hint": "use one of: " + ", ".join(sorted(canonical)[:10]) + ", …",
                    })

    elif aspect == "image":
        # Picture-slot presence + style consistency check. Surfaces slides
        # that declare an image_style but have no picture slot, or
        # vice-versa. LLM verification of actual style match is a separate
        # subagent step (needs_llm=True).
        if args.design_brief:
            try:
                db = _json.loads(Path(args.design_brief).read_text(encoding="utf-8"))
                out["image_style"] = db.get("image_style")
            except (OSError, _json.JSONDecodeError):
                pass
        out["needs_llm"] = True
        for i, spec in enumerate(slides):
            c = spec.get("content") or {}
            has_pic = any(k in c for k in ("hero_image", "picture", "image", "photo"))
            if has_pic and not out.get("image_style"):
                out["findings"].append({
                    "slide": i + 1, "kind": "image-style-undeclared",
                    "severity": "warn",
                    "message": "Slide uses a picture slot but design_brief.image_style "
                               "is unset.",
                    "hint": "set image_style in design_brief.json",
                })

    elif aspect == "content":
        # Title-body coherence + filler-word lint. Deterministic part:
        # scan body slots for filler words. LLM part: claim/proof match.
        FILLER = {"basically", "actually", "really", "very", "just", "simply", "in order to"}
        out["needs_llm"] = True
        for i, spec in enumerate(slides):
            c = spec.get("content") or {}
            for key in ("body", "supporting_body", "so_what"):
                val = c.get(key) or ""
                if not isinstance(val, str):
                    continue
                low = val.lower()
                hits = [w for w in FILLER if w in low]
                if hits:
                    out["findings"].append({
                        "slide": i + 1, "kind": "filler-word",
                        "severity": "warn",
                        "message": f"{key} contains filler: {', '.join(hits)}",
                        "hint": "Cut the filler — strong sentences don't need it.",
                    })

    elif aspect == "notes-coherence":
        # Pair the deck's red_line against each slide's (claim, notes).
        # The orchestrator LLM judges whether the spoken delivery tracks
        # the arc: drift / contradiction / off-arc tangents → dirty.
        _require_builder("deck verify-aspect notes-coherence")
        from feinschliff_builder.verify.deck.notes_coherence import (
            SlideForCoherence,
            render_contact_sheet as _render_notes_sheet,
        )
        out["needs_llm"] = True
        red_line = ""
        if args.design_brief:
            db_path = Path(args.design_brief).resolve()
            if db_path.is_file():
                try:
                    db = _json.loads(db_path.read_text(encoding="utf-8"))
                    red_line = db.get("red_line", "") or ""
                    out["red_line"] = red_line
                    out["design_brief"] = db
                except _json.JSONDecodeError:
                    pass
        coherence_slides: list[SlideForCoherence] = []
        for i, spec in enumerate(slides):
            c = spec.get("content") or {}
            coherence_slides.append(SlideForCoherence(
                index=i,
                role=spec.get("_meta", {}).get("role", ""),
                claim=c.get("title") or c.get("action_title") or "",
                notes=spec.get("notes"),
            ))
        # Cheap deterministic pre-flag: hook slide missing notes.
        # The LLM judge handles drift / off-arc semantics.
        if coherence_slides and not (coherence_slides[0].notes or "").strip():
            out["findings"].append({
                "slide": 1, "kind": "hook-notes-missing",
                "severity": "warn",
                "message": "Hook slide has no speaker notes; expected the "
                           "deck-level storyline articulating the red_line.",
                "hint": "Author the full red_line arc into slide 1's notes.",
            })
        out["contact_sheet"] = _render_notes_sheet(red_line, coherence_slides)

    out["iteration_check_ms"] = int(
        (time.perf_counter() - _t0) * 1000
    )
    out["summary"] = (
        f"{len(out['findings'])} finding(s) "
        f"({'needs LLM' if out['needs_llm'] else 'deterministic only'})"
    )

    out_path = Path(args.output).resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(_json.dumps(out, indent=2, ensure_ascii=False),
                        encoding="utf-8")
    log_event(plan_path.parent, f"verify-aspect:{aspect}", "tick",
              elapsed_ms=out["iteration_check_ms"],
              findings=len(out["findings"]))
    print(f"wrote {out_path} — {out['summary']}")
    return 0


def cmd_verify_collate(args) -> int:
    """`feinschliff deck verify-collate --aspect a.json --aspect b.json -o report.md`
    Merge per-aspect verify outputs into a unified verify_report.md.
    """
    import json as _json
    plan_path = Path(args.plan).resolve()
    plan = yaml.safe_load(plan_path.read_text(encoding="utf-8")) or {}
    slides = plan.get("slides") or []

    # Aggregate findings grouped by slide index. Each finding carries its
    # source aspect so the report makes the parallel checks visible.
    by_slide: dict[int, list[dict]] = {}
    aspects_seen: list[str] = []
    needs_llm: list[str] = []
    total_findings = 0

    for path_str in args.aspect:
        p = Path(path_str).resolve()
        if not p.is_file():
            print(f"deck verify-collate: aspect file not found: {p}",
                  file=sys.stderr)
            continue
        try:
            data = _json.loads(p.read_text(encoding="utf-8"))
        except _json.JSONDecodeError as e:
            print(f"deck verify-collate: {p}: {e}", file=sys.stderr)
            continue
        aspect = data.get("aspect", p.stem)
        aspects_seen.append(aspect)
        if data.get("needs_llm"):
            needs_llm.append(aspect)
        for f in data.get("findings", []) or []:
            f = dict(f)
            f["aspect"] = aspect
            by_slide.setdefault(int(f.get("slide", 0)), []).append(f)
            total_findings += 1

    # Decide verdict: clean iff no findings + no LLM aspects waiting.
    fatal = any(
        f.get("severity") == "fatal"
        for fs in by_slide.values() for f in fs
    )
    dirty = total_findings > 0 or fatal
    verdict = "dirty" if dirty else ("pending-llm" if needs_llm else "clean")

    lines: list[str] = []
    lines.append(f"# Verify Report — {plan_path.name}")
    lines.append("")
    lines.append(f"- **Iteration:** {args.iteration} of {args.budget}")
    lines.append(f"- **Verdict:** {verdict}"
                 + (f" — {total_findings} finding(s) across "
                    f"{len(by_slide)} slide(s)" if dirty else ""))
    lines.append(f"- **Aspects checked:** {', '.join(aspects_seen) or '(none)'}")
    if needs_llm:
        lines.append(f"- **Pending LLM verdict from:** {', '.join(needs_llm)}")
    if args.png_dir:
        lines.append(f"- **Rendered PNGs:** `{args.png_dir}`")
    lines.append("")
    lines.append("---")
    lines.append("")

    for i, spec in enumerate(slides, start=1):
        title = (spec.get("content") or {}).get("title") \
                or (spec.get("content") or {}).get("action_title") \
                or spec.get("_meta", {}).get("title", "(untitled)")
        layout_name = Path(spec.get("layout", "?")).name.replace(".slide.dsl", "")
        findings = by_slide.get(i, [])
        if not findings:
            lines.append(f"## Slide {i} — {title!r} ({layout_name}) ✅")
            lines.append("")
            lines.append("_No defects._")
        else:
            lines.append(f"## Slide {i} — {title!r} ({layout_name})")
            lines.append("")
            for f in findings:
                lines.append(
                    f"- **[{f['aspect']}/{f.get('kind', '?')}]** "
                    f"{f.get('message', '')}"
                )
                if f.get("hint"):
                    lines.append(f"  → {f['hint']}")
        lines.append("")

    out_path = Path(args.output).resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text("\n".join(lines), encoding="utf-8")
    log_event(plan_path.parent, "verify-collate", "tick",
              aspects=len(aspects_seen), findings=total_findings,
              verdict=verdict, iteration=args.iteration)
    print(f"wrote {out_path} — verdict: {verdict}, "
          f"{total_findings} finding(s) across {len(by_slide)} slide(s)")
    return 0


def cmd_verify_static(args) -> int:
    """`feinschliff deck verify-static <plan.yaml>` — pre-render static check.

    Inspects a plan.yaml for geometry defects that can be detected from the
    DSL + populated content without rendering (slot-overflow,
    empty-placeholder). Cheaper than a full build: catches class of defect
    in ~10-50 ms vs ~3 s/slide for a render-based check.

    Exit codes:
      0 — clean (no defects)
      1 — one or more defects found
      2 — plumbing error (plan not found, brand resolution failure, etc.)
    """
    _require_builder("deck verify-static")
    import json as _json
    from feinschliff_builder.verify.static import validate as _validate_static

    plan_path = Path(args.plan).resolve()
    if not plan_path.is_file():
        print(f"deck verify-static: plan not found: {plan_path}", file=sys.stderr)
        return 2

    try:
        plan = yaml.safe_load(plan_path.read_text(encoding="utf-8")) or {}
    except (OSError, yaml.YAMLError) as exc:
        print(f"deck verify-static: could not load plan: {exc}", file=sys.stderr)
        return 2

    brand_name = getattr(args, "brand", None) or plan.get("brand") or "feinschliff"
    try:
        brand_obj = find_brand(brand_name)
    except ValueError as exc:
        print(f"deck verify-static: {exc}", file=sys.stderr)
        return 2

    try:
        bag = _validate_static(plan, brand=brand_obj, plan_dir=plan_path.parent)
    except Exception as exc:  # noqa: BLE001
        print(f"deck verify-static: unexpected error: {exc}", file=sys.stderr)
        return 2

    if getattr(args, "json", False):
        # Produce a backward-compatible schema so `deck apply-fixes --defects`
        # can consume this output: {slide_index, kind, severity, message, meta}.
        out = []
        for d in bag:
            extra = d.extra or {}
            entry = {
                "slide_index": extra.get("slide_index", 0),
                "kind": d.kind.value,
                "severity": d.severity.value,
                "message": d.message,
                "meta": {k: v for k, v in extra.items() if k != "slide_index"},
            }
            if d.location is not None:
                entry["location"] = d.location
            out.append(entry)
        print(_json.dumps(out, indent=2, ensure_ascii=False))
    else:
        if bag:
            for d in bag:
                _loc = d.location or "slide ?"
                print(
                    f"{_loc}: [{d.severity.value.upper()}] {d.kind.value} — {d.message}"
                )
        else:
            print("verify-static: clean — no defects found")

    return 1 if bag else 0


def cmd_apply_fixes(args) -> int:
    """`feinschliff deck apply-fixes <plan.yaml> --defects <defects.json> [-o out.yaml]`

    Read the defects JSON (flat list OR collated shape), translate to
    deterministic FixPatch objects, apply them to the plan, and write the
    result.  Prints a markdown diff summary to stdout.

    Defects JSON shapes supported:
      - Flat list:  [{slide_index, kind, severity, message, meta}, ...]
        (output of `deck verify-static --json`)
      - Collated:   {"defects": {"1": [...], "2": [...]}, ...}
        (output shape used by cli/verify.py)

    Exit codes:
      0 — at least one patch was applied
      1 — no patches applied (defects present but none mechanically fixable,
          OR no defects at all)
      2 — plumbing error
    """
    _require_builder("deck apply-fixes")
    import json as _json
    from feinschliff_builder.verify.autofix import plan_fixes, apply_fixes, diff_summary
    from feinschliff.defects import Defect, DefectKind, Severity

    plan_path = Path(args.plan).resolve()
    if not plan_path.is_file():
        print(f"deck apply-fixes: plan not found: {plan_path}", file=sys.stderr)
        return 2

    try:
        plan = yaml.safe_load(plan_path.read_text(encoding="utf-8")) or {}
    except (OSError, yaml.YAMLError) as exc:
        print(f"deck apply-fixes: could not load plan: {exc}", file=sys.stderr)
        return 2

    defects_path = Path(args.defects).resolve()
    if not defects_path.is_file():
        print(f"deck apply-fixes: defects file not found: {defects_path}", file=sys.stderr)
        return 2

    try:
        raw = _json.loads(defects_path.read_text(encoding="utf-8"))
    except (_json.JSONDecodeError, OSError) as exc:
        print(f"deck apply-fixes: could not load defects: {exc}", file=sys.stderr)
        return 2

    # Normalise both JSON shapes to a flat list of Defect objects.
    defect_dicts: list[dict] = []
    if isinstance(raw, list):
        # Flat list shape: [{slide_index, kind, severity, message, meta}, ...]
        defect_dicts = raw
    elif isinstance(raw, dict):
        # Collated shape: {"defects": {slide_idx: [...]}, ...}
        nested = raw.get("defects") or {}
        for _slide_defects in nested.values():
            if isinstance(_slide_defects, list):
                defect_dicts.extend(_slide_defects)
    else:
        print("deck apply-fixes: unrecognised defects JSON shape", file=sys.stderr)
        return 2

    defects: list[Defect] = []
    for dd in defect_dicts:
        try:
            defects.append(Defect(
                slide_index=int(dd["slide_index"]),
                kind=DefectKind(dd["kind"]),
                severity=Severity(dd["severity"]),
                message=str(dd.get("message", "")),
                meta=dict(dd.get("meta") or {}),
            ))
        except (KeyError, ValueError) as exc:
            print(
                f"deck apply-fixes: skipping malformed defect entry {dd!r}: {exc}",
                file=sys.stderr,
            )

    if not defects:
        print("deck apply-fixes: no defects to process — nothing to do")
        return 1

    brand_name = getattr(args, "brand", None) or plan.get("brand") or "feinschliff"
    try:
        brand_obj = find_brand(brand_name)
    except ValueError as exc:
        print(f"deck apply-fixes: {exc}", file=sys.stderr)
        return 2

    patches = plan_fixes(defects, plan, brand_obj.root)
    if not patches:
        print("deck apply-fixes: no mechanical fixes available for these defects")
        print("Auto-fix passes: 0")
        return 1

    fixed_plan = apply_fixes(plan, patches)
    summary = diff_summary(plan, fixed_plan)

    out_path = Path(args.output).resolve() if args.output else plan_path
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(
        yaml.safe_dump(fixed_plan, sort_keys=False, allow_unicode=True),
        encoding="utf-8",
    )

    # Print diff summary to stdout.
    if summary:
        print(summary)
    print(f"Auto-fix passes: {len(patches)}")
    print(f"wrote {out_path} ({len(patches)} patch(es) applied)")
    return 0
