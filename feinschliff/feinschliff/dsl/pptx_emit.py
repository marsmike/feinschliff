"""Primitive-AST → python-pptx emitter.

Consumes the post-expansion node list (primitives only) plus a Tokens
bundle, builds a single-slide Presentation, returns it.

Coordinate system: DSL uses design pixels in a 1920×1080 frame, matching
the existing brand-pack convention. We map 1 design-px → 6350 EMU
(same as gs-ramspau's `px()` helper).

Primitives implemented:
  canvas WxH                 — set slide size
  theme <brand>              — no-op here; brand was already picked at load
  text X,Y "label" style:S align:left|right|center maxwidth:W
  rect X,Y WxH fill:role stroke:role stroke-width:N
  line X,Y X2,Y2 stroke:role stroke-width:N
  picture X,Y WxH path:PATH cover:true
  picture X,Y WxH query:"…" — resolve at build time via image_provider
  picture X,Y WxH path:PATH query:"…" — layered: file if path resolves,
                               else provider search with the explicit query
"""
from __future__ import annotations

import functools
import io
import os
import re
import sys
import tempfile
import warnings
from dataclasses import dataclass, field
from pathlib import Path
from typing import TYPE_CHECKING

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt

from .. import textfit
from feinschmiede.dsl.tokens import Tokens

# Provider search + asset-lock pinning + HTTP materialise live in
# feinschliff.io.image_materialise. Accessed via module attributes
# (`_img_mat._materialise(...)`) — NOT `from … import` — so tests can
# patch `feinschliff.io.image_materialise.X` once and affect the emitter.
from ..io import image_materialise as _img_mat
from .parser import DSLNode, parse_xy, parse_wh
from .polish import normalize_text

if TYPE_CHECKING:
    from ..io.image_provider import ImageProvider
    from feinschmiede.brand import BrandPack
    from feinschmiede.dsl.ast import Document, Element


class DSLError(Exception):
    """Raised when a DSL primitive carries a parse-or-emit-time error
    that the brand author or layout author can fix by editing source
    (e.g. mutually-exclusive keywords, missing required wiring).

    Distinct from ``ValueError`` / ``SyntaxError`` so callers can branch
    on "the deck source is malformed" vs. "the parser hit an internal
    inconsistency".
    """


# Single source of truth for the legacy 13.33"×7.5" widescreen baseline used
# by brand packs that don't declare an explicit physical slide size. Source-
# decompiled brand packs (brand_decompile_all.py writes `slide.width_emu` /
# `slide.height_emu` to tokens.json) override this via `_configure_slide_scale`.
# 13.333" × 914400 EMU/in = 12192000 EMU wide.
_LEGACY_SLIDE_WIDTH_EMU = 12_192_000
_LEGACY_CANVAS_W = 1920


@functools.lru_cache(maxsize=1)
def _installed_fonts() -> set[str]:
    """Set of font family names installed on this system, lowercased.

    Sources, in order:
      1. `fc-list` (Linux/macOS via homebrew fontconfig)
      2. macOS Microsoft Office app bundles (Calibri + family ships inside
         /Applications/Microsoft *.app/Contents/Resources/DFonts/ — those
         fonts are scoped to the host app and don't register with
         fontconfig, so a Calibri-using source PPTX appears uninstalled).

    Falls back to an empty set on systems without fontconfig — in which
    case `_assert_font_available` becomes a no-op (we can't block the
    build on a check we can't run reliably).
    """
    import shutil
    import subprocess
    fams: set[str] = set()
    fc = shutil.which("fc-list")
    if fc:
        try:
            out = subprocess.check_output([fc, ":", "family"], timeout=5).decode("utf-8", "ignore")
            for line in out.splitlines():
                for fam in line.split(","):
                    fams.add(fam.strip().lower())
        except (subprocess.SubprocessError, OSError):
            pass
    fams |= _office_bundle_fonts()
    return fams


def _office_bundle_fonts() -> set[str]:
    """Family names bundled with macOS Microsoft Office (Calibri etc).

    These ship inside the .app bundle and don't register with fontconfig,
    but `python-pptx` writes them as the font name and LibreOffice can
    pick them up via the shared font cache when they've been registered
    once. We surface them to the availability check so the build doesn't
    spuriously hard-fail on `Calibri` when Office is installed.
    """
    from pathlib import Path
    fams: set[str] = set()
    candidates = [
        "/Applications/Microsoft PowerPoint.app/Contents/Resources/DFonts",
        "/Applications/Microsoft Word.app/Contents/Resources/DFonts",
        "/Applications/Microsoft Excel.app/Contents/Resources/DFonts",
    ]
    # Filename → family-name map for the Office font set. PowerPoint stores
    # Calibri across multiple files (Calibri.ttf, calibril.ttf for Light,
    # Calibrib.ttf for Bold, etc); collapse all variants to the family name
    # that the source theme would reference.
    name_for: dict[str, str] = {
        "calibri.ttf": "Calibri", "calibrib.ttf": "Calibri",
        "calibrii.ttf": "Calibri", "calibriz.ttf": "Calibri",
        "calibril.ttf": "Calibri Light", "calibrili.ttf": "Calibri Light",
        "cambria.ttc": "Cambria", "cambriab.ttf": "Cambria",
        "candara.ttf": "Candara", "consola.ttf": "Consolas",
        "constan.ttf": "Constantia", "corbel.ttf": "Corbel",
    }
    for c in candidates:
        d = Path(c)
        if not d.is_dir():
            continue
        for f in d.iterdir():
            fam = name_for.get(f.name.lower())
            if fam:
                fams.add(fam.lower())
    return fams


def _assert_font_available(family: str, brand_name: str) -> None:
    """Check whether the brand's primary font is installed on the system.

    Behaviour gated by env var ``FEINSCHLIFF_STRICT_FONTS``:

    - ``=1`` → hard-fail with DSLError (source-faithful rendering: a missing
      font means glyph metrics + stroke weights will differ from the
      authoring tool's render, invalidating any pixel-level comparison).
    - unset / ``=0`` → soft-warn to stderr but let the build proceed (lets
      CI and toolkit-internal tests build against brand packs whose
      canonical fonts may not be present on every runner).

    Decompile sets the brand's tokens.json font-family to the source theme
    font, so source-faithful pipelines should run with the strict env set.
    """
    installed = _installed_fonts()
    if not installed:
        return  # fc-list unavailable — can't verify; defer to the renderer.
    if family.lower() in installed:
        return
    msg = (
        f"brand '{brand_name}': required font '{family}' is not installed. "
        f"Install it (`brew install --cask font-{family.lower().replace(' ', '-')}` "
        f"on macOS, or copy the .ttf into ~/.fonts on Linux + run `fc-cache -f`) "
        f"then re-run. Substituting another family would invalidate source-"
        f"matched rendering."
    )
    if os.environ.get("FEINSCHLIFF_STRICT_FONTS") == "1":
        raise DSLError(msg)
    print(f"WARN: {msg}", file=sys.stderr)

EMU_PER_PT = 12700           # PowerPoint standard: 1pt = 12700 EMU (914400 / 72).
_EMU_PER_PX = _LEGACY_SLIDE_WIDTH_EMU / _LEGACY_CANVAS_W   # 6350 — default fallback
_PX_TO_PT = _EMU_PER_PX / EMU_PER_PT                       # 0.5  — default fallback
_STROKE_PX_TO_PT = 0.75      # CSS px → pt for stroke widths (96/72 inverse rounded)


def _configure_slide_scale(tokens: "Tokens", canvas_w: int) -> None:
    """Recompute EMU_PER_PX + PX_TO_PT from tokens.json's slide.width_emu.

    Idempotent + global — called once per build before any shape is emitted.
    When tokens.slide.width_emu is absent, falls back to the legacy 13.33"
    baseline. Sourced-from-PPTX brand packs (decompile writes width_emu)
    automatically render at the source's physical slide size.
    """
    global _EMU_PER_PX, _PX_TO_PT
    try:
        width_emu = tokens.slide("width_emu") or _LEGACY_SLIDE_WIDTH_EMU
    except Exception:
        width_emu = _LEGACY_SLIDE_WIDTH_EMU
    cw = canvas_w or _LEGACY_CANVAS_W
    _EMU_PER_PX = width_emu / cw
    _PX_TO_PT = _EMU_PER_PX / EMU_PER_PT
_DEFAULT_PADDING_X = 100     # fallback right/left margin if brand has no slide.padding-x token

# Defense-in-depth: if interpolation somehow left a `{{ … }}` placeholder in
# an `if:` condition, treat it as falsy so the node is suppressed (never
# leaked into the rendered slide as literal text).
_RESIDUAL_PLACEHOLDER = re.compile(r"\{\{[^}]*\}\}")

# Tabular-numeral activation: a run is "numeric" if its content (after polish)
# is purely digits, separators, and numeric signs. Whitespace and Unicode
# minus (U+2212) count as numeric chrome.
_NUMERIC_CONTENT_RE = re.compile(r"^[\d,\.\-−%+\s]+$")

# Glyphs that, when leading a run, hang into the left margin. The bearing
# fraction is approximate (fraction of em — i.e. of font size).
_HANG_SIDE_BEARING_EM: dict[str, float] = {
    "•": 0.45,  # • bullet
    "“": 0.20,  # “ left double quote (en)
    "„": 0.20,  # „ low-9 (de)
    "«": 0.25,  # « left guillemet (fr)
    "—": 0.30,  # — em-dash
    "–": 0.25,  # – en-dash
    "…": 0.15,  # … ellipsis
    "\"":     0.20,  # ASCII straight double (legacy; rare after polish)
    "-":      0.15,  # ASCII hyphen
}


# 3-channel hierarchy stepping for indent levels. Each level beyond 0 steps:
#   size  ×= 0.85   (round to 0.5pt)
#   weight -= 100   (clamped at 300)
#   color  walks ink → graphite → fog (clamped at fog)
_HIERARCHY_COLOR_WALK = ["ink", "graphite", "fog"]


def _step_hierarchy(
    size_px: float, weight: int, color_role: str, *, level: int,
) -> tuple[float, int, str]:
    """Step all three channels for `level` indent levels. level<=0 is a no-op."""
    if level <= 0:
        return size_px, weight, color_role
    new_size = size_px
    new_weight = weight
    new_color = color_role
    for _ in range(level):
        new_size *= 0.85
        new_weight = max(300, new_weight - 100)
        if new_color in _HIERARCHY_COLOR_WALK:
            idx = _HIERARCHY_COLOR_WALK.index(new_color)
            new_color = _HIERARCHY_COLOR_WALK[min(idx + 1, len(_HIERARCHY_COLOR_WALK) - 1)]
    # Round size_pt to nearest 0.5pt. size_px ↔ pt: 2 design-px per pt.
    pt = new_size * _PX_TO_PT
    pt = round(pt * 2) / 2
    new_size = pt / _PX_TO_PT
    return new_size, new_weight, new_color


def _leading_hang_offset_px(text: str, size_pt: float) -> float:
    """Return leftward offset (design-px) for a textbox whose first glyph
    hangs into the margin. Zero if no hang glyph at position 0."""
    if not text:
        return 0.0
    bearing_em = _HANG_SIDE_BEARING_EM.get(text[0], 0.0)
    if bearing_em == 0.0:
        return 0.0
    # 1pt = 2 design-px (the canvas-to-PPT mapping). Em ≈ font size.
    return bearing_em * size_pt / _PX_TO_PT


def _px(n: float) -> Emu:
    return Emu(int(n * _EMU_PER_PX))


@dataclass
class EmitContext:
    tokens: Tokens
    canvas_w: float = 1920.0
    canvas_h: float = 1080.0
    asset_root: Path | None = None        # for resolving picture paths
    # Plugin-level fallback root walked when the per-brand `asset_root` does
    # not contain the requested relative path. Lets shared assets (universal
    # placeholder illustrations, common icons) live once at the plugin root
    # while brands keep the freedom to override per-asset by putting a file
    # at the same relative path under their own `assets/` dir.
    asset_root_fallback: Path | None = None
    # Required-asset accumulator (Review #7). _emit_picture appends an entry
    # for every picture node whose `path` is unset or points at a missing
    # file AND that does not carry `optional:true`. Callers consult this
    # post-build to decide whether to abort.
    missing_assets: list[dict] = field(default_factory=list)
    # Active image provider for resolving `picture query:"..."` nodes
    # (Task 7). String forward-ref keeps module import cheap and avoids any
    # circular-import risk between `feinschliff.dsl.pptx_emit` and `feinschliff.io.image_provider`.
    # Defaults to None so existing callers that never set it keep working.
    image_provider: "ImageProvider | None" = None
    # Deck output directory — used by the picture-query path to write the
    # `asset_lock.json` manifest and the `.cache/` of downloaded images
    # alongside the produced `.pptx`. Defaults to None; the picture-query
    # branch falls back to a temp dir / disabled cache when unset.
    deck_dir: Path | None = None


def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _hex_to_rgb_tuple(hex_color: str) -> tuple[int, int, int]:
    """Same as `_hex_to_rgb` but returns a plain `(r, g, b)` tuple for use
    with PIL APIs that don't accept python-pptx's `RGBColor` subclass.
    """
    h = hex_color.lstrip("#")
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _luminance(rgb: tuple[int, int, int]) -> float:
    """Rec. 601 luma (perceptual brightness 0..255)."""
    r, g, b = rgb
    return 0.299 * r + 0.587 * g + 0.114 * b


def _duotone_endpoints_for_brand(tokens) -> tuple[tuple[int, int, int], tuple[int, int, int]] | None:
    """Pick `(dark, light)` RGB endpoints for the plugin-fallback duotone
    treatment from a brand's tokens. Strategy:

    - Try ink + paper, picking the darker as shadow and lighter as highlight
      (auto-swap so dark-canvas brands still get a readable mapping).
    - If their luminance gap is too narrow to read (`< 80` on 0..255),
      substitute the lighter endpoint with the brand's accent color so the
      duotone keeps tonal range and brand identity.

    Returns None if ink/paper tokens are missing — caller falls back to the
    image's untreated form.
    """
    try:
        a = _hex_to_rgb_tuple(tokens.color("ink"))
        b = _hex_to_rgb_tuple(tokens.color("paper"))
    except (KeyError, ValueError):
        return None
    dark, light = (a, b) if _luminance(a) <= _luminance(b) else (b, a)
    if _luminance(light) - _luminance(dark) < 80.0:
        try:
            accent = _hex_to_rgb_tuple(tokens.color("accent"))
            if _luminance(accent) > _luminance(dark):
                light = accent
        except (KeyError, ValueError):
            pass
    return dark, light


def _align_pp(name: str) -> PP_ALIGN:
    return {
        "left":   PP_ALIGN.LEFT,
        "right":  PP_ALIGN.RIGHT,
        "center": PP_ALIGN.CENTER,
    }.get(name, PP_ALIGN.LEFT)


# ---------------------------------------------------------------------------
# Primitive handlers
# ---------------------------------------------------------------------------

def _emit_text(slide, node: DSLNode, ctx: EmitContext) -> None:
    """text X,Y "label" style:S align:A maxwidth:W maxheight:H color:T autoshrink:true lang:de_DE italic:true

    `autoshrink:true` — shrink font from style size down to a 10pt floor until
        the label fits the `maxwidth × maxheight` box. Off by default; preserves
        byte-identical output for layouts that don't opt in.
    `lang:LOCALE` — pyphen locale (e.g. `de_DE`). Inserts U+00AD soft hyphens
        into the label before emission so python-pptx can break long compound
        words at syllable boundaries. Applied BEFORE autoshrink so the shrink
        sees the hyphenated string.
    `padding:L,T,R,B` / `padding:N` (px) — explicit text-frame insets that
        mirror the source PPTX bodyPr lIns/tIns/rIns/bIns values. Affects
        BOTH the rendered text-frame margins (tf.margin_*) AND the fit budgets
        (autoshrink width/height clamp); omitting it uses PowerPoint's
        built-in defaults (~19 px left/right, ~9 px top/bottom).
    `bullet:true` / `bullet:"CHAR"` — add native PPTX bullets (`<a:buChar>`)
        with a hanging indent to every paragraph in the text frame. `true`
        uses the standard “•” bullet; any other string is used verbatim as the
        bullet character (e.g. `bullet:"–"` for an en-dash list). Opt-in only:
        omitting the kwarg leaves the output byte-identical. When using
        `bullet:`, the label text MUST NOT start with a literal bullet glyph —
        the `_leading_hang_offset_px` side-bearing hack is intentionally left
        alone; since the label won’t start with “•” the hang logic no-ops naturally.
    """
    if not node.pos_args:
        raise ValueError(f"text at line {node.line_no}: expected 'X,Y' positional")
    x, y = parse_xy(node.pos_args[0])
    style_name = node.kw_args.get("style", "body")
    style = ctx.tokens.resolve_style(style_name)
    color_override = node.kw_args.get("color")
    if color_override:
        from dataclasses import replace as _replace
        style = _replace(style, color_hex=ctx.tokens.color(color_override),
                         color_role=color_override)
    # `weight:<token>` overrides the style's default font weight without
    # forcing the author to switch style bundles. Lets a single layout pair
    # display-size with bold (or huge with regular), which the predefined
    # bundles don't express (huge/display are light-only, title-l is
    # bold-only). Token must exist in tokens.json.font-weight.
    weight_override = node.kw_args.get("weight")
    if weight_override:
        from dataclasses import replace as _replace
        style = _replace(style, weight=ctx.tokens.font_weight(weight_override))
    # `size:<N>px` or `size:<N>pt` or `size:<token-name>` lets a single text
    # primitive escape its style bundle's fixed size. Critical for matching
    # source decks whose pt sizes fall between the bundle steps (16/26/44/80
    # /120/160 px) — without it, the decompiler has to round to the nearest
    # bundle and a 42pt source title renders at the 44px sub bundle ≈ 33pt,
    # noticeably small. Numeric forms accepted: "32pt", "56px", or bare int
    # treated as px.
    size_override = node.kw_args.get("size")
    if size_override:
        from dataclasses import replace as _replace
        raw = size_override.strip().lower()
        if raw.endswith("pt"):
            # `pt` → design-px uses the SAME conversion the emitter rounds-
            # trip with (Pt(_px_to_pt(size_px)) downstream). Using the CSS
            # convention (pt × 4/3) here bakes a 96-DPI assumption and
            # halves the rendered font when the slide is sized for a
            # different DPI — e.g. 42pt → 56px → 21pt on a 10" slide.
            size_px = float(raw[:-2]) / _PX_TO_PT
        elif raw.endswith("px"):
            size_px = float(raw[:-2])
        else:
            try:
                size_px = float(raw)
            except ValueError:
                size_px = ctx.tokens.font_size_px(raw)
        style = _replace(style, size_px=size_px)
    # Hierarchy stepping: indent:N steps size/weight/color N times.
    try:
        indent_level = int(node.kw_args.get("indent", "0"))
    except (TypeError, ValueError):
        indent_level = 0
    if indent_level > 0:
        from dataclasses import replace as _replace
        new_size, new_weight, new_color_role = _step_hierarchy(
            style.size_px, style.weight, style.color_role, level=indent_level,
        )
        try:
            new_color_hex = ctx.tokens.color(new_color_role)
        except KeyError:
            new_color_hex = style.color_hex
        style = _replace(
            style,
            size_px=new_size, weight=new_weight,
            color_hex=new_color_hex, color_role=new_color_role,
        )
    if str(node.kw_args.get("italic", "")).lower() == "true":
        from dataclasses import replace as _replace
        style = _replace(style, italic=True)
    align = _align_pp(node.kw_args.get("align", "left"))
    # Default right margin = brand's slide.padding-x token; fall back to a
    # sensible canvas-relative inset if the brand omits it.
    try:
        padding_x = ctx.tokens.slide("padding-x")
    except Exception:
        padding_x = _DEFAULT_PADDING_X
    if not padding_x:
        padding_x = _DEFAULT_PADDING_X
    maxw = float(node.kw_args.get("maxwidth", ctx.canvas_w - x - padding_x))
    # Default height: ~1.5× font size (in design-px) to give room for one line.
    h = float(node.kw_args.get("maxheight", style.size_px * 1.5))

    label_text = node.label or ""
    if label_text:
        label_text = normalize_text(label_text, locale=ctx.tokens.locale)
    lang = node.kw_args.get("lang")
    if lang:
        label_text = textfit.hyphenate(label_text, lang=lang)

    # Text-frame insets: source PPTX bodyPr carries `lIns/tIns/rIns/bIns`
    # (EMU). When the decompiler captures them they ride through as
    # `padding:` kwarg (px, applied to tf.margin_* once the frame exists
    # below). Parsed here, before the fit math, because the insets eat into
    # the wrap width/height — the fit predictor must see the same usable
    # area the renderer sees. Without explicit padding, PowerPoint's
    # defaults apply: lIns/rIns 91440 EMU, tIns/bIns 45720 EMU.
    padding = node.kw_args.get("padding")
    pad_left = pad_top = pad_right = pad_bottom = None
    if padding is not None:
        # Accept `padding:L,T,R,B` (px) or `padding:N` (uniform px).
        parts = [p.strip() for p in str(padding).split(",")]
        if len(parts) == 1:
            pad_left = pad_top = pad_right = pad_bottom = float(parts[0])
        elif len(parts) == 4:
            pad_left, pad_top, pad_right, pad_bottom = (float(p) for p in parts)
        else:
            pad_left = pad_top = pad_right = pad_bottom = 0.0
    if pad_left is not None:
        inset_w_emu = int(pad_left * _EMU_PER_PX) + int(pad_right * _EMU_PER_PX)
        inset_h_emu = int(pad_top * _EMU_PER_PX) + int(pad_bottom * _EMU_PER_PX)
    else:
        # PowerPoint OOXML defaults, already in EMU — no scale conversion
        inset_w_emu = 91440 + 91440
        inset_h_emu = 45720 + 45720

    autoshrink = str(node.kw_args.get("autoshrink", "")).lower() == "true"
    # Resolve font face once for all fit/shrink paths (autoshrink, orphan
    # control, autofit gating). style.weight and style.font_family are
    # frozen at this point; only size_px changes below, which does not
    # affect face selection. _style_run resolves its own face independently.
    fit_face, fit_bold = _resolve_face(style.font_family[0], style.weight)
    if autoshrink and label_text:
        from dataclasses import replace as _replace
        # Use the longest single paragraph as the worst-case for width fit; the
        # height check below works on the joined text.
        max_pt = _px_to_pt(style.size_px)
        fitted_pt = textfit.autoshrink_size(
            label_text,
            font=fit_face,
            max_size_pt=max_pt,
            min_size_pt=10,
            bold=fit_bold,
            width_emu=max(1, int(maxw * _EMU_PER_PX) - inset_w_emu),
            height_emu=max(1, int(h * _EMU_PER_PX) - inset_h_emu),
            line_height=style.line_height,
        )
        if fitted_pt < max_pt:
            # Convert pt back to design-px (2 design-px per pt).
            style = _replace(style, size_px=fitted_pt * 2.0)

    # Orphan control: per paragraph, if the greedy wrap would orphan one
    # word on the final line, replace its preceding space with NBSP so
    # the pair wraps together. Use the final (post-autoshrink) size.
    if label_text:
        paragraphs = label_text.split("\n")
        paragraphs = [
            textfit.prevent_orphan(
                p,
                font=fit_face,
                size_pt=_px_to_pt(style.size_px),
                bold=fit_bold,
                width_emu=max(1, int(maxw * _EMU_PER_PX) - inset_w_emu),
            )
            for p in paragraphs
        ]
        label_text = "\n".join(paragraphs)

    # Hanging punctuation: shift the textbox left so the leading glyph
    # hangs into the margin (its body-letter optical edge sits at x).
    first_line = label_text.split("\n", 1)[0] if label_text else ""
    hang_offset_px = _leading_hang_offset_px(first_line, _px_to_pt(style.size_px))
    box = slide.shapes.add_textbox(
        _px(x - hang_offset_px), _px(y), _px(maxw), _px(h)
    )
    # `rotate:N` — rotate the textbox N degrees clockwise around its
    # geometric center. Use -90 for a y-axis label reading bottom-to-top.
    # Bbox semantics: x,y,maxwidth,maxheight describe the *pre-rotation*
    # rect; PPTX rotates around the rect's center on render. To place a
    # vertical label at a chart's left edge, set x to the chart edge minus
    # half the natural width, and y so the rect's vertical center matches
    # the chart's vertical midpoint.
    rotate_str = node.kw_args.get("rotate")
    if rotate_str:
        try:
            box.rotation = float(rotate_str)
        except (TypeError, ValueError):
            pass
    # Tag title text-shapes so downstream verify helpers
    # (extract_titles_from_pptx, extract_slide_title_and_body) can identify
    # them. Feinschliff layouts don't use python-pptx placeholder types, so
    # slide.shapes.title is always None — the name prefix is how we recover
    # the semantic role post-render. Prefix is chosen to NOT collide with the
    # chrome-shape prefixes in lib/verify/chrome.py.
    #
    # The 5 styles tagged here cover every title-role text in the shared
    # layouts catalog:
    #   title, act-title — standard content-slide titles
    #   sub             — used only in agenda.slide.dsl for the slide headline
    #   display-xl      — used only in end.slide.dsl for the closer headline
    #   quote           — used in full-bleed-cover + quote.slide.dsl (both
    #                     are "headline-shaped" usages, not body decoration)
    if style_name in ("title", "act-title", "sub", "display-xl", "quote"):
        box.name = f"feinschliff-title-{style_name}"
    tf = box.text_frame
    tf.word_wrap = True
    # Apply the explicit `padding:` (parsed above, before the fit math) to
    # the text-frame margins. Zeroing when nothing is set would bake "render
    # at frame edge" which mismatches PowerPoint's default 91440/45720
    # insets (~19px / ~9px at the source slide scale) and shifts text by
    # exactly those amounts versus source. Honour explicit padding;
    # otherwise leave python-pptx defaults (which mirror PowerPoint's
    # authoring defaults).
    if pad_left is not None:
        tf.margin_left = _px(pad_left)
        tf.margin_right = _px(pad_right)
        tf.margin_top = _px(pad_top)
        tf.margin_bottom = _px(pad_bottom)
    valign = node.kw_args.get("valign", "top")
    if valign == "middle":
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif valign == "bottom":
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    if autoshrink:
        if not textfit.has_real_metrics(fit_face, fit_bold):
            # Heuristic pre-shrink only: keep PPT's native shrink-to-fit as
            # the last line of defense. With real measured metrics the
            # computed size is authoritative — writing scale-less autofit
            # would let PowerPoint and LibreOffice re-derive different sizes.
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        else:
            # Drop the new-textbox default <a:spAutoFit/> too: the box was
            # sized for the computed fit; no renderer should regrow it.
            # Trade-off: if real-metrics textfit ever over-estimates usable
            # area, the box overflows visually — verify catches it post-render;
            # there is no renderer-side fallback on this path.
            tf.auto_size = None
    lines = label_text.split("\n")
    p0 = tf.paragraphs[0]
    p0.alignment = align
    p0.line_spacing = style.line_height
    _style_run(p0.add_run(), lines[0], style, tokens=ctx.tokens)
    for extra in lines[1:]:
        p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = style.line_height
        # Inter-paragraph gap defaults to a fraction of font-size; zero it so
        # the visual gap matches CSS line-height applied across all lines.
        p.space_before = Pt(0)
        _style_run(p.add_run(), extra, style, tokens=ctx.tokens)

    # Native bullets (opt-in: only when `bullet:` kwarg is present and not
    # explicitly disabled). Adds <a:buChar> + hanging indent to every
    # paragraph. Existing layouts that omit the kwarg are byte-identical
    # (early return). `bullet:false` (case-insensitive) and empty string are
    # treated as "no bullet" — mirrors the autoshrink idiom so that
    # `bullet:false` in DSL source doesn't yield a literal "false" glyph.
    bullet_raw = node.kw_args.get("bullet")
    if bullet_raw and str(bullet_raw).lower() not in ("false", ""):
        from pptx.oxml.ns import qn
        char = "•" if str(bullet_raw).lower() == "true" else str(bullet_raw)
        if len(char) != 1:
            raise DSLError(
                f"text at line {node.line_no}: bullet: value must be a single "
                f"Unicode character, got {char!r}"
            )
        # Hanging indent ≈ 1.4em so wrapped lines align under the first
        # character, not under the bullet. marL/indent are in EMU.
        mar_l = int(_px_to_pt(style.size_px) * 1.4 * EMU_PER_PT)
        for p in tf.paragraphs:
            pPr = p._p.get_or_add_pPr()
            pPr.set("marL", str(mar_l))
            pPr.set("indent", str(-mar_l))
            # CT_TextParagraphProperties child order: lnSpc, spcBef, spcAft,
            # buClr*, buSz*, buFont, buChar/buNone/buAutoNum, tabLst, defRPr,
            # extLst. Insert before defRPr/extLst when present, else append.
            anchor = pPr.find(qn("a:defRPr"))
            if anchor is None:
                anchor = pPr.find(qn("a:extLst"))
            bu_font = pPr.makeelement(qn("a:buFont"), {"typeface": fit_face})
            bu_char = pPr.makeelement(qn("a:buChar"), {"char": char})
            if anchor is not None:
                anchor.addprevious(bu_font)
                anchor.addprevious(bu_char)
            else:
                pPr.append(bu_font)
                pPr.append(bu_char)


def _is_numeric_run(text: str) -> bool:
    """True if every character in text is a digit, separator, or numeric sign."""
    return bool(text) and bool(_NUMERIC_CONTENT_RE.match(text))


def _tracking_for_size(size_pt: float, curve: dict[int, float]) -> float:
    """Step-function lookup: return curve[max-key ≤ size_pt], or 0.0 if no
    curve key is at or below size_pt. The curve is a brand-level default
    for display-tier tracking; per-style `letter_spacing` overrides it.
    """
    if not curve:
        return 0.0
    sorted_keys = sorted(curve.keys())
    if size_pt < sorted_keys[0]:
        return 0.0
    chosen = sorted_keys[0]
    for k in sorted_keys:
        if k <= size_pt:
            chosen = k
        else:
            break
    return curve[chosen]


def _style_run(run, text: str, style, *, tokens: Tokens | None = None) -> None:
    if style.transform == "upper":
        text = text.upper()
    run.text = text
    f = run.font
    # Tabular-numeral override: when the brand exposes a complete tnum_font
    # face name (e.g. "Inter Tabular") and the run is a pure number, use
    # that face verbatim; bold comes from the style weight so emphasis is
    # preserved. Otherwise resolve face from family + weight as usual.
    if tokens is not None and tokens.tnum_font and _is_numeric_run(text):
        face = tokens.tnum_font
        bold = style.weight >= 600
    else:
        face, bold = _resolve_face(style.font_family[0], style.weight)
    f.name = face
    f.size = Pt(_px_to_pt(style.size_px))
    f.bold = bold
    if style.italic:
        f.italic = True
    color = style.color_hex
    if style.opacity < 1.0:
        # python-pptx has no native alpha on text runs; approximate by
        # pre-blending toward the brand paper color. Canonical .pgmeta
        # opacity 0.7 on ink over white paper → ~#4A586F; over dark paper
        # the blend recedes toward the background (not toward white).
        try:
            paper = tokens.color("paper") if tokens is not None else "#FFFFFF"
        except KeyError:
            paper = "#FFFFFF"
        color = _blend_toward(color, paper, 1.0 - style.opacity)
    f.color.rgb = _hex_to_rgb(color)
    # Tracking: per-style letter_spacing wins; if the style declares none,
    # fall back to the brand-level display_tracking_curve (step function).
    size_pt = _px_to_pt(style.size_px)
    tracking_em = style.letter_spacing
    if tracking_em == 0.0 and tokens is not None:
        tracking_em = _tracking_for_size(size_pt, tokens.display_tracking_curve)
    if tracking_em != 0.0:
        # OOXML `spc` is 100ths of a point. Negative tightens; positive loosens.
        spc = int(tracking_em * size_pt * 100)
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(spc))


# Map design weight → face-name suffix. Noto Sans ships every weight as a
# separate face with the weight encoded in the family name (e.g. "Noto Sans
# Light"). soffice resolves family-by-name, so we pick the right face here.
# Weight 700 is the conventional "bold" face — use regular family + bold flag.
_WEIGHT_SUFFIX = [
    (200, "Thin"),
    (300, "Light"),
    (400, None),        # Regular
    (500, "Medium"),
    (600, "SemiBold"),
    (700, "__bold__"),  # signal: use regular family + bold flag
    (800, "ExtraBold"),
    (1000, "Black"),
]


def _resolve_face(family: str, weight: int) -> tuple[str, bool]:
    for thresh, suffix in _WEIGHT_SUFFIX:
        if weight <= thresh:
            if suffix == "__bold__":
                return family, True
            return (f"{family} {suffix}" if suffix else family), False
    return family, False


def _blend_toward(hex_color: str, target_hex: str, t: float) -> str:
    """Linearly blend `hex_color` toward `target_hex` by fraction `t` (0..1).
    Approximates run alpha over a solid background — python-pptx has no
    native alpha on text runs, so we pre-blend toward the slide paper."""
    h = hex_color.lstrip("#"); g = target_hex.lstrip("#")
    out = []
    for i in (0, 2, 4):
        c = int(h[i:i+2], 16); tg = int(g[i:i+2], 16)
        out.append(round(c + (tg - c) * t))
    return "#{:02X}{:02X}{:02X}".format(*out)


def _px_to_pt(px: float) -> float:
    """Design-px (1920×1080 baseline) → PowerPoint points.

    A standard 16:9 PPT slide is 13.33×7.5in = 960×540pt. The 1920-wide
    design baseline maps 2 design-px per point (960/1920 = 0.5).
    """
    return px * _PX_TO_PT


def _emit_rect(slide, node: DSLNode, ctx: EmitContext) -> None:
    """rect X,Y WxH fill:role [radius:N] [stroke:role stroke-width:N dash:preset]

    `radius:N` (design-px) switches the shape primitive to a rounded
    rectangle and sets the corner adjustment to N as a fraction of the
    shortest side. Decompiled brand packs use this for source PPTX
    `<a:prstGeom prst="roundRect">` shapes — without it, source rounded
    cards render as sharp rectangles.

    `dash:preset` accepts the PowerPoint preset names (`solid`, `dash`,
    `dot`, `sysDash`, `sysDot`, `dashDot`, `lgDashDot`, …) — passes
    through to python-pptx's `line.dash_style` via `MSO_LINE_DASH_STYLE`.
    """
    x, y = parse_xy(node.pos_args[0])
    w, h = parse_wh(node.pos_args[1])
    radius_attr = node.kw_args.get("radius")
    radius_px = float(radius_attr) if radius_attr else 0.0
    if radius_px > 0:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, _px(x), _px(y), _px(w), _px(h))
        # Adjustment value is fraction of shortest side. python-pptx exposes
        # this via shape.adjustments — index 0 is the corner-radius knob on
        # a rounded rect, valid range 0.0..0.5.
        shortest = min(w, h)
        if shortest > 0:
            adj = max(0.0, min(0.5, radius_px / shortest))
            try:
                shape.adjustments[0] = adj
            except (IndexError, AttributeError):
                pass
    else:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, _px(x), _px(y), _px(w), _px(h))
    shape.shadow.inherit = False
    _strip_theme_style(shape)

    gradient_attr = node.kw_args.get("gradient")
    fill_name = node.kw_args.get("fill")
    if gradient_attr:
        _apply_gradient_fill(shape, gradient_attr, ctx.tokens)
    elif fill_name:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(ctx.tokens.color(fill_name))
    else:
        shape.fill.background()

    stroke_name = node.kw_args.get("stroke")
    if stroke_name:
        shape.line.color.rgb = _hex_to_rgb(ctx.tokens.color(stroke_name))
        sw = float(node.kw_args.get("stroke-width", 1))
        shape.line.width = Pt(sw * _STROKE_PX_TO_PT)
        dash_name = node.kw_args.get("dash")
        if dash_name:
            _apply_line_dash(shape.line, dash_name)
    else:
        shape.line.fill.background()

    shadow_attr = node.kw_args.get("shadow")
    if shadow_attr:
        _apply_drop_shadow(shape, shadow_attr, ctx.tokens)


def _apply_gradient_fill(shape, descriptor: str, tokens) -> None:
    """Replace the shape's solidFill with a `<a:gradFill>` block.

    Descriptor format: `angle=Ddeg;0.00=token;1.00=token` — matches what
    the decompiler emits. Stops are positioned 0..1. Replaces any existing
    solid/no-fill on `<p:spPr>` and survives `sanitize_chrome` (which by
    default replaces gradFill with the first stop's solid colour) via the
    `_EFFECT_ALLOW_ATTR` opt-in on the parent `<p:sp>`.
    """
    from lxml import etree
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    parts = [p.strip() for p in descriptor.split(";") if p.strip()]
    angle_deg = 0.0
    stops: list[tuple[float, str]] = []
    for p in parts:
        if "=" not in p:
            continue
        k, v = p.split("=", 1)
        if k == "angle":
            try:
                angle_deg = float(v.replace("deg", ""))
            except ValueError:
                pass
        else:
            try:
                stops.append((float(k), v))
            except ValueError:
                continue
    if not stops:
        return
    shape._element.set(_EFFECT_ALLOW_ATTR, "1")
    spPr = shape._element.spPr
    # Remove any existing fill on the spPr (solidFill / noFill / gradFill).
    for tag in ("solidFill", "noFill", "gradFill", "blipFill", "pattFill"):
        for old in spPr.findall(f"{{{A}}}{tag}"):
            spPr.remove(old)
    grad = etree.SubElement(spPr, f"{{{A}}}gradFill")
    grad.set("flip", "none")
    grad.set("rotWithShape", "1")
    gs_lst = etree.SubElement(grad, f"{{{A}}}gsLst")
    for pos, color_ref in stops:
        gs = etree.SubElement(gs_lst, f"{{{A}}}gs")
        gs.set("pos", str(int(round(pos * 100000))))
        try:
            color_hex = tokens.color(color_ref).lstrip("#")
        except KeyError:
            color_hex = color_ref.lstrip("#") if color_ref.startswith("#") else "808080"
        srgb = etree.SubElement(gs, f"{{{A}}}srgbClr")
        srgb.set("val", color_hex.upper())
    lin = etree.SubElement(grad, f"{{{A}}}lin")
    lin.set("ang", str(int(round(angle_deg * 60000))))
    lin.set("scaled", "0")


def _apply_drop_shadow(shape, descriptor: str, tokens) -> None:
    """Attach an `<a:effectLst><a:outerShdw>` block to the shape's spPr.

    Descriptor format: `blur:Npx,dist:Npx,angle:Ddeg,color:T,alpha:0.30` —
    matches what the decompiler emits. python-pptx's `shape.shadow` API
    can't write outerShdw effects, so we build the XML directly. PowerPoint
    units: blur/dist in EMU, angle in 1/60000ths of a degree, alpha in
    0..100000 (≡ 0..100%).
    """
    fields = dict(kv.split(":", 1) for kv in descriptor.split(",") if ":" in kv)
    try:
        blur_emu = int(float(fields.get("blur", "0")) * _EMU_PER_PX)
        dist_emu = int(float(fields.get("dist", "0")) * _EMU_PER_PX)
        angle = int(float(fields.get("angle", "0")) * 60000)
        alpha = int(float(fields.get("alpha", "1.0")) * 100000)
    except ValueError:
        return
    color_token = fields.get("color", "black")
    try:
        color_hex = tokens.color(color_token).lstrip("#")
    except KeyError:
        color_hex = color_token.lstrip("#") if color_token.startswith("#") else "000000"
    from lxml import etree
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    spPr = shape._element.spPr
    # Opt this shape into surviving sanitize_chrome's effectLst strip
    # (which is otherwise scoped to flat templated brands that don't want
    # PowerPoint's default shadow ride-alongs). Mark via the shared
    # opt-in attribute so a decompile-derived shadow rebuilt from source
    # XML isn't stripped after we've just attached it.
    shape._element.set(_EFFECT_ALLOW_ATTR, "1")
    # Strip any existing effectLst so a re-emit doesn't stack shadows.
    for old in spPr.findall(f"{{{A}}}effectLst"):
        spPr.remove(old)
    effect_lst = etree.SubElement(spPr, f"{{{A}}}effectLst")
    shdw = etree.SubElement(effect_lst, f"{{{A}}}outerShdw")
    shdw.set("blurRad", str(blur_emu))
    shdw.set("dist", str(dist_emu))
    shdw.set("dir", str(angle))
    shdw.set("algn", "bl")
    shdw.set("rotWithShape", "0")
    srgb = etree.SubElement(shdw, f"{{{A}}}srgbClr")
    srgb.set("val", color_hex.upper())
    alpha_el = etree.SubElement(srgb, f"{{{A}}}alpha")
    alpha_el.set("val", str(alpha))


# PowerPoint preset-dash names → python-pptx MSO_LINE_DASH_STYLE. Keys are
# the literal `prstDash val=` strings the decompiler emits; misses fall
# through to the renderer default (solid).
_DASH_PRESETS = {
    "solid": "SOLID", "dash": "DASH", "dot": "ROUND_DOT",
    "sysDash": "SQUARE_DOT", "sysDot": "ROUND_DOT",
    "dashDot": "DASH_DOT", "lgDash": "LONG_DASH",
    "lgDashDot": "LONG_DASH_DOT", "lgDashDotDot": "LONG_DASH_DOT_DOT",
    "sysDashDot": "DASH_DOT", "sysDashDotDot": "DASH_DOT_DOT",
}


def _apply_line_dash(line, dash_name: str) -> None:
    """Map a `prstDash val=` string onto MSO_LINE_DASH_STYLE."""
    try:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
    except ImportError:
        return
    enum_name = _DASH_PRESETS.get(dash_name)
    if enum_name and hasattr(MSO_LINE_DASH_STYLE, enum_name):
        line.dash_style = getattr(MSO_LINE_DASH_STYLE, enum_name)


def _emit_line(slide, node: DSLNode, ctx: EmitContext) -> None:
    """line X,Y X2,Y2 stroke:role stroke-width:N"""
    x1, y1 = parse_xy(node.pos_args[0])
    x2, y2 = parse_xy(node.pos_args[1])
    line = slide.shapes.add_connector(1, _px(x1), _px(y1), _px(x2), _px(y2))
    stroke = node.kw_args.get("stroke", "fog")
    line.line.color.rgb = _hex_to_rgb(ctx.tokens.color(stroke))
    sw = float(node.kw_args.get("stroke-width", 1))
    line.line.width = Pt(sw * _STROKE_PX_TO_PT)
    _strip_theme_style(line)


def _emit_polyline(slide, node: DSLNode, ctx: EmitContext) -> None:
    """polyline X1,Y1 X2,Y2 X3,Y3 ... stroke:role stroke-width:N

    Open multi-segment line via python-pptx's freeform builder. N≥2 points.
    Stroke styled like `line`; no fill (open path).
    """
    pts = [parse_xy(p) for p in node.pos_args]
    if len(pts) < 2:
        raise ValueError(f"polyline: needs ≥2 points, got {len(pts)}")
    x0, y0 = pts[0]
    ff = slide.shapes.build_freeform(_px(x0), _px(y0), scale=1.0)
    ff.add_line_segments(
        [(_px(x), _px(y)) for x, y in pts[1:]],
        close=False,
    )
    poly = ff.convert_to_shape()
    poly.fill.background()
    stroke = node.kw_args.get("stroke", "ink")
    poly.line.color.rgb = _hex_to_rgb(ctx.tokens.color(stroke))
    sw = float(node.kw_args.get("stroke-width", 2))
    poly.line.width = Pt(sw * _STROKE_PX_TO_PT)
    _strip_theme_style(poly)


def _apply_picture_treatment(
    pil_image: Image.Image,
    treatment: str,
    *,
    duotone_dark: tuple[int, int, int] | None = None,
    duotone_light: tuple[int, int, int] | None = None,
) -> Image.Image:
    """Apply tokens.picture_treatment to a PIL image. Pure function — returns
    a new image. Supported treatments:
      - "none"          — no change
      - "desat(<frac>)" — blend toward grayscale by `frac` (0..1)
      - "duotone"       — convert to grayscale, then map luminance to a
        gradient from `duotone_dark` (shadow) to `duotone_light` (highlight).
        Caller must supply both RGB tuples; without them the image returns
        unchanged.
    """
    if not treatment or treatment == "none":
        return pil_image
    if treatment.startswith("desat(") and treatment.endswith(")"):
        try:
            frac = float(treatment[len("desat("):-1])
        except ValueError:
            return pil_image
        has_alpha = pil_image.mode in ("RGBA", "LA") or "transparency" in pil_image.info
        rgb = pil_image.convert("RGB")
        gray = rgb.convert("L").convert("RGB")
        blended = Image.blend(rgb, gray, frac)
        if has_alpha:
            alpha = pil_image.convert("RGBA").split()[3]
            r, g, b = blended.split()
            return Image.merge("RGBA", (r, g, b, alpha))
        return blended
    if treatment == "duotone":
        if duotone_dark is None or duotone_light is None:
            return pil_image
        gray = pil_image.convert("L")
        # Build a 256-entry RGB lookup from dark (luminance 0) to light (255).
        dr, dg, db = duotone_dark
        lr, lg, lb = duotone_light
        r_lut = bytes(round(dr + (lr - dr) * i / 255) for i in range(256))
        g_lut = bytes(round(dg + (lg - dg) * i / 255) for i in range(256))
        b_lut = bytes(round(db + (lb - db) * i / 255) for i in range(256))
        r = gray.point(r_lut)
        g = gray.point(g_lut)
        b = gray.point(b_lut)
        return Image.merge("RGB", (r, g, b))
    return pil_image


# Shared gem illustration used as the fallback when a picture asset can't be
# resolved. Lives at the feinschliff project root `assets/illustrations/`
# (parents[2] of this module: dsl -> feinschliff pkg -> project).
_PLACEHOLDER_ILLUSTRATION = (
    Path(__file__).resolve().parents[2] / "assets" / "illustrations" / "placeholder.jpg"
)


def _placeholder_image_path(ctx: EmitContext) -> "Path | None":
    """Locate the shared gem placeholder illustration, if available.

    Prefers the plugin asset-fallback root (so an installed wheel resolves it
    from packaged assets); falls back to the source-tree project assets.
    """
    if ctx.asset_root_fallback is not None:
        cand = ctx.asset_root_fallback / "illustrations" / "placeholder.jpg"
        if cand.is_file():
            return cand
    if _PLACEHOLDER_ILLUSTRATION.is_file():
        return _PLACEHOLDER_ILLUSTRATION
    return None


def _emit_picture_placeholder(
    slide,
    *,
    pos_xy: str,
    pos_wh: str,
    node: DSLNode,
    ctx: EmitContext,
) -> None:
    """Fallback for an unresolvable picture (missing file, unresolvable query,
    failed fetch, or unset slot).

    Places the shared gem illustration (`assets/illustrations/placeholder.jpg`)
    cropped to the slot so a missing image reads as intentional brand chrome
    rather than a grey box. Falls back to a brand-tonal rect only when the
    illustration can't be located or fails to load.

    Used exclusively by ``_emit_picture`` for all four fallback branches.
    """
    placeholder = _placeholder_image_path(ctx)
    if placeholder is not None:
        try:
            x, y = parse_xy(pos_xy)
            w, h = parse_wh(pos_wh)
            bytes_io = _prepare_picture_bytes(
                placeholder, target_aspect=(w / h),
                treatment=ctx.tokens.picture_treatment,
            )
            slide.shapes.add_picture(
                bytes_io, _px(x), _px(y), width=_px(w), height=_px(h)
            )
            return
        except Exception:
            pass  # any load/crop failure -> brand-tonal rect below
    _emit_rect(slide, DSLNode(
        kind="rect", pos_args=[pos_xy, pos_wh],
        kw_args={"fill": "paper-2", "stroke": "fog"},
        label=None, line_no=node.line_no, source=node.source,
    ), ctx)


def _resolve_provider_image(
    ctx: EmitContext,
    query: str,
    slot_id: str,
    *,
    slide,
    node: DSLNode,
    pos_xy: str,
    pos_wh: str,
) -> "Path | None":
    """Search the active image provider for *query*, materialise the hit, and
    return the local Path.  On any failure appends to ctx.missing_assets,
    emits a placeholder rect, and returns None.  Caller must not emit a
    further placeholder when None is returned."""
    hit = _img_mat.lookup_lock_then_search(
        ctx.image_provider, ctx.deck_dir, slot_id, query,
    )
    if hit is None or isinstance(hit, _img_mat._SearchError):
        entry: dict = {
            "kind": "search-error" if isinstance(hit, _img_mat._SearchError) else "no-hit",
            "query": query,
            "slot_id": slot_id,
            "provider": ctx.image_provider.name,  # type: ignore[union-attr]
            "line_no": node.line_no,
            "source": node.source,
        }
        if isinstance(hit, _img_mat._SearchError):
            entry["exc_type"] = hit.exc_type.__name__
        ctx.missing_assets.append(entry)
        _emit_picture_placeholder(slide, pos_xy=pos_xy, pos_wh=pos_wh, node=node, ctx=ctx)
        return None
    cache_dir = (ctx.deck_dir / ".cache") if ctx.deck_dir else None
    if cache_dir is None:
        cache_dir = Path(tempfile.mkdtemp(prefix="feinschliff-imgcache-"))
        _img_mat._THROWAWAY_CACHE_DIRS.append(cache_dir)
        _img_mat._register_throwaway_cache_cleanup()
        warnings.warn(
            "EmitContext.deck_dir is unset; HTTP image materialise will "
            "use a throwaway tempdir cache (no rebuild reuse). Wire "
            "deck_dir on the EmitContext (or pass it to "
            "build_presentation/build_multi_slide) to persist cached "
            "downloads in <deck_dir>/.cache/.",
            RuntimeWarning,
            stacklevel=2,
        )
    materialised, fetch_err = _img_mat._materialise(hit, cache_dir)
    if materialised is None:
        fail_entry: dict = {
            "kind": "fetch-failed",
            "query": query,
            "slot_id": slot_id,
            "url": hit.url,
            "provider": ctx.image_provider.name,  # type: ignore[union-attr]
            "line_no": node.line_no,
            "source": node.source,
        }
        if fetch_err is not None:
            fail_entry["error"] = f"{type(fetch_err).__name__}: {fetch_err}"
        ctx.missing_assets.append(fail_entry)
        _emit_picture_placeholder(slide, pos_xy=pos_xy, pos_wh=pos_wh, node=node, ctx=ctx)
        return None
    return materialised


def _contain_geometry(
    img_w: int, img_h: int, x: float, y: float, w: float, h: float,
) -> tuple[float, float, float, float]:
    """Scale img to fit w×h preserving aspect; center in the box. Design-px in/out."""
    if img_w <= 0 or img_h <= 0:
        return x, y, w, h
    if w <= 0 or h <= 0:
        return x, y, w, h
    src_aspect = img_w / img_h
    box_aspect = w / h
    if abs(src_aspect - box_aspect) < 1e-3:
        return x, y, w, h
    if src_aspect > box_aspect:
        new_h = w / src_aspect
        return x, y + (h - new_h) / 2.0, w, new_h
    new_w = h * src_aspect
    return x + (w - new_w) / 2.0, y, new_w, h


def _emit_picture(slide, node: DSLNode, ctx: EmitContext) -> None:
    """picture X,Y WxH path:PATH cover:true

    `path` is the resolved image location — either a literal in the layout
    or interpolated from a `{{ slot }}` placeholder by the expander. If
    `path` is missing, the node is skipped silently. If `path` does not
    resolve to a local file but an image provider is active, the node falls
    back to a provider search: an explicit ``query:`` kwarg takes
    precedence, otherwise the unresolved path value itself is used as the
    query (e.g. Unsplash) so plan authors can write
    ``image: "regensburg medieval bridge"`` without changing layouts. A
    ``query:`` is ignored whenever the path resolves to a real file.
    `cover:true` center-crops the source image to the box aspect ratio
    (default behaviour is contain).

    Diagram-emitted picture nodes (produced by ``expand_diagram_blocks``)
    store geometry in ``kw_args`` (x, y, w, h as ints) and carry a
    ``_diagram_meta`` sentinel.  Slide-authored nodes use ``pos_args``.
    """
    # Diagram-emitted picture: geometry in kw_args, src in kw_args["src"],
    # sentinel in kw_args["_diagram_meta"].  (Slide-authored picture uses
    # pos_args for position and size.)
    if "_diagram_meta" in node.kw_args:
        x = int(node.kw_args["x"])
        y = int(node.kw_args["y"])
        w = int(node.kw_args["w"])
        h = int(node.kw_args["h"])
        path = node.kw_args.get("src")
        _pos_xy = f"{x},{y}"
        _pos_wh = f"{w}x{h}"
        query: str | None = None
    else:
        x, y = parse_xy(node.pos_args[0])
        w, h = parse_wh(node.pos_args[1])
        path = node.kw_args.get("path")
        query = node.kw_args.get("query")
        _pos_xy = node.pos_args[0]
        _pos_wh = node.pos_args[1]

    # `query:` alongside `path:` is a layered fallback: the file wins when
    # `path` resolves; the explicit `query:` is consulted only when the
    # path misses (see the provider branch below), taking precedence over
    # the synthesized use-the-path-as-query fallback. A bare `query:`
    # (no path) resolves through the provider directly — and fails loud
    # below when no provider is wired, because that combination can never
    # produce an image.
    if query and not path:
        if not ctx.image_provider:
            # The brand author wrote `query:` in a layout but forgot to
            # wire `$image_provider` in tokens.json. Silent-fallback to a
            # placeholder rect would mask the misconfiguration — fail
            # loud so the next build run surfaces the problem.
            raise DSLError(
                f"picture at line {node.line_no}: `query:{query!r}` requires "
                f"an image_provider on the EmitContext, but ctx.image_provider "
                f"is None. Add `$image_provider` to your brand's tokens.json "
                f"(or your `extends` ancestor) so the build can resolve it."
            )
        slot_id = node.label or _img_mat._slot_id_from_query(query)
        materialised = _resolve_provider_image(
            ctx, query, slot_id,
            slide=slide, node=node, pos_xy=_pos_xy, pos_wh=_pos_wh,
        )
        if materialised is None:
            return
        path = str(materialised)

    optional = str(node.kw_args.get("optional", "false")).lower() == "true"
    if not path:
        # No image bound — emit a placeholder rect so the slot's location is
        # visible in dev renders. Required pictures append to ctx.missing_assets
        # so callers (build / deck build) can decide whether to abort. Mark
        # the picture optional via `optional:true` if a missing slot is OK.
        if not optional:
            ctx.missing_assets.append({
                "kind": "unset",
                "path": None,
                "line_no": node.line_no,
                "source": node.source,
            })
        _emit_picture_placeholder(slide, pos_xy=_pos_xy, pos_wh=_pos_wh, node=node, ctx=ctx)
        return
    p = Path(path)
    from_fallback = False
    if not p.is_absolute() and ctx.asset_root:
        primary = ctx.asset_root / p
        if primary.is_file():
            p = primary
        elif ctx.asset_root_fallback:
            fallback = ctx.asset_root_fallback / p
            if fallback.is_file():
                p = fallback
                from_fallback = True
            else:
                p = primary
        else:
            p = primary
    if not p.is_file():
        # When an image provider is active, fall back to a provider search
        # instead of failing. An explicit `query:` kwarg takes precedence;
        # otherwise the unresolved path value itself is treated as the
        # search query. This lets plan authors write
        # image: "regensburg aerial" and have it resolve through e.g.
        # Unsplash without requiring query: in every layout DSL file.
        if ctx.image_provider:
            provider_query = query or path
            slot_id = node.label or _img_mat._slot_id_from_query(provider_query)
            materialised = _resolve_provider_image(
                ctx, provider_query, slot_id,
                slide=slide, node=node, pos_xy=_pos_xy, pos_wh=_pos_wh,
            )
            if materialised is None:
                return
            p = materialised
        else:
            if not optional:
                ctx.missing_assets.append({
                    "kind": "missing-file",
                    "path": str(p),
                    "line_no": node.line_no,
                    "source": node.source,
                })
            _emit_picture_placeholder(slide, pos_xy=_pos_xy, pos_wh=_pos_wh, node=node, ctx=ctx)
            return

    cover = node.kw_args.get("cover", "false").lower() == "true"
    treatment = ctx.tokens.picture_treatment
    duotone_dark = duotone_light = None
    # Plugin-fallback images (generic placeholders shared across brands) get
    # duotoned to the brand palette so they read as native to each brand
    # pack instead of always showing the same warm-tone source photo. Brand-
    # specific overrides bypass this — brands choose their own treatment.
    if from_fallback:
        endpoints = _duotone_endpoints_for_brand(ctx.tokens)
        if endpoints is not None:
            treatment = "duotone"
            duotone_dark, duotone_light = endpoints
    if cover:
        bytes_io = _prepare_picture_bytes(p, target_aspect=(w / h),
                                          treatment=treatment,
                                          duotone_dark=duotone_dark,
                                          duotone_light=duotone_light)
        slide.shapes.add_picture(bytes_io, _px(x), _px(y), width=_px(w), height=_px(h))
    else:
        # Contain: scale to fit the box preserving aspect ratio, center in the slot.
        try:
            # Header-only open — PIL defers pixel decode, so this costs a stat+header read
            # even though _prepare_picture_bytes opens the file again.
            with Image.open(p) as _im:
                img_w, img_h = _im.size
            cx, cy, cw, ch = _contain_geometry(img_w, img_h, x, y, w, h)
        except Exception:
            cx, cy, cw, ch = x, y, w, h
        if treatment != "none":
            bytes_io = _prepare_picture_bytes(p, target_aspect=None,
                                              treatment=treatment,
                                              duotone_dark=duotone_dark,
                                              duotone_light=duotone_light)
            slide.shapes.add_picture(bytes_io, _px(cx), _px(cy), width=_px(cw), height=_px(ch))
        else:
            # Fast path: no treatment — let python-pptx read the file directly.
            slide.shapes.add_picture(str(p), _px(cx), _px(cy), width=_px(cw), height=_px(ch))


def _prepare_picture_bytes(
    image_path: Path, *, target_aspect: float | None, treatment: str,
    duotone_dark: tuple[int, int, int] | None = None,
    duotone_light: tuple[int, int, int] | None = None,
) -> io.BytesIO:
    """Load → optional center-crop → optional treatment → JPEG/PNG bytes."""
    with Image.open(image_path) as im:
        im.load()
        if target_aspect is not None:
            sw, sh = im.size
            src_aspect = sw / sh
            if abs(src_aspect - target_aspect) >= 1e-3:
                if src_aspect > target_aspect:
                    new_w = max(1, int(round(sh * target_aspect)))
                    offset = (sw - new_w) // 2
                    im = im.crop((offset, 0, offset + new_w, sh))
                else:
                    new_h = max(1, int(round(sw / target_aspect)))
                    offset = (sh - new_h) // 2
                    im = im.crop((0, offset, sw, offset + new_h))
        if treatment and treatment != "none":
            im = _apply_picture_treatment(
                im, treatment,
                duotone_dark=duotone_dark, duotone_light=duotone_light,
            )
        out = io.BytesIO()
        fmt = (im.format or "PNG").upper()
        if fmt in ("JPG", "JPEG") and im.mode not in ("RGB", "L"):
            im = im.convert("RGB")
        # Treated images may have lost the source-format hint; pick by mode.
        if fmt in ("JPG", "JPEG"):
            im.save(out, format="JPEG", quality=92, optimize=True)
        else:
            im.save(out, format="PNG", optimize=True)
        return io.BytesIO(out.getvalue())


# ---------------------------------------------------------------------------
# Entry
# ---------------------------------------------------------------------------

_SHAPE_KIND = {
    "rect":          MSO_SHAPE.RECTANGLE,
    "rectangle":     MSO_SHAPE.RECTANGLE,
    "oval":          MSO_SHAPE.OVAL,
    "ellipse":       MSO_SHAPE.OVAL,
    "circle":        MSO_SHAPE.OVAL,
    "triangle":      MSO_SHAPE.ISOSCELES_TRIANGLE,
    "triangle-down": MSO_SHAPE.ISOSCELES_TRIANGLE,   # rotated 180 at emit
    "triangle-left": MSO_SHAPE.RIGHT_TRIANGLE,
    "triangle-right":MSO_SHAPE.RIGHT_TRIANGLE,
    "right-triangle":MSO_SHAPE.RIGHT_TRIANGLE,
    "chevron":       MSO_SHAPE.CHEVRON,
    "right-arrow":   MSO_SHAPE.RIGHT_ARROW,
    "left-arrow":    MSO_SHAPE.LEFT_ARROW,
    "diamond":       MSO_SHAPE.DIAMOND,
    "trapezoid":     MSO_SHAPE.TRAPEZOID,
    "pie":           MSO_SHAPE.PIE,          # 270° wedge; rotate to position the missing 90°
    "pie-wedge":     MSO_SHAPE.PIE_WEDGE,    # 90° wedge
    "arc":           MSO_SHAPE.ARC,
    "block-arc":     MSO_SHAPE.BLOCK_ARC,
}


def _emit_shape(slide, node: DSLNode, ctx: EmitContext) -> None:
    """shape X,Y WxH kind:NAME fill:role stroke:role stroke-width:N rotate:DEG

    Generic shape emitter — wraps python-pptx MSO_SHAPE so layouts can
    place ovals, triangles, chevrons, etc. Fill / stroke behaviour matches
    `rect`. `rotate` is in degrees clockwise.
    """
    x, y = parse_xy(node.pos_args[0])
    w, h = parse_wh(node.pos_args[1])
    kind = node.kw_args.get("kind", "rect")
    mso_kind = _SHAPE_KIND.get(kind)
    if mso_kind is None:
        raise ValueError(f"shape: unknown kind '{kind}' (known: {sorted(_SHAPE_KIND)})")
    shape = slide.shapes.add_shape(mso_kind, _px(x), _px(y), _px(w), _px(h))
    shape.shadow.inherit = False
    _strip_theme_style(shape)
    rot = node.kw_args.get("rotate")
    if rot is None and kind == "triangle-down":
        rot = "180"
    if rot is None and kind == "triangle-left":
        rot = "270"
    if rot is None and kind == "triangle-right":
        rot = "90"
    if rot is not None:
        shape.rotation = float(rot)

    # Preset-shape adjustment handle (e.g. parallelogram top-edge skew).
    # Maps to <a:gd name="adj" fmla="val …"/>. Range 0..1; python-pptx stores
    # as fraction. Silently no-op for shapes without an adj handle.
    adj1 = node.kw_args.get("adj1")
    if adj1 is not None:
        try:
            shape.adjustments[0] = float(adj1)
        except (IndexError, ValueError):
            pass

    fill_name = node.kw_args.get("fill")
    if fill_name:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(ctx.tokens.color(fill_name))
        opacity = node.kw_args.get("fill-opacity")
        if opacity is not None:
            _set_fill_alpha(shape, float(opacity))
    else:
        shape.fill.background()

    stroke_name = node.kw_args.get("stroke")
    if stroke_name:
        shape.line.color.rgb = _hex_to_rgb(ctx.tokens.color(stroke_name))
        sw = float(node.kw_args.get("stroke-width", 1))
        shape.line.width = Pt(sw * _STROKE_PX_TO_PT)
    else:
        shape.line.fill.background()


# ---------------------------------------------------------------------------
# Chrome sanitation — post-build XML pass (Layer 1 task 10)
# ---------------------------------------------------------------------------

_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

# Custom attribute the emitter writes on a <p:sp> when its source DSL used
# `effect:allow` to opt out of sanitation.
# Namespaced form (Clark notation for lxml) — the sanctioned OOXML extension
# mechanism. Writers always use this; readers also accept the legacy bare form.
_FS_NS = "urn:feinschliff:emit"
_EFFECT_ALLOW_ATTR = f"{{{_FS_NS}}}effect-allow"   # Clark notation for lxml
_EFFECT_ALLOW_LEGACY = "effect-allow"              # accepted on read (old files)

# Outline clamp: hairline at 0.5pt = 6350 EMU; anything above 1pt (12700 EMU)
# is clamped down to hairline.
_OUTLINE_CLAMP_THRESHOLD_EMU = 12700
_OUTLINE_CLAMP_TARGET_EMU = 6350


def _effect_allowed(sp) -> bool:
    """Return True if a <p:sp> element carries the effect opt-in marker.

    The caller is responsible for passing the ``<p:sp>`` element itself —
    no ancestor walk is performed inside this helper.  Accepts both the
    current namespaced form (``{urn:feinschliff:emit}effect-allow``) and
    the legacy bare attribute so that old decks keep their opt-in after
    upgrading.
    """
    return (
        sp.get(_EFFECT_ALLOW_ATTR) == "1"
        or sp.get(_EFFECT_ALLOW_LEGACY) == "1"
    )


def sanitize_chrome(slide_xml) -> None:
    """Walk a slide's <p:spTree> and remove PowerPoint defaults that read
    as "AI-templated": drop-shadow / glow / soft-edge effects, gradient
    fills (replaced with solid), and outlines wider than 1pt.

    Scoped to <p:spTree> so master/layout inheritance is preserved.
    Mutates `slide_xml` in place. Idempotent.
    """
    from lxml import etree

    sptree = slide_xml.find(f".//{{{_NS_P}}}spTree")
    if sptree is None:
        return

    # 1. effectLst — strip unless parent <p:sp> opts in.
    for effect_lst in list(sptree.iter(f"{{{_NS_A}}}effectLst")):
        parent = effect_lst.getparent()
        while parent is not None and not parent.tag.endswith("}sp"):
            parent = parent.getparent()
        if parent is not None and _effect_allowed(parent):
            continue
        effect_lst.getparent().remove(effect_lst)

    # 2. gradFill → solidFill (with the first stop's color), unless the
    #    parent <p:sp> opts in via _effect_allowed() (same gate as the
    #    effectLst case above; accepts both namespaced and legacy bare
    #    forms — used by decompile-derived shapes that deliberately carry
    #    a source-faithful gradient).
    ns = {"a": _NS_A}
    for grad in list(sptree.iter(f"{{{_NS_A}}}gradFill")):
        sp_anc = grad.getparent()
        while sp_anc is not None and not sp_anc.tag.endswith("}sp"):
            sp_anc = sp_anc.getparent()
        if sp_anc is not None and _effect_allowed(sp_anc):
            continue
        grad_parent = grad.getparent()
        first_clr = grad.find(".//a:gs[1]/a:srgbClr", ns)
        idx = list(grad_parent).index(grad)
        grad_parent.remove(grad)
        if first_clr is None:
            continue
        color_val = first_clr.get("val") or "808080"
        solid = etree.SubElement(grad_parent, f"{{{_NS_A}}}solidFill")
        clr = etree.SubElement(solid, f"{{{_NS_A}}}srgbClr")
        clr.set("val", color_val)
        grad_parent.remove(solid)
        grad_parent.insert(idx, solid)

    # 3. Clamp outline widths > 1pt down to a 0.5pt hairline.
    for ln in sptree.iter(f"{{{_NS_A}}}ln"):
        w_str = ln.get("w")
        if w_str is None:
            continue
        try:
            w = int(w_str)
        except ValueError:
            continue
        if w > _OUTLINE_CLAMP_THRESHOLD_EMU:
            ln.set("w", str(_OUTLINE_CLAMP_TARGET_EMU))


def _set_fill_alpha(shape, opacity: float) -> None:
    """Set alpha on a solid fill via OOXML. opacity 0..1."""
    from lxml import etree
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    spPr = shape.fill._xPr
    solidFill = spPr.find("a:solidFill", nsmap)
    if solidFill is None:
        return
    srgbClr = solidFill.find("a:srgbClr", nsmap)
    if srgbClr is None:
        return
    alpha = etree.SubElement(srgbClr, "{%s}alpha" % nsmap["a"])
    alpha.set("val", str(int(opacity * 100000)))


def _strip_theme_style(shape) -> None:
    """Drop the `<p:style>` block from a newly-added shape.

    python-pptx's `add_shape` writes a default `<p:style>` element that
    references the theme's effect/fill/line slots (e.g. `<a:effectRef
    idx="2">` points at the second `<a:outerShdw>` entry in
    `theme1.xml`). LibreOffice and PowerPoint both honour that effect
    reference at render time *even though* `shape.shadow.inherit = False`
    writes an explicit `<a:effectLst/>` on the shape's `<p:spPr>`.

    Result before the strip: a subtle drop-shadow under every primitive
    in rendered PNGs. Flat-shape brand packs (the default, feinschliff,
    …) want the source to match the render; removing the `<p:style>`
    block does that.
    """
    nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    el = shape._element.find("p:style", nsmap)
    if el is not None:
        shape._element.remove(el)


def _merge_table_style(slide, style_b64: str, node: DSLNode) -> None:
    """Merge a carried `<a:tblStyle>` into the deck's tableStyles part.

    No-op when the deck already has a style with that styleId (dedupes the
    same style carried by several table slides) or ships no tableStyles part.
    """
    import base64

    from lxml import etree as _etree

    try:
        el = _etree.fromstring(base64.b64decode(style_b64))
    except Exception as exc:
        raise DSLError(
            f"native table (line {node.line_no}): unparseable carried table style — {exc}"
        )
    style_id = el.get("styleId")
    if not style_id:
        return
    for part in slide.part.package.iter_parts():
        if str(part.partname) != "/ppt/tableStyles.xml":
            continue
        root = _etree.fromstring(part.blob)
        if any(ch.get("styleId") == style_id for ch in root):
            return
        root.append(el)
        part._blob = _etree.tostring(
            root, xml_declaration=True, encoding="UTF-8", standalone=True
        )
        return


def _splice_native_parts(slide, frame_el, parts_b64: str, node: DSLNode) -> None:
    """Re-create a native-carried CHART or DIAGRAM external part-graph in this deck.

    `parts_b64` is base64 of a JSON list (emitted by the decompiler's
    `_emit_graphic_frame` chart / diagram branches). Each entry is either a full
    part {partname, content_type, blob (base64 bytes), reltype, parent, src_rid}
    or — only for a diagram's media image shared by two parents (data + drawing
    both point at the same imageN.png) — a ref entry {ref (old partname of an
    already-materialised part), parent, src_rid} that wires a SECOND relationship
    to the existing part without re-materialising it.

    The carrier references point at SOURCE rIds that are dead in this deck:
      * chart: the frame's `<c:chart r:id>` + each chart part's `<c:externalData
        r:id>`;
      * diagram: the frame's `<dgm:relIds r:dm/r:lo/r:qs/r:cs>` (slide-level) +
        the dataN part's `<dsp:dataModelExt relId>` (a slide rel, but stored
        INSIDE the data part's XML) + any data/drawing → image `r:embed`.

    We:
      1. Materialise every full part as a fresh `Part` (new partname from the
         package so it can't collide with parts already in the output deck).
      2. Wire relationships bottom-up: each child part is `relate_to`'d from its
         parent (a chart/data/drawing part for the leaves) or from `slide.part`
         (chart part + the four dgm parts + the drawing part). `relate_to`
         returns the NEW rId. Ref entries wire an extra rel to an existing part.
      3. Rewrite every r:id / r:embed / relId reference from the relevant src→new
         map — the frame's chart/diagram refs from the slide-level map; each
         part's own XML from that part's map; PLUS the dataN dataModelExt relId
         from the slide-level map (it's a slide rel living inside the data part)
         — then push the rewritten XML back onto each part's blob.

    `[Content_Types].xml` + save-reachability are automatic once each part has
    the right content_type and is related into the slide's part graph.
    """
    import base64
    import json
    from pptx.opc.package import Part
    from pptx.oxml.ns import qn

    try:
        entries = json.loads(base64.b64decode(parts_b64).decode("utf-8"))
    except Exception as exc:
        raise DSLError(
            f"native chart (line {node.line_no}): unparseable carried parts — {exc}"
        )
    # A native TABLE carries its referenced <a:tblStyle> (the tableStyleId in
    # the spliced <a:tbl> is dead in this deck's tableStyles.xml otherwise —
    # the renderer would fall back to a default style with the wrong header
    # fill / band colours / borders).
    style_entries = [e for e in entries if "table_style" in e]
    entries = [e for e in entries if "table_style" not in e]
    for e in style_entries:
        _merge_table_style(slide, e["table_style"], node)
    if not entries:
        return

    pkg = slide.part.package
    # Per content-type partname template (so a re-created xlsx lands in
    # /ppt/embeddings/, styles in /ppt/charts/, etc.). Falls back to a generic
    # chart-part slot for anything unforeseen — partname uniqueness is what
    # matters; the directory is cosmetic.
    _CT = {
        "application/vnd.openxmlformats-officedocument.drawingml.chart+xml":
            "/ppt/charts/chart%d.xml",
        "application/vnd.ms-office.chartstyle+xml":
            "/ppt/charts/style%d.xml",
        "application/vnd.ms-office.chartcolorstyle+xml":
            "/ppt/charts/colors%d.xml",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            "/ppt/embeddings/Microsoft_Excel_Worksheet%d.xlsx",
        # SmartArt / diagram parts — land them in /ppt/diagrams/ so the package
        # mirrors a PowerPoint-authored deck (cosmetic; uniqueness is what counts).
        "application/vnd.openxmlformats-officedocument.drawingml.diagramData+xml":
            "/ppt/diagrams/data%d.xml",
        "application/vnd.openxmlformats-officedocument.drawingml.diagramLayout+xml":
            "/ppt/diagrams/layout%d.xml",
        "application/vnd.openxmlformats-officedocument.drawingml.diagramStyle+xml":
            "/ppt/diagrams/quickStyle%d.xml",
        "application/vnd.openxmlformats-officedocument.drawingml.diagramColors+xml":
            "/ppt/diagrams/colors%d.xml",
        "application/vnd.ms-office.drawingml.diagramDrawing+xml":
            "/ppt/diagrams/drawing%d.xml",
        "image/png": "/ppt/media/image%d.png",
        "image/jpeg": "/ppt/media/image%d.jpeg",
    }

    # 1. Materialise every FULL part (no rels yet). `ref` entries (a diagram's
    #    shared image, already materialised under its first parent) carry no blob
    #    — skip them here; step 2 wires the extra relationship.
    new_part: dict[str, Part] = {}
    for e in entries:
        if "ref" in e:
            continue
        tmpl = _CT.get(e["content_type"], "/ppt/charts/chart%d.xml")
        pn = pkg.next_partname(tmpl)
        p = Part(pn, e["content_type"], pkg, blob=base64.b64decode(e["blob"]))
        new_part[e["partname"]] = p

    # 2. Wire rels bottom-up + build a per-PARENT src_rid → new_rid map so we can
    #    rewrite references inside each parent's XML afterwards. Keyed by the
    #    parent's OLD partname (chart/data/drawing parts) or the literal "slide"
    #    sentinel (chart part + the four dgm parts + the drawing part).
    rid_map: dict[str, dict[str, str]] = {}
    for e in entries:
        # A ref entry points at an already-materialised part (its OLD partname);
        # a full entry materialised one keyed by its own OLD partname.
        child = new_part.get(e["ref"]) if "ref" in e else new_part.get(e["partname"])
        if child is None:
            continue
        parent_key = e["parent"]
        if parent_key == "slide":
            parent_part = slide.part
        else:
            parent_part = new_part.get(parent_key)
            if parent_part is None:
                # Parent wasn't carried (shouldn't happen — capture is rooted at
                # the chart / dgm parts) — skip wiring this orphan rather than crash.
                continue
        new_rid = parent_part.relate_to(child, e["reltype"])
        src_rid = e.get("src_rid")
        if src_rid:
            rid_map.setdefault(parent_key, {})[src_rid] = new_rid

    # 3a. Rewrite the FRAME's references (parent == "slide"): chart `<c:chart
    #     r:id>` AND diagram `<dgm:relIds r:dm/r:lo/r:qs/r:cs>`.
    slide_map = rid_map.get("slide", {})
    if slide_map:
        for cref in frame_el.iter(qn("c:chart")):
            old = cref.get(qn("r:id"))
            if old in slide_map:
                cref.set(qn("r:id"), slide_map[old])
        # `dgm` isn't a prefix python-pptx's `qn` knows, so match relIds via its
        # Clark-notation namespace literal; the r:dm/lo/qs/cs attrs resolve fine.
        _DGM = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
        for relids in frame_el.iter("{%s}relIds" % _DGM):
            for _a in ("r:dm", "r:lo", "r:qs", "r:cs"):
                old = relids.get(qn(_a))
                if old in slide_map:
                    relids.set(qn(_a), slide_map[old])
        # A native-carried `<p:grpSp>` of decorative chrome may embed grouped
        # `<p:pic>` images whose `<a:blip r:embed>` points at a dead SOURCE rId.
        # Their media is carried as image parts (parent == "slide"), so rewrite
        # every blip on the frame from the slide map. (Chart/diagram frames carry
        # no slide-rel-referencing blip on the FRAME — their image refs live
        # inside the chart/data/drawing PARTS and are fixed by 3b's part_map — so
        # this is a no-op for them.)
        for _blip in frame_el.iter(qn("a:blip")):
            old = _blip.get(qn("r:embed"))
            if old in slide_map:
                _blip.set(qn("r:embed"), slide_map[old])
            old_link = _blip.get(qn("r:link"))
            if old_link in slide_map:
                _blip.set(qn("r:link"), slide_map[old_link])

    # 3b. Rewrite each carried part's own XML and push the rewritten bytes back.
    #     Each part rewrites r:id / r:embed from its OWN map (chart →
    #     <c:externalData r:id>; data/drawing → image <…r:embed>). The dataN part
    #     ALSO carries a `<dsp:dataModelExt relId>` that is a SLIDE rel (it points
    #     at the drawing part) — rewrite that one from the slide_map. We therefore
    #     visit every part that has either its own map OR a dataModelExt to fix.
    from lxml import etree as _etree
    _R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    _DSP = "http://schemas.microsoft.com/office/drawing/2008/diagram"
    for old_pn, part in new_part.items():
        part_map = rid_map.get(old_pn)
        try:
            root = _etree.fromstring(part.blob)
        except Exception:
            continue
        changed = False
        # The dataN dataModelExt relId is a SLIDE-level rel stored inside the data
        # part. Rewrite it from slide_map (NOT this part's map).
        if slide_map:
            for dme in root.iter("{%s}dataModelExt" % _DSP):
                old = dme.get("relId")
                if old in slide_map:
                    dme.set("relId", slide_map[old])
                    changed = True
        if part_map:
            for el2 in root.iter():
                for attr, val in list(el2.attrib.items()):
                    if attr.startswith(f"{{{_R}}}") and val in part_map:
                        el2.set(attr, part_map[val])
                        changed = True
        if changed:
            part.blob = _etree.tostring(root, xml_declaration=True,
                                        encoding="UTF-8", standalone=True)


def _native_sidecar_bytes(ref: str, ctx: EmitContext, node: DSLNode,
                          what: str) -> bytes:
    """Resolve a `*_file:` native-payload reference against the brand pack's
    asset root (then the plugin fallback root), mirroring picture-path
    resolution. Sidecar files keep multi-MB carried fragments out of the DSL
    text (a 33 MB carried group made a 44 MB .slide.dsl when inlined)."""
    p = Path(ref)
    if not p.is_absolute():
        for root in (ctx.asset_root, ctx.asset_root_fallback):
            if root is not None and (root / p).is_file():
                p = root / p
                break
    if not p.is_file():
        raise DSLError(
            f"native {what} (line {node.line_no}): sidecar file not found — {ref!r}"
            + (f" (searched under {ctx.asset_root})" if ctx.asset_root else "")
        )
    return p.read_bytes()


def _remap_colliding_shape_ids(slide, el) -> None:
    """Native fragments keep their SOURCE `p:cNvPr` ids. Two fragments carried
    from different sources (slide chrome vs layout template image) — or a
    fragment vs a python-pptx-generated shape — can claim the same id, and
    slide-wide duplicate shape ids make PowerPoint open the deck with the
    "needs repair" dialog. Remap every id in `el` that is already used in this
    slide to a fresh unique one, and rewrite in-fragment connector references
    (`a:stCxn`/`a:endCxn id=`) that pointed at a remapped shape. Scope is
    per-fragment on purpose: a connector's target id only means anything
    within the id space of the source the fragment came from."""
    def _cnvprs(root):
        for e in root.iter():
            if isinstance(e.tag, str) and e.tag.endswith("}cNvPr"):
                yield e

    used: set[int] = set()
    for cnv in _cnvprs(slide.shapes._spTree):
        try:
            used.add(int(cnv.get("id")))
        except (TypeError, ValueError):
            pass
    frag = [(cnv, int(cnv.get("id"))) for cnv in _cnvprs(el)
            if (cnv.get("id") or "").isdigit()]
    # Fresh ids start past everything in the slide AND the fragment, so a
    # remapped id can never collide with a fragment id we keep.
    next_id = max(used | {i for _c, i in frag}, default=0) + 1
    remap: dict[int, int] = {}
    for cnv, old in frag:
        if old in used:
            remap[old] = next_id
            cnv.set("id", str(next_id))
            used.add(next_id)
            next_id += 1
        else:
            used.add(old)
    if not remap:
        return
    for e in el.iter():
        if isinstance(e.tag, str) and (e.tag.endswith("}stCxn")
                                       or e.tag.endswith("}endCxn")):
            try:
                old = int(e.get("id"))
            except (TypeError, ValueError):
                continue
            if old in remap:
                e.set("id", str(remap[old]))


def _emit_native(slide, node: DSLNode, ctx: EmitContext) -> None:
    """native <id> b64:"<base64 element>" [media:"<base64 image>"] [parts:"<b64 json>"]
                   xml_file:"<rel path>"  [media_file:"<rel path>"] [parts_file:"<rel path>"]

    Splice a native PPTX element carried verbatim from the source deck. Small
    payloads ride inline as base64 in the DSL; large ones live as sidecar files
    under the brand pack's assets dir (`*_file:` refs, resolved like picture
    paths) so .slide.dsl files stay readable. Used for fixed corporate-design
    chrome: complex custGeom shapes (`<p:sp>`), template images (`<p:pic>`,
    `media:`/`media_file:` carries the embedded bytes), and charts
    (`<p:graphicFrame>` with `parts:`/`parts_file:` carrying the external chart
    part-graph). The element stays real + EDITABLE in the output — no
    rasterisation, no picture "cheat" — preserving geometry + colour.
    """
    import base64
    from pptx.oxml import parse_xml
    blob = node.kw_args.get("b64")
    if blob:
        xml_bytes = base64.b64decode(blob)
    else:
        ref = node.kw_args.get("xml_file")
        if not ref:
            return
        xml_bytes = _native_sidecar_bytes(ref, ctx, node, "shape")
    try:
        el = parse_xml(xml_bytes.decode("utf-8"))
    except Exception as exc:
        raise DSLError(
            f"native shape (line {node.line_no}): unparseable embedded element — {exc}"
        )
    parts_b64 = node.kw_args.get("parts")
    if not parts_b64 and node.kw_args.get("parts_file"):
        parts_b64 = base64.b64encode(
            _native_sidecar_bytes(node.kw_args["parts_file"], ctx, node, "parts")
        ).decode("ascii")
    if parts_b64:
        _splice_native_parts(slide, el, parts_b64, node)
    media_b64 = node.kw_args.get("media")
    if media_b64:
        media_bytes = base64.b64decode(media_b64)
    elif node.kw_args.get("media_file"):
        media_bytes = _native_sidecar_bytes(node.kw_args["media_file"], ctx, node,
                                            "media")
    else:
        media_bytes = None
    if media_bytes is not None:
        # Carried <p:pic>: the source rId is meaningless in this deck, so re-embed
        # the image here and re-point every <a:blip> to the fresh relationship.
        import io
        from pptx.oxml.ns import qn
        _part, rid = slide.part.get_or_add_image_part(io.BytesIO(media_bytes))
        # Drop the Microsoft svgBlip sidecar: we carry only the raster media, so
        # its stale rId would dangle / collide with another shape's relationship
        # (some renderers prefer the svgBlip and then show the WRONG image). The
        # re-embedded raster blip below is the faithful fallback.
        _SVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
        for _svg in list(el.iter("{%s}svgBlip" % _SVG_NS)):
            _ext = _svg.getparent()
            if _ext is not None and _ext.getparent() is not None:
                _ext.getparent().remove(_ext)
        for blip in el.iter(qn("a:blip")):
            if blip.get(qn("r:embed")) is not None:
                blip.set(qn("r:embed"), rid)
    _remap_colliding_shape_ids(slide, el)
    slide.shapes._spTree.append(el)


_EMITTERS = {
    "text":     _emit_text,
    "rect":     _emit_rect,
    "line":     _emit_line,
    "polyline": _emit_polyline,
    "picture":  _emit_picture,
    "shape":    _emit_shape,
    "native":   _emit_native,
}


def _slide_canvas(nodes: list[DSLNode]) -> tuple[float, float]:
    """Read the `canvas WxH` directive if present; otherwise default 1920×1080."""
    for n in nodes:
        if n.kind == "canvas" and n.pos_args:
            return parse_wh(n.pos_args[0])
    return 1920.0, 1080.0


def _write_speaker_notes(slide, notes: str) -> None:
    """Write `notes` into the slide's PPTX notes pane. python-pptx lazily
    creates the notes slide on first access. Empty/whitespace input is a no-op
    so we never materialise an empty notes slide just to hold blanks."""
    if not notes or not notes.strip():
        return
    slide.notes_slide.notes_text_frame.text = notes


def _append_slide(prs: Presentation, nodes: list[DSLNode], tokens: Tokens, *,
                  asset_root: Path | None,
                  asset_root_fallback: Path | None = None,
                  missing_assets: list[dict] | None = None,
                  image_provider: "ImageProvider | None" = None,
                  deck_dir: Path | None = None,
                  notes: str | None = None) -> None:
    """Append one slide built from `nodes` to `prs`, using the tokens for fills."""
    cw, ch = _slide_canvas(nodes)
    slide = prs.slides.add_slide(prs.slide_layouts[6])    # blank
    try:
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = _hex_to_rgb(tokens.color("paper"))
    except KeyError:
        pass

    ctx = EmitContext(tokens=tokens, canvas_w=cw, canvas_h=ch,
                      asset_root=asset_root, asset_root_fallback=asset_root_fallback,
                      image_provider=image_provider, deck_dir=deck_dir)
    for n in nodes:
        cond = n.kw_args.get("if")
        if cond is not None:
            s = cond.strip()
            if not s or s.lower() == "false" or _RESIDUAL_PLACEHOLDER.search(s):
                continue
        if n.kind in ("canvas", "theme"):
            continue
        emit = _EMITTERS.get(n.kind)
        if emit is None:
            raise DSLError(
                f"no emitter for primitive '{n.kind}' (line {n.line_no}) — "
                f"add it to _EMITTERS or remove it from the DSL"
            )
        emit(slide, n, ctx)
    if notes is not None:
        _write_speaker_notes(slide, notes)
    if missing_assets is not None and ctx.missing_assets:
        missing_assets.extend(ctx.missing_assets)


def _set_theme_fonts(prs, major: str | None, minor: str | None) -> None:
    """Align the deck theme's fontScheme major/minor latin typeface to the brand fonts.

    python-pptx's default template ships a **Calibri** theme. Decompiled `style:` runs
    carry explicit faces, but any text that INHERITS the theme font (`+mj-lt`/`+mn-lt`) —
    notably NATIVE-CARRIED table cells and non-styled runs — would otherwise render in
    Calibri, a deck-wide mismatch vs the source's brand theme. The theme part loads as a
    generic `Part` (no `_element`), so edit its `.blob`.
    """
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    from lxml import etree as _etree
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    seen: set[int] = set()
    for master in prs.slide_masters:
        try:
            tp = master.part.part_related_by(RT.THEME)
        except KeyError:
            continue
        if id(tp) in seen:
            continue
        seen.add(id(tp))
        try:
            root = _etree.fromstring(tp.blob)
        except Exception:
            continue
        changed = False
        for tag, fam in (("majorFont", major), ("minorFont", minor)):
            if not fam:
                continue
            latin = root.find(f".//{{{A}}}fontScheme/{{{A}}}{tag}/{{{A}}}latin")
            if latin is not None and latin.get("typeface") != fam:
                latin.set("typeface", fam)
                changed = True
        if changed:
            tp._blob = _etree.tostring(root, xml_declaration=True,
                                       encoding="UTF-8", standalone=True)


def build_presentation(nodes: list[DSLNode], tokens: Tokens, *,
                       asset_root: Path | None = None,
                       asset_root_fallback: Path | None = None,
                       image_provider: "ImageProvider | None" = None,
                       deck_dir: Path | None = None,
                       notes: str | None = None) -> Presentation:
    """Walk primitive nodes, build a Presentation with one filled slide.

    The returned Presentation carries `missing_assets: list[dict]` as an
    attribute — empty when every required picture resolved, populated
    otherwise. Callers can branch on `prs.missing_assets` to abort.

    `asset_root_fallback` is the plugin-level shared assets dir; it is
    walked when the per-brand `asset_root` does not contain the requested
    relative path. Brand-specific files always win over plugin defaults.

    `image_provider` and `deck_dir` are the optional pair for the picture
    ``query:`` branch (Task 7). Both default to ``None``; when unset, any
    ``query:`` keyword on a picture node raises :class:`DSLError`.

    `notes` populates the slide's PPTX speaker-notes pane. None / empty
    leaves the slide without a notes slide.
    """
    cw, ch = _slide_canvas(nodes)
    _configure_slide_scale(tokens, cw)
    # Source-faithful rendering requires the brand's primary fonts to be
    # installed; substitution by LibreOffice/PowerPoint changes glyph metrics
    # and weights, which the user has called out as breaking exact match.
    brand_fonts: dict[str, str] = {}
    for fam_key in ("display", "body"):
        try:
            family = tokens.font_family(fam_key)[0]
        except (KeyError, IndexError):
            continue
        brand_fonts[fam_key] = family
        _assert_font_available(family, tokens.brand_name)
    prs = Presentation()
    prs.slide_width  = _px(cw)
    prs.slide_height = _px(ch)
    # The default python-pptx template theme is Calibri; align its fontScheme to the
    # brand so theme-inheriting text (native-carried table cells, non-styled runs)
    # renders in the brand font instead of falling back to Calibri.
    _set_theme_fonts(prs, brand_fonts.get("display"), brand_fonts.get("body"))
    missing: list[dict] = []
    _append_slide(prs, nodes, tokens, asset_root=asset_root,
                  asset_root_fallback=asset_root_fallback,
                  missing_assets=missing,
                  image_provider=image_provider,
                  deck_dir=deck_dir,
                  notes=notes)
    for slide in prs.slides:
        sanitize_chrome(slide._element)
    prs.missing_assets = missing
    return prs


# A slide payload is either a 3-tuple `(nodes, tokens, asset_root)` or a
# 4-tuple `(nodes, tokens, asset_root, notes)`. 3-tuples remain valid so
# existing callers don't need to know about speaker notes.
def _unpack_slide_payload(
    entry: tuple,
) -> tuple[list[DSLNode], Tokens, Path | None, str | None]:
    if len(entry) == 4:
        nodes, tokens, asset_root, notes = entry
        return nodes, tokens, asset_root, notes
    if len(entry) == 3:
        nodes, tokens, asset_root = entry
        return nodes, tokens, asset_root, None
    raise ValueError(
        f"build_multi_slide: slide payload must be a 3- or 4-tuple "
        f"(nodes, tokens, asset_root[, notes]); got len={len(entry)}"
    )


def emit_pptx_from_document(
    doc: "Document",
    pack: "BrandPack",
    out_path: Path,
    *,
    image_provider: "ImageProvider | None" = None,
) -> Path:
    """Emit a PPTX file from a typed :class:`~feinschmiede.dsl.ast.Document`.

    Delegates to :func:`build_presentation` (single-slide) or
    :func:`build_multi_slide` (multi-slide) via DSLNode reconstruction.

    Parameters
    ----------
    doc:
        Typed Document AST. Each slide is emitted as one PPTX slide.
    pack:
        BrandPack used to load tokens and resolve asset paths.
    out_path:
        Destination ``.pptx`` path (created or overwritten).
    image_provider:
        Optional image provider for ``query:`` picture slots.

    Returns
    -------
    Path
        The written ``out_path``.

    Note: deck_dir is not exposed via this typed entry point. Decks needing
    asset_lock.json persistence or a build cache should use build_presentation /
    build_multi_slide directly with an explicit deck_dir argument.
    """
    from feinschmiede.dsl.tokens import load_tokens

    tokens = load_tokens(pack.root)
    asset_root = pack.root / "assets" if (pack.root / "assets").is_dir() else None

    # Convert Document slides → DSLNode lists for the existing emitter.
    slide_payloads: list[tuple] = []
    for slide in doc.slides:
        nodes = _elements_to_nodes(slide.elements)
        slide_payloads.append((nodes, tokens, asset_root, slide.notes))

    if not slide_payloads:
        raise ValueError("emit_pptx_from_document: document has no slides")

    if len(slide_payloads) == 1:
        nodes, tok, a_root, notes = slide_payloads[0]
        prs = build_presentation(
            nodes, tok,
            asset_root=a_root,
            image_provider=image_provider,
            notes=notes,
        )
    else:
        prs = build_multi_slide(
            slide_payloads,
            image_provider=image_provider,
        )

    prs.save(str(out_path))
    return out_path


def _elements_to_nodes(elements: "list[Element]") -> list[DSLNode]:
    """Reconstruct DSLNode objects from typed Element AST.

    Each Element stores the original DSL kind in ``props['_dsl_kind']``.
    Position args and keyword args are recovered from props.
    """
    nodes: list[DSLNode] = []
    for el in elements:
        props = el.props
        kind = props.get("_dsl_kind") or el.kind.value
        pos_args = list(props.get("pos_args") or [])
        kw: dict = {
            k: v for k, v in props.items()
            if k not in ("pos_args", "label", "source", "line_no",
                         "_dsl_kind", "compound_name")
        }
        label = props.get("label")
        source = props.get("source")
        line_no = int(props.get("line_no") or 0)
        children_nodes = _elements_to_nodes(el.children) if el.children else None
        nodes.append(DSLNode(
            kind=kind,
            pos_args=pos_args,
            kw_args=kw,
            label=label,
            line_no=line_no,
            source=source,
            body=children_nodes,
        ))
    return nodes


def build_multi_slide(
    slides: list[tuple],
    *,
    asset_root_fallback: Path | None = None,
    image_provider: "ImageProvider | None" = None,
    deck_dir: Path | None = None,
) -> Presentation:
    """Build a Presentation with N slides. Each slide entry is
    `(nodes, tokens, asset_root)` or `(nodes, tokens, asset_root, notes)`.
    The slide deck's canvas comes from the first slide's `canvas` directive;
    remaining slides reuse it.

    `asset_root_fallback` applies to every slide; see
    :func:`build_presentation` for semantics.

    `image_provider` and `deck_dir` are the optional pair for the picture
    ``query:`` branch (Task 7).

    The optional 4th element of each slide tuple is speaker-notes text
    written into the PPTX notes pane for that slide.
    """
    if not slides:
        raise ValueError("build_multi_slide: no slides provided")
    first_nodes, first_tokens, _, _ = _unpack_slide_payload(slides[0])
    cw, ch = _slide_canvas(first_nodes)
    _configure_slide_scale(first_tokens, cw)
    prs = Presentation()
    prs.slide_width  = _px(cw)
    prs.slide_height = _px(ch)
    # Align the deck theme fontScheme to the brand (see build_presentation) so
    # theme-inheriting / native-carried text doesn't fall back to Calibri.
    try:
        _maj = first_tokens.font_family("display")[0]
    except (KeyError, IndexError):
        _maj = None
    try:
        _min = first_tokens.font_family("body")[0]
    except (KeyError, IndexError):
        _min = None
    _set_theme_fonts(prs, _maj, _min)
    missing: list[dict] = []
    try:
        first_w_emu = first_tokens.slide("width_emu")
    except Exception:
        first_w_emu = 0
    for slide_idx, entry in enumerate(slides, start=1):
        nodes, tokens, asset_root, notes = _unpack_slide_payload(entry)
        try:
            slide_w_emu = tokens.slide("width_emu")
        except Exception:
            slide_w_emu = 0
        if slide_w_emu and first_w_emu and slide_w_emu != first_w_emu:
            print(
                f"WARN: build_multi_slide: slide {slide_idx} tokens declare "
                f"width_emu={slide_w_emu} but the deck renders at "
                f"{first_w_emu} (first slide wins; geometry will be scaled "
                f"to the deck size)", file=sys.stderr,
            )
        per_slide: list[dict] = []
        _append_slide(prs, nodes, tokens, asset_root=asset_root,
                      asset_root_fallback=asset_root_fallback,
                      missing_assets=per_slide,
                      image_provider=image_provider,
                      deck_dir=deck_dir,
                      notes=notes)
        for entry_d in per_slide:
            entry_d.setdefault("slide_index", slide_idx)
            missing.append(entry_d)
    for slide in prs.slides:
        sanitize_chrome(slide._element)
    prs.missing_assets = missing
    return prs
