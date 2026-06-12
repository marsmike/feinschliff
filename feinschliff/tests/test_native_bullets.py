"""Native PPTX bullets: `bullet:true` / `bullet:"CHAR"` kwarg on text nodes.

Adds `<a:buChar>` + hanging indent to every paragraph in the text frame.
Existing layouts that do not use `bullet:` must stay byte-identical.
"""
import io

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

from feinschliff.dsl.pptx_emit import DSLError

from feinschmiede.dsl.tokens import Tokens
from feinschliff.dsl.parser import parse_lines
from feinschliff.dsl.pptx_emit import build_presentation


def _minimal_raw():
    return {
        "color": {"ink": "#000000", "paper": "#FFFFFF"},
        "font-family": {"display": ["Open Sans"], "body": ["Open Sans"]},
        "font-size": {"body": "32px"},
        "font-weight": {"regular": 400},
        "style": {
            "body": {
                "font": "body",
                "size": "body",
                "weight": "regular",
                "color": "ink",
            },
        },
    }


def _build_tb(dsl_source: str):
    """Parse DSL, build a presentation, return the first text-frame shape."""
    nodes, _ = parse_lines(dsl_source, source="<test>")
    prs = build_presentation(nodes, Tokens.from_dict(_minimal_raw(), brand_name="t"))
    shapes_with_tf = [s for s in prs.slides[0].shapes if s.has_text_frame]
    assert shapes_with_tf, "No text-frame shapes found on slide"
    return shapes_with_tf[0]


# ---------------------------------------------------------------------------
# Tests
# ---------------------------------------------------------------------------


def test_bullet_true_adds_buchar_per_paragraph():
    """`bullet:true` gives every paragraph an <a:buChar char="•"> with hanging indent."""
    # \\n inside a quoted label → real newline (parser's _unquote converts it)
    tb = _build_tb('text 100,100 "First\\nSecond" style:body maxwidth:800 bullet:true')
    paras = tb.text_frame.paragraphs
    assert len(paras) == 2, f"Expected 2 paragraphs, got {len(paras)}"
    for p in paras:
        pPr = p._p.find(qn("a:pPr"))
        assert pPr is not None, "pPr missing on paragraph"
        bu = pPr.find(qn("a:buChar"))
        assert bu is not None, "a:buChar missing"
        assert bu.get("char") == "•", f"Expected bullet char '•', got {bu.get('char')!r}"
        mar_l = int(pPr.get("marL"))
        indent = int(pPr.get("indent"))
        assert mar_l > 0, f"marL should be positive, got {mar_l}"
        assert indent == -mar_l, (
            f"indent ({indent}) should equal -marL ({-mar_l}) for hanging indent"
        )


def test_custom_bullet_char():
    """`bullet:"–"` uses the specified character as the bullet glyph."""
    tb = _build_tb('text 100,100 "Item one\\nItem two" style:body maxwidth:800 bullet:"–"')
    paras = tb.text_frame.paragraphs
    assert len(paras) == 2
    for p in paras:
        pPr = p._p.find(qn("a:pPr"))
        assert pPr is not None
        bu = pPr.find(qn("a:buChar"))
        assert bu is not None
        assert bu.get("char") == "–", f"Expected '–', got {bu.get('char')!r}"


def test_no_bullet_kwarg_emits_no_buchar():
    """Plain text without `bullet:` → no <a:buChar> in any paragraph (byte-identical legacy)."""
    tb = _build_tb('text 100,100 "First\\nSecond" style:body maxwidth:800')
    for p in tb.text_frame.paragraphs:
        pPr = p._p.find(qn("a:pPr"))
        if pPr is not None:
            bu = pPr.find(qn("a:buChar"))
            assert bu is None, "Legacy path must not add a:buChar"


def test_bullet_font_matches_run_face():
    """`a:buFont typeface` equals the resolved typeface of the style's font."""
    tb = _build_tb('text 100,100 "Hello" style:body maxwidth:800 bullet:true')
    para = tb.text_frame.paragraphs[0]
    pPr = para._p.find(qn("a:pPr"))
    assert pPr is not None
    bu_font = pPr.find(qn("a:buFont"))
    assert bu_font is not None, "a:buFont element missing"
    # The body style uses "Open Sans" at weight 400 → regular face → "Open Sans"
    assert bu_font.get("typeface") == "Open Sans", (
        f"Expected 'Open Sans', got {bu_font.get('typeface')!r}"
    )


def test_buFont_precedes_buChar_in_pPr():
    """Schema order: <a:buFont> must immediately precede <a:buChar> inside pPr."""
    tb = _build_tb('text 100,100 "Hello" style:body maxwidth:800 bullet:true')
    para = tb.text_frame.paragraphs[0]
    pPr = para._p.find(qn("a:pPr"))
    assert pPr is not None
    children = list(pPr)
    bu_font_idx = next(
        (i for i, c in enumerate(children) if c.tag == qn("a:buFont")), None
    )
    bu_char_idx = next(
        (i for i, c in enumerate(children) if c.tag == qn("a:buChar")), None
    )
    assert bu_font_idx is not None, "a:buFont not found"
    assert bu_char_idx is not None, "a:buChar not found"
    assert bu_font_idx < bu_char_idx, (
        f"a:buFont (index {bu_font_idx}) must precede a:buChar (index {bu_char_idx})"
    )


def test_bullet_presentation_roundtrip():
    """A slide with bullet:true can be saved and re-opened by python-pptx."""
    nodes, _ = parse_lines(
        'text 100,100 "Alpha\\nBeta\\nGamma" style:body maxwidth:800 bullet:true',
        source="<test>",
    )
    prs = build_presentation(nodes, Tokens.from_dict(_minimal_raw(), brand_name="t"))
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    prs2 = Presentation(buf)
    shapes = [s for s in prs2.slides[0].shapes if s.has_text_frame]
    assert shapes, "No text shapes after roundtrip"
    paras = shapes[0].text_frame.paragraphs
    assert len(paras) == 3
    for p in paras:
        pPr = p._p.find(qn("a:pPr"))
        assert pPr is not None
        bu = pPr.find(qn("a:buChar"))
        assert bu is not None and bu.get("char") == "•"


def test_single_line_bullet():
    """`bullet:true` also works when the label has a single paragraph (no \\n)."""
    tb = _build_tb('text 100,100 "One line" style:body maxwidth:800 bullet:true')
    paras = tb.text_frame.paragraphs
    assert len(paras) == 1
    pPr = paras[0]._p.find(qn("a:pPr"))
    assert pPr is not None
    bu = pPr.find(qn("a:buChar"))
    assert bu is not None and bu.get("char") == "•"


def test_bullet_false_emits_no_buchar():
    """`bullet:false` must suppress bullets — identical output to omitting the kwarg."""
    tb = _build_tb('text 100,100 "First\\nSecond" style:body maxwidth:800 bullet:false')
    for p in tb.text_frame.paragraphs:
        pPr = p._p.find(qn("a:pPr"))
        if pPr is not None:
            bu = pPr.find(qn("a:buChar"))
            assert bu is None, (
                "bullet:false must not emit <a:buChar> "
                f"— got char={bu.get('char')!r}"
            )


def test_multichar_bullet_raises():
    """`bullet:"ab"` (two chars) must raise DSLError mentioning the value."""
    with pytest.raises(DSLError, match="ab"):
        _build_tb('text 100,100 "Item" style:body maxwidth:800 bullet:"ab"')
