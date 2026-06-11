"""DSL text-emission feature parity with v1: per-slot `autoshrink:` and
language-aware `lang:` soft hyphenation on the `text` primitive.

Both kw-args are opt-in — when absent, output is byte-identical to before.
The tests target `_emit_text` directly to keep them cheap.
"""
from __future__ import annotations

import shutil
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Pt

from feinschliff import textfit
from feinschliff.dsl.parser import DSLNode
from feinschliff.dsl.pptx_emit import EmitContext, _emit_text, _px_to_pt, _px
from feinschmiede.dsl.tokens import load_tokens
from feinschmiede.text import measure as _measure


REPO_ROOT = Path(__file__).resolve().parents[1]
BRANDS_DIR = REPO_ROOT / "brands"


@pytest.fixture()
def heuristic_metrics(monkeypatch):
    """Force the heuristic (non-real-metrics) code path for the duration of the
    test, regardless of which fonts happen to be installed on this machine.
    Clears resolution caches before yielding and again on teardown so no
    warm-cache state leaks between tests.
    """
    monkeypatch.setenv("FEINSCHMIEDE_NO_REAL_METRICS", "1")
    _measure.clear_caches()
    yield
    _measure.clear_caches()


def _fresh_slide(canvas_w: float = 1920.0, canvas_h: float = 1080.0):
    """Build a fresh single-slide Presentation and return (slide, ctx)."""
    tokens = load_tokens(BRANDS_DIR / "feinschliff", brands_dir=BRANDS_DIR)
    prs = Presentation()
    prs.slide_width = _px(canvas_w)
    prs.slide_height = _px(canvas_h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])    # blank
    ctx = EmitContext(tokens=tokens, canvas_w=canvas_w, canvas_h=canvas_h)
    return slide, ctx


def _only_textbox(slide):
    """Return the single textbox shape on the slide."""
    boxes = [s for s in slide.shapes if s.has_text_frame]
    assert len(boxes) == 1, f"expected exactly 1 textbox, got {len(boxes)}"
    return boxes[0]


# ---------------------------------------------------------------------------
# 1. No opt-in flags → unchanged output
# ---------------------------------------------------------------------------

def test_text_without_optins_is_unchanged(heuristic_metrics):
    """A plain `text` node (no autoshrink, no lang) emits the same shape as
    before: original text intact, font size = style.size_px → pt."""
    slide, ctx = _fresh_slide()
    node = DSLNode(
        kind="text",
        pos_args=["100,100"],
        kw_args={"style": "body", "maxwidth": "800", "maxheight": "40"},
        label="Donaudampfschifffahrtsgesellschaft is a long compound word.",
        line_no=1,
    )
    _emit_text(slide, node, ctx)

    box = _only_textbox(slide)
    run = box.text_frame.paragraphs[0].runs[0]
    style = ctx.tokens.resolve_style("body")
    expected_pt = _px_to_pt(style.size_px)
    assert run.font.size == Pt(expected_pt), (
        f"expected unshrunk size {expected_pt}pt, got {run.font.size}"
    )
    # No hyphens injected — raw text passes through verbatim.
    assert run.text == "Donaudampfschifffahrtsgesellschaft is a long compound word."
    assert "­" not in run.text


# ---------------------------------------------------------------------------
# 2. autoshrink:true on overflow → font shrinks below requested size
# ---------------------------------------------------------------------------

def test_autoshrink_shrinks_font_when_text_overflows(heuristic_metrics):
    """A deliberately oversized title in a thin/short box should shrink. The
    emitted run's pt size must drop below the input size_px → pt."""
    slide, ctx = _fresh_slide()
    style = ctx.tokens.resolve_style("display")    # 160px = 80pt
    requested_pt = _px_to_pt(style.size_px)
    # Cram a long paragraph into a tiny box (400x80 design-px ≈ 200x40pt).
    long_text = (
        "The quick brown fox jumps over the lazy dog, and then continues "
        "running for many more lines just to be sure it overflows."
    )
    node = DSLNode(
        kind="text",
        pos_args=["100,100"],
        kw_args={
            "style": "display",
            "maxwidth": "400",
            "maxheight": "80",
            "autoshrink": "true",
        },
        label=long_text,
        line_no=2,
    )
    _emit_text(slide, node, ctx)

    run = _only_textbox(slide).text_frame.paragraphs[0].runs[0]
    assert run.font.size is not None
    # Compare as raw EMU (Pt is just an int-EMU wrapper) — must be strictly less.
    assert run.font.size < Pt(requested_pt), (
        f"autoshrink failed: size {run.font.size} not < {Pt(requested_pt)}"
    )
    # And must not have collapsed below the 10pt floor.
    assert run.font.size >= Pt(10)


# ---------------------------------------------------------------------------
# 3. lang:de_DE inserts U+00AD soft hyphens into long German compounds
# ---------------------------------------------------------------------------

def test_lang_de_inserts_soft_hyphens(heuristic_metrics):
    """`lang:de_DE` should produce a hyphenated string containing U+00AD."""
    slide, ctx = _fresh_slide()
    node = DSLNode(
        kind="text",
        pos_args=["100,100"],
        kw_args={
            "style": "body",
            "maxwidth": "800",
            "maxheight": "40",
            "lang": "de_DE",
        },
        label="Donaudampfschifffahrtsgesellschaft",
        line_no=3,
    )
    _emit_text(slide, node, ctx)

    run = _only_textbox(slide).text_frame.paragraphs[0].runs[0]
    assert "­" in run.text, (
        f"expected U+00AD soft hyphen in hyphenated output, got {run.text!r}"
    )
    # The visible characters (sans soft-hyphens) are still the original word.
    assert run.text.replace("­", "") == "Donaudampfschifffahrtsgesellschaft"


# ---------------------------------------------------------------------------
# 4. prevent_orphan: NBSP retry to keep the last two words on one line
# ---------------------------------------------------------------------------

def test_prevent_orphan_replaces_space_with_nbsp(heuristic_metrics):
    """Given text that would wrap to a final line containing one word, the
    space between the last two words is replaced with U+00A0 (NBSP) so the
    pair wraps together."""
    text = "This is a sentence that ends with one orphan word"
    # Narrow box → "word" gets pushed alone onto the final line.
    width_emu = 5_200_000
    result = textfit.prevent_orphan(
        text,
        font="Open Sans",
        size_pt=18,
        bold=False,
        width_emu=width_emu,
    )
    assert "orphan word" in result, (
        f"expected NBSP between last two words, got {result!r}"
    )
    # No regular space between "orphan" and "word".
    assert "orphan word" not in result


def test_prevent_orphan_returns_original_when_no_orphan(heuristic_metrics):
    short = "Two words"
    width_emu = 20_000_000
    out = textfit.prevent_orphan(
        short, font="Open Sans", size_pt=18, bold=False, width_emu=width_emu
    )
    assert out == short


def test_prevent_orphan_idempotent(heuristic_metrics):
    """Re-applying produces the same string (already-NBSP'd text doesn't re-loop)."""
    text = "This is a sentence that ends with one orphan word"
    width_emu = 5_200_000
    once = textfit.prevent_orphan(text, font="Open Sans", size_pt=18, bold=False, width_emu=width_emu)
    twice = textfit.prevent_orphan(once, font="Open Sans", size_pt=18, bold=False, width_emu=width_emu)
    assert once == twice


# ---------------------------------------------------------------------------
# 5. Real font metrics (feinschmiede.text.measure) wired into textfit
# ---------------------------------------------------------------------------

def _real_face(face="DejaVu Sans"):
    if shutil.which("fc-match") is None or _measure.find_font_file(face) is None:
        pytest.skip(f"{face} not resolvable")
    return face


def test_real_metrics_beat_default_table():
    face = _real_face()
    ratio = _measure.avg_char_width_ratio(face)
    w = textfit._avg_char_width_emu(face, 18, False)
    assert abs(w - ratio * 18 * 12700) < 1.0     # measured, not table default 0.52


def test_registered_metrics_win_over_measurement():
    face = _real_face()
    # Snapshot any builtin table entry so cleanup can restore (never delete) it.
    builtin = textfit._FONT_WIDTH_RATIO.get(face)
    textfit.register_font_metrics(face, normal=0.99, bold=0.99)
    try:
        w = textfit._avg_char_width_emu(face, 18, False)
        assert abs(w - 0.99 * 18 * 12700) < 1.0
    finally:
        if builtin is None:
            textfit._FONT_WIDTH_RATIO.pop(face, None)
        else:
            textfit._FONT_WIDTH_RATIO[face] = builtin
        textfit._REGISTERED.discard(face)


def test_measure_height_real_wrap_counts_lines():
    face = _real_face()
    one_line_pt = _measure.line_width_pt("hello world", face, 18)
    h1 = textfit.measure_height_emu("hello world", font=face, size_pt=18,
                                    width_emu=int(one_line_pt * 12700) + 200)
    h2 = textfit.measure_height_emu("hello world", font=face, size_pt=18,
                                    width_emu=int(one_line_pt * 12700 * 0.6))
    assert h2 == 2 * h1


def test_real_wrap_counts_overwide_word_as_multiple_lines():
    face = _real_face()
    word = "Donaudampfschifffahrtsgesellschaftskapitaen"
    w_pt = _measure.line_width_pt(word, face, 18)
    half = int(w_pt * 12700 * 0.5)
    h_half = textfit.measure_height_emu(word, font=face, size_pt=18, width_emu=half)
    h_full = textfit.measure_height_emu(word, font=face, size_pt=18,
                                        width_emu=int(w_pt * 12700) + 200)
    assert h_half >= 2 * h_full      # mid-word breaking modeled, not 1 line


def test_has_real_metrics_false_for_unknown():
    assert textfit.has_real_metrics("No Such Font Family XYZ") is False


def test_has_real_metrics_false_when_registered():
    """Registered ratios override measurement, so predictions for a registered
    family are NOT based on real metrics — even when the font resolves."""
    face = _real_face()
    textfit.register_font_metrics(face, normal=0.99, bold=0.99)
    try:
        assert textfit.has_real_metrics(face) is False
    finally:
        textfit._FONT_WIDTH_RATIO.pop(face, None)
        textfit._REGISTERED.discard(face)
    assert textfit.has_real_metrics(face) is True


def test_kill_switch_forces_heuristics(monkeypatch):
    face = _real_face()
    monkeypatch.setenv("FEINSCHMIEDE_NO_REAL_METRICS", "1")
    _measure.clear_caches()
    w = textfit._avg_char_width_emu(face, 18, False)
    # falls back to the builtin/default table ratio
    table = textfit._FONT_WIDTH_RATIO.get(face, textfit._FONT_WIDTH_RATIO["default"])
    assert abs(w - table["normal"] * 18 * 12700) < 1.0
