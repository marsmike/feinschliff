"""Default (non-cover) picture placement must contain, not stretch."""
import pytest
from PIL import Image, UnidentifiedImageError
from pptx.enum.shapes import MSO_SHAPE_TYPE

from feinschmiede.dsl.tokens import Tokens
from feinschliff.dsl.parser import parse_lines
from feinschliff.dsl.pptx_emit import build_presentation

RAW = {
    "color": {
        "ink": "#000000",
        "paper": "#FFFFFF",
        "paper-2": "#EEEEEE",
        "fog": "#CCCCCC",
        "graphite": "#444444",
    },
    "font-family": {"display": ["Open Sans"], "body": ["Open Sans"]},
    "font-size": {"body": "32px"},
    "font-weight": {"regular": 400},
}


def _build(tmp_path, dsl_line):
    img = tmp_path / "wide.png"
    Image.new("RGB", (200, 100), "red").save(img)
    nodes, _ = parse_lines(dsl_line.format(img=img), source="<test>")
    tokens = Tokens.from_dict(RAW, brand_name="t")
    prs = build_presentation(nodes, tokens)
    return prs.slides[0].shapes


def test_contain_preserves_aspect_and_centers(tmp_path):
    shapes = _build(tmp_path, 'picture 100,100 400x400 path:"{img}"')
    pic = [s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE][0]
    # 2:1 source in 1:1 box -> width 400px, height 200px, y-centered
    assert pic.width == int(400 * 6350)
    assert pic.height == int(200 * 6350)
    assert pic.top == int((100 + 100) * 6350)   # 100 + (400-200)/2


def test_cover_still_fills_box(tmp_path):
    shapes = _build(tmp_path, 'picture 100,100 400x400 path:"{img}" cover:true')
    pic = [s for s in shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE][0]
    assert pic.width == int(400 * 6350)
    assert pic.height == int(400 * 6350)


def test_contain_tall_image_x_centers(tmp_path):
    """1:2 (tall) image in a 1:1 box → letterboxed on the sides, x-centered."""
    img = tmp_path / "tall.png"
    Image.new("RGB", (100, 200), "blue").save(img)
    nodes, _ = parse_lines(f'picture 100,100 400x400 path:"{img}"', source="<test>")
    tokens = Tokens.from_dict(RAW, brand_name="t")
    prs = build_presentation(nodes, tokens)
    pic = [s for s in prs.slides[0].shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE][0]
    # 1:2 source in 1:1 box → height 400px, width 200px, x-centered
    assert pic.width == int(200 * 6350)
    assert pic.height == int(400 * 6350)
    assert pic.left == int((100 + 100) * 6350)   # 100 + (400-200)/2
    assert pic.top == int(100 * 6350)


def test_unreadable_image_falls_back_to_box_geometry(tmp_path):
    """A file with garbage bytes that PIL can't open causes add_picture to raise.

    Behavior: _emit_picture's except-block falls back to stretched box geometry
    (cx,cy,cw,ch = x,y,w,h), then the fast-path (treatment=="none") calls
    slide.shapes.add_picture(str(p), ...) with the unreadable file. python-pptx
    itself tries to identify the image format and raises UnidentifiedImageError.
    This is the current unhandled-crash behavior — pinned here so a future fix
    (wrapping bad-file into the placeholder flow) is a deliberate choice.
    """
    img = tmp_path / "bad.png"
    img.write_bytes(b"not an image")
    nodes, _ = parse_lines(f'picture 100,100 400x400 path:"{img}"', source="<test>")
    tokens = Tokens.from_dict(RAW, brand_name="t")
    with pytest.raises(UnidentifiedImageError):
        build_presentation(nodes, tokens)
