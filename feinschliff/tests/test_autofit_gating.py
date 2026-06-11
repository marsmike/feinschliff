"""Native PPT autofit is gated on metrics quality, and fit budgets are
inset-aware.

(A) With REAL measured glyph metrics the textfit pre-shrink is authoritative —
    writing scale-less `TEXT_TO_FIT_SHAPE` autofit would let PowerPoint and
    LibreOffice re-derive different sizes. Native autofit stays ON only for
    the heuristic (ratio-table) path, as the last line of defense.

(B) PowerPoint text frames carry default insets (lIns/rIns 91440 EMU,
    tIns/bIns 45720 EMU) unless `padding:` is set. The fit predictor must
    see the same usable area the renderer sees, so the emitter subtracts
    the effective insets from the EMU budgets passed to textfit.
"""
from __future__ import annotations

import pytest
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

from feinschliff import textfit
from feinschliff.dsl.parser import parse_lines
from feinschliff.dsl.pptx_emit import _resolve_face, build_presentation
from feinschmiede.dsl.tokens import Tokens
from feinschmiede.text import measure as _measure


# Default scale for token sets without slide.width_emu: 12192000 EMU / 1920 px.
_EMU_PER_PX = 6350
# PowerPoint default text-frame insets (EMU): lIns/rIns and tIns/bIns.
_DEFAULT_INSET_W = 91440 + 91440
_DEFAULT_INSET_H = 45720 + 45720

# 40px style → 20pt max; long enough to force a shrink in the boxes below.
_LONG_TEXT = " ".join(["Measured metrics make the predictor honest"] * 3)


def _tokens() -> Tokens:
    raw = {
        "color": {"ink": "#000000", "paper": "#FFFFFF", "graphite": "#444444"},
        "font-family": {"display": ["DejaVu Sans"], "body": ["DejaVu Sans"]},
        "font-size": {"body": "40px"},
        "font-weight": {"regular": 400},
        "style": {"body": {"font": "body", "size": "body",
                           "weight": "regular", "color": "ink"}},
    }
    return Tokens.from_dict(raw, brand_name="t")


def _emit(maxwidth: int, maxheight: int, *, extra: str = ""):
    """Build a one-slide presentation with a single autoshrink text node and
    return its textbox shape."""
    line = (f'text 100,100 "{_LONG_TEXT}" style:body '
            f'maxwidth:{maxwidth} maxheight:{maxheight} autoshrink:true{extra}')
    nodes, _ = parse_lines(line, source="<test>")
    prs = build_presentation(nodes, _tokens())
    return [s for s in prs.slides[0].shapes if s.has_text_frame][0]


def _require_dejavu():
    if _measure.find_font_file("DejaVu Sans") is None:
        pytest.skip("DejaVu Sans not resolvable on this machine")


@pytest.fixture()
def heuristic_metrics(monkeypatch):
    """Force the heuristic (non-real-metrics) path; clear resolution caches
    before and after so no warm-cache state leaks between tests."""
    monkeypatch.setenv("FEINSCHMIEDE_NO_REAL_METRICS", "1")
    _measure.clear_caches()
    yield
    _measure.clear_caches()


# ---------------------------------------------------------------------------
# (A) autofit gating
# ---------------------------------------------------------------------------

def test_autofit_off_when_real_metrics():
    """Real measured metrics → the computed size is authoritative; native
    shrink-to-fit must NOT be written (it carries no fontScale and lets
    each renderer re-derive its own shrink)."""
    _require_dejavu()
    tb = _emit(460, 80)
    assert tb.text_frame.auto_size is None


def test_autofit_on_when_heuristic(heuristic_metrics):
    """Heuristic ratio-table pre-shrink → keep PPT's native shrink-to-fit
    as the last line of defense."""
    tb = _emit(460, 80)
    assert tb.text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE


# ---------------------------------------------------------------------------
# (B) inset-aware fit budgets
# ---------------------------------------------------------------------------

def _expected_sizes(maxwidth: int, maxheight: int) -> tuple[float, float]:
    """(full-budget size, inset-aware size) computed with the exact inputs
    the emitter uses — the test recomputes rather than pinning numbers, so
    it stays correct across font versions."""
    tokens = _tokens()
    style = tokens.resolve_style("body")
    face, bold = _resolve_face(style.font_family[0], style.weight)
    max_pt = style.size_px * 0.5    # default scale: 2 design-px per pt
    kwargs = dict(font=face, max_size_pt=max_pt, min_size_pt=10, bold=bold,
                  line_height=style.line_height)
    wf = int(maxwidth * _EMU_PER_PX)
    hf = int(maxheight * _EMU_PER_PX)
    full = textfit.autoshrink_size(_LONG_TEXT, width_emu=wf, height_emu=hf,
                                   **kwargs)
    inset = textfit.autoshrink_size(_LONG_TEXT,
                                    width_emu=max(1, wf - _DEFAULT_INSET_W),
                                    height_emu=max(1, hf - _DEFAULT_INSET_H),
                                    **kwargs)
    return full, inset


def test_insets_tighten_fit_budget():
    """Default insets (no `padding:`) shrink the usable area; the emitted
    size must match the inset-aware budget, not the full box."""
    _require_dejavu()
    full, inset = _expected_sizes(820, 140)
    assert full != inset, (
        "probe geometry no longer discriminates on this host's DejaVu Sans; "
        f"full={full} inset={inset} — pick a new (text, width, height)"
    )
    run = _emit(820, 140).text_frame.paragraphs[0].runs[0]
    assert run.font.size == Pt(inset), (
        f"expected inset-aware size {inset}pt, got {run.font.size.pt}pt "
        f"(full-budget size would be {full}pt)"
    )


def test_explicit_padding_overrides_default_insets():
    """`padding:0` zeroes the text-frame margins, so the budgets use the
    full box — the emitted size must match the full-width expectation."""
    _require_dejavu()
    full, inset = _expected_sizes(820, 140)
    assert full != inset
    run = _emit(820, 140, extra=" padding:0").text_frame.paragraphs[0].runs[0]
    assert run.font.size == Pt(full), (
        f"expected full-budget size {full}pt with padding:0, "
        f"got {run.font.size.pt}pt (inset-aware size would be {inset}pt)"
    )
