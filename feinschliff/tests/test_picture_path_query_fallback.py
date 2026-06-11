"""Tests for the layered ``picture path:"…" query:"…"`` fallback.

A decompiled brand layout emits picture lines like::

    picture 0,0 1920x1080 path:"decompile/cover/image.png" \
        query:"soccer stadium crowd" cover:true

At build time the emitter resolves them in layers:

  1. ``path:`` resolves to an existing file → use the file. The ``query:``
     kwarg is IGNORED and the provider is never consulted.
  2. ``path:`` misses + ``query:`` present + provider wired → resolve via
     the provider with the EXPLICIT query string (precedence over the
     synthesized use-the-path-as-query fallback from #26).
  3. ``path:`` misses + no provider → existing fallback chain (gem
     placeholder illustration); the build completes without crashing.

Stub-provider style mirrors ``tests/test_picture_query.py``
(``MagicMock(spec=ImageProvider)`` with ``.name`` + ``.search``).
"""
from __future__ import annotations

from pathlib import Path
from unittest.mock import MagicMock

import pytest
from PIL import Image

from feinschliff.io import image_provider
from feinschliff.dsl.parser import parse_lines
from feinschliff.dsl.pptx_emit import build_presentation
from feinschliff.io.image_provider import ImageHit, ImageProvider
from tests.test_emitter_restraint import _minimal_tokens


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture(autouse=True)
def _isolate_registry(monkeypatch):
    """Per-test registry isolation — fakes registered here don't leak."""
    monkeypatch.setattr(image_provider, "_REGISTRY", {})
    monkeypatch.setattr(image_provider, "_DISCOVERED", False)
    yield


@pytest.fixture
def tiny_png(tmp_path) -> Path:
    """Generate a 10×10 red PNG and return its path."""
    p = tmp_path / "fixture.png"
    Image.new("RGB", (10, 10), color=(255, 0, 0)).save(p, "PNG")
    return p


@pytest.fixture
def provider_png(tmp_path) -> Path:
    """A second PNG, returned by the fake provider (distinct from tiny_png)."""
    p = tmp_path / "provider.png"
    Image.new("RGB", (10, 10), color=(0, 255, 0)).save(p, "PNG")
    return p


def _fake_provider(hit_path: Path) -> MagicMock:
    """Stub provider matching the ImageProvider interface used in pptx_emit."""
    hit = ImageHit(
        url=f"file://{hit_path}",
        license="Test License",
        attribution="Jane Doe",
        width=10,
        height=10,
        mime="image/png",
    )
    provider = MagicMock(spec=ImageProvider)
    provider.name = "fakeprov"
    provider.search.return_value = [hit]
    return provider


def _build(dsl: str, *, deck_dir: Path, provider: ImageProvider | None = None):
    """Render one DSL slide through build_presentation."""
    nodes, _ = parse_lines(dsl)
    return build_presentation(
        nodes,
        _minimal_tokens(),
        image_provider=provider,
        deck_dir=deck_dir,
    )


# ---------------------------------------------------------------------------
# (a) path resolves → file used, provider NOT called, query ignored
# ---------------------------------------------------------------------------

def test_path_resolves_file_used_provider_not_called(tmp_path, tiny_png, provider_png):
    provider = _fake_provider(provider_png)

    dsl = (
        'canvas 1920x1080\n'
        'theme test\n'
        f'picture 100,100 200x200 path:"{tiny_png}" query:"soccer stadium crowd" cover:true'
    )
    prs = _build(dsl, deck_dir=tmp_path, provider=provider)

    # The provider must never be consulted when the path resolves.
    provider.search.assert_not_called()

    # A picture shape was emitted and nothing landed in missing_assets.
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shapes = list(prs.slides[0].shapes)
    assert len(shapes) == 1
    assert shapes[0].shape_type == MSO_SHAPE_TYPE.PICTURE
    assert prs.missing_assets == []

    # No provider pin was written — the file path was authoritative.
    assert not (tmp_path / "asset_lock.json").is_file(), (
        "asset_lock.json must not be written when path: resolves"
    )


# ---------------------------------------------------------------------------
# (b) path misses + query + provider → provider called with the query string
# ---------------------------------------------------------------------------

def test_path_misses_provider_called_with_explicit_query(tmp_path, provider_png):
    provider = _fake_provider(provider_png)
    missing = tmp_path / "decompile" / "cover" / "image.png"  # never created

    dsl = (
        'canvas 1920x1080\n'
        'theme test\n'
        f'picture 0,0 1920x1080 path:"{missing}" query:"soccer stadium crowd" cover:true'
    )
    prs = _build(dsl, deck_dir=tmp_path, provider=provider)

    # The EXPLICIT query string drives the search — not a query synthesized
    # from the missing path (the #26 behaviour for path-only nodes).
    provider.search.assert_called_once_with("soccer stadium crowd", count=1)

    # Provider hit materialised into a real picture shape.
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shapes = list(prs.slides[0].shapes)
    assert len(shapes) == 1
    assert shapes[0].shape_type == MSO_SHAPE_TYPE.PICTURE
    assert prs.missing_assets == []

    # The lock slot id is derived from the query, not the path.
    import json
    lock = json.loads((tmp_path / "asset_lock.json").read_text())
    assert "soccer_stadium_crowd" in lock["slots"]
    assert lock["slots"]["soccer_stadium_crowd"]["query"] == "soccer stadium crowd"


def test_path_only_miss_still_synthesizes_query_from_path(tmp_path, provider_png):
    """Regression guard for the pre-existing #26 behaviour: WITHOUT an
    explicit ``query:``, a missing path is itself used as the search
    query."""
    provider = _fake_provider(provider_png)

    dsl = (
        'canvas 1920x1080\n'
        'theme test\n'
        'picture 100,100 200x200 path:"regensburg aerial"'
    )
    _build(dsl, deck_dir=tmp_path, provider=provider)

    provider.search.assert_called_once_with("regensburg aerial", count=1)


# ---------------------------------------------------------------------------
# (c) path misses + no provider → existing fallback, no crash
# ---------------------------------------------------------------------------

def test_path_misses_no_provider_falls_back_without_crash(tmp_path):
    missing = tmp_path / "decompile" / "cover" / "image.png"  # never created

    dsl = (
        'canvas 1920x1080\n'
        'theme test\n'
        f'picture 0,0 1920x1080 path:"{missing}" query:"soccer stadium crowd" cover:true'
    )
    prs = _build(dsl, deck_dir=tmp_path, provider=None)  # must not raise

    # Existing fallback chain: missing-file recorded, gem placeholder
    # illustration emitted in the slot.
    missing_files = [e for e in prs.missing_assets if e.get("kind") == "missing-file"]
    assert len(missing_files) == 1
    assert missing_files[0]["path"] == str(missing)

    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shapes = list(prs.slides[0].shapes)
    assert len(shapes) == 1
    assert shapes[0].shape_type == MSO_SHAPE_TYPE.PICTURE, (
        "gem placeholder illustration expected in the unresolved slot"
    )
