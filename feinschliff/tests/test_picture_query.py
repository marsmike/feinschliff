"""Tests for the ``picture query:`` provider-resolved branch (Task 7).

Covers the 6 cases enumerated in the image-provider-framework plan:

  1. ``query:`` + ``path:`` together → ``DSLError``.
  2. ``query:`` with no provider wired → loud error.
  3. Happy path: a fake provider returns one hit; the slide gets a
     picture shape and ``asset_lock.json`` is written.
  4. Lock re-use: the second build with the same slot id does NOT call
     ``provider.search`` again.
  5. Stale lock (``file://`` URL points at a deleted file): the
     provider is re-called and the lock is overwritten.
  6. Provider returns ``[]``: ``missing_assets`` gets a ``no-hit`` entry
     and a placeholder rect is emitted in place of the picture.

The test isolates the provider registry per-test so other tests' fakes
don't leak. The PNG fixture is generated via Pillow at the start of
each test.
"""
from __future__ import annotations

import shutil
import tempfile
from pathlib import Path
from unittest.mock import MagicMock

import pytest
from PIL import Image

from lib import image_provider
from lib.dsl.parser import parse_lines
from lib.dsl.pptx_emit import DSLError, build_presentation
from lib.image_provider import ImageHit, ImageProvider
from tests.test_emitter_restraint import _minimal_tokens


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture(autouse=True)
def _isolate_registry(monkeypatch):
    """Per-test registry isolation — fakes registered here don't leak."""
    monkeypatch.setattr(image_provider, "_REGISTRY", {})
    monkeypatch.setattr(image_provider, "_DISCOVERED", False)
    yield


@pytest.fixture
def tiny_png(tmp_path) -> Path:
    """Generate a 10×10 red PNG and return its path."""
    p = tmp_path / "fixture.png"
    Image.new("RGB", (10, 10), color=(255, 0, 0)).save(p, "PNG")
    return p


def _build(dsl: str, *, deck_dir: Path, provider: ImageProvider | None = None):
    """Render one DSL slide through build_presentation with the new kwargs."""
    nodes, _ = parse_lines(dsl)
    return build_presentation(
        nodes,
        _minimal_tokens(),
        image_provider=provider,
        deck_dir=deck_dir,
    )


# ---------------------------------------------------------------------------
# Case 1 — mutex: query: and path: cannot coexist
# ---------------------------------------------------------------------------

def test_picture_query_and_path_together_raises_dsl_error(tmp_path):
    dsl = (
        'canvas 1920x1080\n'
        'theme test\n'
        'picture 100,100 200x200 query:"kitchen" path:"/tmp/x.png"'
    )
    with pytest.raises(DSLError) as exc:
        _build(dsl, deck_dir=tmp_path)
    assert "mutually exclusive" in str(exc.value).lower()


# ---------------------------------------------------------------------------
# Case 2 — query: without ctx.image_provider raises loudly
# ---------------------------------------------------------------------------

def test_picture_query_without_provider_raises_dsl_error(tmp_path):
    """The plan is explicit: failing loud is correct when the brand author
    forgot to wire ``$image_provider`` but layouts still reference
    ``query:``. We do NOT silently fall through to a placeholder rect —
    that would mask a brand-pack misconfiguration."""
    dsl = (
        'canvas 1920x1080\n'
        'theme test\n'
        'picture 100,100 200x200 query:"kitchen"'
    )
    with pytest.raises(DSLError) as exc:
        _build(dsl, deck_dir=tmp_path, provider=None)
    msg = str(exc.value).lower()
    assert "provider" in msg
    assert "query" in msg


# ---------------------------------------------------------------------------
# Case 3 — happy path: provider returns a hit, picture emitted, lock written
# ---------------------------------------------------------------------------

def test_picture_query_happy_path_writes_picture_and_lock(tmp_path, tiny_png):
    import json

    hit = ImageHit(
        url=f"file://{tiny_png}",
        license="Test License",
        attribution="Jane Doe",
        width=10,
        height=10,
        mime="image/png",
    )
    provider = MagicMock(spec=ImageProvider)
    provider.name = "fakeprov"
    provider.search.return_value = [hit]

    dsl = (
        'canvas 1920x1080\n'
        'theme test\n'
        'picture 100,100 200x200 query:"kitchen morning light"'
    )
    prs = _build(dsl, deck_dir=tmp_path, provider=provider)

    # Picture shape on the slide (not a placeholder rect).
    shapes = list(prs.slides[0].shapes)
    assert len(shapes) == 1
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    assert shapes[0].shape_type == MSO_SHAPE_TYPE.PICTURE, (
        f"expected a PICTURE shape, got {shapes[0].shape_type!r}"
    )

    # Lock file written with the hit data.
    lock_path = tmp_path / "asset_lock.json"
    assert lock_path.is_file(), "asset_lock.json should have been written"
    lock = json.loads(lock_path.read_text())
    assert lock["version"] == 1
    assert lock["provider"] == "fakeprov"
    assert "slots" in lock
    # Slot id derived from query (deterministic).
    assert len(lock["slots"]) == 1
    slot_id = next(iter(lock["slots"]))
    entry = lock["slots"][slot_id]
    assert entry["query"] == "kitchen morning light"
    assert entry["url"] == f"file://{tiny_png}"
    assert entry["license"] == "Test License"
    assert entry["attribution"] == "Jane Doe"
    assert entry["mime"] == "image/png"
    assert entry["pinned_at"].endswith("Z")  # UTC ISO with Z suffix

    # No missing_assets — we got a real picture.
    assert prs.missing_assets == []


# ---------------------------------------------------------------------------
# Case 4 — second build with same slot id reuses lock, does NOT re-search
# ---------------------------------------------------------------------------

def test_picture_query_lock_reuse_skips_provider_search(tmp_path, tiny_png):
    hit = ImageHit(
        url=f"file://{tiny_png}",
        license="Test License",
        attribution="Jane Doe",
        width=10,
        height=10,
        mime="image/png",
    )
    provider = MagicMock(spec=ImageProvider)
    provider.name = "fakeprov"
    provider.search.return_value = [hit]

    dsl = (
        'canvas 1920x1080\n'
        'theme test\n'
        'picture 100,100 200x200 query:"reuse me"'
    )

    # First build: provider.search is called once, lock is written.
    _build(dsl, deck_dir=tmp_path, provider=provider)
    assert provider.search.call_count == 1

    # Second build: lock is hit, provider.search is NOT called again.
    provider.search.reset_mock()
    _build(dsl, deck_dir=tmp_path, provider=provider)
    provider.search.assert_not_called()


# ---------------------------------------------------------------------------
# Case 5 — stale lock (file:// URL no longer exists) → re-search, overwrite
# ---------------------------------------------------------------------------

def test_picture_query_stale_lock_triggers_research(tmp_path):
    import json

    # Pre-stage a lock file pointing at a non-existent fixture.
    stale_path = tmp_path / "stale.png"  # never created
    fresh_path = tmp_path / "fresh.png"
    Image.new("RGB", (10, 10), color=(0, 255, 0)).save(fresh_path, "PNG")

    # Derive the slot id manually — re-emit must overwrite this slot.
    # Slot id derivation: lowercase, non-alnum → '_', trim, truncate.
    # "stale check" → "stale_check".
    lock = {
        "version": 1,
        "provider": "fakeprov",
        "slots": {
            "stale_check": {
                "query": "stale check",
                "url": f"file://{stale_path}",
                "license": "Old License",
                "attribution": "Old Author",
                "mime": "image/png",
                "pinned_at": "2020-01-01T00:00:00Z",
            }
        },
    }
    (tmp_path / "asset_lock.json").write_text(json.dumps(lock))

    fresh_hit = ImageHit(
        url=f"file://{fresh_path}",
        license="Fresh License",
        attribution="Fresh Author",
        width=10,
        height=10,
        mime="image/png",
    )
    provider = MagicMock(spec=ImageProvider)
    provider.name = "fakeprov"
    provider.search.return_value = [fresh_hit]

    dsl = (
        'canvas 1920x1080\n'
        'theme test\n'
        'picture 100,100 200x200 query:"stale check"'
    )
    prs = _build(dsl, deck_dir=tmp_path, provider=provider)

    # Provider must have been re-called because the pinned file is gone.
    assert provider.search.call_count == 1

    # Lock entry overwritten with the fresh URL.
    new_lock = json.loads((tmp_path / "asset_lock.json").read_text())
    entry = new_lock["slots"]["stale_check"]
    assert entry["url"] == f"file://{fresh_path}"
    assert entry["license"] == "Fresh License"
    assert entry["pinned_at"] != "2020-01-01T00:00:00Z"

    # Picture shape rendered from the fresh fixture.
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    assert list(prs.slides[0].shapes)[0].shape_type == MSO_SHAPE_TYPE.PICTURE


# ---------------------------------------------------------------------------
# Case 6 — provider returns [] → no-hit missing_assets entry + placeholder
# ---------------------------------------------------------------------------

def test_picture_query_no_hit_records_missing_and_emits_placeholder(tmp_path):
    provider = MagicMock(spec=ImageProvider)
    provider.name = "fakeprov"
    provider.search.return_value = []

    dsl = (
        'canvas 1920x1080\n'
        'theme test\n'
        'picture 100,100 200x200 query:"nothing matches"'
    )
    prs = _build(dsl, deck_dir=tmp_path, provider=provider)

    # missing_assets gets a no-hit entry.
    assert prs.missing_assets, "expected a missing_assets entry"
    no_hit = [e for e in prs.missing_assets if e.get("kind") == "no-hit"]
    assert len(no_hit) == 1
    assert no_hit[0]["query"] == "nothing matches"

    # Placeholder rect was emitted (not a picture).
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shapes = list(prs.slides[0].shapes)
    assert len(shapes) == 1
    # AUTO_SHAPE for a rect; definitely not PICTURE.
    assert shapes[0].shape_type != MSO_SHAPE_TYPE.PICTURE

    # No lock entry should have been written for a failed search — that
    # would pin the failure permanently.
    lock_path = tmp_path / "asset_lock.json"
    if lock_path.is_file():
        import json
        lock = json.loads(lock_path.read_text())
        assert "nothing_matches" not in lock.get("slots", {}), (
            "failed search must not be pinned in asset_lock.json"
        )


# ---------------------------------------------------------------------------
# Extras — additional smoke coverage
# ---------------------------------------------------------------------------

def test_slot_id_from_query_is_deterministic(tmp_path, tiny_png):
    """Same query string → same slot id across builds, regardless of provider."""
    import json

    hit = ImageHit(
        url=f"file://{tiny_png}",
        license="L", attribution="A",
        width=10, height=10, mime="image/png",
    )

    def _new_provider():
        p = MagicMock(spec=ImageProvider)
        p.name = "fakeprov"
        p.search.return_value = [hit]
        return p

    dsl = (
        'canvas 1920x1080\n'
        'theme test\n'
        'picture 100,100 200x200 query:"Kitchen morning light!"'
    )
    _build(dsl, deck_dir=tmp_path, provider=_new_provider())
    lock = json.loads((tmp_path / "asset_lock.json").read_text())
    slot_id_first = next(iter(lock["slots"]))

    # Wipe lock; build again — slot id must match.
    (tmp_path / "asset_lock.json").unlink()
    _build(dsl, deck_dir=tmp_path, provider=_new_provider())
    lock_again = json.loads((tmp_path / "asset_lock.json").read_text())
    slot_id_second = next(iter(lock_again["slots"]))

    assert slot_id_first == slot_id_second
    # Sanity: no leading/trailing underscore, lowercase, no punctuation.
    assert slot_id_first == "kitchen_morning_light"


def test_brand_switch_invalidates_lock(tmp_path, tiny_png):
    """A lock written under provider A is ignored when provider B is active.

    Different providers have different URL schemes / licensing — the
    pin must be re-acquired against the new provider.
    """
    import json

    hit_a = ImageHit(
        url=f"file://{tiny_png}", license="A", attribution="AA",
        width=10, height=10, mime="image/png",
    )
    prov_a = MagicMock(spec=ImageProvider)
    prov_a.name = "provider_a"
    prov_a.search.return_value = [hit_a]

    dsl = (
        'canvas 1920x1080\n'
        'theme test\n'
        'picture 100,100 200x200 query:"shared query"'
    )
    _build(dsl, deck_dir=tmp_path, provider=prov_a)
    assert prov_a.search.call_count == 1

    # Now build with provider_b — search must be called because the lock
    # is scoped to provider_a.
    prov_b = MagicMock(spec=ImageProvider)
    prov_b.name = "provider_b"
    prov_b.search.return_value = [hit_a]
    _build(dsl, deck_dir=tmp_path, provider=prov_b)
    assert prov_b.search.call_count == 1, (
        "provider switch should invalidate the pin and trigger re-search"
    )

    # Lock now reflects provider_b.
    lock = json.loads((tmp_path / "asset_lock.json").read_text())
    assert lock["provider"] == "provider_b"


def test_picture_query_search_exception_marked_as_search_error(tmp_path):
    """``provider.search()`` raising must be caught at the
    ``_emit_picture`` boundary so the build still completes with a
    placeholder rect (per spec §"Failure modes"), BUT the failure must
    be distinguishable from a legitimate empty-result miss:

    - A ``RuntimeWarning`` is emitted naming the provider, query, and
      exception type so the operator notices in the build log.
    - The ``missing_assets`` entry is ``kind="search-error"`` (not
      ``"no-hit"``) and carries an ``exc_type`` field.

    Without these signals a transient provider crash (e.g. expired API
    token, network blip) would look identical to "we searched and the
    provider has nothing matching this query" — masking real bugs.
    """
    provider = MagicMock(spec=ImageProvider)
    provider.name = "fakeprov"
    provider.search.side_effect = RuntimeError("upstream API exploded")

    dsl = (
        'canvas 1920x1080\n'
        'theme test\n'
        'picture 100,100 200x200 query:"explodes"'
    )
    # No exception should bubble out of the build, but a RuntimeWarning
    # must fire identifying the provider crash.
    with pytest.warns(RuntimeWarning, match="raised on search"):
        prs = _build(dsl, deck_dir=tmp_path, provider=provider)

    # missing_assets gets a search-error entry (distinct from no-hit).
    assert prs.missing_assets, "expected a missing_assets entry"
    errors = [e for e in prs.missing_assets if e.get("kind") == "search-error"]
    assert len(errors) == 1
    assert errors[0]["query"] == "explodes"
    assert errors[0]["provider"] == "fakeprov"
    assert errors[0]["exc_type"] == "RuntimeError"

    # And specifically NOT a no-hit entry — those are reserved for
    # provider returning ``[]``.
    no_hits = [e for e in prs.missing_assets if e.get("kind") == "no-hit"]
    assert no_hits == []

    # Placeholder rect emitted (not a picture).
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shapes = list(prs.slides[0].shapes)
    assert len(shapes) == 1
    assert shapes[0].shape_type != MSO_SHAPE_TYPE.PICTURE

    # Failed search must NOT be pinned in asset_lock.json — a stale
    # "no results" entry would block the slot forever.
    lock_path = tmp_path / "asset_lock.json"
    if lock_path.is_file():
        import json
        lock = json.loads(lock_path.read_text())
        assert "explodes" not in lock.get("slots", {}), (
            "exception-raising search must not be pinned in asset_lock.json"
        )


def test_picture_query_warns_when_deck_dir_unset_for_http_hit(tmp_path):
    """When ``ctx.deck_dir`` is ``None`` and the resolved hit needs HTTP
    materialise, a ``RuntimeWarning`` must fire so library callers who
    forgot to wire ``deck_dir`` notice (no rebuild cache reuse). The
    fallback itself still completes the build via a throwaway tempdir.

    We mock ``urllib.request.urlopen`` so no real network call happens.
    """
    import urllib.request as _urlreq

    # Minimal valid PNG payload — pptx Picture insertion calls Pillow to
    # introspect dimensions, so the bytes need to actually be parseable.
    from io import BytesIO
    buf = BytesIO()
    Image.new("RGB", (10, 10), color=(0, 0, 255)).save(buf, "PNG")
    png_bytes = buf.getvalue()

    class _FakeResp:
        def __init__(self, data):
            self._data = data

        def read(self):
            return self._data

        def __enter__(self):
            return self

        def __exit__(self, *_a):
            return False

    def _fake_urlopen(url, timeout=None):  # noqa: ARG001
        return _FakeResp(png_bytes)

    hit = ImageHit(
        url="https://example.invalid/test.png",
        license="L", attribution="A",
        width=10, height=10, mime="image/png",
    )
    provider = MagicMock(spec=ImageProvider)
    provider.name = "fakeprov"
    provider.search.return_value = [hit]

    dsl = (
        'canvas 1920x1080\n'
        'theme test\n'
        'picture 100,100 200x200 query:"needs http"'
    )

    # deck_dir=None forces the tempdir fallback path. We still need a
    # tmp_path for any other I/O the build might do, but pass None to
    # the EmitContext via the helper signature.
    nodes, _ = parse_lines(dsl)
    from lib.dsl.pptx_emit import build_presentation as _bp

    import unittest.mock as _um
    with _um.patch.object(_urlreq, "urlopen", _fake_urlopen):
        with pytest.warns(RuntimeWarning, match="deck_dir is unset"):
            prs = _bp(
                nodes,
                _minimal_tokens(),
                image_provider=provider,
                deck_dir=None,
            )

    # Build still produces a picture shape (tempdir fallback works).
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shapes = list(prs.slides[0].shapes)
    assert shapes and shapes[0].shape_type == MSO_SHAPE_TYPE.PICTURE


def test_picture_query_fetch_failed_carries_error_diagnostic(tmp_path):
    """When the HTTP materialise step fails on every retry, the
    ``missing_assets`` entry for ``kind="fetch-failed"`` must carry an
    ``error`` field naming the underlying exception type and message.
    Without that diagnostic, a build that silently degraded a slot to a
    placeholder rect looks identical to a slot that hit ``[]`` from the
    provider — operators have no way to triage transient HTTP failures.
    """
    import urllib.error
    import urllib.request as _urlreq

    hit = ImageHit(
        url="https://example.invalid/never-resolves.png",
        license="L", attribution="A",
        width=10, height=10, mime="image/png",
    )
    provider = MagicMock(spec=ImageProvider)
    provider.name = "fakeprov"
    provider.search.return_value = [hit]

    def _raising_urlopen(url, timeout=None):  # noqa: ARG001
        raise urllib.error.HTTPError(
            url, 404, "Not Found", hdrs=None, fp=None,  # type: ignore[arg-type]
        )

    dsl = (
        'canvas 1920x1080\n'
        'theme test\n'
        'picture 100,100 200x200 query:"http boom"'
    )

    import unittest.mock as _um
    with _um.patch.object(_urlreq, "urlopen", _raising_urlopen):
        prs = _build(dsl, deck_dir=tmp_path, provider=provider)

    fetch_failed = [e for e in prs.missing_assets if e.get("kind") == "fetch-failed"]
    assert len(fetch_failed) == 1, (
        f"expected exactly one fetch-failed entry, got {prs.missing_assets}"
    )
    entry = fetch_failed[0]
    assert entry["query"] == "http boom"
    assert entry["url"] == "https://example.invalid/never-resolves.png"
    assert entry["provider"] == "fakeprov"
    # Error context is the whole point of this test.
    assert "error" in entry, (
        f"fetch-failed entry missing 'error' diagnostic: {entry!r}"
    )
    assert "HTTPError" in entry["error"], (
        f"expected exception type in error string, got {entry['error']!r}"
    )

    # Placeholder rect emitted in place of the un-fetchable picture.
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shapes = list(prs.slides[0].shapes)
    assert len(shapes) == 1
    assert shapes[0].shape_type != MSO_SHAPE_TYPE.PICTURE


def test_picture_query_label_used_as_slot_id(tmp_path, tiny_png):
    """When the picture node carries a quoted label, that label becomes the
    slot id (overriding the query-derived slug)."""
    import json

    hit = ImageHit(
        url=f"file://{tiny_png}", license="L", attribution="A",
        width=10, height=10, mime="image/png",
    )
    provider = MagicMock(spec=ImageProvider)
    provider.name = "fakeprov"
    provider.search.return_value = [hit]

    dsl = (
        'canvas 1920x1080\n'
        'theme test\n'
        'picture 100,100 200x200 query:"any" "hero_image"'
    )
    _build(dsl, deck_dir=tmp_path, provider=provider)
    lock = json.loads((tmp_path / "asset_lock.json").read_text())
    assert "hero_image" in lock["slots"]


# ---------------------------------------------------------------------------
# Findings #2 + #3 — _materialise trusts HTTP Content-Type at download time
# ---------------------------------------------------------------------------

class _FakeHTTPResponse:
    """Minimal stand-in for an ``http.client.HTTPResponse`` good enough for
    ``_materialise``'s consumption: it reads bytes and exposes ``.headers``
    that lookups ``Content-Type`` from."""

    def __init__(self, data: bytes, content_type: str | None = None):
        self._data = data
        # ``email.message.Message`` is what urllib hands us in real life;
        # a plain dict also works because we only call ``.get()``.
        self.headers = {"Content-Type": content_type} if content_type else {}

    def read(self) -> bytes:
        return self._data

    def __enter__(self):
        return self

    def __exit__(self, *_a):
        return False


def test_materialise_trusts_content_type_over_hit_mime(tmp_path):
    """When the HTTP response declares ``Content-Type: image/webp`` but the
    provider's ``ImageHit.mime`` says ``image/jpeg``, the cached file
    must take the ``.webp`` extension. This is the core Finding #2/#3
    behaviour: the response is authoritative, the hint is a fallback."""
    import urllib.request as _urlreq
    import unittest.mock as _um

    from io import BytesIO
    buf = BytesIO()
    Image.new("RGB", (10, 10), color=(0, 255, 0)).save(buf, "WEBP")
    webp_bytes = buf.getvalue()

    hit = ImageHit(
        url="https://example.invalid/photo.jpg",
        license="L", attribution="A",
        width=10, height=10, mime="image/jpeg",  # hint says jpeg
    )

    def _fake_urlopen(url, timeout=None):  # noqa: ARG001
        # Response says webp — that wins.
        return _FakeHTTPResponse(webp_bytes, content_type="image/webp")

    from lib.dsl.pptx_emit import _materialise as _mat
    cache_dir = tmp_path / "cache"
    with _um.patch.object(_urlreq, "urlopen", _fake_urlopen):
        result, err = _mat(hit, cache_dir)

    assert err is None, f"unexpected error: {err!r}"
    assert result is not None, "materialise returned None for a happy fetch"
    assert result.suffix == ".webp", (
        f"expected .webp extension from Content-Type, got {result.suffix!r}"
    )
    assert result.is_file()


def test_materialise_strips_charset_suffix_from_content_type(tmp_path):
    """``Content-Type: image/png; charset=binary`` is a real-world variant.
    The mime token must be parsed out before the lookup."""
    import urllib.request as _urlreq
    import unittest.mock as _um

    from io import BytesIO
    buf = BytesIO()
    Image.new("RGB", (10, 10), color=(0, 0, 255)).save(buf, "PNG")
    png_bytes = buf.getvalue()

    hit = ImageHit(
        url="https://example.invalid/img",  # no extension in URL
        license="L", attribution="A",
        width=10, height=10, mime="",  # no hint either
    )

    def _fake_urlopen(url, timeout=None):  # noqa: ARG001
        return _FakeHTTPResponse(png_bytes, content_type="image/png; charset=binary")

    from lib.dsl.pptx_emit import _materialise as _mat
    cache_dir = tmp_path / "cache"
    with _um.patch.object(_urlreq, "urlopen", _fake_urlopen):
        result, err = _mat(hit, cache_dir)

    assert err is None
    assert result is not None
    assert result.suffix == ".png", (
        f"expected .png after stripping charset suffix, got {result.suffix!r}"
    )


def test_materialise_falls_back_to_hit_mime_when_content_type_missing(tmp_path):
    """No ``Content-Type`` header — the cache filename must take its
    extension from ``hit.mime`` (the pre-fix code path), so existing
    callers that never set Content-Type behave unchanged."""
    import urllib.request as _urlreq
    import unittest.mock as _um

    from io import BytesIO
    buf = BytesIO()
    Image.new("RGB", (10, 10), color=(255, 0, 0)).save(buf, "PNG")
    png_bytes = buf.getvalue()

    hit = ImageHit(
        url="https://example.invalid/img",  # no extension hint in URL
        license="L", attribution="A",
        width=10, height=10, mime="image/png",
    )

    def _fake_urlopen(url, timeout=None):  # noqa: ARG001
        return _FakeHTTPResponse(png_bytes, content_type=None)

    from lib.dsl.pptx_emit import _materialise as _mat
    cache_dir = tmp_path / "cache"
    with _um.patch.object(_urlreq, "urlopen", _fake_urlopen):
        result, err = _mat(hit, cache_dir)

    assert err is None
    assert result is not None
    assert result.suffix == ".png", (
        f"expected .png from hit.mime fallback, got {result.suffix!r}"
    )


def test_materialise_defaults_to_bin_when_nothing_resolvable(tmp_path):
    """Pin the last-resort behaviour: no Content-Type header, no known
    ``hit.mime``, no extension in the URL path → ``.bin``. This is rare
    after the Content-Type trust fix but the fallback must stay intact."""
    import urllib.request as _urlreq
    import unittest.mock as _um

    hit = ImageHit(
        url="https://example.invalid/img",  # no extension
        license="L", attribution="A",
        width=10, height=10, mime="application/octet-stream",
    )

    def _fake_urlopen(url, timeout=None):  # noqa: ARG001
        # No Content-Type header, opaque bytes.
        return _FakeHTTPResponse(b"\x00\x01\x02\x03", content_type=None)

    from lib.dsl.pptx_emit import _materialise as _mat
    cache_dir = tmp_path / "cache"
    with _um.patch.object(_urlreq, "urlopen", _fake_urlopen):
        result, err = _mat(hit, cache_dir)

    assert err is None
    assert result is not None
    assert result.suffix == ".bin", (
        f"expected .bin last-resort fallback, got {result.suffix!r}"
    )


def test_mime_to_ext_recognises_extended_image_mimes():
    """Direct lookup sanity check on the extended ``_MIME_TO_EXT`` map.
    These are the new entries added with the Content-Type trust change."""
    from lib.dsl.pptx_emit import _MIME_TO_EXT

    assert _MIME_TO_EXT["image/webp"] == ".webp"
    assert _MIME_TO_EXT["image/gif"] == ".gif"
    assert _MIME_TO_EXT["image/avif"] == ".avif"
    assert _MIME_TO_EXT["image/svg+xml"] == ".svg"
    assert _MIME_TO_EXT["image/bmp"] == ".bmp"
    assert _MIME_TO_EXT["image/tiff"] == ".tiff"
    # The pre-existing entries must still be there.
    assert _MIME_TO_EXT["image/jpeg"] == ".jpg"
    assert _MIME_TO_EXT["image/jpg"] == ".jpg"  # tolerated variant
    assert _MIME_TO_EXT["image/png"] == ".png"


# ---------------------------------------------------------------------------
# Finding #1 — atexit cleanup for throwaway tempdir caches
# ---------------------------------------------------------------------------

def test_throwaway_cache_registry_single_atexit_registration(monkeypatch):
    """The throwaway-cache cleanup must register with ``atexit`` exactly
    once, even across many fallback invocations. We don't drive the
    process-exit path (atexit is too lifecycle-specific to unit-test);
    we assert the registry shape instead — the handler is present, the
    guard fires once, and tempdirs accumulate in the registry list."""
    import urllib.request as _urlreq
    import unittest.mock as _um

    from io import BytesIO
    buf = BytesIO()
    Image.new("RGB", (10, 10), color=(128, 0, 128)).save(buf, "PNG")
    png_bytes = buf.getvalue()

    def _fake_urlopen(url, timeout=None):  # noqa: ARG001
        return _FakeHTTPResponse(png_bytes, content_type="image/png")

    # Reset the module-level registry state so this test is hermetic.
    # ``monkeypatch.setattr`` restores the originals on teardown.
    from lib.dsl import pptx_emit as _pe
    monkeypatch.setattr(_pe, "_THROWAWAY_CACHE_DIRS", [])
    monkeypatch.setattr(_pe, "_THROWAWAY_CLEANUP_REGISTERED", False)

    hit = ImageHit(
        url="https://example.invalid/throwaway.png",
        license="L", attribution="A",
        width=10, height=10, mime="image/png",
    )
    provider = MagicMock(spec=ImageProvider)
    provider.name = "fakeprov"
    provider.search.return_value = [hit]

    dsl_a = (
        'canvas 1920x1080\n'
        'theme test\n'
        'picture 100,100 200x200 query:"throwaway one"'
    )
    dsl_b = (
        'canvas 1920x1080\n'
        'theme test\n'
        'picture 100,100 200x200 query:"throwaway two"'
    )

    from lib.dsl.pptx_emit import build_presentation as _bp

    # Track atexit.register calls so we can assert single registration.
    mock_register = MagicMock()
    monkeypatch.setattr(_pe.atexit, "register", mock_register)

    with _um.patch.object(_urlreq, "urlopen", _fake_urlopen):
        # Build A — first fallback, expect atexit.register fired once.
        with pytest.warns(RuntimeWarning, match="deck_dir is unset"):
            nodes_a, _ = parse_lines(dsl_a)
            _bp(nodes_a, _minimal_tokens(),
                image_provider=provider, deck_dir=None)
        # Build B — second fallback, expect NO additional registration.
        with pytest.warns(RuntimeWarning, match="deck_dir is unset"):
            nodes_b, _ = parse_lines(dsl_b)
            _bp(nodes_b, _minimal_tokens(),
                image_provider=provider, deck_dir=None)

    # The guard ensures register fires exactly once across both builds.
    assert mock_register.call_count == 1, (
        f"expected atexit.register to fire exactly once across two "
        f"fallback invocations, got {mock_register.call_count}"
    )
    # The handler the guard registered is our cleanup function.
    registered_callable = mock_register.call_args.args[0]
    assert registered_callable is _pe._cleanup_throwaway_caches

    # Both throwaway dirs accumulated in the registry — the single atexit
    # handler walks the full list on exit.
    assert len(_pe._THROWAWAY_CACHE_DIRS) == 2, (
        f"expected two registered tempdirs, got {_pe._THROWAWAY_CACHE_DIRS!r}"
    )
    # And the guard reflects the registration happened.
    assert _pe._THROWAWAY_CLEANUP_REGISTERED is True


def test_cleanup_throwaway_caches_tolerates_missing_dirs(tmp_path):
    """The atexit handler must be safe to invoke even when a registered
    tempdir was already removed by something else (a second handler, the
    OS wiping ``/tmp``, an explicit teardown). ``shutil.rmtree`` with
    ``ignore_errors=True`` swallows the FileNotFoundError."""
    import unittest.mock as _um

    from lib.dsl import pptx_emit as _pe

    # Use real tempdirs so the rmtree call is meaningful.
    d1 = Path(tempfile.mkdtemp(prefix="feinschliff-cleanup-test-"))
    d2 = Path(tempfile.mkdtemp(prefix="feinschliff-cleanup-test-"))
    # Pre-remove d2 so the cleanup hits a "missing" entry.
    shutil.rmtree(d2)
    assert not d2.exists()

    with _um.patch.object(_pe, "_THROWAWAY_CACHE_DIRS", [d1, d2]):
        # Must not raise — ignore_errors=True swallows the missing dir.
        _pe._cleanup_throwaway_caches()

    # d1 was actually cleaned.
    assert not d1.exists(), "live dir should have been removed by cleanup"
