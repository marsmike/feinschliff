from __future__ import annotations

import tempfile
from pathlib import Path

import pytest
import yaml
from pptx import Presentation
from pptx.util import Inches


def _make_3slide_pptx(path: Path) -> None:
    prs = Presentation()
    tl = prs.slide_layouts[0]  # title slide
    bl = prs.slide_layouts[6]  # blank

    s1 = prs.slides.add_slide(tl)
    s1.shapes.title.text = "Cover Title"

    s2 = prs.slides.add_slide(bl)
    tb = s2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    tb.text_frame.text = "Content Slide"
    tb2 = s2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(3))
    tb2.text_frame.text = "Point one"
    p2 = tb2.text_frame.add_paragraph()
    p2.text = "Point two"

    s3 = prs.slides.add_slide(bl)
    tb3 = s3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    tb3.text_frame.text = "Thank You"

    prs.save(str(path))


def _find_brand_available() -> bool:
    try:
        from feinschmiede.brand_discovery import find_brand
        find_brand("feinschliff")
        return True
    except (ImportError, ValueError):
        return False


@pytest.fixture(scope="module")
def brand_available():
    if not _find_brand_available():
        pytest.skip("feinschliff brand pack not found on this checkout")


def test_cosmetic_polish_returns_report(brand_available):
    from feinschliff.polish import cosmetic_polish, CosmeticReport

    with tempfile.TemporaryDirectory() as tmpdir:
        src = Path(tmpdir) / "src.pptx"
        out = Path(tmpdir) / "out.pptx"
        _make_3slide_pptx(src)

        report = cosmetic_polish(src, "feinschliff", out)

    assert isinstance(report, CosmeticReport)
    assert report.slides_preserved == 3
    assert report.slides_dropped == 0


def test_cosmetic_plan_yaml_exists(brand_available):
    from feinschliff.polish import cosmetic_polish

    with tempfile.TemporaryDirectory() as tmpdir:
        src = Path(tmpdir) / "src.pptx"
        out = Path(tmpdir) / "out.pptx"
        _make_3slide_pptx(src)

        report = cosmetic_polish(src, "feinschliff", out)

        assert report.plan_path.is_file(), "plan YAML not written"

        plan = yaml.safe_load(report.plan_path.read_text(encoding="utf-8"))

    assert "slides" in plan
    assert len(plan["slides"]) == 3
    assert plan["brand"] == "feinschliff"


def test_cosmetic_titles_preserved(brand_available):
    from feinschliff.polish import cosmetic_polish

    with tempfile.TemporaryDirectory() as tmpdir:
        src = Path(tmpdir) / "src.pptx"
        out = Path(tmpdir) / "out.pptx"
        _make_3slide_pptx(src)

        report = cosmetic_polish(src, "feinschliff", out)

        plan = yaml.safe_load(report.plan_path.read_text(encoding="utf-8"))

    titles = [s["content_inline"]["title"] for s in plan["slides"]]
    assert titles[0] == "Cover Title"
    assert titles[1] == "Content Slide"
    assert titles[2] == "Thank You"


def test_cosmetic_cover_layout(brand_available):
    from feinschliff.polish import cosmetic_polish

    with tempfile.TemporaryDirectory() as tmpdir:
        src = Path(tmpdir) / "src.pptx"
        out = Path(tmpdir) / "out.pptx"
        _make_3slide_pptx(src)

        report = cosmetic_polish(src, "feinschliff", out)

        plan = yaml.safe_load(report.plan_path.read_text(encoding="utf-8"))

    cover_layout = plan["slides"][0]["layout"].lower()
    assert "cover" in cover_layout or "title" in cover_layout, (
        f"Expected cover/title layout for slide 1, got: {cover_layout}"
    )


def test_cosmetic_closer_layout(brand_available):
    from feinschliff.polish import cosmetic_polish

    with tempfile.TemporaryDirectory() as tmpdir:
        src = Path(tmpdir) / "src.pptx"
        out = Path(tmpdir) / "out.pptx"
        _make_3slide_pptx(src)

        report = cosmetic_polish(src, "feinschliff", out)

        plan = yaml.safe_load(report.plan_path.read_text(encoding="utf-8"))

    closer_layout = plan["slides"][-1]["layout"].lower()
    assert "closer" in closer_layout or "end" in closer_layout, (
        f"Expected closer/end layout for last slide, got: {closer_layout}"
    )
