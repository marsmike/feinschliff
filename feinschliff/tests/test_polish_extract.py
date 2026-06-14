from __future__ import annotations

import io
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches


def _make_test_pptx(path: Path) -> None:
    """Build a 3-slide PPTX for testing: title, content+bullets, content+notes."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank
    title_layout = prs.slide_layouts[0]  # title slide

    # Slide 1: title slide
    s1 = prs.slides.add_slide(title_layout)
    s1.shapes.title.text = "Hello World"

    # Slide 2: blank with a text-frame for title + bullets
    s2 = prs.slides.add_slide(blank_layout)
    txb_title = s2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    txb_title.text_frame.text = "Slide Two Title"

    txb_body = s2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(4))
    tf = txb_body.text_frame
    tf.text = "First bullet"
    p2 = tf.add_paragraph()
    p2.text = "Second bullet"
    p3 = tf.add_paragraph()
    p3.text = ""  # empty paragraph — should be skipped
    p4 = tf.add_paragraph()
    p4.text = "Third bullet"

    # Slide 3: blank with notes
    s3 = prs.slides.add_slide(blank_layout)
    txb3 = s3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    txb3.text_frame.text = "Slide Three"
    notes_frame = s3.notes_slide.notes_text_frame
    notes_frame.text = "These are the speaker notes."

    prs.save(str(path))


def test_extract_three_slides():
    from feinschliff.polish import extract_text_from_pptx

    with tempfile.TemporaryDirectory() as tmpdir:
        p = Path(tmpdir) / "test.pptx"
        _make_test_pptx(p)
        slides = extract_text_from_pptx(p)

    assert len(slides) == 3

    s1, s2, s3 = slides
    assert s1.index == 1
    assert s1.title == "Hello World"

    assert s2.index == 2
    assert s2.title != ""  # some title found
    assert len(s2.bullets) == 3  # empty paragraph skipped
    assert "First bullet" in s2.bullets
    assert "Second bullet" in s2.bullets
    assert "Third bullet" in s2.bullets
    assert s2.notes == ""

    assert s3.index == 3
    assert s3.title == "Slide Three"
    assert s3.notes == "These are the speaker notes."


def test_extract_image_detection():
    from feinschliff.polish import extract_text_from_pptx

    with tempfile.TemporaryDirectory() as tmpdir:
        p = Path(tmpdir) / "img_test.pptx"
        prs = Presentation()
        blank = prs.slide_layouts[6]

        # Slide with only a rectangle (auto-shape) — NOT an image
        s1 = prs.slides.add_slide(blank)
        s1.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE value
            Inches(1), Inches(1), Inches(2), Inches(1),
        )

        # Slide with a real picture
        s2 = prs.slides.add_slide(blank)
        # Build a minimal valid 1x1 PNG in memory
        import struct
        import zlib
        def _tiny_png() -> bytes:
            def chunk(name, data):
                c = struct.pack(">I", len(data)) + name + data
                return c + struct.pack(">I", zlib.crc32(name + data) & 0xFFFFFFFF)
            ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
            idat = zlib.compress(b"\x00\xff\xff\xff")
            return b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", ihdr) + chunk(b"IDAT", idat) + chunk(b"IEND", b"")

        png_bytes = _tiny_png()
        img_stream = io.BytesIO(png_bytes)
        s2.shapes.add_picture(img_stream, Inches(1), Inches(1), Inches(2), Inches(2))

        prs.save(str(p))
        slides = extract_text_from_pptx(p)

    assert len(slides) == 2
    assert slides[0].has_image is False
    assert slides[1].has_image is True


def test_no_chart_no_table_by_default():
    from feinschliff.polish import extract_text_from_pptx

    with tempfile.TemporaryDirectory() as tmpdir:
        p = Path(tmpdir) / "plain.pptx"
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(str(p))
        slides = extract_text_from_pptx(p)

    assert len(slides) == 1
    assert slides[0].has_chart is False
    assert slides[0].has_table is False
