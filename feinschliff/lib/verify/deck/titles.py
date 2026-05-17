"""Title extraction helpers for deck-level verify.

Two source shapes:
1. **Planned form** — a JSON file with a top-level `slides` list. Each
   slide carries either `slot_values.title` (deck_plan.json shape, post-
   layout-pick) or a raw `title` (content_plan.json shape, pre-layout-
   pick). Returns one string per slide; missing/blank → "".
2. **Built form** — a `.pptx` file. We read each slide's title placeholder
   via python-pptx. Slides without a title placeholder → "".

Both return `list[str]` of titles in slide order.
"""
from __future__ import annotations

import json
from pathlib import Path


def extract_titles_from_plan(path: Path) -> list[str]:
    """Read titles from a deck plan file.

    Supports three plan shapes, auto-detected:
    - `deck_plan.json` — `slides[].slot_values.title`
    - `content_plan.json` — `slides[].title` (bare)
    - `<deck>.yaml` (build-CLI deck plan) — `slides[].content.title` or
      `slides[].content.action_title`

    File format detected by extension (`.yaml`/`.yml` → YAML, otherwise JSON).
    """
    if not path.is_file():
        raise FileNotFoundError(f"plan not found: {path}")
    text = path.read_text()
    if path.suffix.lower() in (".yaml", ".yml"):
        import yaml
        plan = yaml.safe_load(text)
    else:
        plan = json.loads(text)
    if not isinstance(plan, dict):
        raise ValueError(
            f"plan root is not a mapping: {path} "
            f"(expected deck_plan / content_plan / deck-YAML shape)"
        )
    slides = plan.get("slides")
    if not isinstance(slides, list):
        raise ValueError(
            f"plan has no top-level `slides` list: {path} "
            f"(expected deck_plan.json, content_plan.json, or deck-YAML shape)"
        )
    out: list[str] = []
    for s in slides:
        if not isinstance(s, dict):
            out.append("")
            continue
        # deck_plan.json: slot_values.title
        title = (s.get("slot_values") or {}).get("title")
        if title is None:
            # content_plan.json: bare `title` field. Accept `title_draft`
            # too — that's the natural step-1 field name and the AI agent
            # gravitates toward it when authoring fresh briefs.
            title = s.get("title") or s.get("title_draft")
        if title is None:
            # deck-YAML (build CLI): content.action_title preferred,
            # falls back to content.title.
            content = s.get("content") or {}
            if isinstance(content, dict):
                title = content.get("action_title") or content.get("title")
        out.append(str(title).strip() if title else "")
    return out


def extract_titles_from_pptx(path: Path) -> list[str]:
    """Read titles from a built .pptx by walking each slide's title shape.

    Discovery order per slide:
    1. Feinschliff-tagged shape (name starts with "feinschliff-title-") —
       emitter sets this for any text primitive with style:title or
       style:act-title. Always preferred when present.
    2. python-pptx title placeholder (slide.shapes.title) — fallback for
       vanilla .pptx decks not built by the Feinschliff emitter.
    3. Empty string — neither found.
    """
    from pptx import Presentation

    prs = Presentation(str(path))
    out: list[str] = []
    for slide in prs.slides:
        title_text = ""
        # 1. Feinschliff-tagged title shape.
        for shape in slide.shapes:
            name = (shape.name or "")
            if name.startswith("feinschliff-title-") and shape.has_text_frame:
                title_text = shape.text_frame.text.strip()
                if title_text:
                    break
        # 2. Fallback to python-pptx title placeholder.
        if not title_text and slide.shapes.title is not None:
            if slide.shapes.title.has_text_frame:
                title_text = slide.shapes.title.text_frame.text.strip()
        out.append(title_text)
    return out
