"""Title-body coherence check — McKinsey hard rule.

Each slide's title makes a claim; the body must prove that claim AND
nothing more. The LLM judgment runs at step 4 verify (see
`skills/deck/references/iteration-loop.md` defect class title-body-coherence,
#21).

This module ships the deterministic helpers the orchestrating LLM
invokes: title + body text extraction from a built .pptx slide.
"""
from __future__ import annotations

from pathlib import Path


def extract_slide_title_and_body(
    pptx_path: Path,
    slide_index: int,
) -> tuple[str, str]:
    """Return (title, body) text from the slide at 1-based slide_index.

    title is the title placeholder text (empty string if no title shape).
    body is the concatenation of all non-title text-frame shape contents,
    joined with "\\n\\n". Slides with no body text → empty body string.
    """
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    slides = list(prs.slides)
    if not 1 <= slide_index <= len(slides):
        raise IndexError(
            f"slide_index {slide_index} out of range "
            f"(deck has {len(slides)} slide(s))"
        )
    slide = slides[slide_index - 1]

    # Identify the title shape: prefer a Feinschliff-tagged shape (emitter
    # sets name="feinschliff-title-{title|act-title}" on title primitives),
    # fall back to python-pptx's title placeholder for non-Feinschliff decks.
    title_shape = None
    for shape in slide.shapes:
        if (shape.name or "").startswith("feinschliff-title-") and shape.has_text_frame:
            title_shape = shape
            break
    if title_shape is None and slide.shapes.title is not None:
        title_shape = slide.shapes.title

    title_text = ""
    if title_shape is not None and title_shape.has_text_frame:
        title_text = title_shape.text_frame.text.strip()

    body_parts: list[str] = []
    for shape in slide.shapes:
        if shape == title_shape:
            continue
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if text:
            body_parts.append(text)

    return title_text, "\n\n".join(body_parts)
