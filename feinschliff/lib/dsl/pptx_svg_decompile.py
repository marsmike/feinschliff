"""Hybrid PPTX+SVG → Feinschliff DSL decompiler — brand-agnostic.

A higher-fidelity alternative to `lib/dsl/pptx_decompile.py`. Uses PPTX
XML as the canonical source for shape semantics and (optionally) SVG —
rendered from the slide's PDF via pdf2svg — as a secondary source for
custGeom geometry that PPTX xfrm inherits from the layout.

Sources of truth:

  PPTX (`ppt/slides/slideN.xml` + slideLayouts + slideMasters + theme1)
    - canonical for shape semantics: placeholder type/idx, group structure,
      text content + per-run style, schemeClr → token resolution, z-order,
      footer / page-number / wordmark placeholders inherited from master.

  SVG (pdf2svg of the slide's PDF page; optional)
    - canonical for final rendered geometry: bounding boxes of custGeom
      shapes, fall-back bbox when PPTX xfrm is inherited from layout.
    - NOT used for text classification or color matching (PPTX wins).

  Picture handling
    - every <p:pic> shape and every placeholder with type="pic" is
      emitted as a `picture` statement pointing at a configurable
      placeholder image (default: `assets/illustrations/placeholder.jpg`,
      the feinschliff convention).

Coordinate systems:
  - PPTX EMU: 914400 EMU/inch, slide size from presentation.xml::sldSz.
  - SVG: PDF points (1 pt = 1/72 inch). pdf2svg width/height map 1:1 to
    the slide's printable area, so EMU and pt match through inch.
  - DSL canvas: 1920×1080 px by default (override via canvas_w/h).

Brand-specific knobs (all defaulted to feinschliff baseline):
  - `theme_name` — name of the brand to emit on the `theme` directive
  - `tokens_path` — brand's tokens.json (for nearest-color matching)
  - `placeholder_rel` — DSL-relative path for picture placeholders
  - Footer-region text is emitted as plain `text` primitives; brands
    that ship a `footer(...)` compound can post-process the output to
    collapse the four footer lines into a single compound call.

Usage (programmatic, preferred):
  from lib.dsl.pptx_svg_decompile import derive
  dsl = derive(pptx_path, slide_idx=1, theme_name="acme",
               tokens_path=Path("brands/acme/tokens.json"),
               layout_name="cover-orange")

Usage (CLI smoke test):
  uv run python lib/dsl/pptx_svg_decompile.py SOURCE.pptx --slide N \\
      --theme <brand> --brand-tokens brands/<brand>/tokens.json > out.dsl
"""
from __future__ import annotations

import argparse
import json
import math
import os
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation

from lib.dsl.tokens import STYLE_BUNDLES

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "svg": "http://www.w3.org/2000/svg",
}

EMU_PER_PT = 12700
PLACEHOLDER_REL = "assets/illustrations/placeholder.jpg"


# ---------------------------------------------------------------------------
# Data
# ---------------------------------------------------------------------------


@dataclass
class Shape:
    kind: str            # rect | line | oval | pic | text | table
    x: float             # px in canvas
    y: float
    w: float
    h: float
    fill: str | None = None       # token name, e.g. "accent"
    stroke: str | None = None
    text_runs: list["TextRun"] = field(default_factory=list)
    is_picture: bool = False      # True for <p:pic> or ph type="pic"
    ph_type: str | None = None    # 'title','body','subTitle','pic','ftr','sldNum',...
    ph_idx: str | None = None
    # Border width in design-px. Captured from `<a:ln w="...">` (EMU); when
    # None the emitter uses its default hairline width.
    stroke_width: float | None = None
    # Dashed-line preset name from `<a:ln><a:prstDash val="...">` (e.g.
    # "dash", "dot", "sysDash", "lgDashDot"). None = solid stroke.
    stroke_dash: str | None = None
    # Corner radius in design-px for `roundRect` shapes. Captured from the
    # `<a:gd name="adj" fmla="val N">` adjustment value where N is N/100000
    # of the shape's shortest side. None when the source uses sharp corners.
    corner_radius: float | None = None
    # Drop shadow descriptor from `<a:effectLst><a:outerShdw>`. Tuple
    # (blur_px, dist_px, angle_deg, color_token, alpha) so the decompiler
    # can emit a compact `shadow:` kwarg and the emitter can rebuild the
    # `<a:effectLst>` XML at build time. None = no shadow.
    shadow: tuple[float, float, float, str, float] | None = None
    # Gradient fill descriptor from `<a:gradFill>`. List of (position, color)
    # stops (position 0..1, color = token or hex) plus the linear angle in
    # degrees. None = solid fill (use `Shape.fill` instead). When set, the
    # emitter writes a `gradFill` XML block onto the shape's spPr; the
    # decompiler emits a `gradient:angle=Ddeg;0=token;1=token` kwarg.
    gradient: tuple[list[tuple[float, str]], float] | None = None
    # Source bodyPr `anchor` attribute — controls text vertical position
    # within the shape bbox. "ctr" centers (PowerPoint default for many
    # text frames); "b" bottoms; "t" or absent = top. Without this the DSL
    # always emits top-anchored text, so source content that's vertically
    # centered in its frame renders shifted up by half the frame height.
    valign: str | None = None
    # Text-frame internal insets as (l, t, r, b) in design-px. Source carries
    # these on `<a:bodyPr lIns="..." tIns="..." rIns="..." bIns="...">` in EMU.
    # When absent on source, defaults to PowerPoint's published 91440 / 45720
    # EMU (= ~7.2px / ~3.6px at 13.33" canvas scale). Captured per-text so
    # decompiled DSL renders at the exact source text position; without this,
    # my emitter zeroed all four and shifted every text by ~9-19 px versus
    # source — visible as the persistent blue/red ghost offsets in the redline.
    padding: tuple[float, float, float, float] | None = None
    # When image_extract_dir is passed to derive(), pictures are extracted
    # from the source PPTX and this holds the brand-pack-relative path the
    # DSL's `default:` should resolve to. None falls back to the generic
    # placeholder (genericised brand-template behaviour).
    media_path: str | None = None
    media_rid: str | None = None  # rId of the <a:blip r:embed=...>, for extraction
    # Resolved python-pptx Part for the embedded image, captured at walk
    # time so the rId is looked up against the part that actually owns it
    # (slide vs. layout vs. master). Without this, layout-inherited
    # pictures (most corporate-template chrome) fail extraction because
    # their rId only resolves on the layout part, not the slide.
    media_part: Any = None
    # For custGeom shapes: the full pathLst converted to an SVG-`d` string
    # in *canvas pixels* (already scaled from the path's path-local space).
    # When set, emit_dsl emits the shape as an `svg { path … }` block
    # instead of the lossy bbox-rect fallback. Lets the renderer reproduce
    # rings, arrows, callouts, and other vector decoration that the PDF
    # pipeline would otherwise rasterise.
    svg_path_d: str | None = None


@dataclass
class TextRun:
    text: str
    pt: float            # font size in points
    bold: bool = False
    italic: bool = False
    color: str | None = None  # token name
    align: str | None = None  # paragraph alignment: "center" | "right" | "justify" | None


# ---------------------------------------------------------------------------
# Palette + color resolution
# ---------------------------------------------------------------------------


def load_palette(tokens_path: Path) -> dict[str, tuple[int, int, int]]:
    # Walk the brand-pack `extends:` chain when DESIGN.md declares one, so a
    # child pack that only carries local overrides still gets the parent's
    # palette for nearest-colour matching. Falls back to reading just the
    # immediate file when there's no DESIGN.md or the parent can't be
    # resolved (e.g. out-of-tree packs whose parent lives elsewhere).
    brand_root = tokens_path.parent
    data = None
    if (brand_root / "DESIGN.md").is_file():
        from lib.dsl.tokens import load_tokens
        # Try sibling-located parent first (the default), then fall back
        # to the toolkit's bundled brands/ dir. Out-of-tree packs (e.g.
        # `.debug/brands/<name>` or `~/customer-brands/<name>`) declare
        # `extends: feinschliff` but their sibling dir isn't the toolkit
        # repo, so without this fallback the parent palette never loads
        # and nearest_token() degrades to raw hex emission for every
        # shape — visible in the decompiled DSL as `fill:#ffed00` instead
        # of `fill:accent` and `fill:neutral` on every custGeom because
        # _svg_color_token() then sees an unknown brand token.
        candidate_dirs = [brand_root.parent]
        toolkit_brands = Path(__file__).resolve().parents[2] / "brands"
        if toolkit_brands.is_dir():
            candidate_dirs.append(toolkit_brands)
        env_paths = os.environ.get("FEINSCHLIFF_BRAND_PATH", "")
        for ep in env_paths.split(os.pathsep):
            if ep and Path(ep).is_dir():
                candidate_dirs.append(Path(ep))
        seen: set[Path] = set()
        for cd in candidate_dirs:
            cd = cd.resolve()
            if cd in seen:
                continue
            seen.add(cd)
            try:
                data = load_tokens(brand_root, brands_dir=cd).raw
                break
            except (FileNotFoundError, ValueError):
                continue
    if data is None:
        data = json.loads(tokens_path.read_text(encoding="utf-8"))
    palette: dict[str, tuple[int, int, int]] = {}
    colors = data.get("color") or data.get("colors") or {}

    def _hex_of(entry):
        if isinstance(entry, dict):
            return entry.get("$value") or entry.get("value")
        if isinstance(entry, str):
            return entry
        return None

    for name, entry in colors.items():
        if name.startswith("$"):
            continue
        v = _hex_of(entry)
        if not isinstance(v, str):
            continue
        v = v.strip()
        if v.startswith("{") and v.endswith("}"):
            ref = v[1:-1].split(".")[-1]
            if ref in colors:
                v = _hex_of(colors[ref]) or ""
        if v.startswith("#") and len(v) == 7:
            palette[name] = (
                int(v[1:3], 16), int(v[3:5], 16), int(v[5:7], 16),
            )
    return palette


def nearest_token(rgb: tuple[int, int, int], palette: dict[str, tuple[int, int, int]]) -> str:
    # No palette → emit the raw hex literal so the DSL preserves the source
    # color verbatim. `derive()` documents this fallback for callers that
    # don't pass a tokens.json.
    if not palette:
        return "#{:02x}{:02x}{:02x}".format(*rgb)
    best = None
    best_d = math.inf
    for name, prgb in palette.items():
        d = sum((a - b) ** 2 for a, b in zip(rgb, prgb))
        if d < best_d:
            best_d = d
            best = name
    # Source-fidelity guard. Squared-euclidean threshold ≈ 25 per channel
    # (3 * 25^2 = 1875). When the closest brand token is further than this,
    # the source colour isn't really represented in the palette — emit the
    # raw hex literal instead of approximating to a token that renders as
    # a visibly different colour (e.g. the Sartorius source's #FFED00
    # yellow shouldn't squash to the feinschliff parent's gold #C9A24A
    # accent token just because it's the closest of 30 mostly-cool tokens).
    if best is not None and best_d <= _NEAREST_TOKEN_THRESHOLD_SQ:
        return best
    return "#{:02x}{:02x}{:02x}".format(*rgb)


_NEAREST_TOKEN_THRESHOLD_SQ = 1875


def load_theme_scheme(pres: Presentation) -> dict[str, str]:
    """Map theme scheme keys (accent1..6, dk1, lt1, hlink, folHlink) to #RRGGBB.

    Falls back to empty dict if theme can't be reached.
    """
    out: dict[str, str] = {}
    # Direct XML approach: locate ppt/theme/theme1.xml inside the zip.
    try:
        import zipfile
        with zipfile.ZipFile(pres.part.package._path_to_part_for_uri) as _:
            pass
    except Exception:
        pass
    # Simpler: pres.part.package iter_parts for theme parts
    try:
        for part in pres.part.package.iter_parts():
            if part.partname.endswith("theme1.xml"):
                root = etree.fromstring(part.blob)
                scheme = root.find(".//a:clrScheme", NS)
                if scheme is None:
                    continue
                for child in scheme:
                    key = etree.QName(child).localname  # dk1, lt1, accent1, ...
                    srgb = child.find("a:srgbClr", NS)
                    sys_ = child.find("a:sysClr", NS)
                    if srgb is not None:
                        out[key] = "#" + srgb.get("val").upper()
                    elif sys_ is not None:
                        out[key] = "#" + (sys_.get("lastClr") or "000000").upper()
                # PowerPoint default clrMap aliases — these slots are
                # always present and resolve to the matching scheme entry.
                # Without them, a shape fill of `schemeClr val="bg2"` (used
                # widely in corporate templates that put the slide bg in a
                # layout rect rather than `<p:bg>`) falls through unmapped.
                for alias, real in (("bg1", "lt1"), ("bg2", "lt2"),
                                    ("tx1", "dk1"), ("tx2", "dk2")):
                    if alias not in out and real in out:
                        out[alias] = out[real]
                break
    except Exception:
        pass
    return out


# ---------------------------------------------------------------------------
# Geometry conversion
# ---------------------------------------------------------------------------


class CanvasMap:
    def __init__(self, slide_cx_emu: int, slide_cy_emu: int, canvas_w: int, canvas_h: int):
        self.cx = slide_cx_emu
        self.cy = slide_cy_emu
        self.cw = canvas_w
        self.ch = canvas_h
        self.sx = canvas_w / slide_cx_emu
        self.sy = canvas_h / slide_cy_emu
        # SVG renders at 1pt per pt; emu→pt = 1/12700. Conversion to SVG units:
        self.emu_to_pt = 1 / EMU_PER_PT

    def x(self, emu: float) -> int:
        return round(emu * self.sx)

    def y(self, emu: float) -> int:
        return round(emu * self.sy)

    def w(self, emu: float) -> int:
        return max(1, round(emu * self.sx))

    def h(self, emu: float) -> int:
        return max(1, round(emu * self.sy))

    def pt_to_px(self, pt: float) -> float:
        # 1 pt in slide-y → (cy/emu_per_inch / canvas_h) ... easier:
        # canvas_h px corresponds to slide height in pt = cy/12700
        return pt * (self.ch / (self.cy / EMU_PER_PT))


# ---------------------------------------------------------------------------
# Shape walking
# ---------------------------------------------------------------------------


def _split_runs_by_color(runs: list["TextRun"]) -> list[list["TextRun"]]:
    """Group runs into consecutive same-color blocks (paragraphs).

    Newline markers (`text="\\n"`) act as paragraph separators — they don't
    own a colour, so they're attached to the *preceding* block. Returns a
    single-element list when all content shares one colour (the caller can
    cheap-check len(blocks) and skip splitting in the common case).
    """
    blocks: list[list[TextRun]] = []
    current: list[TextRun] = []
    current_color: str | None = None
    for r in runs:
        rc = r.color if (r.text and r.text != "\n") else None
        if rc is None:
            # Newline marker or coloured-less run — append to current block.
            if current:
                current.append(r)
            continue
        if current_color is None:
            current_color = rc
            current.append(r)
        elif rc == current_color:
            current.append(r)
        else:
            blocks.append(current)
            current = [r]
            current_color = rc
    if current:
        blocks.append(current)
    return blocks


def _resolve_gradient(spPr: etree._Element, theme: dict[str, str],
                       palette: dict[str, tuple[int, int, int]]
                       ) -> tuple[list[tuple[float, str]], float] | None:
    """Extract `<a:gradFill>` as ([(pos, token)], angle_deg) — or None."""
    if spPr is None:
        return None
    grad = spPr.find("a:gradFill", NS)
    if grad is None:
        return None
    stops: list[tuple[float, str]] = []
    for gs in grad.findall("a:gsLst/a:gs", NS):
        try:
            pos = int(gs.get("pos") or 0) / 100000
        except ValueError:
            continue
        srgb = gs.find("a:srgbClr", NS)
        if srgb is not None and srgb.get("val"):
            hx = srgb.get("val")
            try:
                rgb = (int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))
                color = nearest_token(rgb, palette) if palette else f"#{hx}"
            except ValueError:
                continue
            stops.append((pos, color))
    if not stops:
        return None
    angle_deg = 0.0
    lin = grad.find("a:lin", NS)
    if lin is not None and lin.get("ang"):
        try:
            angle_deg = int(lin.get("ang")) / 60000.0
        except ValueError:
            pass
    return (stops, angle_deg)


def _resolve_fill(spPr: etree._Element, theme: dict[str, str], palette: dict[str, tuple[int, int, int]]) -> str | None:
    """Return a token name, or None if no fill.

    Handles `<a:grpFill/>` by walking up to the nearest `<p:grpSp>` ancestor
    and resolving its `grpSpPr/solidFill`. Without this, custGeom shapes that
    declare `<a:grpFill/>` (a common pattern for vector logo glyph bundles
    on slide masters) render unfilled instead of inheriting the group's
    solid colour.
    """
    if spPr is None:
        return None
    gf = spPr.find("a:grpFill", NS)
    if gf is not None:
        # Walk up to find the enclosing <p:grpSp> and resolve its fill.
        anc = spPr.getparent()
        while anc is not None:
            tag = etree.QName(anc).localname
            if tag == "grpSp":
                grpSpPr = anc.find("p:grpSpPr", NS)
                if grpSpPr is not None:
                    grp_color = _resolve_fill(grpSpPr, theme, palette)
                    if grp_color:
                        return grp_color
                anc = anc.getparent()
                continue
            anc = anc.getparent()
    sf = spPr.find("a:solidFill", NS)
    if sf is None:
        return None
    srgb = sf.find("a:srgbClr", NS)
    if srgb is not None:
        hx = srgb.get("val")
        rgb = (int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))
        return nearest_token(rgb, palette)
    scheme = sf.find("a:schemeClr", NS)
    if scheme is not None:
        key = scheme.get("val")
        hex_str = theme.get(key)
        if hex_str:
            rgb = (int(hex_str[1:3], 16), int(hex_str[3:5], 16), int(hex_str[5:7], 16))
            # Apply lumMod/lumOff tints/shades crudely.
            lumMod = scheme.find("a:lumMod", NS)
            lumOff = scheme.find("a:lumOff", NS)
            if lumMod is not None or lumOff is not None:
                mod = int(lumMod.get("val")) / 100000 if lumMod is not None else 1.0
                off = int(lumOff.get("val")) / 100000 if lumOff is not None else 0.0
                rgb = tuple(
                    max(0, min(255, int(c * mod + 255 * off))) for c in rgb
                )
            return nearest_token(rgb, palette)
    return None


def _get_xfrm(spPr: etree._Element) -> tuple[int, int, int, int] | None:
    if spPr is None:
        return None
    xfrm = spPr.find("a:xfrm", NS)
    if xfrm is None:
        return None
    off = xfrm.find("a:off", NS)
    ext = xfrm.find("a:ext", NS)
    if off is None or ext is None:
        return None
    return int(off.get("x")), int(off.get("y")), int(ext.get("cx")), int(ext.get("cy"))


def _placeholder_info(node: etree._Element) -> tuple[str | None, str | None]:
    ph = node.find(".//p:nvSpPr/p:nvPr/p:ph", NS)
    if ph is None:
        ph = node.find(".//p:nvPicPr/p:nvPr/p:ph", NS)
    if ph is None:
        ph = node.find(".//p:nvCxnSpPr/p:nvPr/p:ph", NS)
    if ph is None:
        return None, None
    return ph.get("type"), ph.get("idx")


def _layout_placeholder_xfrm(slide, ph_type: str | None, ph_idx: str | None) -> tuple[int, int, int, int] | None:
    """Walk slide layout + master to resolve an inherited placeholder bbox.

    Accepts a Slide, SlideLayout, or SlideMaster — when the caller is
    already a layout/master (because walk_slide now recurses through the
    inheritance chain), there's no further parent to walk so this just
    no-ops out.
    """
    layout = getattr(slide, "slide_layout", None)
    master = getattr(layout, "slide_master", None) if layout is not None else None
    parents = [p for p in (layout, master) if p is not None]
    for parent in parents:
        root = parent.element
        for sp in root.iter("{%s}sp" % NS["p"]):
            ph = sp.find(".//p:nvSpPr/p:nvPr/p:ph", NS)
            if ph is None:
                continue
            if (ph_type and ph.get("type") == ph_type) or (
                ph_idx and ph.get("idx") == ph_idx
            ):
                xfrm = _get_xfrm(sp.find("p:spPr", NS))
                if xfrm:
                    return xfrm
    return None


def _text_runs(node: etree._Element, theme: dict[str, str], palette: dict[str, tuple[int, int, int]]) -> list[TextRun]:
    runs: list[TextRun] = []
    txBody = node.find(".//p:txBody", NS)
    if txBody is None:
        txBody = node.find(".//a:txBody", NS)
    if txBody is None:
        return runs
    # Body-level lstStyle/lvl1pPr/defRPr sz acts as the cascade default for any
    # run that does not set its own sz. Without this we miss titles whose size
    # is defined only at the body level (common in master-driven decks).
    body_default_sz: int | None = None
    lstStyle = txBody.find("a:lstStyle", NS)
    if lstStyle is not None:
        lvl1 = lstStyle.find("a:lvl1pPr", NS)
        if lvl1 is not None:
            d = lvl1.find("a:defRPr", NS)
            if d is not None and d.get("sz"):
                body_default_sz = int(d.get("sz"))
    for para in txBody.findall("a:p", NS):
        para_runs: list[TextRun] = []
        # Pick up para-level defRPr or pPr/defRPr for sz fallback.
        pPr = para.find("a:pPr", NS)
        default_sz = body_default_sz if body_default_sz is not None else 1800
        if pPr is not None:
            d = pPr.find("a:defRPr", NS)
            if d is not None and d.get("sz"):
                default_sz = int(d.get("sz"))
        # Paragraph-level alignment (`<a:pPr algn="ctr|r|just">`). Source
        # frequently centers KPI numbers + their labels within a card; the
        # emitter's default `align:left` shifts them left of source. Stored
        # on each TextRun in this paragraph so emit_dsl can lift the
        # majority-vote into a per-text `align:` kwarg.
        para_align = None
        if pPr is not None:
            algn = pPr.get("algn")
            if algn == "ctr":
                para_align = "center"
            elif algn == "r":
                para_align = "right"
            elif algn == "just":
                para_align = "justify"
        for r in para.findall("a:r", NS):
            rPr = r.find("a:rPr", NS)
            t = r.find("a:t", NS)
            if t is None or t.text is None:
                continue
            sz = default_sz
            bold = False
            italic = False
            color = None
            text = t.text
            if rPr is not None:
                if rPr.get("sz"):
                    sz = int(rPr.get("sz"))
                bold = rPr.get("b") == "1"
                italic = rPr.get("i") == "1"
                sf = rPr.find("a:solidFill", NS)
                if sf is not None:
                    color = _resolve_fill(rPr, theme, palette) or _resolve_solid(sf, theme, palette)
                # PPTX `cap="all"` is a render-time text-transform: the run's
                # stored text stays mixed-case but draws uppercase. Bake the
                # transform into the emitted DSL since downstream layouts
                # carry the literal text, not a `text-transform` directive.
                if rPr.get("cap") == "all":
                    text = text.upper()
            para_runs.append(TextRun(text=text, pt=sz / 100, bold=bold, italic=italic, color=color, align=para_align))
        if para_runs:
            # Insert a newline marker between paragraphs so emit_dsl can preserve
            # line breaks. Without this, "Headline" + "Lorem ipsum…" paragraphs
            # collapse into "HeadlineLorem ipsum…" and ruin the body diff score.
            if runs:
                runs.append(TextRun(text="\n", pt=default_sz / 100))
            runs.extend(para_runs)
        else:
            # No <a:r> — could be just <a:fld> (page number, date). Capture as one run.
            for fld in para.findall("a:fld", NS):
                t = fld.find("a:t", NS)
                if t is not None and t.text:
                    runs.append(TextRun(text=t.text, pt=default_sz / 100))
    return runs


def _resolve_solid(sf: etree._Element, theme: dict[str, str], palette: dict[str, tuple[int, int, int]]) -> str | None:
    srgb = sf.find("a:srgbClr", NS)
    if srgb is not None:
        hx = srgb.get("val")
        return nearest_token((int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16)), palette)
    scheme = sf.find("a:schemeClr", NS)
    if scheme is not None:
        key = scheme.get("val")
        hex_str = theme.get(key)
        if hex_str:
            return nearest_token(
                (int(hex_str[1:3], 16), int(hex_str[3:5], 16), int(hex_str[5:7], 16)),
                palette,
            )
    return None


# Map brand-pack tokens (the full feinschliff vocabulary) onto the SVG
# DSL's 17-name semantic vocabulary (defined in skills/svg/references/
# dsl-reference.md, resolved through lib.diagrams.brand_bridge). Tokens
# that have no direct counterpart fall back to the nearest neutral so
# the SVG block builds — at worst we lose a shade of grey, never the
# shape itself.
_BRAND_TO_SVG_COLOR: dict[str, str] = {
    "accent":         "accent",
    "accent-hover":   "accent",
    "highlight":      "accent",
    "ink":            "ink",
    "black":          "ink",
    "white":          "paper",
    "graphite":       "neutral-strong",
    "steel":          "neutral",
    "silver":         "neutral-soft",
    "fog":            "surface-2",
    "paper":          "paper",
    "paper-2":        "surface-2",
    "off-white":      "paper",
    "off-white-2":    "surface-2",
    "rule-dark":      "neutral-strong",
    "accent-2":       "primary",
    "accent-3":       "secondary",
    "severity-low":   "success",
    "severity-medium":"warning",
    "severity-high":  "danger",
    "status-done":    "status-on",
    "status-current": "status-pending",
    # Chart-series ramp — identity-mapped because chart-series-N is BOTH
    # a brand-pack token and a valid SVG semantic name (added together
    # in the chart-decompile feature so pie/bar slice fills survive the
    # round-trip from XML → DSL → SVG block).
    "chart-series-1": "chart-series-1",
    "chart-series-2": "chart-series-2",
    "chart-series-3": "chart-series-3",
    "chart-series-4": "chart-series-4",
    "chart-series-5": "chart-series-5",
    "chart-series-6": "chart-series-6",
    "status-next":    "status-off",
}


def _svg_color_token(brand_token: str | None, *, default: str = "neutral") -> str:
    """Best-effort map of brand-pack color token → SVG DSL semantic name.

    Inline `#rrggbb` literals (produced by nearest_token when the source
    colour is too far from any palette entry) pass through unchanged so
    the SVG block renders the source colour verbatim instead of
    collapsing to the default neutral. Without this, the source-fidelity
    guard in nearest_token would buy nothing for custGeom paths — they'd
    still render as a generic grey because `_BRAND_TO_SVG_COLOR` only
    knows the named feinschliff vocabulary.
    """
    if not brand_token:
        return default
    if brand_token.startswith("#"):
        return brand_token
    return _BRAND_TO_SVG_COLOR.get(brand_token, default)


def _pptx_path_to_svg_d(
    path_el: etree._Element,
    path_w: float, path_h: float,
    target_w: float, target_h: float,
) -> str:
    """Convert one `<a:path>` element to an SVG `d` string in target coords.

    PPTX path commands map onto SVG as follows:
      <a:moveTo>     → M x,y
      <a:lnTo>       → L x,y
      <a:cubicBezTo> → C x1,y1 x2,y2 x,y   (3 pts)
      <a:quadBezTo>  → Q x1,y1 x,y         (2 pts)
      <a:arcTo>      → A rx,ry 0 large,sweep x,y  (computed from wR/hR/stAng/swAng)
      <a:close/>     → Z

    PPTX path-local coordinates are integers in the path's own
    (w, h) box. We scale linearly to (target_w, target_h) in px so the
    resulting SVG `d` lives in the same coordinate space as the
    surrounding `svg <id> X,Y WxH { … }` block (which already places
    the origin at the shape's bbox top-left).

    PPTX arcs are documented angle-based (start angle, sweep angle, both
    in 60000ths of a degree, measured from the +x axis going clockwise).
    We compute the arc endpoint manually and emit SVG's endpoint-form
    arc command. PPTX `sweep > 0` is clockwise; SVG sweep flag `1` is
    clockwise — they line up.
    """
    sx = target_w / path_w if path_w else 1.0
    sy = target_h / path_h if path_h else 1.0
    out: list[str] = []
    cx_cur = cy_cur = 0.0

    def _xy(pt: etree._Element) -> tuple[float, float]:
        return float(pt.get("x")) * sx, float(pt.get("y")) * sy

    for cmd in path_el:
        tag = etree.QName(cmd).localname
        if tag == "moveTo":
            pt = cmd.find("a:pt", NS)
            x, y = _xy(pt)
            out.append(f"M {x:.2f},{y:.2f}")
            cx_cur, cy_cur = x, y
        elif tag == "lnTo":
            pt = cmd.find("a:pt", NS)
            x, y = _xy(pt)
            out.append(f"L {x:.2f},{y:.2f}")
            cx_cur, cy_cur = x, y
        elif tag == "cubicBezTo":
            pts = cmd.findall("a:pt", NS)
            if len(pts) == 3:
                (x1, y1), (x2, y2), (x, y) = _xy(pts[0]), _xy(pts[1]), _xy(pts[2])
                out.append(f"C {x1:.2f},{y1:.2f} {x2:.2f},{y2:.2f} {x:.2f},{y:.2f}")
                cx_cur, cy_cur = x, y
        elif tag == "quadBezTo":
            pts = cmd.findall("a:pt", NS)
            if len(pts) == 2:
                (x1, y1), (x, y) = _xy(pts[0]), _xy(pts[1])
                out.append(f"Q {x1:.2f},{y1:.2f} {x:.2f},{y:.2f}")
                cx_cur, cy_cur = x, y
        elif tag == "arcTo":
            wR = float(cmd.get("wR")) * sx
            hR = float(cmd.get("hR")) * sy
            stAng = float(cmd.get("stAng")) / 60000.0     # degrees
            swAng = float(cmd.get("swAng")) / 60000.0
            # PPTX arc convention: the arc is on an ellipse whose CENTER
            # is at (cx_cur - wR*cos(stAng), cy_cur - hR*sin(stAng))
            # — i.e. the current point IS the arc's start; the angle
            # tells us where the start lives on the ellipse. Compute the
            # endpoint by adding the sweep.
            import math
            st_rad = math.radians(stAng)
            end_rad = math.radians(stAng + swAng)
            centre_x = cx_cur - wR * math.cos(st_rad)
            centre_y = cy_cur - hR * math.sin(st_rad)
            end_x = centre_x + wR * math.cos(end_rad)
            end_y = centre_y + hR * math.sin(end_rad)
            large_arc = 1 if abs(swAng) > 180 else 0
            sweep_flag = 1 if swAng > 0 else 0
            out.append(
                f"A {wR:.2f},{hR:.2f} 0 {large_arc} {sweep_flag} {end_x:.2f},{end_y:.2f}"
            )
            cx_cur, cy_cur = end_x, end_y
        elif tag == "close":
            out.append("Z")
    return " ".join(out)


def _custgeom_svg_d(spPr: etree._Element, target_w: float, target_h: float) -> str | None:
    """Walk a custGeom's pathLst and concatenate the SVG `d` strings.

    Multiple `<a:path>` siblings inside `<a:pathLst>` become subpaths in
    a single `d` — each starts with its own `M`. Returns None when the
    spPr is not a custGeom or there's no usable path.
    """
    if spPr is None:
        return None
    cg = spPr.find("a:custGeom", NS)
    if cg is None:
        return None
    pathLst = cg.find("a:pathLst", NS)
    if pathLst is None:
        return None
    parts: list[str] = []
    for path_el in pathLst.findall("a:path", NS):
        pw = float(path_el.get("w") or 0)
        ph = float(path_el.get("h") or 0)
        if pw <= 0 or ph <= 0:
            continue
        d = _pptx_path_to_svg_d(path_el, pw, ph, target_w, target_h)
        if d:
            parts.append(d)
    return " ".join(parts) if parts else None


def _shape_geometry_kind(spPr: etree._Element) -> str:
    """Classify a sp by its geometry: rect, oval, line, or 'shape' (custGeom)."""
    if spPr is None:
        return "rect"
    pg = spPr.find("a:prstGeom", NS)
    if pg is not None:
        preset = pg.get("prst")
        if preset in ("ellipse",):
            return "oval"
        if preset in ("line", "straightConnector1"):
            return "line"
        if preset in ("rect", "roundRect"):
            return "rect"
        return "rect"
    if spPr.find("a:custGeom", NS) is not None:
        return "shape"
    return "rect"


# ---------------------------------------------------------------------------
# Tree walk
# ---------------------------------------------------------------------------


def walk_slide(slide, cmap: CanvasMap, theme: dict[str, str], palette: dict[str, tuple[int, int, int]]) -> list[Shape]:
    """Walk the slide's spTree; also collect shapes inherited from the
    slide layout + master that aren't already represented on the slide.

    A typical PowerPoint slide carries only its unique content (a title
    placeholder, the body text). All decorative chrome — corporate logo,
    page-number bar, branded background blocks — lives on the slide's
    layout and master. Without this inheritance walk, the decompile of
    a corporate-template deck emits ~10% of what the source renders.

    Layout/master shapes that are placeholder fills already provided by
    the slide itself (same ph_idx) are skipped — slide content wins.
    Everything else is added at the front of the shape list so it draws
    behind slide-level content.
    """
    shapes: list[Shape] = []
    spTree = slide.element.find(".//p:cSld/p:spTree", NS)
    _walk(spTree, (0, 0), shapes, slide, cmap, theme, palette)
    # A slide-level placeholder only "owns" its idx when it actually
    # carries content (text, fill, or geometry) — empty placeholders
    # (common pattern: <p:ph idx="N"/> + empty <p:spPr/>) inherit
    # EVERYTHING from the layout, including the layout's fill/size.
    # Filtering by content here lets the layout walk re-add the rich
    # version of those placeholders below.
    def _has_content(sh: Shape) -> bool:
        if sh.fill or sh.stroke or sh.gradient or sh.svg_path_d:
            return True
        if sh.is_picture:
            return True
        if any(r.text and r.text.strip() and r.text != "\n" for r in sh.text_runs):
            return True
        return False

    slide_ph_idxs = {s.ph_idx for s in shapes if s.ph_idx and _has_content(s)}
    # Filter out the empty placeholders themselves — the layout version
    # will provide the actual content.
    shapes = [s for s in shapes if not (s.ph_idx and not _has_content(s))]
    inherited: list[Shape] = []
    layout_master_chain = _layout_master_chain(slide)
    for src in layout_master_chain:
        chain_spTree = src.element.find(".//p:cSld/p:spTree", NS)
        if chain_spTree is None:
            continue
        chain_shapes: list[Shape] = []
        _walk(chain_spTree, (0, 0), chain_shapes, src, cmap, theme, palette)
        for s in chain_shapes:
            # Skip placeholder shapes the slide already owns with content.
            if s.ph_idx and s.ph_idx in slide_ph_idxs:
                continue
            # Skip pure page-number / footer placeholders by type when the
            # slide hasn't overridden them — these are pp:fld things that
            # the source-deck renderer fills at slide time; decompiling
            # them from the master emits literal "<#>" tokens that pollute
            # the output. (sldNum/ftr/dt placeholders all carry ph_type.)
            if s.ph_type in ("sldNum", "ftr", "dt") and not s.text_runs:
                continue
            # Inherited PICTURE placeholders with no real media binary
            # (no media_rid) are shells the slide should fill but didn't.
            # PowerPoint renders them as empty; the build's missing-asset
            # fallback paints them as white rectangles which inflate the
            # diff visibly. Skip them so the rendered output matches
            # PowerPoint's "absent" behaviour.
            if s.is_picture and not s.media_rid:
                continue
            # Skip layout placeholders whose only text is template
            # prompt copy ("Hier Zitat einfügen.", "Überschrift 1, TT
            # Norms Pro, 28 pt", etc) — these are author hints that
            # PowerPoint suppresses at render time when the slide's
            # version is empty. We can't easily distinguish prompt from
            # real content, but layouts mark prompts with
            # `<p:nvPr hasCustomPrompt="1">` — drop just the text on
            # those, keep the fill/geometry so the layout still emits
            # its visual frame.
            # Inherited PLACEHOLDER text is template prompt copy that
            # PowerPoint never renders ("Überschrift 1, TT Norms Pro",
            # "Diagramm durch Klicken …", "Click to edit Master title").
            # The slide's placeholder either overrides it (filtered above
            # via slide_ph_idxs) or is empty — in which case PowerPoint
            # shows nothing. Drop the text but keep fill / geometry so
            # the layout's visual frame (yellow rect, black square)
            # still emits. Non-placeholder layout/master shapes (logo
            # glyphs, decorative rects) pass through untouched.
            if s.ph_type and s.text_runs:
                s.text_runs = []
                if not _has_content(s):
                    continue
            inherited.append(s)
            if s.ph_idx:
                slide_ph_idxs.add(s.ph_idx)
    # Inherited chrome draws behind slide content.
    return inherited + shapes


def _is_custom_prompt(src, ph_idx: str | None) -> bool:
    """True when the layout's placeholder for `ph_idx` carries the
    `hasCustomPrompt="1"` marker. Those text bodies hold the template
    instruction shown in PowerPoint while the slide placeholder is empty;
    they're suppressed at render time and must not be emitted into the
    DSL or they leak as visible template copy in derived slides.
    """
    if ph_idx is None:
        return False
    for sp in src.element.iter("{%s}sp" % NS["p"]):
        ph = sp.find(".//p:nvSpPr/p:nvPr/p:ph", NS)
        if ph is None or ph.get("idx") != ph_idx:
            continue
        nvPr = sp.find(".//p:nvSpPr/p:nvPr", NS)
        if nvPr is not None and ph.get("hasCustomPrompt") == "1":
            return True
    return False


def _layout_master_chain(slide) -> list:
    """Return [layout, master] for a slide (best-effort, never raises)."""
    chain = []
    try:
        layout = slide.slide_layout
        if layout is not None:
            chain.append(layout)
            master = layout.slide_master
            if master is not None:
                chain.append(master)
    except Exception:
        pass
    return chain


def extract_slide_bg_fill(slide, theme: dict[str, str],
                          palette: dict[str, tuple[int, int, int]]) -> str | None:
    """Return the slide's background solid-fill colour as a token / hex.

    Walks the inheritance chain: slide → layout → master. Each level can
    carry an explicit `<p:cSld><p:bg>` (solidFill or bgRef→theme); the
    first one found wins. Without the layout/master fallback, decks whose
    slides have empty bg (the common case for branded corporate templates)
    would render with the brand default paper colour even when the master
    declares a solid yellow / navy background.
    """
    for src in [slide, *_layout_master_chain(slide)]:
        bg = src.element.find(".//p:cSld/p:bg", NS)
        if bg is None:
            continue
        bgPr = bg.find("p:bgPr", NS)
        if bgPr is not None:
            color = _resolve_fill(bgPr, theme, palette)
            if color:
                return color
        # <p:bgRef idx="N"><a:schemeClr val="bg1"/></p:bgRef> — referenced
        # background-fill style from theme1.xml's bgFillStyleLst. The
        # schemeClr fills the `phClr` placeholder in the referenced style.
        bgRef = bg.find("p:bgRef", NS)
        if bgRef is not None:
            color = _resolve_bg_ref(bgRef, theme, palette)
            if color:
                return color
    return None


def _resolve_bg_ref(bgRef, theme: dict[str, str],
                    palette: dict[str, tuple[int, int, int]]) -> str | None:
    """Resolve a `<p:bgRef idx="N">` reference against the theme's
    `bgFillStyleLst`, filling `<a:schemeClr val="phClr"/>` with the
    schemeClr that the bgRef carries.

    PowerPoint idx encoding: 1001 = bgFillStyleLst[0], 1002 = [1], etc.
    Only `<a:solidFill>` entries are inlined; gradient/blip refs fall
    through (no DSL primitive for those yet at the bg level).
    """
    scheme = bgRef.find("a:schemeClr", NS)
    if scheme is None or not scheme.get("val"):
        return None
    scheme_key = scheme.get("val")
    # Resolve scheme → hex via the theme dict captured by load_theme_scheme.
    # Some keys are aliases: `bg1`→`lt1`, `bg2`→`lt2`, `tx1`→`dk1`, `tx2`→`dk2`.
    alias = {"bg1": "lt1", "bg2": "lt2", "tx1": "dk1", "tx2": "dk2"}
    hex_val = theme.get(scheme_key) or theme.get(alias.get(scheme_key, scheme_key))
    if not hex_val:
        return None
    try:
        rgb = (int(hex_val[1:3], 16), int(hex_val[3:5], 16), int(hex_val[5:7], 16))
    except (ValueError, IndexError):
        return None
    return nearest_token(rgb, palette) if palette else f"#{hex_val[1:].upper()}"


def _walk(node, offset, shapes, slide, cmap, theme, palette):
    # offset is either a 2-tuple (ox, oy) or an 8-tuple carrying a scaled-
    # group affine — the latter is unpacked by `_shape_bbox`. We only need
    # ox/oy here to forward to nested non-scaled groups.
    ox, oy = offset[:2]
    for ch in node:
        tag = etree.QName(ch).localname
        if tag == "sp":
            _emit_sp(ch, offset, shapes, slide, cmap, theme, palette)
        elif tag == "pic":
            _emit_pic(ch, offset, shapes, slide, cmap, theme, palette)
        elif tag == "cxnSp":
            _emit_cxn(ch, offset, shapes, cmap, theme, palette)
        elif tag == "graphicFrame":
            _emit_graphic_frame(ch, offset, shapes, slide, cmap, theme, palette)
        elif tag == "grpSp":
            # Walk children with the group's offset added. Scaled groups
            # (ext != chExt — typical for master-level logo bundles
            # dropped into smaller slots) are skipped because the walker
            # only carries a pure translation offset; emitting their
            # children at unscaled coords would land them off-canvas.
            grp_xfrm = ch.find("p:grpSpPr/a:xfrm", NS)
            child_off = (ox, oy)
            if grp_xfrm is not None:
                off = grp_xfrm.find("a:off", NS)
                ext = grp_xfrm.find("a:ext", NS)
                chOff = grp_xfrm.find("a:chOff", NS)
                chExt = grp_xfrm.find("a:chExt", NS)
                if (off is not None and ext is not None
                        and chOff is not None and chExt is not None):
                    try:
                        ox_emu = int(off.get("x"))
                        oy_emu = int(off.get("y"))
                        cx = int(ext.get("cx"))
                        cy = int(ext.get("cy"))
                        chcx = int(chExt.get("cx"))
                        chcy = int(chExt.get("cy"))
                        chox = int(chOff.get("x"))
                        choy = int(chOff.get("y"))
                        if chcx > 0 and chcy > 0:
                            sx = cx / chcx
                            sy = cy / chcy
                            # Scaled group: thread a 6-tuple offset that
                            # _shape_bbox unpacks and applies as an EMU-level
                            # affine. Translation-only groups stay as 2-tuple
                            # for backward compat.
                            if abs(sx - 1.0) > 0.001 or abs(sy - 1.0) > 0.001:
                                child_off = (ox, oy, ox_emu, oy_emu, chox, choy, sx, sy)
                                _walk(ch, child_off, shapes, slide, cmap, theme, palette)
                                continue
                        # Pure translation fallthrough.
                        child_off = (ox + ox_emu - chox, oy + oy_emu - choy)
                    except (ValueError, TypeError):
                        pass
                elif off is not None and chOff is not None:
                    try:
                        dx = int(off.get("x")) - int(chOff.get("x"))
                        dy = int(off.get("y")) - int(chOff.get("y"))
                        child_off = (ox + dx, oy + dy)
                    except (ValueError, TypeError):
                        pass
            _walk(ch, child_off, shapes, slide, cmap, theme, palette)


def _shape_bbox(ch, offset, slide):
    spPr = ch.find("p:spPr", NS)
    xfrm = _get_xfrm(spPr)
    if xfrm is None:
        ph_type, ph_idx = _placeholder_info(ch)
        xfrm = _layout_placeholder_xfrm(slide, ph_type, ph_idx)
    if xfrm is None:
        return None
    x, y, w, h = xfrm
    # Offset shapes:
    #   2-tuple (ox, oy)      — translation-only ancestor group(s)
    #   8-tuple (ox, oy, ax, ay, chox, choy, sx, sy)
    #                         — current shape lives inside a scaled group;
    #                           apply the EMU-level affine before adding
    #                           any outer translation.
    if len(offset) == 2:
        ox, oy = offset
        return x + ox, y + oy, w, h
    ox, oy, ax, ay, chox, choy, sx, sy = offset
    x_emu = ax + (x - chox) * sx
    y_emu = ay + (y - choy) * sy
    w_emu = w * sx
    h_emu = h * sy
    return x_emu + ox, y_emu + oy, w_emu, h_emu


def _emit_sp(ch, offset, shapes, slide, cmap, theme, palette):
    spPr = ch.find("p:spPr", NS)
    bbox = _shape_bbox(ch, offset, slide)
    if bbox is None:
        return
    x, y, w, h = bbox
    ph_type, ph_idx = _placeholder_info(ch)
    runs = _text_runs(ch, theme, palette)
    fill = _resolve_fill(spPr, theme, palette)
    gradient = _resolve_gradient(spPr, theme, palette)
    kind = _shape_geometry_kind(spPr)
    # Vertical anchor — `<a:bodyPr anchor="ctr">` / `b` / `t`. Without
    # this the rendered text lands at frame-top even when source centers
    # it, which is the dominant cause of the redline "two ghost positions"
    # pattern: source content at frame center, render content at frame top.
    valign: str | None = None
    padding_emu: tuple[int, int, int, int] | None = None
    txBody = ch.find(".//p:txBody", NS) or ch.find(".//a:txBody", NS)
    if txBody is not None:
        bodyPr = txBody.find("a:bodyPr", NS)
        if bodyPr is not None:
            anc = bodyPr.get("anchor")
            if anc == "ctr":
                valign = "middle"
            elif anc == "b":
                valign = "bottom"
            # Insets — l/t/r/b. Source omits = PowerPoint defaults
            # (91440 / 45720 EMU). Capture whatever's there so decompile
            # preserves exact text position, including the default insets.
            left = int(bodyPr.get("lIns") or 91440)
            top = int(bodyPr.get("tIns") or 45720)
            right = int(bodyPr.get("rIns") or 91440)
            bottom = int(bodyPr.get("bIns") or 45720)
            padding_emu = (left, top, right, bottom)
    # Convert insets EMU → design-px for the Shape (CanvasMap-relative).
    padding_px: tuple[float, float, float, float] | None = None
    if padding_emu is not None:
        left, top, right, bottom = padding_emu
        padding_px = (cmap.w(left), cmap.h(top), cmap.w(right), cmap.h(bottom))
    # Stroke (line) colour + width + dash. PowerPoint stores stroke width
    # in EMU on `<a:ln w="...">` (default ~9525 EMU = 0.75pt = 1px hairline);
    # dash preset on the optional `<a:prstDash val="...">` child.
    stroke = None
    stroke_width: float | None = None
    stroke_dash: str | None = None
    ln = spPr.find("a:ln", NS) if spPr is not None else None
    if ln is not None:
        sf = ln.find("a:solidFill", NS)
        if sf is not None:
            stroke = _resolve_solid(sf, theme, palette)
        w_attr = ln.get("w")
        if w_attr:
            try:
                w_emu = int(w_attr)
                # Width is uniform in EMU (1pt = 12700); converting via the
                # horizontal scale keeps it consistent with the rest of the
                # design-px coordinate system on this canvas.
                stroke_width = cmap.w(w_emu)
            except (ValueError, TypeError):
                pass
        dash = ln.find("a:prstDash", NS)
        if dash is not None and dash.get("val"):
            stroke_dash = dash.get("val")
    # Corner radius — captured from `prstGeom prst="roundRect"` with an
    # `<a:gd name="adj" fmla="val N">`. N is N/100000ths of the shape's
    # shortest side. PowerPoint defaults to 0.10 when adj is absent.
    corner_radius: float | None = None
    if spPr is not None:
        pg = spPr.find("a:prstGeom", NS)
        if pg is not None and pg.get("prst") == "roundRect":
            gd = pg.find(".//a:gd[@name='adj']", NS)
            adj_frac = 0.10  # PowerPoint default when omitted
            if gd is not None and gd.get("fmla"):
                m = re.search(r"val (\d+)", gd.get("fmla"))
                if m:
                    adj_frac = int(m.group(1)) / 100000
            corner_radius = cmap.w(min(w, h)) * adj_frac
    # Drop shadow — `<a:effectLst><a:outerShdw>`. Standard PowerPoint card
    # shadows use blurRad/dist in EMU, dir in 1/60000ths of a degree, and a
    # solid colour with an `<a:alpha val="...">` modifier (0-100000 = 0-100%).
    shadow: tuple[float, float, float, str, float] | None = None
    if spPr is not None:
        eff = spPr.find(".//a:effectLst/a:outerShdw", NS)
        if eff is not None:
            blur_emu = int(eff.get("blurRad") or 0)
            dist_emu = int(eff.get("dist") or 0)
            dir_60k = int(eff.get("dir") or 0)
            blur_px = cmap.w(blur_emu)
            dist_px = cmap.w(dist_emu)
            angle_deg = dir_60k / 60000.0
            sh_color = "black"
            sh_alpha = 1.0
            srgb = eff.find("a:srgbClr", NS)
            if srgb is not None and srgb.get("val"):
                hx = srgb.get("val")
                try:
                    rgb = (int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))
                    sh_color = nearest_token(rgb, palette) if palette else f"#{hx}"
                except ValueError:
                    pass
                alpha = srgb.find("a:alpha", NS)
                if alpha is not None and alpha.get("val"):
                    sh_alpha = int(alpha.get("val")) / 100000.0
            shadow = (blur_px, dist_px, angle_deg, sh_color, sh_alpha)

    # Picture-typed placeholder → picture shape (no actual <p:pic>).
    if ph_type == "pic":
        shapes.append(Shape(
            kind="pic", x=cmap.x(x), y=cmap.y(y), w=cmap.w(w), h=cmap.h(h),
            is_picture=True, ph_type=ph_type, ph_idx=ph_idx,
        ))
        return

    # Pure-text shape (placeholder, label, etc.) — no rect, just text.
    if runs and fill is None and kind == "rect":
        # Multi-paragraph colour split: when a text shape carries paragraphs
        # in different colours (typical of a cyan headline above silver body
        # bullets), emit one DSL `text` primitive per consecutive same-colour
        # block at calculated y-offsets, instead of collapsing them into a
        # single primitive whose first-run colour overrides the rest.
        blocks = _split_runs_by_color(runs)
        if len(blocks) > 1:
            frame_x_px = cmap.x(x)
            frame_y_px = cmap.y(y)
            frame_w_px = cmap.w(w)
            # Cursor-style y placement: each block consumes its own height
            # based on its primary pt (line-height factor 1.25 captures
            # standard slide leading without over-spacing tight headlines).
            cursor = frame_y_px
            for block_runs in blocks:
                block_pts = [r.pt for r in block_runs if r.text and r.text != "\n"]
                primary_pt = max(block_pts) if block_pts else 12
                line_count = sum(1 for r in block_runs if r.text and r.text != "\n")
                block_h_px = max(int(round(primary_pt * (4 / 3) * 1.25 * max(line_count, 1))), 24)
                # Split blocks always anchor top so the next block starts at
                # the bottom of the previous one — using the parent shape's
                # valign:middle would re-center each split block and overlap
                # neighbours.
                shapes.append(Shape(
                    kind="text", x=frame_x_px, y=cursor,
                    w=frame_w_px, h=block_h_px,
                    text_runs=block_runs, ph_type=ph_type, ph_idx=ph_idx,
                    valign=None, padding=padding_px,
                ))
                cursor += block_h_px
            return
        shapes.append(Shape(
            kind="text", x=cmap.x(x), y=cmap.y(y), w=cmap.w(w), h=cmap.h(h),
            text_runs=runs, ph_type=ph_type, ph_idx=ph_idx, valign=valign,
            padding=padding_px,
        ))
        return

    # custGeom paths convert directly to SVG path `d`. Build it in
    # canvas-pixel space so the surrounding svg-block can simply
    # `path "<d>"` without further transforms.
    svg_d = None
    if kind == "shape":
        svg_d = _custgeom_svg_d(spPr, cmap.w(w), cmap.h(h))

    # Geometry shape (rect / oval / shape). May also carry text.
    shapes.append(Shape(
        kind=kind, x=cmap.x(x), y=cmap.y(y), w=cmap.w(w), h=cmap.h(h),
        fill=fill, stroke=stroke, stroke_width=stroke_width,
        stroke_dash=stroke_dash, corner_radius=corner_radius, shadow=shadow,
        gradient=gradient,
        text_runs=runs, ph_type=ph_type, ph_idx=ph_idx, svg_path_d=svg_d,
    ))


def _emit_pic(ch, offset, shapes, slide, cmap, theme, palette):
    bbox = _shape_bbox(ch, offset, slide)
    if bbox is None:
        return
    x, y, w, h = bbox
    ph_type, ph_idx = _placeholder_info(ch)
    # Capture the embedded media rId so derive() can extract the binary
    # when image_extract_dir is set (pipeline-optimization mode).
    rid = None
    media_part = None
    blip = ch.find(".//a:blip", NS)
    if blip is not None:
        rid = blip.get(f"{{{NS['r']}}}embed")
        # Resolve the related part NOW against the source object that
        # actually owns the rId (the slide, layout, or master `slide`
        # param of this call). Storing only the rId and re-resolving on
        # `slide.part` later breaks for layout-inherited pictures —
        # their rId is scoped to the layout's relationships, not the
        # slide's.
        if rid is not None:
            try:
                media_part = slide.part.related_part(rid)
            except (KeyError, AttributeError):
                media_part = None
    shapes.append(Shape(
        kind="pic", x=cmap.x(x), y=cmap.y(y), w=cmap.w(w), h=cmap.h(h),
        is_picture=True, ph_type=ph_type, ph_idx=ph_idx, media_rid=rid,
        media_part=media_part,
    ))


def _emit_cxn(ch, offset, shapes, cmap, theme, palette):
    spPr = ch.find("p:spPr", NS)
    xfrm = _get_xfrm(spPr)
    if xfrm is None:
        return
    x, y, w, h = xfrm
    ox, oy = offset
    x += ox
    y += oy
    # Stroke color from line/solidFill.
    ln = spPr.find("a:ln", NS)
    stroke = None
    if ln is not None:
        sf = ln.find("a:solidFill", NS)
        if sf is not None:
            stroke = _resolve_solid(sf, theme, palette)
    shapes.append(Shape(
        kind="line", x=cmap.x(x), y=cmap.y(y), w=cmap.w(w), h=cmap.h(h),
        stroke=stroke or "fog",
    ))


CHART_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"
RELS_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _emit_graphic_frame(ch, offset, shapes, slide, cmap, theme, palette):
    """Tables and charts both arrive as <p:graphicFrame>. Dispatch by inner kind."""
    xfrm = ch.find("p:xfrm", NS)
    if xfrm is None:
        return
    off = xfrm.find("a:off", NS)
    ext = xfrm.find("a:ext", NS)
    if off is None or ext is None:
        return
    ox_local, oy_local = offset
    x0 = int(off.get("x")) + ox_local
    y0 = int(off.get("y")) + oy_local
    fw = int(ext.get("cx"))
    fh = int(ext.get("cy"))

    tbl = ch.find(".//a:tbl", NS)
    if tbl is not None:
        _emit_table(tbl, x0, y0, shapes, cmap, theme, palette)
        return

    # <c:chart r:id="..."/> inside graphicData → resolve chart part via slide rels.
    chart_ref = ch.find(f".//{{{CHART_NS}}}chart")
    if chart_ref is not None:
        rid = chart_ref.get(f"{{{RELS_NS}}}id")
        if rid:
            try:
                chart_part = slide.part.related_part(rid)
            except Exception:
                chart_part = None
            if chart_part is not None:
                _emit_chart(chart_part, x0, y0, fw, fh, shapes, cmap, theme, palette)


def _emit_table(tbl, x0, y0, shapes, cmap, theme, palette):
    grid = tbl.find("a:tblGrid", NS)
    if grid is None:
        return
    col_widths = [int(c.get("w")) for c in grid.findall("a:gridCol", NS)]

    # Auto-fit h=0 rows: PowerPoint expands them to fit text. Use a heuristic
    # of 350000 EMU (~31 px) per text line so subsequent rows shift down and
    # don't overlap the header. Without this the header sits visually inside
    # row 1.
    EMU_PER_LINE = 350000

    rows = list(tbl.findall("a:tr", NS))
    effective_h = []
    for tr in rows:
        h = int(tr.get("h", "0"))
        if h == 0:
            # estimate from text content
            text = "".join(t.text or "" for t in tr.findall(".//a:t", NS))
            lines = max(1, len(text.split("\n"))) if text else 1
            h = EMU_PER_LINE * lines
        effective_h.append(h)

    y_cursor = y0
    for ri, tr in enumerate(rows):
        row_h = effective_h[ri]
        x_cursor = x0
        cells = tr.findall("a:tc", NS)
        for ci, tc in enumerate(cells):
            cw = col_widths[ci] if ci < len(col_widths) else 0
            tcPr = tc.find("a:tcPr", NS)
            fill = _resolve_fill(tcPr, theme, palette) if tcPr is not None else None
            runs = _text_runs(tc, theme, palette)
            if fill is not None:
                shapes.append(Shape(
                    kind="rect", x=cmap.x(x_cursor), y=cmap.y(y_cursor),
                    w=cmap.w(cw), h=cmap.h(row_h), fill=fill,
                ))
            # Cell bottom border → emit as line (orange separators under headers).
            if tcPr is not None:
                lnB = tcPr.find("a:lnB", NS)
                if lnB is not None:
                    sf = lnB.find("a:solidFill", NS)
                    stroke = _resolve_solid(sf, theme, palette) if sf is not None else None
                    if stroke is not None:
                        y_b = y_cursor + row_h
                        shapes.append(Shape(
                            kind="line", x=cmap.x(x_cursor), y=cmap.y(y_b),
                            w=cmap.w(cw), h=0, stroke=stroke,
                        ))
            if runs:
                shapes.append(Shape(
                    kind="text", x=cmap.x(x_cursor + cw // 20), y=cmap.y(y_cursor + row_h // 4),
                    w=cmap.w(cw - cw // 10), h=cmap.h(row_h),
                    text_runs=runs,
                ))
            x_cursor += cw
        y_cursor += row_h


def _emit_pie_chart(pie_el, x0, y0, fw, fh, shapes, cmap):
    """Extract pie/doughnut chart geometry and emit one svg{} arc path per slice.

    Each slice becomes a Shape with kind='shape', svg_path_d set to an SVG
    arc path of the form 'M cx,cy L x1,y1 A r,r 0 large,sweep x2,y2 Z',
    fill mapped to chart-series-N via slice index (using the brand's
    chart-series ramp — for BSH that's accent → accent-80 → ... → accent-10).
    emit_dsl() converts each Shape into a standalone `svg{}` block; the
    blocks share the same bbox (the chart frame) so slices overlay into
    one unified pie at render time.

    Percentage labels are emitted as `text` Shapes positioned just outside
    the slice's arc bisector when <c:dLbls><c:showPercent val="1"> is set
    in the source — matches PowerPoint's default external-label placement.

    Per-slice colors in the source XML are intentionally ignored in favour
    of the brand's chart-series ramp: source decks authored against the
    brand already use the same hue sequence, so index-based mapping gives
    fidelity AND brand-correctness in one pass. Source decks authored
    off-brand will use the brand's hues here — that is by design (a
    brand-pack decompile is brand-conforming output, not pixel mimicry).
    """
    ser = pie_el.find(f"{{{CHART_NS}}}ser")
    if ser is None:
        return
    val_els = ser.findall(
        f".//{{{CHART_NS}}}val//{{{CHART_NS}}}pt/{{{CHART_NS}}}v"
    )
    try:
        values = [float(v.text) for v in val_els if v.text]
    except (TypeError, ValueError):
        return
    if not values or sum(values) <= 0:
        return
    total = sum(values)

    cat_els = ser.findall(
        f".//{{{CHART_NS}}}cat//{{{CHART_NS}}}pt/{{{CHART_NS}}}v"
    )
    categories = [c.text or "" for c in cat_els]

    # Legend position: source PowerPoint convention is r/l/t/b. We honour
    # only r/l (column-style legend with one row per category — the
    # dominant case for pies); t/b are rare on small pies and fall back
    # to right-side. Search at chart-space level since legend lives on
    # the chart root, not inside pie_el.
    chart_root = pie_el
    while chart_root is not None and chart_root.tag != f"{{{CHART_NS}}}chartSpace":
        parent = chart_root.getparent()
        if parent is None:
            break
        chart_root = parent
    legend_pos = "r"
    if chart_root is not None:
        lp = chart_root.find(f".//{{{CHART_NS}}}legend/{{{CHART_NS}}}legendPos")
        if lp is not None and lp.get("val") in ("l", "r", "t", "b"):
            legend_pos = lp.get("val")

    # Data label flags — <c:dLbls><c:showPercent val="1"/> or
    # <c:dLbls><c:showVal val="1"/>. Mutually exclusive in practice;
    # percent wins when both are set, matching PowerPoint behaviour.
    show_percent = False
    show_val = False
    sp = pie_el.find(
        f".//{{{CHART_NS}}}dLbls/{{{CHART_NS}}}showPercent"
    )
    if sp is not None and sp.get("val") == "1":
        show_percent = True
    sv = pie_el.find(
        f".//{{{CHART_NS}}}dLbls/{{{CHART_NS}}}showVal"
    )
    if sv is not None and sv.get("val") == "1":
        show_val = True

    # svg-block-local pixel coords for slice paths. The block's outer
    # bbox is the chart frame; coords inside are 0..bbox_w_px by
    # 0..bbox_h_px. min(w,h) keeps pies circular in non-square frames.
    # Pie-area fraction adapts to chart-frame aspect ratio: wide frames
    # (multi-pie-in-column layouts) keep ~60% pie area so pies fill the
    # narrow column adequately; square-ish frames (single-big-pie
    # layouts) shrink to 45% so the pie + adjacent legend mirror
    # PowerPoint's left-edge placement.
    bbox_w_px = cmap.w(fw)
    bbox_h_px = cmap.h(fh)
    frame_aspect = bbox_w_px / bbox_h_px if bbox_h_px else 1.0
    if categories and legend_pos in ("l", "r"):
        pie_w_frac = 0.60 if frame_aspect > 1.4 else 0.50
    else:
        pie_w_frac = 1.0
    pie_w_px = bbox_w_px * pie_w_frac
    pie_h_px = bbox_h_px
    pie_off_x = (bbox_w_px - pie_w_px) if legend_pos == "l" else 0
    cx_px = pie_off_x + pie_w_px / 2
    cy_px = pie_h_px / 2
    # 0.36 of pie-area min dimension leaves margin for external percentage
    # labels around the circumference.
    r_px = min(pie_w_px, pie_h_px) * 0.36

    # Start at 12 o'clock (-π/2), sweep clockwise. PowerPoint pies follow
    # this convention; matching it preserves slice-to-color correspondence
    # against the source.
    angle_start = -math.pi / 2

    for i, v in enumerate(values):
        if v <= 0:
            angle_start += 0
            continue
        sweep = (v / total) * 2 * math.pi
        angle_end = angle_start + sweep
        x1 = cx_px + r_px * math.cos(angle_start)
        y1 = cy_px + r_px * math.sin(angle_start)
        x2 = cx_px + r_px * math.cos(angle_end)
        y2 = cy_px + r_px * math.sin(angle_end)
        large_arc = 1 if sweep > math.pi else 0
        # SVG arc: A rx,ry x-axis-rotation large-arc-flag sweep-flag x,y
        # sweep-flag=1 (clockwise) matches our positive-angle direction in
        # screen-space (y grows down).
        d = (
            f"M {cx_px:.1f},{cy_px:.1f} "
            f"L {x1:.1f},{y1:.1f} "
            f"A {r_px:.1f},{r_px:.1f} 0 {large_arc},1 {x2:.1f},{y2:.1f} Z"
        )
        fill_token = f"chart-series-{(i % 6) + 1}"
        shapes.append(Shape(
            kind="shape",
            x=cmap.x(x0), y=cmap.y(y0),
            w=bbox_w_px, h=bbox_h_px,
            fill=fill_token,
            svg_path_d=d,
        ))

        if show_percent or show_val:
            mid_angle = angle_start + sweep / 2
            # showPercent labels sit ~15% outside circumference (PPT
            # default external placement on small pies). showVal labels
            # sit inside the slice at ~65% radius — PowerPoint's default
            # internal-label position.
            label_r = r_px * 1.15 if show_percent else r_px * 0.65
            lx = cx_px + label_r * math.cos(mid_angle)
            ly = cy_px + label_r * math.sin(mid_angle)
            if show_percent:
                label_text = f"{round((v / total) * 100)} %"
            else:
                # Format like the source: 8.2 → "8,2", 1 → "1".
                if v == int(v):
                    label_text = str(int(v))
                else:
                    label_text = f"{v:.1f}".replace(".", ",")
            shapes.append(Shape(
                kind="text",
                x=cmap.x(x0) + int(lx) - 20,
                y=cmap.y(y0) + int(ly) - 10,
                w=40, h=20,
                text_runs=[TextRun(text=label_text, pt=10)],
            ))

        angle_start = angle_end

    # Legend (categories + colour swatches). Position: right or left of pie.
    # Stack one row per category centred vertically against the pie. Skip
    # entirely when no categories — labelled-from-percentage slices alone
    # carry enough signal for the small pies common in showcase decks.
    if categories and legend_pos in ("l", "r"):
        swatch_w = 14
        swatch_h = 14
        row_gap = 32  # vertical pitch between legend rows
        text_gap = 8  # swatch-to-label horizontal gap
        legend_w_px = bbox_w_px - pie_w_px
        # Legend bbox top-left within the chart frame.
        legend_x_local = 0 if legend_pos == "l" else pie_w_px
        # Center the legend block vertically against the pie.
        block_h = len(categories) * row_gap
        legend_y_local = max(0, (bbox_h_px - block_h) / 2)
        # Insets keep legend out of frame edges.
        legend_x_local += 12
        for i, cat in enumerate(categories):
            row_y = legend_y_local + i * row_gap
            color = f"chart-series-{(i % 6) + 1}"
            # Swatch as a rect — survives outside svg{} blocks and accepts
            # the brand vocab directly.
            shapes.append(Shape(
                kind="rect",
                x=cmap.x(x0) + int(legend_x_local),
                y=cmap.y(y0) + int(row_y),
                w=swatch_w, h=swatch_h,
                fill=color,
            ))
            shapes.append(Shape(
                kind="text",
                x=cmap.x(x0) + int(legend_x_local) + swatch_w + text_gap,
                y=cmap.y(y0) + int(row_y) - 4,
                w=int(legend_w_px) - swatch_w - text_gap - 12,
                h=24,
                text_runs=[TextRun(text=cat, pt=10)],
            ))


def _emit_chart(chart_part, x0, y0, fw, fh, shapes, cmap, theme, palette):
    """Extract chart geometry from a c:chartSpace part and emit DSL primitives.

    Dispatches by chart type:
      * c:pieChart / c:doughnutChart → _emit_pie_chart (arc paths + labels)
      * c:barChart                   → _emit_bar_chart (rects + axis + labels)
    Other chart types (line, area, scatter, etc.) fall through unhandled —
    the decompile-as-rects fallback in the hybrid SVG pass still gives a
    rough first-pass for those, and improve-brand sub-agents refine.
    """
    try:
        root = etree.fromstring(chart_part.blob)
    except Exception:
        return

    pie = (root.find(f".//{{{CHART_NS}}}pieChart")
           or root.find(f".//{{{CHART_NS}}}doughnutChart"))
    if pie is not None:
        _emit_pie_chart(pie, x0, y0, fw, fh, shapes, cmap)
        return

    bar = root.find(f".//{{{CHART_NS}}}barChart")
    if bar is None:
        return

    series = []
    for ser in bar.findall(f"{{{CHART_NS}}}ser"):
        name_el = ser.find(f".//{{{CHART_NS}}}tx//{{{CHART_NS}}}v")
        name = name_el.text if name_el is not None else "?"
        vals = [float(v.text) for v in ser.findall(f".//{{{CHART_NS}}}val//{{{CHART_NS}}}pt/{{{CHART_NS}}}v")]
        cats = [v.text for v in ser.findall(f".//{{{CHART_NS}}}cat//{{{CHART_NS}}}pt/{{{CHART_NS}}}v")]
        # Per-series fill colour from <c:ser><c:spPr><a:solidFill>. Falls
        # back to None so the caller knows to use the chart-series-N ramp
        # by index instead. Reads either scheme or srgb fills via the
        # existing solid-fill resolver.
        sp_pr = ser.find(f"{{{CHART_NS}}}spPr")
        ser_color = _resolve_fill(sp_pr, theme, palette) if sp_pr is not None else None
        series.append((name, vals, cats, ser_color))
    if not series:
        return

    n_cats = max(len(s[1]) for s in series)
    n_series = len(series)
    cats = series[0][2] if series[0][2] else [f"Cat {i+1}" for i in range(n_cats)]
    data_max = max((max(s[1]) for s in series if s[1]), default=0)
    # Round axis max up to the next integer above data_max.
    # LibreOffice's auto-axis (the source-PNG ground truth in the verify
    # loop) adds one major-unit of headroom over data_max, so matching its
    # tick count beats the more semantically-correct ceil(data_max) — the
    # struct_diff_ratio improves when ticks line up with the source-PNG
    # rasterisation, not with what PowerPoint would have drawn.
    axis_max = math.ceil(data_max + 0.5) if data_max > 0 else 5

    # Plot-area extents inside the frame (EMU). Match PowerPoint defaults:
    # ~7% left for y-axis labels, ~12% top for cat labels, ~22% bottom for legend.
    plot_x = x0 + int(fw * 0.07)
    plot_y = y0 + int(fh * 0.12)
    plot_w = int(fw * 0.91)
    plot_h = int(fh * 0.66)

    # Y-axis numeric labels on the left.
    n_ticks = axis_max + 1
    for i in range(n_ticks):
        v = axis_max - i  # top→bottom
        ty = plot_y + int(plot_h * i / axis_max)
        shapes.append(Shape(
            kind="text",
            x=cmap.x(x0 + int(fw * 0.005)),
            y=cmap.y(ty - 180000),
            w=cmap.w(int(fw * 0.05)),
            h=cmap.h(360000),
            text_runs=[TextRun(text=str(v), pt=14)],
        ))

    # Skip gridlines: source has barely-visible hairlines; rendering them at
    # 0.75pt over a 1240px-wide plot adds heavy diff pixels and emit_dsl
    # orders lines after rects, so they paint OVER the bars producing stripes.

    # Category labels above each group.
    cat_w = plot_w // n_cats if n_cats else plot_w
    for ci in range(n_cats):
        cx = plot_x + ci * cat_w + cat_w // 4
        shapes.append(Shape(
            kind="text",
            x=cmap.x(cx),
            y=cmap.y(plot_y - int(fh * 0.07)),
            w=cmap.w(cat_w // 2),
            h=cmap.h(int(fh * 0.06)),
            text_runs=[TextRun(text=cats[ci] if ci < len(cats) else "", pt=14)],
        ))

    # Bars: each category has n_series side-by-side bars. Source bars are slim
    # (~8% of category width) and adjacent (no inter-bar gap). Wider/spaced
    # bars produced visibly different pixels and inflated struct_diff.
    # Series colour: prefer the per-series fill resolved from <c:ser><c:spPr>
    # (gives brand-accurate orange/peach/grey progression). Fall back to the
    # chart-series ramp by index when the source omits per-series colour.
    bar_w = int(cat_w * 0.085)
    group_w = bar_w * n_series
    group_inset = (cat_w - group_w) // 2
    for si, (name, vals, _, ser_color) in enumerate(series):
        color = ser_color or f"chart-series-{(si % 6) + 1}"
        for ci, v in enumerate(vals):
            bx = plot_x + ci * cat_w + group_inset + si * bar_w
            bh = int(plot_h * v / axis_max) if axis_max > 0 else 0
            by = plot_y + plot_h - bh
            shapes.append(Shape(
                kind="rect",
                x=cmap.x(bx), y=cmap.y(by),
                w=cmap.w(bar_w), h=cmap.h(bh),
                fill=color,
            ))
            # Value label above the bar.
            label = str(v).rstrip("0").rstrip(".") if "." in str(v) else str(v)
            label = label.replace(".", ",")
            shapes.append(Shape(
                kind="text",
                x=cmap.x(bx - bar_w // 2),
                y=cmap.y(by - 400000),
                w=cmap.w(bar_w * 2),
                h=cmap.h(360000),
                text_runs=[TextRun(text=label, pt=14)],
            ))

    # Legend at bottom-left.
    legend_y = y0 + fh - int(fh * 0.12)
    legend_x = plot_x + int(fw * 0.02)
    swatch_w = int(fw * 0.012)
    swatch_h = int(fh * 0.025)
    # Read chart title from c:title//c:tx//c:rich//a:p//a:r//a:t (or cached
    # strRef). Omit the title primitive entirely when the chart has no title.
    title_el = root.find(f".//{{{CHART_NS}}}title")
    title_text = ""
    if title_el is not None:
        for t_el in title_el.iterfind(f".//{{{NS['a']}}}t"):
            if t_el.text:
                title_text += t_el.text
    if title_text:
        shapes.append(Shape(
            kind="text",
            x=cmap.x(legend_x), y=cmap.y(legend_y),
            w=cmap.w(int(fw * 0.22)), h=cmap.h(int(fh * 0.04)),
            text_runs=[TextRun(text=title_text, pt=14)],
        ))
    lx = legend_x + int(fw * 0.18)
    for si, (name, _, _, ser_color) in enumerate(series):
        color = ser_color or f"chart-series-{(si % 6) + 1}"
        shapes.append(Shape(
            kind="rect",
            x=cmap.x(lx), y=cmap.y(legend_y + (swatch_h // 4)),
            w=cmap.w(swatch_w), h=cmap.h(swatch_h),
            fill=color,
        ))
        shapes.append(Shape(
            kind="text",
            x=cmap.x(lx + swatch_w + 50000),
            y=cmap.y(legend_y),
            # Legend label width must fit "Data N" at 14pt without
            # wrapping; the previous 4% gave ~31 design-px which forced
            # multi-line "Da\nta\n1" — visibly wrong on every chart
            # legend. 10% fits the common case comfortably + leaves
            # room for short multi-word series names.
            w=cmap.w(int(fw * 0.10)),
            h=cmap.h(int(fh * 0.04)),
            text_runs=[TextRun(text=name or "", pt=14)],
        ))
        lx += int(fw * 0.12)


# ---------------------------------------------------------------------------
# Style mapping (PPTX pt size → DSL style token)
# ---------------------------------------------------------------------------


_NUM_RE = re.compile(r"^\s*\d{1,2}\.\s*$")


def _style_for(pt: float, text: str, is_footer: bool) -> str:
    # Map source pt → nearest available style bundle (by emitted px). The
    # standard feinschliff bundles emit these px sizes:
    #   body-sm=16  ·  body=26  ·  sub=44  ·  title-l=80  ·  huge=120  ·  display=160
    # A DSL `text style:sub` round-trips through python-pptx as ≈33pt on the
    # rendered slide (px*0.75). So matching SOURCE pt to BUNDLE px directly
    # (1pt ≈ 1.333px) lands closer to source than the prior bucketing,
    # which mapped any pt ≥ 28 to title-l and over-sized 32pt slide titles
    # by almost 2×. Boundaries below are midpoints between adjacent bundle
    # px values, expressed as source pt.
    if _NUM_RE.match(text):
        return "agenda-num"
    if is_footer:
        return "footer"
    px = pt * 1.333
    if px >= 140:                              # 140+ px (≈ 105pt+) → display
        return "display"
    if px >= 100:                              # 100-140 px (≈ 75-105pt) → huge
        return "huge"
    if px >= 62:                               # 62-100 px (≈ 47-75pt) → title-l
        return "title-l"
    if px >= 35:                               # 35-62 px (≈ 26-47pt) → sub
        return "sub"
    if px >= 21:                               # 21-35 px (≈ 16-26pt) → body
        return "body"
    return "body-sm"                           # <21 px → body-sm (16 px)


# ---------------------------------------------------------------------------
# Emission
# ---------------------------------------------------------------------------


def emit_dsl(shapes: list[Shape], cmap: CanvasMap, layout_name: str,
             theme_name: str = "feinschliff",
             placeholder_rel: str = PLACEHOLDER_REL,
             bg_fill: str | None = None) -> str:
    out: list[str] = [
        "# auto-derived from PPTX+SVG hybrid — review before use",
        f"# layout: {layout_name}",
        f"canvas {cmap.cw}x{cmap.ch}",
        f"theme {theme_name}",
        "",
    ]

    # Slide-level background fill (from <p:cSld><p:bg>) emits first as a
    # full-canvas rect. PowerPoint draws this *under* every shape, so the
    # DSL ordering matches the source z-order. Without this, dark slides
    # rebuild on the brand's paper default and white text disappears.
    if bg_fill:
        out.append(f"rect 0,0 {cmap.cw}x{cmap.ch} fill:{bg_fill}")

    # Pre-split into geometry-first, text-last, footer-last.
    rects = [s for s in shapes if s.kind == "rect" and s.fill]
    ovals = [s for s in shapes if s.kind == "oval"]
    custs = [s for s in shapes if s.kind == "shape"]
    pics = [s for s in shapes if s.kind == "pic" or s.is_picture]
    lines = [s for s in shapes if s.kind == "line"]
    texts: list[Shape] = []
    for s in shapes:
        if s.kind == "text" and s.text_runs:
            texts.append(s)
        elif s.kind in ("rect", "oval", "shape") and s.text_runs:
            # geometry shape that also carries text — emit shape now, defer text
            texts.append(Shape(
                kind="text", x=s.x, y=s.y, w=s.w, h=s.h, text_runs=s.text_runs,
                ph_type=s.ph_type, ph_idx=s.ph_idx, valign=s.valign,
                padding=s.padding,
            ))

    footer_y_threshold = int(cmap.ch * 0.92)

    # Backgrounds first (large area). Append stroke / stroke-width / dash /
    # radius captured from the source PowerPoint shape so framed cards,
    # rounded rects, and dashed dividers survive the round-trip.
    for r in sorted(rects, key=lambda s: -(s.w * s.h)):
        line = f"rect {r.x},{r.y} {r.w}x{r.h} fill:{r.fill}"
        if r.gradient is not None:
            stops, angle = r.gradient
            stops_str = ";".join(f"{p:.2f}={c}" for p, c in stops)
            line += f" gradient:angle={angle:g};{stops_str}"
        if r.corner_radius is not None and r.corner_radius > 0:
            line += f" radius:{r.corner_radius:g}"
        if r.stroke:
            line += f" stroke:{r.stroke}"
            if r.stroke_width is not None and r.stroke_width > 0:
                line += f" stroke-width:{r.stroke_width:g}"
            if r.stroke_dash:
                line += f" dash:{r.stroke_dash}"
        if r.shadow is not None:
            blur, dist, angle, color, alpha = r.shadow
            line += (f" shadow:blur:{blur:g},dist:{dist:g},"
                     f"angle:{angle:g},color:{color},alpha:{alpha:.2f}")
        out.append(line)

    # Custom shapes (puzzle pieces, parallelograms, border paths, ring
    # sectors, etc.). When we recovered an SVG `d` string from the source
    # `<a:custGeom>`, emit an `svg { path "<d>" … }` block so the
    # renderer reproduces the actual vector geometry — this is the
    # difference between "blue donut with the right arc" and "grey bbox
    # rect where the donut should be." Stroke-only paths emit
    # `stroke:<token>` with no fill; otherwise fill (or fog fallback).
    # When no path data is available we keep the lossy bbox-rect fallback
    # so the DSL still builds.
    for i, s in enumerate(custs, 1):
        if s.svg_path_d:
            # `none` is not in the SVG DSL's 17-name semantic colour
            # vocabulary, so we omit `fill:` entirely when the source has
            # no solid fill — the path primitive defaults to stroke-only
            # rendering. With a fill, map the brand token onto an SVG
            # vocabulary name.
            attrs = []
            if s.fill:
                attrs.append(f"fill:{_svg_color_token(s.fill)}")
            if s.stroke:
                attrs.append(f"stroke:{_svg_color_token(s.stroke)}")
            attr_str = (" " + " ".join(attrs)) if attrs else ""
            out.append(f"svg shape{i} {s.x},{s.y} {s.w}x{s.h} {{")
            # Path coordinates are already in svg-block-local pixels (the
            # converter scales from path-local space to the shape's bbox).
            out.append(f"  path p \"{s.svg_path_d}\"{attr_str}")
            out.append("}")
        elif s.fill is None and s.stroke:
            out.append(f"shape {s.x},{s.y} {s.w}x{s.h} kind:rect stroke:{s.stroke}")
        else:
            out.append(f"shape {s.x},{s.y} {s.w}x{s.h} kind:rect fill:{s.fill or 'fog'}")

    # Ovals (circles, decorative dots). Stroke-only ovals (callout
    # circles, annotation marks) emit as stroke without fill; fill-only
    # ovals emit fill; if neither is present, fall back to a muted neutral
    # so the DSL always builds.
    for o in ovals:
        if o.fill is None and o.stroke:
            line = f"shape {o.x},{o.y} {o.w}x{o.h} kind:oval stroke:{o.stroke}"
        else:
            line = f"shape {o.x},{o.y} {o.w}x{o.h} kind:oval fill:{o.fill or 'fog'}"
            if o.stroke:
                line += f" stroke:{o.stroke}"
        if o.stroke_width is not None and o.stroke_width > 0:
            line += f" stroke-width:{o.stroke_width:g}"
        out.append(line)

    # Pictures — default to the brand's generic placeholder (so a derived
    # layout works as a reusable template) OR, when derive() was called
    # with image_extract_dir, to the actual extracted-from-source asset
    # (pipeline-optimization mode: no picture_coverage masking needed,
    # struct_diff_ratio reflects real shape/text mismatch).
    # Clamp bbox to the canvas so that picture-bleed boxes (e.g.
    # 166,-144 2345x1319 on 1920x1080) become canvas-fitted rectangles.
    # PowerPoint crops bleed at slide edges anyway, and the unclamped
    # bbox confuses the visual-diff coverage gate (>90% triggers a
    # struct = total fallback that masks real text deficits).
    for i, p in enumerate(pics, 1):
        slot = "image" if len(pics) == 1 else f"image{i}"
        cx0 = max(0, p.x)
        cy0 = max(0, p.y)
        cx1 = min(cmap.cw, p.x + p.w)
        cy1 = min(cmap.ch, p.y + p.h)
        cw = max(1, cx1 - cx0)
        ch = max(1, cy1 - cy0)
        default_path = p.media_path or placeholder_rel
        # The expander's default-filter grammar requires `default("...")`
        # — parentheses + double quotes (see lib/dsl/expander.py:_DEFAULT_FILTER_RE).
        # The earlier `default:'…'` form silently failed to match, so the
        # slot resolved to empty string and the build fell into the
        # "no image bound" placeholder-rect branch — exactly why pictures
        # rendered as grey rects even when the asset path was correct.
        out.append(
            f'picture {cx0},{cy0} {cw}x{ch} '
            f'path:"{{{{ {slot} | default(\\"{default_path}\\") }}}}" cover:true'
        )

    # Lines.
    for ln in lines:
        x1, y1 = ln.x, ln.y
        x2, y2 = ln.x + ln.w, ln.y + ln.h
        out.append(f"line {x1},{y1} {x2},{y2} stroke:{ln.stroke or 'fog'} stroke-width:1")

    if rects or pics or lines or ovals or custs:
        out.append("")

    # Footer collection: shapes whose y is in the bottom 8%.
    footer_runs: list[tuple[int, int, str]] = []
    body_texts: list[Shape] = []
    for t in texts:
        if t.y >= footer_y_threshold:
            for r in t.text_runs:
                footer_runs.append((t.x, t.y, r.text))
        else:
            body_texts.append(t)

    # Body text in reading order.
    body_texts.sort(key=lambda s: (round(s.y / 5) * 5, s.x))
    for t in body_texts:
        # Concatenate runs verbatim — PPTX emits explicit space-only runs
        # between words, so " ".join would produce double spaces. Collapse
        # any runs of whitespace down to a single space after concat.
        raw = "".join(r.text for r in t.text_runs if r.text)
        # Collapse runs of spaces/tabs but PRESERVE newlines (paragraph breaks).
        full = re.sub(r"[ \t]+", " ", raw).strip()
        # Strip stray spaces on either side of a newline that resulted from
        # space-only runs between paragraphs.
        full = re.sub(r" *\n *", "\n", full)
        if not full:
            continue
        # Exclude paragraph-break marker runs (text == "\n") from the size
        # vote — those runs inherit the body-level default size (often 18pt)
        # which would otherwise drown out the actual text-content size when
        # two 13pt paragraphs share a single shape.
        content_pts = [r.pt for r in t.text_runs if r.text and r.text != "\n"]
        pt = max(content_pts) if content_pts else max((r.pt for r in t.text_runs), default=18)
        style = _style_for(pt, full, is_footer=False)
        text = full.replace("\\", "\\\\").replace("\n", "\\n").replace('"', '\\"')
        mw = max(80, t.w)
        mh = max(24, t.h)
        # Emit `color:` override when the source's text-run colour differs
        # from the chosen style bundle's default. Captured in TextRun.color
        # from `<a:rPr><a:solidFill>`; without this the title that should
        # render in accent-blue lands in ink-grey (or whatever the style
        # bundle's default colour is) and inflates the visual diff.
        run_colors = [r.color for r in t.text_runs if r.color]
        run_color = run_colors[0] if run_colors else None
        style_default = STYLE_BUNDLES.get(style, {}).get("color")
        color_attr = (
            f" color:{run_color}" if run_color and run_color != style_default else ""
        )
        # Per-run weight override. `<a:rPr b="1">` rides on the source run;
        # the size-based classifier picks a bundle whose default weight may
        # not match (e.g. source 60pt bold maps to `huge` which is
        # weight:light by token convention). Without this, the rendered
        # text loses its emphasis even though the source explicitly carried
        # the bold flag. Emit only when source ≠ bundle default to keep
        # decompile output stable for the common case.
        source_bold = any(r.bold for r in t.text_runs)
        bundle_weight = STYLE_BUNDLES.get(style, {}).get("weight")
        weight_attr = ""
        if source_bold and bundle_weight not in ("bold", "semibold", "black"):
            weight_attr = " weight:bold"
        elif (not source_bold) and bundle_weight == "bold":
            # Source author explicitly chose regular against a bold-default
            # bundle — preserve that. `regular` is the most common name in
            # `font-weight` tokens; brands that use a different label can
            # override per-layout.
            weight_attr = " weight:regular"
        # Per-run size override. The classifier rounds source pt to the
        # nearest bundle, but the bundle steps are coarse (16/26/44/80 px)
        # and the emitted pt depends on the brand's slide physical width
        # (via `_PX_TO_PT`). Emit the source pt verbatim so renders match
        # source physical pt regardless of slide scale.
        # exactly. Stable for the common case (no emit when bundle is close).
        # Always emit `size:<pt>pt` from the source pt. The toolkit's style
        # bundles use design-px (`body=26px`, `sub=44px`, ...) which only
        # round-trip to the right rendered pt at the 13.33"-wide slide
        # convention (`_PX_TO_PT=0.5`). When the brand pack inherits a
        # different physical slide size from the source PPTX (`slide.width_emu`
        # in tokens.json), the emitter's `_PX_TO_PT` shifts and the same
        # px-valued bundles render at a different pt — which silently
        # shrinks every untagged text. Locking each run to its source pt
        # makes physical font sizes faithful regardless of slide scale.
        size_attr = f" size:{pt:g}pt"
        valign_attr = f" valign:{t.valign}" if t.valign else ""
        # Horizontal align — pick the first non-None align from the runs.
        # PPTX stores it per-paragraph; for a single emitted `text` primitive
        # the first paragraph's alignment wins (consistent with how we pick
        # color + size from the first content run).
        run_aligns = [r.align for r in t.text_runs if r.align]
        run_align = run_aligns[0] if run_aligns else None
        align_attr = f" align:{run_align}" if run_align else ""
        padding_attr = ""
        if t.padding is not None:
            left, top, right, bottom = t.padding
            # Compact form `padding:N` when all four insets are equal,
            # `padding:L,T,R,B` otherwise — keeps the common-case DSL short.
            if left == right and top == bottom and left == top:
                padding_attr = f" padding:{left:g}"
            else:
                padding_attr = f" padding:{left:g},{top:g},{right:g},{bottom:g}"
        out.append(
            f'text {t.x},{t.y} style:{style}{color_attr}{weight_attr}{size_attr}{valign_attr}{align_attr}{padding_attr} '
            f'maxwidth:{mw} maxheight:{mh} "{text}"'
        )

    # Footer-region text. Anything below `footer_y_threshold` (bottom 8%)
    # is emitted as plain `text` primitives — brand-agnostic. Brands that
    # ship a dedicated `footer(...)` compound can post-process the
    # output to collapse the lines into one compound call (typically via a
    # brand-specific post-pass or plugin emit hook). We deliberately
    # do NOT try to detect master-inherited chrome here: chrome that the
    # slide XML doesn't carry won't appear in this decompile, which is
    # the honest behaviour for a single-slide decompiler. Use the
    # brand's compound + master template at build-time for chrome.
    if footer_runs:
        footer_runs.sort(key=lambda r: (r[1], r[0]))
        out.append("")
        for x, y, raw in footer_runs:
            text = re.sub(r"\s+", " ", raw).strip()
            if not text:
                continue
            escaped = text.replace("\\", "\\\\").replace('"', '\\"')
            out.append(
                f'text {x},{y} style:footer maxwidth:400 maxheight:40 "{escaped}"'
            )

    return "\n".join(out) + "\n"


# ---------------------------------------------------------------------------
# SVG bbox lookup (fallback / verification)
# ---------------------------------------------------------------------------


def svg_path_bboxes(svg_path: Path) -> list[tuple[float, float, float, float]]:
    """Return crude bounding boxes of <path> elements in the SVG, in pt units.
    Used to cross-check PPTX-derived geometry for custGeom shapes.
    """
    try:
        root = etree.parse(str(svg_path)).getroot()
    except Exception:
        return []
    bboxes: list[tuple[float, float, float, float]] = []
    for p in root.iter("{%s}path" % NS["svg"]):
        d = p.get("d") or ""
        nums = re.findall(r"-?\d+(?:\.\d+)?", d)
        if not nums:
            continue
        xs = [float(n) for n in nums[0::2]]
        ys = [float(n) for n in nums[1::2]]
        if not xs or not ys:
            continue
        bboxes.append((min(xs), min(ys), max(xs) - min(xs), max(ys) - min(ys)))
    return bboxes


# ---------------------------------------------------------------------------
# Public derive()
# ---------------------------------------------------------------------------


def derive(
    pptx_path: Path,
    slide_idx: int,
    canvas_w: int = 1920,
    canvas_h: int = 1080,
    tokens_path: Path | None = None,
    layout_name: str = "derived",
    theme_name: str = "feinschliff",
    placeholder_rel: str = PLACEHOLDER_REL,
    pdf_path: Path | None = None,
    image_extract_dir: Path | None = None,
    image_extract_rel: str | None = None,
) -> str:
    """Decompile one slide of `pptx_path` (1-indexed) into a Feinschliff DSL
    string. Brand-agnostic: pass `theme_name` and `tokens_path` to point at
    the target brand pack. `tokens_path` is used only for nearest-color
    matching against brand color tokens; if omitted, raw hex colors land
    in the DSL.

    `pdf_path` is currently unused; reserved for SVG cross-check of
    custGeom bboxes (callers render the slide's PDF page on demand).

    `image_extract_dir` + `image_extract_rel` enable **image carry-over**
    for pipeline-optimization runs: each `<p:pic>` in the source slide
    has its embedded binary written to `image_extract_dir/imageN.<ext>`,
    and the generated DSL's `default:` for that slot points at
    `image_extract_rel/imageN.<ext>` (a brand-pack-relative path the
    build will resolve at render time). Without these args, picture
    statements fall back to the generic placeholder image as before,
    which is the right default for *templating* an out-of-tree brand
    pack (where the source illustrations are not re-used). For *testing
    decompile fidelity* against the source, carry the images over so
    the visual diff measures real shape/text mismatch instead of
    picture-region noise.
    """
    # python-pptx rejects template content-types (.potx / .pptm-template)
    # with a cryptic "is not a PowerPoint file, content type is
    # '…presentationml.template.main+xml'" error. Catch the suffix up front
    # and tell the caller exactly how to convert. (LibreOffice converts
    # cleanly in one pass; renaming the file is NOT enough — the internal
    # [Content_Types].xml carries the template MIME and must be rewritten.)
    if pptx_path.suffix.lower() in (".potx", ".potm"):
        raise ValueError(
            f"{pptx_path.name} is a PowerPoint TEMPLATE (.potx/.potm), not a "
            f"presentation. Convert it first with:\n"
            f"  soffice --headless --convert-to pptx --outdir {pptx_path.parent} "
            f"{pptx_path}\n"
            f"then re-run with the produced .pptx."
        )
    palette: dict[str, tuple[int, int, int]] = {}
    if tokens_path and tokens_path.exists():
        palette = load_palette(tokens_path)
    pres = Presentation(str(pptx_path))
    theme = load_theme_scheme(pres)
    slide = pres.slides[slide_idx - 1]

    cx = pres.slide_width
    cy = pres.slide_height
    cmap = CanvasMap(cx, cy, canvas_w, canvas_h)

    shapes = walk_slide(slide, cmap, theme, palette)
    bg_fill = extract_slide_bg_fill(slide, theme, palette)
    _ = pdf_path  # reserved for SVG cross-check, off by default

    if image_extract_dir is not None:
        if image_extract_rel is None:
            raise ValueError(
                "image_extract_rel is required when image_extract_dir is set "
                "(it goes into the DSL `default:` slot — the build resolves "
                "it relative to the brand pack root)."
            )
        image_extract_dir.mkdir(parents=True, exist_ok=True)
        # Prefer pictures whose Part was resolved at walk-time (handles
        # layout-inherited pics); fall back to slide.part lookup for the
        # in-slide case so callers that pre-date media_part still work.
        pics = [s for s in shapes if s.is_picture and (s.media_part or s.media_rid)]
        for i, p in enumerate(pics, 1):
            part = p.media_part
            if part is None and p.media_rid:
                try:
                    part = slide.part.related_part(p.media_rid)
                except KeyError:
                    continue
            if part is None:
                continue
            blob = getattr(part, "blob", None)
            partname = str(getattr(part, "partname", "/image.bin"))
            ext = partname.rsplit(".", 1)[-1].lower() if "." in partname else "bin"
            stem = "image" if len(pics) == 1 else f"image{i}"
            out_name = f"{stem}.{ext}"
            (image_extract_dir / out_name).write_bytes(blob or b"")
            p.media_path = f"{image_extract_rel.rstrip('/')}/{out_name}"

    return emit_dsl(shapes, cmap, layout_name,
                    theme_name=theme_name,
                    placeholder_rel=placeholder_rel,
                    bg_fill=bg_fill)


def main() -> int:
    ap = argparse.ArgumentParser(description="Decompile one PPTX slide → Feinschliff DSL (hybrid PPTX+SVG)")
    ap.add_argument("pptx", type=Path)
    ap.add_argument("--slide", type=int, default=1)
    ap.add_argument("--canvas", default="1920x1080")
    ap.add_argument("--theme", default="feinschliff",
                    help="Brand name to emit on the `theme` directive (default: feinschliff)")
    ap.add_argument("--brand-tokens", type=Path, default=None,
                    help="Brand tokens.json — used for nearest-color matching against brand palette")
    ap.add_argument("--layout-name", default="derived")
    ap.add_argument("--placeholder", default=PLACEHOLDER_REL,
                    help=f"Picture placeholder path (default: {PLACEHOLDER_REL})")
    args = ap.parse_args()

    if not args.pptx.exists():
        print(f"missing: {args.pptx}", file=sys.stderr)
        return 2
    w, h = (int(x) for x in args.canvas.split("x"))
    print(derive(args.pptx, args.slide, w, h,
                 tokens_path=args.brand_tokens,
                 layout_name=args.layout_name,
                 theme_name=args.theme,
                 placeholder_rel=args.placeholder), end="")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
