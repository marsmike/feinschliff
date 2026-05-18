"""Primitive-AST → python-pptx emitter.

Consumes the post-expansion node list (primitives only) plus a Tokens
bundle, builds a single-slide Presentation, returns it.

Coordinate system: DSL uses design pixels in a 1920×1080 frame, matching
the existing brand-pack convention. We map 1 design-px → 6350 EMU
(same as gs-ramspau's `px()` helper).

Primitives implemented:
  canvas WxH                 — set slide size
  theme <brand>              — no-op here; brand was already picked at load
  text X,Y "label" style:S align:left|right|center maxwidth:W
  rect X,Y WxH fill:role stroke:role stroke-width:N
  line X,Y X2,Y2 stroke:role stroke-width:N
  picture X,Y WxH path:PATH cover:true
  picture X,Y WxH query:"…" — resolve at build time via image_provider
"""
from __future__ import annotations

import hashlib
import io
import json
import os
import re
import tempfile
import urllib.error
import urllib.request
import warnings
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import TYPE_CHECKING

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from .. import textfit
from .parser import DSLNode, parse_xy, parse_wh
from .polish import normalize_text
from .tokens import Tokens

if TYPE_CHECKING:
    from ..image_provider import ImageHit, ImageProvider


class DSLError(Exception):
    """Raised when a DSL primitive carries a parse-or-emit-time error
    that the brand author or layout author can fix by editing source
    (e.g. mutually-exclusive keywords, missing required wiring).

    Distinct from ``ValueError`` / ``SyntaxError`` so callers can branch
    on "the deck source is malformed" vs. "the parser hit an internal
    inconsistency".
    """


# 1920 design-px → 12,192,000 EMU (PowerPoint widescreen). Hence 6350 EMU/px.
_EMU_PER_PX = 6350           # 914400 EMU/in × (1/144 in/px) — for 1920×1080 slides at 96 DPI
_PX_TO_PT = 0.5              # design-px → typographic pt (2 design-px = 1pt at 1920-wide / 960pt PPT slide)
_STROKE_PX_TO_PT = 0.75      # CSS px → pt for stroke widths (96/72 inverse rounded)
_DEFAULT_PADDING_X = 100     # fallback right/left margin if brand has no slide.padding-x token

# Defense-in-depth: if interpolation somehow left a `{{ … }}` placeholder in
# an `if:` condition, treat it as falsy so the node is suppressed (never
# leaked into the rendered slide as literal text).
_RESIDUAL_PLACEHOLDER = re.compile(r"\{\{[^}]*\}\}")

# Tabular-numeral activation: a run is "numeric" if its content (after polish)
# is purely digits, separators, and numeric signs. Whitespace and Unicode
# minus (U+2212) count as numeric chrome.
_NUMERIC_CONTENT_RE = re.compile(r"^[\d,\.\-−%+\s]+$")

# Glyphs that, when leading a run, hang into the left margin. The bearing
# fraction is approximate (fraction of em — i.e. of font size).
_HANG_SIDE_BEARING_EM: dict[str, float] = {
    "•": 0.45,  # • bullet
    "“": 0.20,  # “ left double quote (en)
    "„": 0.20,  # „ low-9 (de)
    "«": 0.25,  # « left guillemet (fr)
    "—": 0.30,  # — em-dash
    "–": 0.25,  # – en-dash
    "…": 0.15,  # … ellipsis
    "\"":     0.20,  # ASCII straight double (legacy; rare after polish)
    "-":      0.15,  # ASCII hyphen
}


# 3-channel hierarchy stepping for indent levels. Each level beyond 0 steps:
#   size  ×= 0.85   (round to 0.5pt)
#   weight -= 100   (clamped at 300)
#   color  walks ink → graphite → fog (clamped at fog)
_HIERARCHY_COLOR_WALK = ["ink", "graphite", "fog"]


def _step_hierarchy(
    size_px: float, weight: int, color_role: str, *, level: int,
) -> tuple[float, int, str]:
    """Step all three channels for `level` indent levels. level<=0 is a no-op."""
    if level <= 0:
        return size_px, weight, color_role
    new_size = size_px
    new_weight = weight
    new_color = color_role
    for _ in range(level):
        new_size *= 0.85
        new_weight = max(300, new_weight - 100)
        if new_color in _HIERARCHY_COLOR_WALK:
            idx = _HIERARCHY_COLOR_WALK.index(new_color)
            new_color = _HIERARCHY_COLOR_WALK[min(idx + 1, len(_HIERARCHY_COLOR_WALK) - 1)]
    # Round size_pt to nearest 0.5pt. size_px ↔ pt: 2 design-px per pt.
    pt = new_size * _PX_TO_PT
    pt = round(pt * 2) / 2
    new_size = pt / _PX_TO_PT
    return new_size, new_weight, new_color


def _leading_hang_offset_px(text: str, size_pt: float) -> float:
    """Return leftward offset (design-px) for a textbox whose first glyph
    hangs into the margin. Zero if no hang glyph at position 0."""
    if not text:
        return 0.0
    bearing_em = _HANG_SIDE_BEARING_EM.get(text[0], 0.0)
    if bearing_em == 0.0:
        return 0.0
    # 1pt = 2 design-px (the canvas-to-PPT mapping). Em ≈ font size.
    return bearing_em * size_pt / _PX_TO_PT


def _px(n: float) -> Emu:
    return Emu(int(n * _EMU_PER_PX))


@dataclass
class EmitContext:
    tokens: Tokens
    canvas_w: float = 1920.0
    canvas_h: float = 1080.0
    asset_root: Path | None = None        # for resolving picture paths
    # Plugin-level fallback root walked when the per-brand `asset_root` does
    # not contain the requested relative path. Lets shared assets (universal
    # placeholder illustrations, common icons) live once at the plugin root
    # while brands keep the freedom to override per-asset by putting a file
    # at the same relative path under their own `assets/` dir.
    asset_root_fallback: Path | None = None
    # Required-asset accumulator (Review #7). _emit_picture appends an entry
    # for every picture node whose `path` is unset or points at a missing
    # file AND that does not carry `optional:true`. Callers consult this
    # post-build to decide whether to abort.
    missing_assets: list[dict] = field(default_factory=list)
    # Active image provider for resolving `picture query:"..."` nodes
    # (Task 7). String forward-ref keeps module import cheap and avoids any
    # circular-import risk between `lib.dsl.pptx_emit` and `lib.image_provider`.
    # Defaults to None so existing callers that never set it keep working.
    image_provider: "ImageProvider | None" = None
    # Deck output directory — used by the picture-query path to write the
    # `asset_lock.json` manifest and the `.cache/` of downloaded images
    # alongside the produced `.pptx`. Defaults to None; the picture-query
    # branch falls back to a temp dir / disabled cache when unset.
    deck_dir: Path | None = None


def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _hex_to_rgb_tuple(hex_color: str) -> tuple[int, int, int]:
    """Same as `_hex_to_rgb` but returns a plain `(r, g, b)` tuple for use
    with PIL APIs that don't accept python-pptx's `RGBColor` subclass.
    """
    h = hex_color.lstrip("#")
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _luminance(rgb: tuple[int, int, int]) -> float:
    """Rec. 601 luma (perceptual brightness 0..255)."""
    r, g, b = rgb
    return 0.299 * r + 0.587 * g + 0.114 * b


def _duotone_endpoints_for_brand(tokens) -> tuple[tuple[int, int, int], tuple[int, int, int]] | None:
    """Pick `(dark, light)` RGB endpoints for the plugin-fallback duotone
    treatment from a brand's tokens. Strategy:

    - Try ink + paper, picking the darker as shadow and lighter as highlight
      (auto-swap so dark-canvas brands still get a readable mapping).
    - If their luminance gap is too narrow to read (`< 80` on 0..255),
      substitute the lighter endpoint with the brand's accent color so the
      duotone keeps tonal range and brand identity.

    Returns None if ink/paper tokens are missing — caller falls back to the
    image's untreated form.
    """
    try:
        a = _hex_to_rgb_tuple(tokens.color("ink"))
        b = _hex_to_rgb_tuple(tokens.color("paper"))
    except (KeyError, ValueError):
        return None
    dark, light = (a, b) if _luminance(a) <= _luminance(b) else (b, a)
    if _luminance(light) - _luminance(dark) < 80.0:
        try:
            accent = _hex_to_rgb_tuple(tokens.color("accent"))
            if _luminance(accent) > _luminance(dark):
                light = accent
        except (KeyError, ValueError):
            pass
    return dark, light


def _align_pp(name: str) -> PP_ALIGN:
    return {
        "left":   PP_ALIGN.LEFT,
        "right":  PP_ALIGN.RIGHT,
        "center": PP_ALIGN.CENTER,
    }.get(name, PP_ALIGN.LEFT)


# ---------------------------------------------------------------------------
# Primitive handlers
# ---------------------------------------------------------------------------

def _emit_text(slide, node: DSLNode, ctx: EmitContext) -> None:
    """text X,Y "label" style:S align:A maxwidth:W maxheight:H color:T autoshrink:true lang:de_DE

    `autoshrink:true` — shrink font from style size down to a 10pt floor until
        the label fits the `maxwidth × maxheight` box. Off by default; preserves
        byte-identical output for layouts that don't opt in.
    `lang:LOCALE` — pyphen locale (e.g. `de_DE`). Inserts U+00AD soft hyphens
        into the label before emission so python-pptx can break long compound
        words at syllable boundaries. Applied BEFORE autoshrink so the shrink
        sees the hyphenated string.
    """
    if not node.pos_args:
        raise ValueError(f"text at line {node.line_no}: expected 'X,Y' positional")
    x, y = parse_xy(node.pos_args[0])
    style_name = node.kw_args.get("style", "body")
    style = ctx.tokens.resolve_style(style_name)
    color_override = node.kw_args.get("color")
    if color_override:
        from dataclasses import replace as _replace
        style = _replace(style, color_hex=ctx.tokens.color(color_override),
                         color_role=color_override)
    # Hierarchy stepping: indent:N steps size/weight/color N times.
    try:
        indent_level = int(node.kw_args.get("indent", "0"))
    except (TypeError, ValueError):
        indent_level = 0
    if indent_level > 0:
        from dataclasses import replace as _replace
        new_size, new_weight, new_color_role = _step_hierarchy(
            style.size_px, style.weight, style.color_role, level=indent_level,
        )
        try:
            new_color_hex = ctx.tokens.color(new_color_role)
        except KeyError:
            new_color_hex = style.color_hex
        style = _replace(
            style,
            size_px=new_size, weight=new_weight,
            color_hex=new_color_hex, color_role=new_color_role,
        )
    align = _align_pp(node.kw_args.get("align", "left"))
    # Default right margin = brand's slide.padding-x token; fall back to a
    # sensible canvas-relative inset if the brand omits it.
    try:
        padding_x = ctx.tokens.slide("padding-x")
    except Exception:
        padding_x = _DEFAULT_PADDING_X
    if not padding_x:
        padding_x = _DEFAULT_PADDING_X
    maxw = float(node.kw_args.get("maxwidth", ctx.canvas_w - x - padding_x))
    # Default height: ~1.5× font size (in design-px) to give room for one line.
    h = float(node.kw_args.get("maxheight", style.size_px * 1.5))

    label_text = node.label or ""
    if label_text:
        label_text = normalize_text(label_text, locale=ctx.tokens.locale)
    lang = node.kw_args.get("lang")
    if lang:
        label_text = textfit.hyphenate(label_text, lang=lang)

    autoshrink = str(node.kw_args.get("autoshrink", "")).lower() == "true"
    if autoshrink and label_text:
        from dataclasses import replace as _replace
        # Use the longest single paragraph as the worst-case for width fit; the
        # height check below works on the joined text.
        face, bold = _resolve_face(style.font_family[0], style.weight)
        max_pt = _px_to_pt(style.size_px)
        fitted_pt = textfit.autoshrink_size(
            label_text,
            font=face,
            max_size_pt=max_pt,
            min_size_pt=10,
            bold=bold,
            width_emu=int(maxw * _EMU_PER_PX),
            height_emu=int(h * _EMU_PER_PX),
            line_height=style.line_height,
        )
        if fitted_pt < max_pt:
            # Convert pt back to design-px (2 design-px per pt).
            style = _replace(style, size_px=fitted_pt * 2.0)

    # Orphan control: per paragraph, if the greedy wrap would orphan one
    # word on the final line, replace its preceding space with NBSP so
    # the pair wraps together. Use the final (post-autoshrink) size.
    if label_text:
        face_for_fit, bold_for_fit = _resolve_face(style.font_family[0], style.weight)
        paragraphs = label_text.split("\n")
        paragraphs = [
            textfit.prevent_orphan(
                p,
                font=face_for_fit,
                size_pt=_px_to_pt(style.size_px),
                bold=bold_for_fit,
                width_emu=int(maxw * _EMU_PER_PX),
            )
            for p in paragraphs
        ]
        label_text = "\n".join(paragraphs)

    # Hanging punctuation: shift the textbox left so the leading glyph
    # hangs into the margin (its body-letter optical edge sits at x).
    first_line = label_text.split("\n", 1)[0] if label_text else ""
    hang_offset_px = _leading_hang_offset_px(first_line, _px_to_pt(style.size_px))
    box = slide.shapes.add_textbox(
        _px(x - hang_offset_px), _px(y), _px(maxw), _px(h)
    )
    # `rotate:N` — rotate the textbox N degrees clockwise around its
    # geometric center. Use -90 for a y-axis label reading bottom-to-top.
    # Bbox semantics: x,y,maxwidth,maxheight describe the *pre-rotation*
    # rect; PPTX rotates around the rect's center on render. To place a
    # vertical label at a chart's left edge, set x to the chart edge minus
    # half the natural width, and y so the rect's vertical center matches
    # the chart's vertical midpoint.
    rotate_str = node.kw_args.get("rotate")
    if rotate_str:
        try:
            box.rotation = float(rotate_str)
        except (TypeError, ValueError):
            pass
    # Tag title text-shapes so downstream verify helpers
    # (extract_titles_from_pptx, extract_slide_title_and_body) can identify
    # them. Feinschliff layouts don't use python-pptx placeholder types, so
    # slide.shapes.title is always None — the name prefix is how we recover
    # the semantic role post-render. Prefix is chosen to NOT collide with the
    # chrome-shape prefixes in lib/verify/chrome.py.
    #
    # The 5 styles tagged here cover every title-role text in the shared
    # layouts catalog:
    #   title, act-title — standard content-slide titles
    #   sub             — used only in agenda.slide.dsl for the slide headline
    #   display-xl      — used only in end.slide.dsl for the closer headline
    #   quote           — used in full-bleed-cover + quote.slide.dsl (both
    #                     are "headline-shaped" usages, not body decoration)
    if style_name in ("title", "act-title", "sub", "display-xl", "quote"):
        box.name = f"feinschliff-title-{style_name}"
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    lines = label_text.split("\n")
    p0 = tf.paragraphs[0]
    p0.alignment = align
    p0.line_spacing = style.line_height
    _style_run(p0.add_run(), lines[0], style, tokens=ctx.tokens)
    for extra in lines[1:]:
        p = tf.add_paragraph()
        p.alignment = align
        p.line_spacing = style.line_height
        # Inter-paragraph gap defaults to a fraction of font-size; zero it so
        # the visual gap matches CSS line-height applied across all lines.
        p.space_before = Pt(0)
        _style_run(p.add_run(), extra, style, tokens=ctx.tokens)


def _is_numeric_run(text: str) -> bool:
    """True if every character in text is a digit, separator, or numeric sign."""
    return bool(text) and bool(_NUMERIC_CONTENT_RE.match(text))


def _tracking_for_size(size_pt: float, curve: dict[int, float]) -> float:
    """Step-function lookup: return curve[max-key ≤ size_pt], or 0.0 if no
    curve key is at or below size_pt. The curve is a brand-level default
    for display-tier tracking; per-style `letter_spacing` overrides it.
    """
    if not curve:
        return 0.0
    sorted_keys = sorted(curve.keys())
    if size_pt < sorted_keys[0]:
        return 0.0
    chosen = sorted_keys[0]
    for k in sorted_keys:
        if k <= size_pt:
            chosen = k
        else:
            break
    return curve[chosen]


def _style_run(run, text: str, style, *, tokens: Tokens | None = None) -> None:
    if style.transform == "upper":
        text = text.upper()
    run.text = text
    f = run.font
    # Tabular-numeral override: when the brand exposes a complete tnum_font
    # face name (e.g. "Inter Tabular") and the run is a pure number, use
    # that face verbatim; bold comes from the style weight so emphasis is
    # preserved. Otherwise resolve face from family + weight as usual.
    if tokens is not None and tokens.tnum_font and _is_numeric_run(text):
        face = tokens.tnum_font
        bold = style.weight >= 600
    else:
        face, bold = _resolve_face(style.font_family[0], style.weight)
    f.name = face
    f.size = Pt(_px_to_pt(style.size_px))
    f.bold = bold
    color = style.color_hex
    if style.opacity < 1.0:
        # python-pptx has no native alpha on text runs; approximate by
        # blending toward white. Canonical .pgmeta opacity 0.7 on ink → ~#4A586F.
        color = _blend_to_white(color, 1.0 - style.opacity)
    f.color.rgb = _hex_to_rgb(color)
    # Tracking: per-style letter_spacing wins; if the style declares none,
    # fall back to the brand-level display_tracking_curve (step function).
    size_pt = _px_to_pt(style.size_px)
    tracking_em = style.letter_spacing
    if tracking_em == 0.0 and tokens is not None:
        tracking_em = _tracking_for_size(size_pt, tokens.display_tracking_curve)
    if tracking_em != 0.0:
        # OOXML `spc` is 100ths of a point. Negative tightens; positive loosens.
        spc = int(tracking_em * size_pt * 100)
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(spc))


# Map design weight → face-name suffix. Noto Sans ships every weight as a
# separate face with the weight encoded in the family name (e.g. "Noto Sans
# Light"). soffice resolves family-by-name, so we pick the right face here.
# Weight 700 is the conventional "bold" face — use regular family + bold flag.
_WEIGHT_SUFFIX = [
    (200, "Thin"),
    (300, "Light"),
    (400, None),        # Regular
    (500, "Medium"),
    (600, "SemiBold"),
    (700, "__bold__"),  # signal: use regular family + bold flag
    (800, "ExtraBold"),
    (1000, "Black"),
]


def _resolve_face(family: str, weight: int) -> tuple[str, bool]:
    for thresh, suffix in _WEIGHT_SUFFIX:
        if weight <= thresh:
            if suffix == "__bold__":
                return family, True
            return (f"{family} {suffix}" if suffix else family), False
    return family, False


def _blend_to_white(hex_color: str, t: float) -> str:
    """Linearly blend `hex_color` toward white by fraction `t` (0..1)."""
    h = hex_color.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = round(r + (255 - r) * t)
    g = round(g + (255 - g) * t)
    b = round(b + (255 - b) * t)
    return f"#{r:02X}{g:02X}{b:02X}"


def _px_to_pt(px: float) -> float:
    """Design-px (1920×1080 baseline) → PowerPoint points.

    A standard 16:9 PPT slide is 13.33×7.5in = 960×540pt. The 1920-wide
    design baseline maps 2 design-px per point (960/1920 = 0.5).
    """
    return px * _PX_TO_PT


def _emit_rect(slide, node: DSLNode, ctx: EmitContext) -> None:
    """rect X,Y WxH fill:role stroke:role stroke-width:N"""
    x, y = parse_xy(node.pos_args[0])
    w, h = parse_wh(node.pos_args[1])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, _px(x), _px(y), _px(w), _px(h))
    shape.shadow.inherit = False
    _strip_theme_style(shape)

    fill_name = node.kw_args.get("fill")
    if fill_name:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(ctx.tokens.color(fill_name))
    else:
        shape.fill.background()

    stroke_name = node.kw_args.get("stroke")
    if stroke_name:
        shape.line.color.rgb = _hex_to_rgb(ctx.tokens.color(stroke_name))
        sw = float(node.kw_args.get("stroke-width", 1))
        shape.line.width = Pt(sw * _STROKE_PX_TO_PT)
    else:
        shape.line.fill.background()


def _emit_line(slide, node: DSLNode, ctx: EmitContext) -> None:
    """line X,Y X2,Y2 stroke:role stroke-width:N"""
    x1, y1 = parse_xy(node.pos_args[0])
    x2, y2 = parse_xy(node.pos_args[1])
    line = slide.shapes.add_connector(1, _px(x1), _px(y1), _px(x2), _px(y2))
    stroke = node.kw_args.get("stroke", "fog")
    line.line.color.rgb = _hex_to_rgb(ctx.tokens.color(stroke))
    sw = float(node.kw_args.get("stroke-width", 1))
    line.line.width = Pt(sw * _STROKE_PX_TO_PT)
    _strip_theme_style(line)


def _emit_polyline(slide, node: DSLNode, ctx: EmitContext) -> None:
    """polyline X1,Y1 X2,Y2 X3,Y3 ... stroke:role stroke-width:N

    Open multi-segment line via python-pptx's freeform builder. N≥2 points.
    Stroke styled like `line`; no fill (open path).
    """
    pts = [parse_xy(p) for p in node.pos_args]
    if len(pts) < 2:
        raise ValueError(f"polyline: needs ≥2 points, got {len(pts)}")
    x0, y0 = pts[0]
    ff = slide.shapes.build_freeform(_px(x0), _px(y0), scale=1.0)
    ff.add_line_segments(
        [(_px(x), _px(y)) for x, y in pts[1:]],
        close=False,
    )
    poly = ff.convert_to_shape()
    poly.fill.background()
    stroke = node.kw_args.get("stroke", "ink")
    poly.line.color.rgb = _hex_to_rgb(ctx.tokens.color(stroke))
    sw = float(node.kw_args.get("stroke-width", 2))
    poly.line.width = Pt(sw * _STROKE_PX_TO_PT)
    _strip_theme_style(poly)


def _apply_picture_treatment(
    pil_image: Image.Image,
    treatment: str,
    *,
    duotone_dark: tuple[int, int, int] | None = None,
    duotone_light: tuple[int, int, int] | None = None,
) -> Image.Image:
    """Apply tokens.picture_treatment to a PIL image. Pure function — returns
    a new image. Supported treatments:
      - "none"          — no change
      - "desat(<frac>)" — blend toward grayscale by `frac` (0..1)
      - "duotone"       — convert to grayscale, then map luminance to a
        gradient from `duotone_dark` (shadow) to `duotone_light` (highlight).
        Caller must supply both RGB tuples; without them the image returns
        unchanged.
    """
    if not treatment or treatment == "none":
        return pil_image
    if treatment.startswith("desat(") and treatment.endswith(")"):
        try:
            frac = float(treatment[len("desat("):-1])
        except ValueError:
            return pil_image
        has_alpha = pil_image.mode in ("RGBA", "LA") or "transparency" in pil_image.info
        rgb = pil_image.convert("RGB")
        gray = rgb.convert("L").convert("RGB")
        blended = Image.blend(rgb, gray, frac)
        if has_alpha:
            alpha = pil_image.convert("RGBA").split()[3]
            r, g, b = blended.split()
            return Image.merge("RGBA", (r, g, b, alpha))
        return blended
    if treatment == "duotone":
        if duotone_dark is None or duotone_light is None:
            return pil_image
        gray = pil_image.convert("L")
        # Build a 256-entry RGB lookup from dark (luminance 0) to light (255).
        dr, dg, db = duotone_dark
        lr, lg, lb = duotone_light
        r_lut = bytes(round(dr + (lr - dr) * i / 255) for i in range(256))
        g_lut = bytes(round(dg + (lg - dg) * i / 255) for i in range(256))
        b_lut = bytes(round(db + (lb - db) * i / 255) for i in range(256))
        r = gray.point(r_lut)
        g = gray.point(g_lut)
        b = gray.point(b_lut)
        return Image.merge("RGB", (r, g, b))
    return pil_image


# ---------------------------------------------------------------------------
# Picture-query helpers (Task 7 — pluggable image provider)
# ---------------------------------------------------------------------------

# Stable slot-id derivation: lowercase, collapse non-alnum to single `_`,
# trim leading/trailing `_`, truncate to 40 chars. Used so re-running a
# build with the same DSL pins the same image again from asset_lock.json.
_SLOT_SLUG_RE = re.compile(r"[^A-Za-z0-9]+")

# MIME → file extension map for materialised cache filenames. Defaults to
# `.bin` for unknown/missing — the renderer falls back to PIL's auto-format
# detection which usually still works.
_MIME_TO_EXT = {
    "image/jpeg": ".jpg",
    "image/jpg":  ".jpg",
    "image/png":  ".png",
    "image/webp": ".webp",
    "image/gif":  ".gif",
    "image/svg+xml": ".svg",
}


def _slot_id_from_query(query: str) -> str:
    """Derive a stable, deterministic slot id from a query string.

    "Kitchen morning light!" → "kitchen_morning_light".
    Empty/non-alnum-only inputs collapse to "asset" so the lock file
    never gets an empty key.
    """
    slug = _SLOT_SLUG_RE.sub("_", query.lower()).strip("_")
    if not slug:
        slug = "asset"
    return slug[:40]


def _read_lock(deck_dir: Path | None) -> dict:
    """Read ``<deck_dir>/asset_lock.json`` or return a fresh empty lock."""
    if deck_dir is None:
        return {"version": 1, "provider": None, "slots": {}}
    lock_path = deck_dir / "asset_lock.json"
    if not lock_path.is_file():
        return {"version": 1, "provider": None, "slots": {}}
    try:
        data = json.loads(lock_path.read_text())
    except (OSError, json.JSONDecodeError):
        return {"version": 1, "provider": None, "slots": {}}
    if not isinstance(data, dict):
        return {"version": 1, "provider": None, "slots": {}}
    data.setdefault("version", 1)
    data.setdefault("provider", None)
    data.setdefault("slots", {})
    return data


def _write_lock(deck_dir: Path | None, lock: dict) -> None:
    """Persist the lock as pretty-printed JSON. No-op if deck_dir is None.

    The write is done via tmp-file + ``os.replace`` so a crashed or
    interrupted build can't leave a half-written ``asset_lock.json`` that
    future runs fail to parse. The tmp file is created in the same
    directory as the lock so the rename is atomic on POSIX (cross-
    filesystem renames aren't atomic).
    """
    if deck_dir is None:
        return
    deck_dir.mkdir(parents=True, exist_ok=True)
    lock_path = deck_dir / "asset_lock.json"
    with tempfile.NamedTemporaryFile(
        mode="w",
        dir=lock_path.parent,
        prefix=".asset_lock.",
        suffix=".tmp",
        delete=False,
    ) as tmp:
        json.dump(lock, tmp, indent=2, sort_keys=True)
        tmp.write("\n")
        tmp_path = Path(tmp.name)
    os.replace(tmp_path, lock_path)


def _hit_from_lock_entry(entry: dict) -> "ImageHit | None":
    """Reconstruct an ``ImageHit`` from a lock-file slot dict. Returns None
    if the entry is missing required fields."""
    from ..image_provider import ImageHit
    try:
        return ImageHit(
            url=entry["url"],
            license=entry.get("license", ""),
            attribution=entry.get("attribution", ""),
            width=entry.get("width"),
            height=entry.get("height"),
            mime=entry.get("mime", ""),
        )
    except KeyError:
        return None


def _url_is_resolvable(url: str) -> bool:
    """For ``file://`` URLs, verify the path exists on disk. For
    ``http(s)://`` URLs we trust the pin — pre-flighting every HEAD on
    every build would defeat the point of the lock cache."""
    if url.startswith("file://"):
        return Path(url[len("file://"):]).is_file()
    if url.startswith(("http://", "https://")):
        return True
    # Bare path — treat as filesystem; resolvable if it exists.
    return Path(url).is_file()


def _utc_now_iso_seconds() -> str:
    """ISO 8601 timestamp in UTC, truncated to seconds, with `Z` suffix.

    `datetime.now(timezone.utc).isoformat()` produces `+00:00` — we
    replace it with `Z` to match the spec example.
    """
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat().replace("+00:00", "Z")


def _entry_from_hit(hit: "ImageHit", query: str) -> dict:
    """Serialise an ImageHit + query into a lock-slot dict."""
    entry: dict = {
        "query": query,
        "url": hit.url,
        "license": hit.license,
        "attribution": hit.attribution,
        "mime": hit.mime,
        "pinned_at": _utc_now_iso_seconds(),
    }
    if hit.width is not None:
        entry["width"] = hit.width
    if hit.height is not None:
        entry["height"] = hit.height
    return entry


class _SearchError:
    """Sentinel returned by ``_lookup_lock_then_search`` when
    ``provider.search`` raised.

    Distinct from ``None`` (legitimate empty-result miss) so the caller
    can mark the ``missing_assets`` entry as ``kind="search-error"``
    rather than ``kind="no-hit"``. Carries the exception type so the
    caller can surface a useful diagnostic without re-raising.
    """
    __slots__ = ("exc_type",)

    def __init__(self, exc_type: type):
        self.exc_type = exc_type


def _lookup_lock_then_search(
    ctx: EmitContext, slot_id: str, query: str,
) -> "ImageHit | _SearchError | None":
    """Return a pinned hit for `slot_id` if available + valid; otherwise
    call ``ctx.image_provider.search(query, count=1)``, pin the first
    result, and return it. Returns ``None`` if the provider returns ``[]``.
    Returns a ``_SearchError`` sentinel if the provider raises — callers
    differentiate provider-crash from legitimate-no-hit on this signal.

    Failed searches are NOT pinned — a stale "no results" entry would
    block the slot from ever resolving, even after the provider's
    backing data improves.
    """
    provider = ctx.image_provider
    assert provider is not None  # caller guards this
    lock = _read_lock(ctx.deck_dir)

    # Lock is scoped to a provider name; a brand switch invalidates the
    # whole file (different URL schemes, different licensing).
    if lock.get("provider") == provider.name:
        slot_entry = lock["slots"].get(slot_id)
        if slot_entry and slot_entry.get("query") == query:
            pinned_url = slot_entry.get("url", "")
            if _url_is_resolvable(pinned_url):
                hit = _hit_from_lock_entry(slot_entry)
                if hit is not None:
                    return hit
            # Stale — fall through to re-search and overwrite below.

    # Either no lock entry for this slot, or it's stale, or the lock
    # belongs to a different provider. Re-search.
    try:
        hits = provider.search(query, count=1)
    except Exception as exc:
        # Provider crashed (network, library bug, bad token, …). Surface
        # via warning + sentinel so the caller writes a search-error
        # entry to missing_assets; the build still completes with a
        # placeholder rect so a single bad slot doesn't block delivery.
        warnings.warn(
            f"image provider {provider.name!r} raised on search({query!r}): "
            f"{type(exc).__name__}: {exc}",
            RuntimeWarning,
            stacklevel=3,
        )
        return _SearchError(type(exc))
    if not hits:
        return None
    hit = hits[0]

    # If the lock belongs to a different provider, blow it away rather
    # than mixing pin sources in one file.
    if lock.get("provider") != provider.name:
        lock = {"version": 1, "provider": provider.name, "slots": {}}
    lock["slots"][slot_id] = _entry_from_hit(hit, query)
    lock["provider"] = provider.name
    _write_lock(ctx.deck_dir, lock)
    return hit


def _materialise(
    hit: "ImageHit", cache_dir: Path,
) -> tuple[Path | None, Exception | None]:
    """Resolve an ``ImageHit`` URL to a local Path.

    - ``file://`` → just check the file exists.
    - ``http(s)://`` → download to ``<cache_dir>/<sha1(url)>.<ext>`` if
      not already present. Single retry. 30s timeout. Returns
      ``(None, last_err)`` on persistent failure.
    - Bare path → treat as a filesystem path.

    Returns ``(path, None)`` on success and ``(None, err_or_None)`` on
    failure. The second element carries the last exception raised on the
    HTTP path so the caller can include error diagnostics in
    ``missing_assets`` entries; it is ``None`` for non-HTTP misses where
    no exception was involved (e.g. a missing ``file://`` target).

    The sha1+ext naming keeps re-runs cheap: repeated builds against the
    same URL skip the network entirely after the first successful fetch.
    """
    url = hit.url
    if url.startswith("file://"):
        p = Path(url[len("file://"):])
        return (p, None) if p.is_file() else (None, None)
    if url.startswith(("http://", "https://")):
        cache_dir.mkdir(parents=True, exist_ok=True)
        ext = _MIME_TO_EXT.get(hit.mime.lower(), "")
        if not ext:
            # Fall back to URL path extension.
            from urllib.parse import urlparse
            url_path = urlparse(url).path
            url_ext = Path(url_path).suffix.lower()
            ext = url_ext if url_ext else ".bin"
        digest = hashlib.sha1(url.encode("utf-8")).hexdigest()
        target = cache_dir / f"{digest}{ext}"
        if target.is_file():
            return (target, None)
        # Single retry, 30s budget. urllib.request.urlretrieve uses the
        # default socket timeout — override globally for this call by
        # building a custom opener with a timeout-aware urlopen.
        last_err: Exception | None = None
        for _ in range(2):
            try:
                with urllib.request.urlopen(url, timeout=30) as resp:  # noqa: S310
                    data = resp.read()
                target.write_bytes(data)
                return (target, None)
            except (urllib.error.URLError, TimeoutError, OSError) as exc:
                last_err = exc
                continue
        return (None, last_err)
    # Bare path fallback.
    p = Path(url)
    return (p, None) if p.is_file() else (None, None)


def _emit_picture(slide, node: DSLNode, ctx: EmitContext) -> None:
    """picture X,Y WxH path:PATH cover:true

    `path` is the resolved image location — either a literal in the layout
    or interpolated from a `{{ slot }}` placeholder by the expander. If
    `path` is missing, the node is skipped silently. If `path` resolves
    to a non-existent file, a placeholder rect is emitted so the absence
    is visible at review time. `cover:true` center-crops the source image
    to the box aspect ratio (default behaviour is contain).

    Diagram-emitted picture nodes (produced by ``expand_diagram_blocks``)
    store geometry in ``kw_args`` (x, y, w, h as ints) and carry a
    ``_diagram_meta`` sentinel.  Slide-authored nodes use ``pos_args``.
    """
    # Diagram-emitted picture: geometry in kw_args, src in kw_args["src"],
    # sentinel in kw_args["_diagram_meta"].  (Slide-authored picture uses
    # pos_args for position and size.)
    if "_diagram_meta" in node.kw_args:
        x = int(node.kw_args["x"])
        y = int(node.kw_args["y"])
        w = int(node.kw_args["w"])
        h = int(node.kw_args["h"])
        path = node.kw_args.get("src")
        _pos_xy = f"{x},{y}"
        _pos_wh = f"{w}x{h}"
        query: str | None = None
    else:
        x, y = parse_xy(node.pos_args[0])
        w, h = parse_wh(node.pos_args[1])
        path = node.kw_args.get("path")
        query = node.kw_args.get("query")
        _pos_xy = node.pos_args[0]
        _pos_wh = node.pos_args[1]

    # `query:` and `path:` are mutually exclusive — a brand author who
    # set both has either copy-pasted half a migration or doesn't know
    # which mode they meant. Fail loud at emit time so the mistake
    # surfaces before the deck ships.
    if query and path:
        raise DSLError(
            f"picture at line {node.line_no}: `query:` and `path:` are "
            f"mutually exclusive (got query={query!r}, path={path!r})"
        )

    if query:
        if not ctx.image_provider:
            # The brand author wrote `query:` in a layout but forgot to
            # wire `$image_provider` in tokens.json. Silent-fallback to a
            # placeholder rect would mask the misconfiguration — fail
            # loud so the next build run surfaces the problem.
            raise DSLError(
                f"picture at line {node.line_no}: `query:{query!r}` requires "
                f"an image_provider on the EmitContext, but ctx.image_provider "
                f"is None. Add `$image_provider` to your brand's tokens.json "
                f"(or your `extends` ancestor) so the build can resolve it."
            )
        slot_id = node.label or _slot_id_from_query(query)
        hit = _lookup_lock_then_search(ctx, slot_id, query)
        if hit is None or isinstance(hit, _SearchError):
            entry: dict = {
                "kind": "search-error" if isinstance(hit, _SearchError) else "no-hit",
                "query": query,
                "slot_id": slot_id,
                "provider": ctx.image_provider.name,
                "line_no": node.line_no,
                "source": node.source,
            }
            if isinstance(hit, _SearchError):
                entry["exc_type"] = hit.exc_type.__name__
            ctx.missing_assets.append(entry)
            _emit_rect(slide, DSLNode(
                kind="rect", pos_args=[_pos_xy, _pos_wh],
                kw_args={"fill": "paper-2", "stroke": "fog"},
                label=None, line_no=node.line_no, source=node.source,
            ), ctx)
            return
        cache_dir = (ctx.deck_dir / ".cache") if ctx.deck_dir else None
        if cache_dir is None:
            # No deck_dir means we can't cache HTTP downloads. Use a
            # process-temp dir so the build still completes — but warn so
            # library callers who forgot to wire deck_dir notice the
            # misconfig (downloads won't be reused across rebuilds).
            import tempfile
            cache_dir = Path(tempfile.mkdtemp(prefix="feinschliff-imgcache-"))
            warnings.warn(
                "EmitContext.deck_dir is unset; HTTP image materialise will "
                "use a throwaway tempdir cache (no rebuild reuse). Wire "
                "deck_dir on the EmitContext (or pass it to "
                "build_presentation/build_multi_slide) to persist cached "
                "downloads in <deck_dir>/.cache/.",
                RuntimeWarning,
                stacklevel=2,
            )
        materialised, fetch_err = _materialise(hit, cache_dir)
        if materialised is None:
            fail_entry: dict = {
                "kind": "fetch-failed",
                "query": query,
                "slot_id": slot_id,
                "url": hit.url,
                "provider": ctx.image_provider.name,
                "line_no": node.line_no,
                "source": node.source,
            }
            if fetch_err is not None:
                fail_entry["error"] = f"{type(fetch_err).__name__}: {fetch_err}"
            ctx.missing_assets.append(fail_entry)
            _emit_rect(slide, DSLNode(
                kind="rect", pos_args=[_pos_xy, _pos_wh],
                kw_args={"fill": "paper-2", "stroke": "fog"},
                label=None, line_no=node.line_no, source=node.source,
            ), ctx)
            return
        path = str(materialised)

    optional = str(node.kw_args.get("optional", "false")).lower() == "true"
    if not path:
        # No image bound — emit a placeholder rect so the slot's location is
        # visible in dev renders. Required pictures append to ctx.missing_assets
        # so callers (build / deck build) can decide whether to abort. Mark
        # the picture optional via `optional:true` if a missing slot is OK.
        if not optional:
            ctx.missing_assets.append({
                "kind": "unset",
                "path": None,
                "line_no": node.line_no,
                "source": node.source,
            })
        _emit_rect(slide, DSLNode(
            kind="rect", pos_args=[_pos_xy, _pos_wh],
            kw_args={"fill": "paper-2", "stroke": "fog"},
            label=None, line_no=node.line_no, source=node.source,
        ), ctx)
        return
    p = Path(path)
    from_fallback = False
    if not p.is_absolute() and ctx.asset_root:
        primary = ctx.asset_root / p
        if primary.is_file():
            p = primary
        elif ctx.asset_root_fallback:
            fallback = ctx.asset_root_fallback / p
            if fallback.is_file():
                p = fallback
                from_fallback = True
            else:
                p = primary
        else:
            p = primary
    if not p.is_file():
        if not optional:
            ctx.missing_assets.append({
                "kind": "missing-file",
                "path": str(p),
                "line_no": node.line_no,
                "source": node.source,
            })
        _emit_rect(slide, DSLNode(
            kind="rect", pos_args=[_pos_xy, _pos_wh],
            kw_args={"fill": "paper-2", "stroke": "fog"},
            label=None, line_no=node.line_no, source=node.source,
        ), ctx)
        return

    cover = node.kw_args.get("cover", "false").lower() == "true"
    treatment = ctx.tokens.picture_treatment
    duotone_dark = duotone_light = None
    # Plugin-fallback images (generic placeholders shared across brands) get
    # duotoned to the brand palette so they read as native to each brand
    # pack instead of always showing the same warm-tone source photo. Brand-
    # specific overrides bypass this — brands choose their own treatment.
    if from_fallback:
        endpoints = _duotone_endpoints_for_brand(ctx.tokens)
        if endpoints is not None:
            treatment = "duotone"
            duotone_dark, duotone_light = endpoints
    if cover or treatment != "none":
        bytes_io = _prepare_picture_bytes(p, target_aspect=(w / h) if cover else None,
                                          treatment=treatment,
                                          duotone_dark=duotone_dark,
                                          duotone_light=duotone_light)
        slide.shapes.add_picture(bytes_io, _px(x), _px(y), width=_px(w), height=_px(h))
    else:
        # Fast path: no crop, no treatment — let python-pptx read the file.
        slide.shapes.add_picture(str(p), _px(x), _px(y), width=_px(w), height=_px(h))


def _prepare_picture_bytes(
    image_path: Path, *, target_aspect: float | None, treatment: str,
    duotone_dark: tuple[int, int, int] | None = None,
    duotone_light: tuple[int, int, int] | None = None,
) -> io.BytesIO:
    """Load → optional center-crop → optional treatment → JPEG/PNG bytes."""
    with Image.open(image_path) as im:
        im.load()
        if target_aspect is not None:
            sw, sh = im.size
            src_aspect = sw / sh
            if abs(src_aspect - target_aspect) >= 1e-3:
                if src_aspect > target_aspect:
                    new_w = max(1, int(round(sh * target_aspect)))
                    offset = (sw - new_w) // 2
                    im = im.crop((offset, 0, offset + new_w, sh))
                else:
                    new_h = max(1, int(round(sw / target_aspect)))
                    offset = (sh - new_h) // 2
                    im = im.crop((0, offset, sw, offset + new_h))
        if treatment and treatment != "none":
            im = _apply_picture_treatment(
                im, treatment,
                duotone_dark=duotone_dark, duotone_light=duotone_light,
            )
        out = io.BytesIO()
        fmt = (im.format or "PNG").upper()
        if fmt in ("JPG", "JPEG") and im.mode not in ("RGB", "L"):
            im = im.convert("RGB")
        # Treated images may have lost the source-format hint; pick by mode.
        if fmt in ("JPG", "JPEG"):
            im.save(out, format="JPEG", quality=92, optimize=True)
        else:
            im.save(out, format="PNG", optimize=True)
        return io.BytesIO(out.getvalue())


# ---------------------------------------------------------------------------
# Entry
# ---------------------------------------------------------------------------

_SHAPE_KIND = {
    "rect":          MSO_SHAPE.RECTANGLE,
    "rectangle":     MSO_SHAPE.RECTANGLE,
    "oval":          MSO_SHAPE.OVAL,
    "ellipse":       MSO_SHAPE.OVAL,
    "circle":        MSO_SHAPE.OVAL,
    "triangle":      MSO_SHAPE.ISOSCELES_TRIANGLE,
    "triangle-down": MSO_SHAPE.ISOSCELES_TRIANGLE,   # rotated 180 at emit
    "triangle-left": MSO_SHAPE.RIGHT_TRIANGLE,
    "triangle-right":MSO_SHAPE.RIGHT_TRIANGLE,
    "right-triangle":MSO_SHAPE.RIGHT_TRIANGLE,
    "chevron":       MSO_SHAPE.CHEVRON,
    "right-arrow":   MSO_SHAPE.RIGHT_ARROW,
    "left-arrow":    MSO_SHAPE.LEFT_ARROW,
    "diamond":       MSO_SHAPE.DIAMOND,
    "trapezoid":     MSO_SHAPE.TRAPEZOID,
    "pie":           MSO_SHAPE.PIE,          # 270° wedge; rotate to position the missing 90°
    "pie-wedge":     MSO_SHAPE.PIE_WEDGE,    # 90° wedge
    "arc":           MSO_SHAPE.ARC,
    "block-arc":     MSO_SHAPE.BLOCK_ARC,
}


def _emit_shape(slide, node: DSLNode, ctx: EmitContext) -> None:
    """shape X,Y WxH kind:NAME fill:role stroke:role stroke-width:N rotate:DEG

    Generic shape emitter — wraps python-pptx MSO_SHAPE so layouts can
    place ovals, triangles, chevrons, etc. Fill / stroke behaviour matches
    `rect`. `rotate` is in degrees clockwise.
    """
    x, y = parse_xy(node.pos_args[0])
    w, h = parse_wh(node.pos_args[1])
    kind = node.kw_args.get("kind", "rect")
    mso_kind = _SHAPE_KIND.get(kind)
    if mso_kind is None:
        raise ValueError(f"shape: unknown kind '{kind}' (known: {sorted(_SHAPE_KIND)})")
    shape = slide.shapes.add_shape(mso_kind, _px(x), _px(y), _px(w), _px(h))
    shape.shadow.inherit = False
    _strip_theme_style(shape)
    rot = node.kw_args.get("rotate")
    if rot is None and kind == "triangle-down":
        rot = "180"
    if rot is None and kind == "triangle-left":
        rot = "270"
    if rot is None and kind == "triangle-right":
        rot = "90"
    if rot is not None:
        shape.rotation = float(rot)

    # Preset-shape adjustment handle (e.g. parallelogram top-edge skew).
    # Maps to <a:gd name="adj" fmla="val …"/>. Range 0..1; python-pptx stores
    # as fraction. Silently no-op for shapes without an adj handle.
    adj1 = node.kw_args.get("adj1")
    if adj1 is not None:
        try:
            shape.adjustments[0] = float(adj1)
        except (IndexError, ValueError):
            pass

    fill_name = node.kw_args.get("fill")
    if fill_name:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(ctx.tokens.color(fill_name))
        opacity = node.kw_args.get("fill-opacity")
        if opacity is not None:
            _set_fill_alpha(shape, float(opacity))
    else:
        shape.fill.background()

    stroke_name = node.kw_args.get("stroke")
    if stroke_name:
        shape.line.color.rgb = _hex_to_rgb(ctx.tokens.color(stroke_name))
        sw = float(node.kw_args.get("stroke-width", 1))
        shape.line.width = Pt(sw * _STROKE_PX_TO_PT)
    else:
        shape.line.fill.background()


# ---------------------------------------------------------------------------
# Chrome sanitation — post-build XML pass (Layer 1 task 10)
# ---------------------------------------------------------------------------

_NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

# Custom attribute the emitter writes on a <p:sp> when its source DSL used
# `effect:allow` to opt out of sanitation.
_EFFECT_ALLOW_ATTR = "effect-allow"

# Outline clamp: hairline at 0.5pt = 6350 EMU; anything above 1pt (12700 EMU)
# is clamped down to hairline.
_OUTLINE_CLAMP_THRESHOLD_EMU = 12700
_OUTLINE_CLAMP_TARGET_EMU = 6350


def sanitize_chrome(slide_xml) -> None:
    """Walk a slide's <p:spTree> and remove PowerPoint defaults that read
    as "AI-templated": drop-shadow / glow / soft-edge effects, gradient
    fills (replaced with solid), and outlines wider than 1pt.

    Scoped to <p:spTree> so master/layout inheritance is preserved.
    Mutates `slide_xml` in place. Idempotent.
    """
    from lxml import etree

    sptree = slide_xml.find(f".//{{{_NS_P}}}spTree")
    if sptree is None:
        return

    # 1. effectLst — strip unless parent <p:sp> opts in.
    for effect_lst in list(sptree.iter(f"{{{_NS_A}}}effectLst")):
        parent = effect_lst.getparent()
        while parent is not None and not parent.tag.endswith("}sp"):
            parent = parent.getparent()
        if parent is not None and parent.get(_EFFECT_ALLOW_ATTR) == "1":
            continue
        effect_lst.getparent().remove(effect_lst)

    # 2. gradFill → solidFill (with the first stop's color).
    ns = {"a": _NS_A}
    for grad in list(sptree.iter(f"{{{_NS_A}}}gradFill")):
        grad_parent = grad.getparent()
        first_clr = grad.find(".//a:gs[1]/a:srgbClr", ns)
        idx = list(grad_parent).index(grad)
        grad_parent.remove(grad)
        if first_clr is None:
            continue
        color_val = first_clr.get("val") or "808080"
        solid = etree.SubElement(grad_parent, f"{{{_NS_A}}}solidFill")
        clr = etree.SubElement(solid, f"{{{_NS_A}}}srgbClr")
        clr.set("val", color_val)
        grad_parent.remove(solid)
        grad_parent.insert(idx, solid)

    # 3. Clamp outline widths > 1pt down to a 0.5pt hairline.
    for ln in sptree.iter(f"{{{_NS_A}}}ln"):
        w_str = ln.get("w")
        if w_str is None:
            continue
        try:
            w = int(w_str)
        except ValueError:
            continue
        if w > _OUTLINE_CLAMP_THRESHOLD_EMU:
            ln.set("w", str(_OUTLINE_CLAMP_TARGET_EMU))


def _set_fill_alpha(shape, opacity: float) -> None:
    """Set alpha on a solid fill via OOXML. opacity 0..1."""
    from lxml import etree
    nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    spPr = shape.fill._xPr
    solidFill = spPr.find("a:solidFill", nsmap)
    if solidFill is None:
        return
    srgbClr = solidFill.find("a:srgbClr", nsmap)
    if srgbClr is None:
        return
    alpha = etree.SubElement(srgbClr, "{%s}alpha" % nsmap["a"])
    alpha.set("val", str(int(opacity * 100000)))


def _strip_theme_style(shape) -> None:
    """Drop the `<p:style>` block from a newly-added shape.

    python-pptx's `add_shape` writes a default `<p:style>` element that
    references the theme's effect/fill/line slots (e.g. `<a:effectRef
    idx="2">` points at the second `<a:outerShdw>` entry in
    `theme1.xml`). LibreOffice and PowerPoint both honour that effect
    reference at render time *even though* `shape.shadow.inherit = False`
    writes an explicit `<a:effectLst/>` on the shape's `<p:spPr>`.

    Result before the strip: a subtle drop-shadow under every primitive
    in rendered PNGs. Flat-shape brand packs (the default, feinschliff,
    …) want the source to match the render; removing the `<p:style>`
    block does that.
    """
    nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    el = shape._element.find("p:style", nsmap)
    if el is not None:
        shape._element.remove(el)


_EMITTERS = {
    "text":     _emit_text,
    "rect":     _emit_rect,
    "line":     _emit_line,
    "polyline": _emit_polyline,
    "picture":  _emit_picture,
    "shape":    _emit_shape,
}


def _slide_canvas(nodes: list[DSLNode]) -> tuple[float, float]:
    """Read the `canvas WxH` directive if present; otherwise default 1920×1080."""
    for n in nodes:
        if n.kind == "canvas" and n.pos_args:
            return parse_wh(n.pos_args[0])
    return 1920.0, 1080.0


def _append_slide(prs: Presentation, nodes: list[DSLNode], tokens: Tokens, *,
                  asset_root: Path | None,
                  asset_root_fallback: Path | None = None,
                  missing_assets: list[dict] | None = None,
                  image_provider: "ImageProvider | None" = None,
                  deck_dir: Path | None = None) -> None:
    """Append one slide built from `nodes` to `prs`, using the tokens for fills."""
    cw, ch = _slide_canvas(nodes)
    slide = prs.slides.add_slide(prs.slide_layouts[6])    # blank
    try:
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = _hex_to_rgb(tokens.color("paper"))
    except KeyError:
        pass

    ctx = EmitContext(tokens=tokens, canvas_w=cw, canvas_h=ch,
                      asset_root=asset_root, asset_root_fallback=asset_root_fallback,
                      image_provider=image_provider, deck_dir=deck_dir)
    for n in nodes:
        cond = n.kw_args.get("if")
        if cond is not None:
            s = cond.strip()
            if not s or s.lower() == "false" or _RESIDUAL_PLACEHOLDER.search(s):
                continue
        if n.kind in ("canvas", "theme"):
            continue
        emit = _EMITTERS.get(n.kind)
        if emit is None:
            import sys
            print(f"WARN: no emitter for primitive '{n.kind}' "
                  f"(line {n.line_no}) — skipping", file=sys.stderr)
            continue
        emit(slide, n, ctx)
    if missing_assets is not None and ctx.missing_assets:
        missing_assets.extend(ctx.missing_assets)


def build_presentation(nodes: list[DSLNode], tokens: Tokens, *,
                       asset_root: Path | None = None,
                       asset_root_fallback: Path | None = None,
                       image_provider: "ImageProvider | None" = None,
                       deck_dir: Path | None = None) -> Presentation:
    """Walk primitive nodes, build a Presentation with one filled slide.

    The returned Presentation carries `missing_assets: list[dict]` as an
    attribute — empty when every required picture resolved, populated
    otherwise. Callers can branch on `prs.missing_assets` to abort.

    `asset_root_fallback` is the plugin-level shared assets dir; it is
    walked when the per-brand `asset_root` does not contain the requested
    relative path. Brand-specific files always win over plugin defaults.

    `image_provider` and `deck_dir` are the optional pair for the picture
    ``query:`` branch (Task 7). Both default to ``None``; when unset, any
    ``query:`` keyword on a picture node raises :class:`DSLError`.
    """
    cw, ch = _slide_canvas(nodes)
    prs = Presentation()
    prs.slide_width  = _px(cw)
    prs.slide_height = _px(ch)
    missing: list[dict] = []
    _append_slide(prs, nodes, tokens, asset_root=asset_root,
                  asset_root_fallback=asset_root_fallback,
                  missing_assets=missing,
                  image_provider=image_provider,
                  deck_dir=deck_dir)
    for slide in prs.slides:
        sanitize_chrome(slide._element)
    prs.missing_assets = missing
    return prs


def build_multi_slide(
    slides: list[tuple[list[DSLNode], Tokens, Path | None]],
    *,
    asset_root_fallback: Path | None = None,
    image_provider: "ImageProvider | None" = None,
    deck_dir: Path | None = None,
) -> Presentation:
    """Build a Presentation with N slides. Each slide entry is
    `(nodes, tokens, asset_root)`. The slide deck's canvas comes from
    the first slide's `canvas` directive; remaining slides reuse it.

    `asset_root_fallback` applies to every slide; see
    :func:`build_presentation` for semantics.

    `image_provider` and `deck_dir` are the optional pair for the picture
    ``query:`` branch (Task 7).
    """
    if not slides:
        raise ValueError("build_multi_slide: no slides provided")
    first_nodes, first_tokens, first_asset_root = slides[0]
    cw, ch = _slide_canvas(first_nodes)
    prs = Presentation()
    prs.slide_width  = _px(cw)
    prs.slide_height = _px(ch)
    missing: list[dict] = []
    for slide_idx, (nodes, tokens, asset_root) in enumerate(slides, start=1):
        per_slide: list[dict] = []
        _append_slide(prs, nodes, tokens, asset_root=asset_root,
                      asset_root_fallback=asset_root_fallback,
                      missing_assets=per_slide,
                      image_provider=image_provider,
                      deck_dir=deck_dir)
        for entry in per_slide:
            entry.setdefault("slide_index", slide_idx)
            missing.append(entry)
    for slide in prs.slides:
        sanitize_chrome(slide._element)
    prs.missing_assets = missing
    return prs
