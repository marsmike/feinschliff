"""`feinschliff deck …` subcommands: multi-slide composer + layout picker.

Subcommands:

  feinschliff deck build <plan.yaml> [-o OUT.pptx]
      Build one .pptx from a plan that lists N slides, each pinning a
      layout and inline content (or a content file).

  feinschliff deck pick <signals.yaml>
      Print a recommended layout id based on structured signals.

  feinschliff deck wireframe <layout.slide.dsl> --brand <id> [-o out.svg]
      Render a DSL layout as an annotated SVG wireframe (bounding boxes for
      text slots, picture slots, rect backgrounds). No PPTX round-trip needed.
      Add --overlay-pptx <file.pptx> to embed the actual rendered slide behind
      the wireframe boxes for pixel-accurate deviation analysis.

  feinschliff deck wireframe-sheet <plan.yaml> [-o sheet.svg]
      Render every slide in a plan as a wireframe and compose them into a
      single SVG contact sheet. Useful as a fast layout-regression baseline:
      store the SVGs in git, regenerate after DSL changes, and diff.

Plan schema (build):

  brand: feinschliff                       # default brand for slides
  out:   deck.pptx                         # output path; --output overrides
  slides:
    - layout: layouts/title-orange.slide.dsl
      content:
        pgmeta: "Q1 2026"
        title:  "..."
    - layout: layouts/quote.slide.dsl
      content_file: examples/v2/quote.yaml  # alternative to inline
    - layout: brands/gs-ramspau/layouts/stundenplan.slide.dsl
      brand:  gs-ramspau                    # per-slide override
      content_file: brands/gs-ramspau/examples/v2/stundenplan.yaml
"""
from __future__ import annotations

import argparse
import sys
import tempfile
from pathlib import Path

import yaml

from lib.dsl.parser import parse_file
from lib.dsl.tokens import load_tokens
from lib.dsl.expander import (
    interpolate_nodes,
    expand_compounds,
    load_compounds_for_brand,
)
from lib.dsl.pptx_emit import build_multi_slide
from lib.content_validator import validate_content, emit_defects_and_abort_message
from lib.slot_budget import compute_slot_budgets
from lib.pipeline import compile_slide
from lib.defects import fatal_kinds, format_defect
from lib.brand_discovery import find_brand
from lib.image_provider import discover_providers, get_provider
from lib.verify.deck.titles import extract_titles_from_plan
from lib.verify.deck.storyline import render_contact_sheet, write_storyline_report
from lib.pipeline_log import (
    log_event, read_events, render_text_report, summarize,
)


REPO_ROOT = Path(__file__).resolve().parents[1]
STD_COMPOUNDS = REPO_ROOT / "compounds"
BRANDS_DIR = REPO_ROOT / "brands"


def register(parser: argparse.ArgumentParser) -> None:
    sub = parser.add_subparsers(dest="deck_command", required=True)

    p_build = sub.add_parser("build", help="Build a multi-slide deck from a plan YAML")
    p_build.add_argument("plan", help="Path to the deck plan YAML")
    p_build.add_argument("-o", "--output", help="Output .pptx (overrides plan.out)")
    p_build.add_argument(
        "--skip-content-lint",
        action="store_true",
        help="Skip pre-render content lints (title-length, action-verb-leading). "
             "For emergency overrides only.",
    )
    p_build.add_argument(
        "--allow-diagram-warnings",
        action="store_true",
        help="Ship even when diagram-overflow or diagram-text-too-small "
             "defects surface. Otherwise these are fatal by default — same "
             "policy as single-slide `feinschliff build`.",
    )
    p_build.add_argument(
        "--allow-missing-assets",
        action="store_true",
        help="Ship even when a picture slot points at a missing file or is "
             "unset. Default: fatal. Mark intentionally-empty slots with "
             "`optional:true` to skip the abort without this flag.",
    )
    p_build.set_defaults(func=cmd_build)

    p_pick = sub.add_parser("pick", help="Recommend a layout for the given signals")
    p_pick.add_argument("signals", help="Path to a signals YAML (or '-' for stdin)")
    p_pick.add_argument("--top-k", type=int, default=3,
                        help="Print the top-K candidates with scores (default 3)")
    p_pick.set_defaults(func=cmd_pick)

    p_storyline = sub.add_parser(
        "storyline",
        help="Emit a title-only contact sheet from a deck_plan.json / "
             "content_plan.json (Phase 2 storyline gate).",
    )
    p_storyline.add_argument("plan", help="Path to deck_plan.json or content_plan.json")
    p_storyline.add_argument(
        "-o", "--output", required=True,
        help="Output path for the storyline_report.md",
    )
    p_storyline.add_argument(
        "--brief-summary", default=None,
        help="Optional one-line summary shown above the contact sheet.",
    )
    p_storyline.set_defaults(func=cmd_storyline)

    p_wf = sub.add_parser(
        "wireframe",
        help="Render a DSL layout as an annotated SVG wireframe.",
    )
    p_wf.add_argument("layout", help="Path to a .slide.dsl file")
    p_wf.add_argument("--brand", required=True, help="Brand id (dir name under brands/)")
    p_wf.add_argument("--content", help="YAML file with slot values (optional)")
    p_wf.add_argument(
        "--show-slots",
        action="store_true",
        help="Preserve {{ slot_name }} labels even when --content is supplied. "
             "Forces skip_interpolation=True so the wireframe shows slot structure.",
    )
    p_wf.add_argument("-o", "--output", required=True, help="Output .svg path")
    p_wf.add_argument(
        "--overlay-pptx",
        help="Path to a .pptx file whose first slide is embedded as background. "
             "Requires LibreOffice on PATH.",
    )
    p_wf.add_argument(
        "--overlay-slide", type=int, default=0,
        help="0-based slide index to use from --overlay-pptx (default 0).",
    )
    p_wf.add_argument(
        "--overlay-opacity", type=float, default=0.55,
        help="Opacity of the background slide image (0.0–1.0, default 0.55).",
    )
    p_wf.set_defaults(func=cmd_wireframe)

    p_wfs = sub.add_parser(
        "wireframe-sheet",
        help="Render all slides in a plan as a SVG wireframe contact sheet.",
    )
    p_wfs.add_argument("plan", help="Path to the deck plan YAML")
    p_wfs.add_argument("-o", "--output", required=True, help="Output .svg path")
    p_wfs.add_argument(
        "--overlay-pptx",
        help="Path to a .pptx file; each slide is embedded behind its wireframe.",
    )
    p_wfs.add_argument(
        "--overlay-opacity", type=float, default=0.55,
        help="Opacity of the background slide images (0.0–1.0, default 0.55).",
    )
    p_wfs.add_argument(
        "--show-slots",
        action="store_true",
        help="Preserve {{ slot_name }} labels in every cell — skip interpolation "
             "even when plans provide inline content or content_file. The contact "
             "sheet is intended as a layout-regression baseline, so showing slot "
             "structure is often more useful than filled content.",
    )
    p_wfs.set_defaults(func=cmd_wireframe_sheet)

    p_polish = sub.add_parser(
        "polish",
        help="Refurbish old PPTX diagrams into brand-perfect DSL artifacts.",
    )
    p_polish.add_argument("input", help=".pptx file to refurbish")
    p_polish.add_argument("-o", "--output", required=True, help="output .pptx path")
    p_polish.add_argument("--brand", default="feinschliff")
    p_polish.add_argument(
        "--refurbish-all",
        action="store_true",
        help="Auto-accept all refurbish proposals (emit DSL for every diagram slide).",
    )
    p_polish.add_argument(
        "--no-refurbish",
        action="store_true",
        help="Skip refurbish; just copy the input to the output path unchanged.",
    )
    p_polish.add_argument(
        "--refurbish-default",
        choices=("excalidraw", "svg"),
        default=None,
        help="Force a specific emitter instead of letting kind_selector choose.",
    )
    p_polish.set_defaults(func=cmd_polish)

    # `deck plan` is a thin alias for `deck storyline` — same CLI surface,
    # same handler. The mode-level work (steps 0 → 1c only, no render) is
    # the skill orchestrator's job; the CLI part is just the storyline
    # report materialization helper. See
    # skills/deck/references/modes.md::plan for the mode semantics.
    p_plan = sub.add_parser(
        "plan",
        help="Emit a title-only storyline report (alias for `deck storyline`). "
             "Implements the CLI surface of the /deck plan mode — steps 0 → 1c "
             "only, no render. See skills/deck/references/modes.md::plan.",
    )
    p_plan.add_argument("plan", help="Path to deck_plan.json or content_plan.json")
    p_plan.add_argument(
        "-o", "--output", required=True,
        help="Output path for the storyline_report.md",
    )
    p_plan.add_argument(
        "--brief-summary", default=None,
        help="Optional one-line summary shown above the contact sheet.",
    )
    p_plan.set_defaults(func=cmd_storyline)

    # ── Pipeline timing log ─────────────────────────────────────────────
    p_log = sub.add_parser(
        "log-event",
        help="Append one event to the deck's timing.jsonl. Used by the "
             "skill orchestrator to mark phase transitions (start/end).",
    )
    p_log.add_argument("phase", help="Phase name, e.g. 'step:2-plan'.")
    p_log.add_argument("status", choices=["start", "end", "tick", "fail"],
                       help="Event status.")
    p_log.add_argument("--dir", required=True,
                       help="Deck working directory (timing.jsonl lives here).")
    p_log.add_argument("--elapsed-ms", type=int, default=None,
                       help="For 'end' events: elapsed milliseconds.")
    p_log.add_argument("--agent", default=None,
                       help="Optional: agent identifier (for parallel fan-out).")
    p_log.add_argument("--slide", type=int, default=None,
                       help="Optional: slide index this event refers to.")
    p_log.add_argument("--note", default=None,
                       help="Optional one-line note for the event.")
    p_log.set_defaults(func=cmd_log_event)

    p_timing = sub.add_parser(
        "timing",
        help="Print stats from a deck's timing.jsonl — total wall clock, "
             "per-phase breakdown, and parallelism speedup.",
    )
    p_timing.add_argument("dir", help="Deck working directory.")
    p_timing.add_argument("--format", choices=["text", "json"], default="text",
                          help="Output format (default text).")
    p_timing.set_defaults(func=cmd_timing)

    # ── Parallel-generation helpers ─────────────────────────────────────
    p_skel = sub.add_parser(
        "plan-skeleton",
        help="Centrally pick a layout per slide from a content_plan.json "
             "and emit a skeleton plan.yaml with empty `content:` blocks. "
             "Fan-out authoring subagents fill the blocks in parallel.",
    )
    p_skel.add_argument("content_plan", help="Path to content_plan.json (or .yaml).")
    p_skel.add_argument("-o", "--output", required=True,
                        help="Output path for the skeleton plan.yaml.")
    p_skel.add_argument("--brand", default=None,
                        help="Override brand. Default: from content_plan or 'feinschliff'.")
    p_skel.add_argument("--out-pptx", default=None,
                        help="Sets the `out:` field of the skeleton (final pptx path).")
    p_skel.set_defaults(func=cmd_plan_skeleton)

    p_merge = sub.add_parser(
        "plan-merge",
        help="Merge per-slide content chunks (from parallel authoring "
             "subagents) into the skeleton plan.yaml.",
    )
    p_merge.add_argument("skeleton", help="Path to the skeleton plan.yaml.")
    p_merge.add_argument(
        "--chunk", action="append", default=[],
        metavar="PATH",
        help="One per slide chunk YAML. Each chunk is "
             "{index: N, content: {...}} OR a list of such entries.",
    )
    p_merge.add_argument("-o", "--output", required=True,
                         help="Output path for the merged plan.yaml.")
    p_merge.set_defaults(func=cmd_plan_merge)

    # ── Parallel-verify aspect helpers ──────────────────────────────────
    p_aspect = sub.add_parser(
        "verify-aspect",
        help="Run one focused verify aspect on a built deck. Each aspect "
             "is independently runnable; spawning multiple aspects in "
             "parallel + collating gives the parallel-verify path.",
    )
    p_aspect.add_argument(
        "aspect",
        choices=["bbox", "font", "narrative", "brand", "image", "content",
                 "notes-coherence"],
        help="bbox = bounding-box / overflow; font = legibility / role "
             "mismatches; narrative = SCQA / claim-title across titles; "
             "brand = color contrast / token discipline; image = picture "
             "slot fit / style consistency; content = title-body coherence; "
             "notes-coherence = per-slide speaker notes track the deck's "
             "red_line.",
    )
    p_aspect.add_argument("--plan", required=True,
                          help="Path to the deck plan.yaml.")
    p_aspect.add_argument("--pptx", default=None,
                          help="Path to the built .pptx (required for bbox/font/brand/image).")
    p_aspect.add_argument("--png-dir", default=None,
                          help="Directory with rendered slide-NN.png files.")
    p_aspect.add_argument("--design-brief", default=None,
                          help="Path to design_brief.json (used by narrative).")
    p_aspect.add_argument("-o", "--output", required=True,
                          help="Output path for verify-<aspect>.json.")
    p_aspect.set_defaults(func=cmd_verify_aspect)

    p_collate = sub.add_parser(
        "verify-collate",
        help="Merge per-aspect verify outputs into a single verify_report.md.",
    )
    p_collate.add_argument(
        "--aspect", action="append", default=[], metavar="PATH",
        help="One per aspect: path to a verify-<aspect>.json file.",
    )
    p_collate.add_argument("--plan", required=True,
                           help="Path to the deck plan.yaml (for slide titles).")
    p_collate.add_argument("--iteration", type=int, default=1,
                           help="Iteration number (header field).")
    p_collate.add_argument("--budget", type=int, default=3,
                           help="Iteration budget (header field).")
    p_collate.add_argument("--png-dir", default=None,
                           help="PNG directory (header reference).")
    p_collate.add_argument("-o", "--output", required=True,
                           help="Output path for verify_report.md.")
    p_collate.set_defaults(func=cmd_verify_collate)


def cmd_build(args) -> int:
    plan_path = Path(args.plan).resolve()
    if not plan_path.is_file():
        print(f"deck: plan not found: {plan_path}", file=sys.stderr)
        return 2
    plan = yaml.safe_load(plan_path.read_text()) or {}

    default_brand = plan.get("brand", "feinschliff")
    slides_spec = plan.get("slides") or []
    if not slides_spec:
        print(f"deck: plan '{plan_path}' has no slides", file=sys.stderr)
        return 2

    # Resolve build-time image provider from the default brand's
    # `$image_provider` (extends-resolved). Per-slide `brand:` overrides
    # still drive tokens/compounds, but the provider used to resolve
    # `picture query:` nodes is deck-wide — `asset_lock.json` lives next
    # to the output deck. Absent → provider is None; brands using only
    # `picture path:` build as before. A typo'd kind raises KeyError with
    # a registry listing — surfaces as the normal CLI traceback.
    discover_providers()
    provider = None
    try:
        default_brand_obj = find_brand(default_brand)
    except ValueError as e:
        print(f"deck: {e}", file=sys.stderr)
        return 2
    if default_brand_obj.image_provider_config:
        cfg = default_brand_obj.image_provider_config
        provider = get_provider(cfg["kind"], cfg.get("config"))

    # Compute the output deck path up front so it can serve as `deck_dir`
    # for `asset_lock.json` + `.cache/` during the build.
    out_path = Path(args.output or plan.get("out", "deck.pptx")).resolve()

    plan_dir = plan_path.parent
    slides_payload: list[tuple[list, object, Path]] = []
    content_defects_by_slide: dict[int, list] | None = (
        {} if not args.skip_content_lint else None
    )
    all_diagram_defects: list = []

    # Top-level timing phase for the build. Per-slide events would push
    # indentation past sanity here; instead we emit one `build:total`
    # event manually around the whole compile loop (see end of function)
    # and emit `build:slide` lightweight events inline via log_event.
    import time as _time
    _build_t0 = _time.perf_counter()
    log_event(plan_dir, "build:total", "start",
              slides=len(slides_spec), brand=default_brand)

    with tempfile.TemporaryDirectory() as _tmp:
        diagrams_out = Path(_tmp) / "diagrams"
        diagrams_out.mkdir()

        for i, spec in enumerate(slides_spec):
            layout_path = (plan_dir / spec["layout"]).resolve()
            if not layout_path.is_file():
                # also try plan-dir-relative + repo-relative.
                alt = (REPO_ROOT / spec["layout"]).resolve()
                if alt.is_file():
                    layout_path = alt
                else:
                    print(f"deck: slide {i}: layout not found: {spec['layout']}", file=sys.stderr)
                    return 2

            brand = spec.get("brand", default_brand)
            try:
                brand_dir = find_brand(brand).root
            except ValueError as e:
                print(f"deck: slide {i}: {e}", file=sys.stderr)
                return 2

            tokens = load_tokens(brand_dir, brands_dir=BRANDS_DIR)
            compounds = load_compounds_for_brand(
                brand_dir, std_dir=STD_COMPOUNDS, brands_dir=BRANDS_DIR
            )

            layout_nodes, layout_compounds = parse_file(layout_path)
            for cd in layout_compounds:
                compounds[cd.name] = cd

            ctx = spec.get("content") or {}
            if not ctx and "content_file" in spec:
                content_path = (plan_dir / spec["content_file"]).resolve()
                if not content_path.is_file():
                    alt = (REPO_ROOT / spec["content_file"]).resolve()
                    if alt.is_file():
                        content_path = alt
                    else:
                        print(f"deck: slide {i}: content_file not found: {spec['content_file']}", file=sys.stderr)
                        return 2
                ctx = yaml.safe_load(content_path.read_text()) or {}

            if content_defects_by_slide is not None:
                slide_index = i + 1
                # Strip `.slide.dsl` (and any leading path) to get the bare
                # layout name — this is what the structural validators key on.
                layout_name = layout_path.name
                if layout_name.endswith(".slide.dsl"):
                    layout_name = layout_name[: -len(".slide.dsl")]
                slot_budgets = compute_slot_budgets(layout_nodes, tokens, compounds=compounds)
                slide_defects = validate_content(
                    ctx, slide_index=slide_index, layout=layout_name,
                    slot_budgets=slot_budgets,
                )
                if slide_defects:
                    content_defects_by_slide[slide_index] = slide_defects

            _slide_t0 = _time.perf_counter()
            log_event(plan_dir, "build:slide", "start", slide=i + 1,
                      layout=Path(spec["layout"]).name)
            slide_result = compile_slide(
                layout_path=layout_path,
                ctx=ctx,
                brand_dir=brand_dir,
                slide_index=i + 1,
                diagrams_out_dir=diagrams_out,
            )
            log_event(
                plan_dir, "build:slide", "end", slide=i + 1,
                layout=Path(spec["layout"]).name,
                elapsed_ms=int((_time.perf_counter() - _slide_t0) * 1000),
                defects=len(slide_result.defects),
            )
            for d in slide_result.defects:
                print(f"deck: slide {i + 1}: {format_defect(d)}", file=sys.stderr)
            all_diagram_defects.extend(slide_result.defects)
            slides_payload.append((slide_result.primitives, slide_result.tokens, brand_dir / "assets"))

        if content_defects_by_slide:
            emit_defects_and_abort_message(content_defects_by_slide, cli_name="deck build")
            return 1

        allowed_to_skip: set[str] = set()
        if getattr(args, "allow_diagram_warnings", False):
            allowed_to_skip |= {"diagram-overflow", "diagram-text-too-small"}
        _fatal_diagram = [
            d for d in all_diagram_defects
            if d.kind.value in fatal_kinds() and d.kind.value not in allowed_to_skip
        ]
        if _fatal_diagram:
            print(
                f"deck build: aborting — {len(_fatal_diagram)} fatal defect(s) "
                f"across {len(slides_spec)} slide(s). Pass "
                f"--allow-diagram-warnings to demote "
                f"diagram-overflow/diagram-text-too-small (if those are the only blockers).",
                file=sys.stderr,
            )
            return 1

        prs = build_multi_slide(
            slides_payload,
            asset_root_fallback=REPO_ROOT / "assets",
            image_provider=provider,
            deck_dir=out_path.parent,
        )
        missing = getattr(prs, "missing_assets", []) or []
        if missing and not getattr(args, "allow_missing_assets", False):
            for entry in missing:
                kind = entry.get("kind", "missing")
                path = entry.get("path") or "(unset)"
                slide_n = entry.get("slide_index", "?")
                line = entry.get("line_no", "?")
                print(
                    f"deck build: slide {slide_n}: missing asset ({kind}) "
                    f"at line {line}: {path}",
                    file=sys.stderr,
                )
            print(
                f"deck build: aborting — {len(missing)} missing required "
                f"asset(s) across {len(slides_spec)} slide(s). Mark optional "
                f"slots with `optional:true` or pass --allow-missing-assets.",
                file=sys.stderr,
            )
            return 1
        out_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out_path))
        print(f"wrote {out_path} ({len(prs.slides)} slides)")
        log_event(
            plan_dir, "build:total", "end",
            elapsed_ms=int((_time.perf_counter() - _build_t0) * 1000),
            slides=len(prs.slides), output=str(out_path),
        )
        return 0


def cmd_pick(args) -> int:
    from lib.layout_picker import pick_layout

    if args.signals == "-":
        raw = sys.stdin.read()
    else:
        raw = Path(args.signals).read_text()
    signals = yaml.safe_load(raw) or {}

    candidates = pick_layout(
        role=signals.get("role"),
        concept_count=signals.get("concept_count"),
        data_quantity=signals.get("data_quantity"),
        comparison=signals.get("comparison"),
        narrative_role=signals.get("narrative_role"),
        narrative_act=signals.get("narrative_act"),
        time_axis_role=signals.get("time_axis_role"),
        audience_mode=signals.get("audience_mode"),
        top_k=args.top_k,
    )
    if not candidates:
        print("deck: no candidate layouts matched the signals", file=sys.stderr)
        return 1
    for c in candidates:
        rationale = c["rationale"]
        if isinstance(rationale, list):
            rationale = ", ".join(rationale)
        print(f"{c['score']:5.2f}  {c['layout']:<24}  {rationale}")
    return 0


def _build_primitives_for_layout(
    layout_path: Path, brand: str, content_path: Path | None,
    *, skip_interpolation: bool = False,
) -> tuple[list, object]:
    """Parse, expand, and return (primitives, tokens) for a single layout.

    When *skip_interpolation* is True the slot-filling pass is skipped so
    ``{{ slot_name }}`` labels are preserved in the primitives. This is the
    correct mode for wireframe rendering, where the slot structure matters
    more than the filled content.
    """
    brand_dir = find_brand(brand).root
    tokens = load_tokens(brand_dir, brands_dir=BRANDS_DIR)
    compounds = load_compounds_for_brand(
        brand_dir, std_dir=STD_COMPOUNDS, brands_dir=BRANDS_DIR
    )
    layout_nodes, layout_compounds = parse_file(layout_path)
    for cd in layout_compounds:
        compounds[cd.name] = cd
    if skip_interpolation:
        source_nodes = layout_nodes
    else:
        ctx: dict = {}
        if content_path and content_path.is_file():
            ctx = yaml.safe_load(content_path.read_text()) or {}
        source_nodes = interpolate_nodes(layout_nodes, ctx)
    primitives, _ = expand_compounds(source_nodes, compounds)
    return primitives, tokens


def cmd_wireframe(args) -> int:
    from lib.dsl.svg_wireframe import render_wireframe

    layout_path = Path(args.layout).resolve()
    if not layout_path.is_file():
        alt = (REPO_ROOT / args.layout).resolve()
        if alt.is_file():
            layout_path = alt
        else:
            print(f"deck wireframe: layout not found: {args.layout}", file=sys.stderr)
            return 2

    if args.content:
        content_path = Path(args.content).resolve()
        if not content_path.is_file():
            print(f"deck wireframe: content file not found: {args.content}", file=sys.stderr)
            return 2
    else:
        content_path = None
    # Skip interpolation when no content is given (preserves slot labels) or
    # when --show-slots forces slot-structure mode even with content provided.
    skip_interp = (content_path is None) or args.show_slots
    try:
        primitives, tokens = _build_primitives_for_layout(
            layout_path, args.brand, content_path,
            skip_interpolation=skip_interp,
        )
    except ValueError as exc:
        print(f"deck wireframe: {exc}", file=sys.stderr)
        return 2

    bg_b64: str | None = None
    if args.overlay_pptx:
        from lib.dsl.pptx_to_png import slide_to_b64
        try:
            bg_b64 = slide_to_b64(
                Path(args.overlay_pptx).resolve(),
                slide_index=args.overlay_slide,
            )
        except (FileNotFoundError, RuntimeError, IndexError) as exc:
            print(f"deck wireframe: overlay failed — {exc}", file=sys.stderr)
            return 2

    title = layout_path.name.replace(".slide.dsl", "")
    svg = render_wireframe(
        primitives, tokens,
        title=title,
        background_png_b64=bg_b64,
        background_opacity=args.overlay_opacity,
    )
    out = Path(args.output).resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(svg, encoding="utf-8")
    mode = "overlay" if bg_b64 else "wireframe"
    print(f"wrote {out} ({mode}, {len(primitives)} primitives)")
    return 0


def cmd_wireframe_sheet(args) -> int:
    from lib.dsl.svg_wireframe import render_wireframe_sheet

    plan_path = Path(args.plan).resolve()
    if not plan_path.is_file():
        print(f"deck wireframe-sheet: plan not found: {plan_path}", file=sys.stderr)
        return 2

    plan = yaml.safe_load(plan_path.read_text()) or {}
    default_brand = plan.get("brand", "feinschliff")
    slides_spec = plan.get("slides") or []
    if not slides_spec:
        print(f"deck wireframe-sheet: plan '{plan_path}' has no slides", file=sys.stderr)
        return 2

    plan_dir = plan_path.parent
    slides_data: list[tuple[list, object, str]] = []

    for i, spec in enumerate(slides_spec):
        layout_rel = spec["layout"]
        layout_path = (plan_dir / layout_rel).resolve()
        if not layout_path.is_file():
            alt = (REPO_ROOT / layout_rel).resolve()
            if alt.is_file():
                layout_path = alt
            else:
                print(f"deck wireframe-sheet: slide {i}: layout not found: {layout_rel}",
                      file=sys.stderr)
                return 2

        brand = spec.get("brand", default_brand)
        ctx_inline = spec.get("content") or {}
        content_path: Path | None = None
        if not ctx_inline and "content_file" in spec:
            cp = (plan_dir / spec["content_file"]).resolve()
            if not cp.is_file():
                cp = (REPO_ROOT / spec["content_file"]).resolve()
            content_path = cp if cp.is_file() else None

        try:
            brand_dir = find_brand(brand).root
            tokens = load_tokens(brand_dir, brands_dir=BRANDS_DIR)
            compounds = load_compounds_for_brand(
                brand_dir, std_dir=STD_COMPOUNDS, brands_dir=BRANDS_DIR
            )
            layout_nodes, layout_compounds = parse_file(layout_path)
            for cd in layout_compounds:
                compounds[cd.name] = cd
            ctx: dict = ctx_inline.copy()
            if not ctx and content_path:
                ctx = yaml.safe_load(content_path.read_text()) or {}
            if args.show_slots:
                # Skip interpolation so {{ slot_name }} labels survive into the cell.
                primitives, _ = expand_compounds(layout_nodes, compounds)
            else:
                interp = interpolate_nodes(layout_nodes, ctx)
                primitives, _ = expand_compounds(interp, compounds)
        except (ValueError, OSError, yaml.YAMLError, KeyError) as exc:
            print(f"deck wireframe-sheet: slide {i}: {exc}", file=sys.stderr)
            return 1

        title = layout_path.name.replace(".slide.dsl", "")
        slides_data.append((primitives, tokens, title))

    bg_list: list[str | None] | None = None
    if args.overlay_pptx:
        from lib.dsl.pptx_to_png import pptx_to_pngs_b64
        try:
            pngs = pptx_to_pngs_b64(Path(args.overlay_pptx).resolve())
        except (FileNotFoundError, RuntimeError) as exc:
            print(f"deck wireframe-sheet: overlay failed — {exc}", file=sys.stderr)
            return 2
        if len(pngs) != len(slides_data):
            if len(pngs) < len(slides_data):
                detail = (
                    f"only {len(pngs)} overlay(s) for {len(slides_data)} slide(s); "
                    "unmatched slides will have no overlay. "
                    "(Some LibreOffice versions only export the first slide.)"
                )
            else:
                detail = (
                    f"{len(pngs)} overlay(s) for {len(slides_data)} slide(s); "
                    "extra overlays will be ignored. Check that --overlay-pptx "
                    "matches the deck plan."
                )
            print(f"deck wireframe-sheet: warning — {detail}", file=sys.stderr)
        # Pad to match slide count; truncate is implicit via index guard in render.
        bg_list = list(pngs) + [None] * max(0, len(slides_data) - len(pngs))

    svg = render_wireframe_sheet(
        slides_data,
        background_pngs_b64=bg_list,
    )
    out = Path(args.output).resolve()
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(svg, encoding="utf-8")
    mode = "overlay sheet" if bg_list else "wireframe sheet"
    print(f"wrote {out} ({mode}, {len(slides_data)} slide(s))")
    return 0


def cmd_polish(args) -> int:
    """Walk a PPTX, extract diagram IR from each slide, emit DSL artifacts,
    and rebuild a brand-perfect polished deck.

    WIRED:
    - Reads input .pptx with python-pptx.
    - Calls extract_vector.extract_from_slide() for every slide.
    - Calls kind_selector.select_kind() (or uses --refurbish-default) to pick
      the emitter target.
    - Emits .exc.dsl (excalidraw) or .svg.dsl (svg) into <output-dir>/refurbished/.
    - Writes refurbish_report.md next to the output file.
    - Rebuilds a polished PPTX from the refurbished diagram slides using the
      excalidraw-diagram or svg-infographic layout with the brand token pack.

    NOTE: Non-diagram slides (no detectable nodes) are skipped in the rebuilt
    deck. The output PPTX contains only refurbished diagram slides. If no
    diagram slides are found, the input is copied to the output as a fallback.

    --refurbish-all is the only interactive-mode variant implemented here;
    interactive per-slide confirmation is not yet wired.
    """
    import shutil

    from pathlib import Path as _Path
    from pptx import Presentation
    from lib.diagrams.refurbish.extract_vector import extract_from_slide
    from lib.diagrams.refurbish.kind_selector import select_kind
    from lib.diagrams.refurbish.emit_excalidraw import emit as emit_excalidraw_dsl
    from lib.diagrams.refurbish.emit_svg import emit as emit_svg_dsl

    def _extract_slide_title(slide) -> str:
        """Pick the best title candidate from a refurbished slide.

        Order: (1) PPTX title placeholder if it has text; (2) the largest
        free-floating textbox (i.e., one whose shape isn't an auto-shape with
        its own label). Never returns a shape-label as the slide title.
        """
        # 1. Slide title placeholder
        try:
            title_shape = slide.shapes.title
            if title_shape and title_shape.has_text_frame and title_shape.text_frame.text.strip():
                return title_shape.text_frame.text.strip().split("\n")[0]
        except (AttributeError, ValueError):
            pass
        # 2. Largest free-floating textbox (skip shapes with their own label)
        best = ""
        best_area = 0
        for sh in slide.shapes:
            try:
                _ = sh.auto_shape_type  # raises if not an auto-shape
                continue  # skip boxes/ellipses (they're nodes, not titles)
            except (ValueError, AttributeError):
                pass
            if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip():
                area = (sh.width or 0) * (sh.height or 0)
                if area > best_area:
                    best = sh.text_frame.text.strip().split("\n")[0]
                    best_area = area
        return best

    src_path = _Path(args.input).resolve()
    if not src_path.is_file():
        print(f"deck polish: input not found: {src_path}", file=sys.stderr)
        return 2

    out_path = _Path(args.output).resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)

    if args.no_refurbish:
        shutil.copy(src_path, out_path)
        print(f"deck polish: --no-refurbish — copied {src_path.name} → {out_path}")
        return 0

    refurbish_dir = out_path.parent / "refurbished"
    refurbish_dir.mkdir(exist_ok=True)
    report_lines = ["# Refurbish Report\n"]

    refurbished_slides: list[dict] = []  # each: {layout, content}

    in_pres = Presentation(str(src_path))
    for idx, slide in enumerate(in_pres.slides, start=1):
        ir = extract_from_slide(slide)
        if not ir.nodes:
            report_lines.append(f"- slide {idx}: no diagram nodes detected — skipped")
            continue

        kind = args.refurbish_default or select_kind(ir)

        if args.refurbish_all or args.refurbish_default:
            if kind == "excalidraw":
                dsl = emit_excalidraw_dsl(ir, 1720, 480)
                ext = ".exc.dsl"
                layout_name = "excalidraw-diagram.slide.dsl"
            else:
                dsl = emit_svg_dsl(ir, 1720, 480)
                ext = ".svg.dsl"
                layout_name = "svg-infographic.slide.dsl"
            artifact = refurbish_dir / f"slide-{idx}{ext}"
            artifact.write_text(dsl)
            report_lines.append(
                f"- slide {idx}: detected {len(ir.nodes)}-node {kind} "
                f"(confidence {ir.confidence:.2f}) → {artifact.name}"
            )

            # Best-effort: pull a title from the source slide.
            # Prefer (1) the slide's title placeholder, (2) the largest free-
            # floating textbox that isn't a shape label. Never use a box's own
            # label as the slide title — that's noise.
            title = _extract_slide_title(slide)

            # Strip the leading "canvas WxH\n" line — the layout template
            # embeds diagram_dsl inside its own canvas block.
            dsl_body = dsl.partition("\n")[2]

            refurbished_slides.append({
                "layout": layout_name,
                "content": {
                    "pgmeta": f"Slide {idx}",
                    "tracker": "Refurbished",
                    "action_title": title or f"Diagram {idx}",
                    "so_what": "",
                    "source": f"Refurbished from input slide {idx}",
                    "diagram_dsl": dsl_body,
                    "footer_left": "Feinschliff",
                    "footer_right": f"Slide {idx}",
                },
            })
        else:
            report_lines.append(
                f"- slide {idx}: detected {len(ir.nodes)}-node {kind} "
                f"(confidence {ir.confidence:.2f}) — no flag set, skipped"
            )

    (out_path.parent / "refurbish_report.md").write_text("\n".join(report_lines))

    if not refurbished_slides:
        # Nothing extracted — fall back to copying input
        shutil.copy(src_path, out_path)
        slide_count = len(in_pres.slides)
        print(
            f"deck polish: processed {slide_count} slide(s), "
            f"0 diagram slides found — input copied to {out_path.name}"
        )
        return 0

    _build_refurbished_deck(
        refurbished_slides,
        brand=args.brand,
        out_path=out_path,
    )

    slide_count = len(in_pres.slides)
    artifact_count = len(refurbished_slides)
    print(
        f"deck polish: processed {slide_count} slide(s), "
        f"{artifact_count} refurbished diagram slide(s) → {out_path.name}"
    )
    return 0


def _build_refurbished_deck(slides_plan: list[dict], brand: str, out_path: Path) -> None:
    """Build a multi-slide PPTX from a plan of {layout, content} entries.

    Uses the existing build_multi_slide pipeline: per-slide nodes are expanded
    and accumulated, then emitted in a single Presentation via build_multi_slide.
    """
    import tempfile
    from lib.dsl.parser import parse_file
    from lib.dsl.expander import (
        expand_compounds, expand_diagram_blocks, interpolate_nodes,
        load_compounds_for_brand,
    )
    from lib.dsl.tokens import load_tokens
    from lib.dsl.pptx_emit import build_multi_slide

    try:
        brand_obj = find_brand(brand)
    except ValueError as e:
        raise ValueError(f"deck polish: {e}") from None
    brand_dir = brand_obj.root

    # Resolve build-time image provider for `picture query:` resolution.
    # Same semantic as `deck build`: provider is brand-scoped, lock file
    # lives next to the polished output deck.
    discover_providers()
    provider = None
    if brand_obj.image_provider_config:
        cfg = brand_obj.image_provider_config
        provider = get_provider(cfg["kind"], cfg.get("config"))

    tokens = load_tokens(brand_dir, brands_dir=BRANDS_DIR)
    compounds = load_compounds_for_brand(
        brand_dir, std_dir=STD_COMPOUNDS, brands_dir=BRANDS_DIR
    )

    slides_payload: list[tuple[list, object, Path]] = []

    with tempfile.TemporaryDirectory() as tmp:
        diagrams_out = Path(tmp) / "diagrams"
        diagrams_out.mkdir()
        for slide_idx, entry in enumerate(slides_plan, start=1):
            layout_path = (REPO_ROOT / "layouts" / entry["layout"]).resolve()
            if not layout_path.is_file():
                raise FileNotFoundError(
                    f"deck polish: layout not found: {layout_path}"
                )

            layout_nodes, layout_compounds = parse_file(layout_path)
            local_compounds = dict(compounds)
            for cd in layout_compounds:
                local_compounds[cd.name] = cd

            interp = interpolate_nodes(layout_nodes, entry["content"])
            interp = expand_diagram_blocks(
                interp,
                brand_dir=brand_dir,
                out_dir=diagrams_out,
                layout_dir=layout_path.parent,
                slide_index=slide_idx,
            )
            primitives, diagnostics = expand_compounds(interp, local_compounds)
            for d in diagnostics:
                print(f"deck polish: {d.format()}", file=sys.stderr)

            slides_payload.append((primitives, tokens, brand_dir / "assets"))

        prs = build_multi_slide(
            slides_payload,
            asset_root_fallback=REPO_ROOT / "assets",
            image_provider=provider,
            deck_dir=out_path.parent,
        )
        out_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out_path))


def cmd_storyline(args) -> int:
    plan_path = Path(args.plan).resolve()
    out_path = Path(args.output).resolve()
    try:
        titles = extract_titles_from_plan(plan_path)
    except FileNotFoundError as exc:
        print(f"deck storyline: {exc}", file=sys.stderr)
        return 2
    except ValueError as exc:
        print(f"deck storyline: {exc}", file=sys.stderr)
        return 2

    contact_sheet = render_contact_sheet(titles, brief_summary=args.brief_summary)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    write_storyline_report(out_path, contact_sheet=contact_sheet)
    print(f"wrote {out_path} ({len([t for t in titles if t])} non-empty title(s) "
          f"across {len(titles)} slide(s))")
    return 0


# ──────────────────────────────────────────────────────────────────────────────
# Timing log + statistics
# ──────────────────────────────────────────────────────────────────────────────

def cmd_log_event(args) -> int:
    """`feinschliff deck log-event <phase> <status> --dir <deck-dir>` — used
    by the skill orchestrator (or any shell script) to append a phase
    transition marker to the deck's timing.jsonl.

    Returns 0 always — logging failures must never abort the pipeline.
    """
    extra: dict = {}
    if args.agent:
        extra["agent"] = args.agent
    if args.slide is not None:
        extra["slide"] = args.slide
    if args.note:
        extra["note"] = args.note
    log_event(
        args.dir, args.phase, args.status,
        elapsed_ms=args.elapsed_ms, **extra,
    )
    return 0


def cmd_timing(args) -> int:
    """`feinschliff deck timing <dir>` — render the timing.jsonl as a
    human-readable phase summary, or emit the structured summary as JSON
    via `--format json`."""
    events = read_events(args.dir)
    if not events:
        print(f"deck timing: no timing.jsonl found in {args.dir}", file=sys.stderr)
        return 1
    summary = summarize(events)
    if args.format == "json":
        import json as _json
        print(_json.dumps({"summary": summary, "events": events},
                          indent=2, ensure_ascii=False))
    else:
        print(render_text_report(events, summary), end="")
    return 0


# ──────────────────────────────────────────────────────────────────────────────
# Parallel generation — plan-skeleton + plan-merge
# ──────────────────────────────────────────────────────────────────────────────

# Maps `diagram_kind` hints from content_plan into the layout picker's
# preferred role. Keeps the centralized pick consistent with what the AI
# would have chosen on a serial pass.
_DIAGRAM_KIND_TO_ROLE: dict[str, str] = {
    "concept": "content-with-visual",
    "chart":   "data-quantity",
    "process": "data-timeline",
    "compare": "data-comparison",
}


def _signals_from_slide(slide: dict) -> dict:
    """Extract the lib.layout_picker / lib.layout_budget kwargs from a
    content-plan slide entry. Centralized here so the per-slide and
    deck-wide pickers agree on which fields feed selection."""
    role = slide.get("role") or slide.get("purpose") or (
        _DIAGRAM_KIND_TO_ROLE.get(str(slide.get("diagram_kind") or ""))
        or "content-columns"
    )
    return {
        "role":            role,
        "concept_count":   slide.get("concept_count"),
        "data_quantity":   slide.get("data_quantity"),
        "comparison":      slide.get("comparison"),
        "narrative_role":  slide.get("narrative_role"),
        "narrative_act":   slide.get("narrative_act"),
        "time_axis_role":  slide.get("time_axis_role"),
        "audience_mode":   slide.get("audience_mode"),
        "diagram_kind":    slide.get("diagram_kind"),
    }


def cmd_plan_skeleton(args) -> int:
    """`feinschliff deck plan-skeleton <content_plan>` — centralized layout
    pick. Reads a content_plan (JSON or YAML) and writes a skeleton
    plan.yaml: one entry per slide, `layout:` filled, `content: {}` left
    empty for parallel authoring subagents to fill in.

    Layout selection runs a two-pass planner (`lib.layout_budget`) that
    re-ranks per-slide picker output with a deck-wide usage budget, so
    eligible-but-overlooked layouts (e.g. `vertical-bullets`,
    `funnel`, `pyramid`) surface instead of the same 2-3 winners
    repeating across the deck."""
    import json as _json
    from lib.layout_budget import plan_deck_layouts

    cp_path = Path(args.content_plan).resolve()
    if not cp_path.is_file():
        print(f"deck plan-skeleton: not found: {cp_path}", file=sys.stderr)
        return 2
    text = cp_path.read_text(encoding="utf-8")
    plan = (yaml.safe_load(text) if cp_path.suffix in (".yaml", ".yml")
            else _json.loads(text))
    if not isinstance(plan, dict) or not isinstance(plan.get("slides"), list):
        print(f"deck plan-skeleton: {cp_path}: missing `slides` list",
              file=sys.stderr)
        return 2

    brand = args.brand or plan.get("brand") or "feinschliff"
    out_pptx = args.out_pptx or "out/deck.pptx"

    signals = [_signals_from_slide(s) for s in plan["slides"]]
    assignments = plan_deck_layouts(signals)

    skeleton_slides: list[dict] = []
    for slide, assignment in zip(plan["slides"], assignments):
        layout = assignment["layout"]
        skeleton_slides.append({
            "layout": f"layouts/{layout}.slide.dsl",
            "content": {},  # left empty for the authoring subagent to fill
            "_meta": {
                "index": slide.get("index"),
                "title": slide.get("title")
                          or slide.get("title_draft")
                          or "(untitled)",
                "role": slide.get("role") or slide.get("purpose"),
                "diagram_kind": slide.get("diagram_kind"),
                "layout_rationale": assignment["rationale"],
            },
        })

    out = {"brand": brand, "out": out_pptx, "slides": skeleton_slides}
    out_path = Path(args.output).resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(
        yaml.safe_dump(out, sort_keys=False, allow_unicode=True),
        encoding="utf-8",
    )
    chosen = [a["layout"] for a in assignments]
    log_event(out_path.parent, "skeleton:write", "tick",
              slides=len(skeleton_slides), brand=brand,
              distinct_layouts=len(set(chosen)))
    print(f"wrote {out_path} ({len(skeleton_slides)} slides, brand={brand})")
    print(f"layouts ({len(set(chosen))} distinct): {', '.join(chosen)}")
    return 0


def cmd_plan_merge(args) -> int:
    """`feinschliff deck plan-merge --skeleton X --chunk Y --chunk Z -o W` —
    merge authored content chunks (from parallel subagents) into the
    skeleton plan.yaml.

    Each chunk is one of:
      {index: N, content: {...}}                — single slide
      [{index: 0, content: {...}}, {index: 1, ...}]  — multiple slides
      {slides: [...]}                            — plan.yaml-style fragment
    """
    skel_path = Path(args.skeleton).resolve()
    if not skel_path.is_file():
        print(f"deck plan-merge: skeleton not found: {skel_path}", file=sys.stderr)
        return 2
    plan = yaml.safe_load(skel_path.read_text(encoding="utf-8")) or {}
    slides = plan.get("slides") or []
    if not slides:
        print(f"deck plan-merge: skeleton has no slides", file=sys.stderr)
        return 2

    merged_count = 0
    for chunk_path in args.chunk:
        cp = Path(chunk_path).resolve()
        if not cp.is_file():
            print(f"deck plan-merge: chunk not found: {cp}", file=sys.stderr)
            return 2
        raw = yaml.safe_load(cp.read_text(encoding="utf-8"))
        entries: list[dict] = []
        if isinstance(raw, dict) and isinstance(raw.get("slides"), list):
            entries = raw["slides"]
        elif isinstance(raw, list):
            entries = raw
        elif isinstance(raw, dict) and "content" in raw:
            entries = [raw]
        else:
            print(f"deck plan-merge: {cp}: unrecognized chunk shape",
                  file=sys.stderr)
            return 2
        for entry in entries:
            idx = entry.get("index")
            content = entry.get("content")
            if idx is None or content is None:
                continue
            if not (0 <= idx < len(slides)):
                print(f"deck plan-merge: chunk index {idx} out of range",
                      file=sys.stderr)
                return 2
            slides[idx]["content"] = content
            # Optional: chunks may override the layout pick.
            if "layout" in entry:
                slides[idx]["layout"] = entry["layout"]
            merged_count += 1

    # Drop the `_meta` annotations from the skeleton — they were only
    # there for the authoring subagents.
    for s in slides:
        s.pop("_meta", None)

    out_path = Path(args.output).resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(
        yaml.safe_dump(plan, sort_keys=False, allow_unicode=True),
        encoding="utf-8",
    )
    empty = [i for i, s in enumerate(slides) if not s.get("content")]
    log_event(out_path.parent, "plan-merge", "tick",
              merged=merged_count, total=len(slides), empty=len(empty))
    print(f"wrote {out_path} ({merged_count} chunk entries → "
          f"{len(slides)} slides; {len(empty)} still empty)")
    if empty:
        print(f"  empty slides: {empty}", file=sys.stderr)
        return 1
    return 0


# ──────────────────────────────────────────────────────────────────────────────
# Parallel verify — focused aspect checks
# ──────────────────────────────────────────────────────────────────────────────

def cmd_verify_aspect(args) -> int:
    """`feinschliff deck verify-aspect <aspect> --plan plan.yaml -o out.json`

    Each aspect runs an independent narrow check. Designed to be spawned
    as one subagent per aspect — N aspects run in parallel, then
    `verify-collate` merges them into a single verify_report.md.

    Implementation note: aspects that need LLM judgment (narrative,
    image-style, content-cohesion) currently emit a stub finding entry
    pointing the orchestrator at the relevant PNGs/files. The deterministic
    aspects (bbox, font, brand) emit actual findings. The orchestrator
    fills LLM-judged aspects via subagent dispatch.
    """
    import json as _json
    aspect = args.aspect
    plan_path = Path(args.plan).resolve()
    plan = yaml.safe_load(plan_path.read_text(encoding="utf-8")) or {}
    slides = plan.get("slides") or []

    out: dict = {
        "aspect": aspect,
        "plan": str(plan_path),
        "iteration_check_ms": None,
        "findings": [],  # list of {slide, kind, message, severity, hint}
        "summary": "",
        "needs_llm": False,
    }

    _t0 = _time.perf_counter() if (_time := __import__("time")) else None

    if aspect == "bbox":
        # Deterministic: re-run feinschliff verify on the plan, surface
        # text-overflow / out-of-bounds / diagram-overflow defects.
        from lib.pipeline import compile_slide
        import tempfile

        with tempfile.TemporaryDirectory() as _tmp:
            diagrams_out = Path(_tmp) / "diagrams"
            diagrams_out.mkdir()
            for i, spec in enumerate(slides):
                layout_path = (plan_path.parent / spec["layout"]).resolve()
                if not layout_path.is_file():
                    layout_path = (REPO_ROOT / spec["layout"]).resolve()
                try:
                    brand_dir = find_brand(spec.get("brand")
                                           or plan.get("brand")
                                           or "feinschliff").root
                except ValueError:
                    continue
                try:
                    r = compile_slide(
                        layout_path=layout_path,
                        ctx=spec.get("content") or {},
                        brand_dir=brand_dir,
                        slide_index=i + 1,
                        diagrams_out_dir=diagrams_out,
                    )
                except Exception as e:
                    out["findings"].append({
                        "slide": i + 1, "kind": "compile-error",
                        "severity": "fatal",
                        "message": str(e)[:200], "hint": "see plan.yaml",
                    })
                    continue
                for d in r.defects:
                    out["findings"].append({
                        "slide": i + 1, "kind": d.kind.value,
                        "severity": d.severity.value,
                        "message": d.message[:240], "hint": "",
                    })

    elif aspect == "font":
        # Deterministic-ish: surface diagram-text-too-small from build,
        # plus a hint pointing the orchestrator at any slide where a text
        # primitive uses role=detail (most fragile per past runs).
        for i, spec in enumerate(slides):
            dsl = (spec.get("content") or {}).get("diagram_dsl") or ""
            if "size:detail" in dsl or " detail " in dsl:
                out["findings"].append({
                    "slide": i + 1, "kind": "detail-role-fragile",
                    "severity": "warn",
                    "message": "Diagram uses role=detail — close to the 10pt floor "
                               "after region/slide scaling; consider promoting to body.",
                    "hint": "search the diagram_dsl for `size:detail` or `detail `",
                })

    elif aspect == "narrative":
        # Reads titles + design_brief + storyline_report (if present).
        # The actual SCQA / claim-title judgment is LLM work; emit a
        # contact-sheet snapshot for the orchestrator subagent.
        titles = [
            (s.get("content") or {}).get("title")
            or (s.get("content") or {}).get("action_title")
            or s.get("_meta", {}).get("title", "?")
            for s in slides
        ]
        out["needs_llm"] = True
        out["contact_sheet"] = list(enumerate(titles, start=1))
        if args.design_brief:
            db = Path(args.design_brief).resolve()
            if db.is_file():
                try:
                    out["design_brief"] = _json.loads(db.read_text(encoding="utf-8"))
                except _json.JSONDecodeError:
                    pass

    elif aspect == "brand":
        # Deterministic-ish: scan diagram_dsl for fill tokens not in the
        # canonical 17-name vocabulary. Brand pack discipline check.
        canonical = {
            "primary", "accent", "secondary", "tertiary", "highlight",
            "ink", "neutral-strong", "neutral-soft", "graphite",
            "paper", "paper-2", "off-white", "surface-2",
            "severity-high", "severity-medium", "severity-low",
            "status-current", "status-next", "status-done",
            "success", "warning", "error", "code", "data", "inactive",
        }
        import re
        for i, spec in enumerate(slides):
            dsl = (spec.get("content") or {}).get("diagram_dsl") or ""
            for m in re.finditer(r"fill:([a-z0-9_-]+)", dsl):
                tok = m.group(1)
                if tok not in canonical:
                    out["findings"].append({
                        "slide": i + 1, "kind": "non-canonical-token",
                        "severity": "warn",
                        "message": f"Diagram uses fill:{tok}, not in the 17-name "
                                   f"semantic vocabulary.",
                        "hint": "use one of: " + ", ".join(sorted(canonical)[:10]) + ", …",
                    })

    elif aspect == "image":
        # Picture-slot presence + style consistency check. Surfaces slides
        # that declare an image_style but have no picture slot, or
        # vice-versa. LLM verification of actual style match is a separate
        # subagent step (needs_llm=True).
        if args.design_brief:
            try:
                db = _json.loads(Path(args.design_brief).read_text(encoding="utf-8"))
                out["image_style"] = db.get("image_style")
            except (OSError, _json.JSONDecodeError):
                pass
        out["needs_llm"] = True
        for i, spec in enumerate(slides):
            c = spec.get("content") or {}
            has_pic = any(k in c for k in ("hero_image", "picture", "image", "photo"))
            if has_pic and not out.get("image_style"):
                out["findings"].append({
                    "slide": i + 1, "kind": "image-style-undeclared",
                    "severity": "warn",
                    "message": "Slide uses a picture slot but design_brief.image_style "
                               "is unset.",
                    "hint": "set image_style in design_brief.json",
                })

    elif aspect == "content":
        # Title-body coherence + filler-word lint. Deterministic part:
        # scan body slots for filler words. LLM part: claim/proof match.
        FILLER = {"basically", "actually", "really", "very", "just", "simply", "in order to"}
        out["needs_llm"] = True
        for i, spec in enumerate(slides):
            c = spec.get("content") or {}
            for key in ("body", "supporting_body", "so_what"):
                val = c.get(key) or ""
                if not isinstance(val, str):
                    continue
                low = val.lower()
                hits = [w for w in FILLER if w in low]
                if hits:
                    out["findings"].append({
                        "slide": i + 1, "kind": "filler-word",
                        "severity": "warn",
                        "message": f"{key} contains filler: {', '.join(hits)}",
                        "hint": "Cut the filler — strong sentences don't need it.",
                    })

    elif aspect == "notes-coherence":
        # Pair the deck's red_line against each slide's (claim, notes).
        # The orchestrator LLM judges whether the spoken delivery tracks
        # the arc: drift / contradiction / off-arc tangents → dirty.
        from lib.verify.deck.notes_coherence import (
            SlideForCoherence,
            render_contact_sheet as _render_notes_sheet,
        )
        out["needs_llm"] = True
        red_line = ""
        if args.design_brief:
            db_path = Path(args.design_brief).resolve()
            if db_path.is_file():
                try:
                    db = _json.loads(db_path.read_text(encoding="utf-8"))
                    red_line = db.get("red_line", "") or ""
                    out["red_line"] = red_line
                    out["design_brief"] = db
                except _json.JSONDecodeError:
                    pass
        coherence_slides: list[SlideForCoherence] = []
        for i, spec in enumerate(slides):
            c = spec.get("content") or {}
            coherence_slides.append(SlideForCoherence(
                index=i,
                role=spec.get("_meta", {}).get("role", ""),
                claim=c.get("title") or c.get("action_title") or "",
                notes=spec.get("notes"),
            ))
        # Cheap deterministic pre-flag: hook slide missing notes.
        # The LLM judge handles drift / off-arc semantics.
        if coherence_slides and not (coherence_slides[0].notes or "").strip():
            out["findings"].append({
                "slide": 1, "kind": "hook-notes-missing",
                "severity": "warn",
                "message": "Hook slide has no speaker notes; expected the "
                           "deck-level storyline articulating the red_line.",
                "hint": "Author the full red_line arc into slide 1's notes.",
            })
        out["contact_sheet"] = _render_notes_sheet(red_line, coherence_slides)

    out["iteration_check_ms"] = int(
        (__import__("time").perf_counter() - (_t0 or 0)) * 1000
    )
    out["summary"] = (
        f"{len(out['findings'])} finding(s) "
        f"({'needs LLM' if out['needs_llm'] else 'deterministic only'})"
    )

    out_path = Path(args.output).resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(_json.dumps(out, indent=2, ensure_ascii=False),
                        encoding="utf-8")
    log_event(plan_path.parent, f"verify-aspect:{aspect}", "tick",
              elapsed_ms=out["iteration_check_ms"],
              findings=len(out["findings"]))
    print(f"wrote {out_path} — {out['summary']}")
    return 0


def cmd_verify_collate(args) -> int:
    """`feinschliff deck verify-collate --aspect a.json --aspect b.json -o report.md`
    Merge per-aspect verify outputs into a unified verify_report.md.
    """
    import json as _json
    plan_path = Path(args.plan).resolve()
    plan = yaml.safe_load(plan_path.read_text(encoding="utf-8")) or {}
    slides = plan.get("slides") or []

    # Aggregate findings grouped by slide index. Each finding carries its
    # source aspect so the report makes the parallel checks visible.
    by_slide: dict[int, list[dict]] = {}
    aspects_seen: list[str] = []
    needs_llm: list[str] = []
    total_findings = 0

    for path_str in args.aspect:
        p = Path(path_str).resolve()
        if not p.is_file():
            print(f"deck verify-collate: aspect file not found: {p}",
                  file=sys.stderr)
            continue
        try:
            data = _json.loads(p.read_text(encoding="utf-8"))
        except _json.JSONDecodeError as e:
            print(f"deck verify-collate: {p}: {e}", file=sys.stderr)
            continue
        aspect = data.get("aspect", p.stem)
        aspects_seen.append(aspect)
        if data.get("needs_llm"):
            needs_llm.append(aspect)
        for f in data.get("findings", []) or []:
            f = dict(f)
            f["aspect"] = aspect
            by_slide.setdefault(int(f.get("slide", 0)), []).append(f)
            total_findings += 1

    # Decide verdict: clean iff no findings + no LLM aspects waiting.
    fatal = any(
        f.get("severity") == "fatal"
        for fs in by_slide.values() for f in fs
    )
    dirty = total_findings > 0 or fatal
    verdict = "dirty" if dirty else ("pending-llm" if needs_llm else "clean")

    lines: list[str] = []
    lines.append(f"# Verify Report — {plan_path.name}")
    lines.append("")
    lines.append(f"- **Iteration:** {args.iteration} of {args.budget}")
    lines.append(f"- **Verdict:** {verdict}"
                 + (f" — {total_findings} finding(s) across "
                    f"{len(by_slide)} slide(s)" if dirty else ""))
    lines.append(f"- **Aspects checked:** {', '.join(aspects_seen) or '(none)'}")
    if needs_llm:
        lines.append(f"- **Pending LLM verdict from:** {', '.join(needs_llm)}")
    if args.png_dir:
        lines.append(f"- **Rendered PNGs:** `{args.png_dir}`")
    lines.append("")
    lines.append("---")
    lines.append("")

    for i, spec in enumerate(slides, start=1):
        title = (spec.get("content") or {}).get("title") \
                or (spec.get("content") or {}).get("action_title") \
                or spec.get("_meta", {}).get("title", "(untitled)")
        layout_name = Path(spec.get("layout", "?")).name.replace(".slide.dsl", "")
        findings = by_slide.get(i, [])
        if not findings:
            lines.append(f"## Slide {i} — {title!r} ({layout_name}) ✅")
            lines.append("")
            lines.append("_No defects._")
        else:
            lines.append(f"## Slide {i} — {title!r} ({layout_name})")
            lines.append("")
            for f in findings:
                lines.append(
                    f"- **[{f['aspect']}/{f.get('kind', '?')}]** "
                    f"{f.get('message', '')}"
                )
                if f.get("hint"):
                    lines.append(f"  → {f['hint']}")
        lines.append("")

    out_path = Path(args.output).resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text("\n".join(lines), encoding="utf-8")
    log_event(plan_path.parent, "verify-collate", "tick",
              aspects=len(aspects_seen), findings=total_findings,
              verdict=verdict, iteration=args.iteration)
    print(f"wrote {out_path} — verdict: {verdict}, "
          f"{total_findings} finding(s) across {len(by_slide)} slide(s)")
    return 0
